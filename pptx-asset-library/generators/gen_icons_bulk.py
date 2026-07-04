# -*- coding: utf-8 -*-
"""ICN 카테고리 대량 확충 — ICN-021 ~ ICN-200 (+180).

- 전부 네이티브 도형 조합(외부 이미지/`<p:pic>` 금지, license self).
- 2색 원칙: 네이비(navy_800/600) 베이스 + 포인트1 액센트. 화이트/그레이는 중립.
- 36개 글리프 × 5색 스윕(cyan/teal/blue/purple + 모노navy) = 180.
- 각 아이콘 = 자체완결 <p:grpSp>(asset:ICN-0NN) — 획 도형 + 하단 한글 라벨을 통째로 그룹화.
- 슬라이드당 8개(4x2), 파일당 60개(≤8슬라이드). 3개 파일로 분산.
- 마지막에 write_fragment('ICN_bulk', entries) 1회.
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE as M

R = c.role
CC = c.C

STROKE = CC["navy_800"]
NAVY = CC["navy_800"]
NAVY6 = CC["navy_600"]
WHITE = CC["white"]
GRAY = CC["gray_300"]
GRAY1 = CC["gray_100"]
LW = 2.0

# ---- 포인트색 스윕 (accent 라벨, RGB, 한글명) ----
COLORS = [
    ("cyan",   CC["cyan_500"],   "시안"),
    ("teal",   CC["teal_500"],   "틸"),
    ("blue",   CC["blue_500"],   "블루"),
    ("purple", CC["purple_600"], "퍼플"),
    ("navy",   CC["navy_600"],   "네이비"),
]


class Painter:
    """96x96 가상 그리드 → 아이콘 박스 내 절대 인치. acc = 이 변형의 포인트색."""
    def __init__(self, slide, ix, iy, S, acc):
        self.s, self.ix, self.iy, self.S, self.acc = slide, ix, iy, S, acc
        self.shapes = []

    def X(self, u): return self.ix + (u / 96.0) * self.S
    def Y(self, u): return self.iy + (u / 96.0) * self.S
    def L(self, u): return (u / 96.0) * self.S

    def box(self, x, y, w, h, fill=None, line=STROKE, lw=LW,
            shape=M.RECTANGLE, rot=None):
        sp = c.add_box(self.s, self.X(x), self.Y(y), self.L(w), self.L(h),
                       fill=fill, line=line, line_w=lw, shape=shape)
        if rot is not None:
            sp.rotation = rot
        self.shapes.append(sp)
        return sp

    def oval(self, x, y, w, h, fill=None, line=STROKE, lw=LW):
        return self.box(x, y, w, h, fill=fill, line=line, lw=lw, shape=M.OVAL)

    def line(self, x1, y1, x2, y2, color=STROKE, w=LW):
        cn = c.connector(self.s, self.X(x1), self.Y(y1), self.X(x2), self.Y(y2),
                         color=color, w=w)
        self.shapes.append(cn)
        return cn

    def glyph_text(self, sp, text, size=20, color=None):
        c.set_shape_text(sp, text, size=size, bold=True,
                         color=color or WHITE, align=PP_ALIGN.CENTER)


# =========================================================
# 글리프 드로잉 (각 함수: painter P, P.acc = 포인트색)
# =========================================================
def g_folder(P):
    a = P.acc
    P.box(20, 26, 26, 10, fill=NAVY6, line=None, shape=M.ROUND_2_SAME_RECTANGLE)  # 탭
    P.box(18, 32, 60, 44, fill=WHITE, line=STROKE, shape=M.ROUNDED_RECTANGLE)     # 본체
    P.box(20, 34, 56, 7, fill=a, line=None)                                        # 라벨줄(액센트)

def g_file(P):
    a = P.acc
    P.box(28, 12, 40, 72, fill=WHITE, line=STROKE, shape=M.FOLDED_CORNER)          # 페이지+접힘
    for yy in (34, 46, 58):
        P.line(36, yy, 60, yy, color=a, w=2.5)                                     # 본문줄(액센트)

def g_piechart(P):
    a = P.acc
    P.oval(20, 20, 56, 56, fill=NAVY6, line=STROKE)                                # 원판
    P.box(48, 22, 28, 26, fill=a, line=None, shape=M.PIE)                          # 조각(액센트)

def g_linechart(P):
    a = P.acc
    P.line(24, 16, 24, 80, color=STROKE, w=2.5)                                    # y축
    P.line(24, 80, 82, 80, color=STROKE, w=2.5)                                    # x축
    P.line(28, 66, 44, 52, color=a, w=3)                                           # 추세선(액센트)
    P.line(44, 52, 58, 60, color=a, w=3)
    P.line(58, 60, 76, 30, color=a, w=3)
    for (px, py) in ((28, 66), (44, 52), (58, 60), (76, 30)):
        P.oval(px - 3, py - 3, 6, 6, fill=a, line=None)                            # 데이터점

def g_arrowdown(P):
    a = P.acc
    P.box(22, 22, 12, 20, fill=NAVY, line=None)                                    # 하강 막대들
    P.box(38, 34, 12, 20, fill=NAVY6, line=None)
    P.box(54, 44, 12, 20, fill=NAVY6, line=None)
    P.box(66, 22, 16, 56, fill=a, line=None, shape=M.DOWN_ARROW)                   # 하향 화살표(액센트)

def g_checklist(P):
    a = P.acc
    for yy in (22, 44, 66):
        P.box(20, yy, 14, 14, fill=WHITE, line=STROKE, lw=1.8)                     # 체크박스
        P.box(40, yy + 3, 36, 8, fill=NAVY6, line=None)                            # 항목줄
    for yy in (22, 44):
        P.line(23, yy + 8, 27, yy + 12, color=a, w=2.5)                            # 체크(액센트)
        P.line(27, yy + 12, 32, yy + 3, color=a, w=2.5)

def g_xcircle(P):
    a = P.acc
    P.oval(20, 20, 56, 56, fill=NAVY, line=None)                                   # 원
    P.line(35, 35, 61, 61, color=a, w=5)                                           # X(액센트)
    P.line(61, 35, 35, 61, color=a, w=5)

def g_warning(P):
    a = P.acc
    P.box(14, 16, 68, 62, fill=a, line=None, shape=M.ISOSCELES_TRIANGLE)           # 삼각(액센트)
    P.box(45, 36, 6, 22, fill=WHITE, line=None)                                    # 느낌표 몸
    P.oval(45, 62, 6, 6, fill=WHITE, line=None)                                    # 느낌표 점

def g_info(P):
    a = P.acc
    sp = P.oval(20, 20, 56, 56, fill=a, line=None)                                 # 원(액센트)
    P.glyph_text(sp, "i", size=26, color=WHITE)                                    # i 글리프

def g_question(P):
    a = P.acc
    sp = P.oval(20, 20, 56, 56, fill=a, line=None)                                 # 원(액센트)
    P.glyph_text(sp, "?", size=26, color=WHITE)                                    # ? 글리프

def g_server(P):
    a = P.acc
    P.box(20, 22, 56, 22, fill=NAVY, line=None, shape=M.ROUNDED_RECTANGLE)         # 랙1
    P.box(20, 52, 56, 22, fill=NAVY6, line=None, shape=M.ROUNDED_RECTANGLE)        # 랙2
    P.oval(28, 29, 8, 8, fill=a, line=None)                                        # LED(액센트)
    P.oval(28, 59, 8, 8, fill=a, line=None)
    P.line(44, 33, 68, 33, color=WHITE, w=2)
    P.line(44, 63, 68, 63, color=WHITE, w=2)

def g_upload(P):
    a = P.acc
    P.box(20, 58, 56, 18, fill=WHITE, line=STROKE, shape=M.ROUNDED_RECTANGLE)      # 트레이
    P.box(41, 18, 14, 36, fill=a, line=None, shape=M.UP_ARROW)                     # 업(액센트)

def g_download(P):
    a = P.acc
    P.box(20, 58, 56, 18, fill=WHITE, line=STROKE, shape=M.ROUNDED_RECTANGLE)      # 트레이
    P.box(41, 18, 14, 36, fill=a, line=None, shape=M.DOWN_ARROW)                   # 다운(액센트)

def g_share(P):
    a = P.acc
    P.line(60, 24, 24, 46, color=STROKE, w=2)                                      # 연결선
    P.line(60, 24, 24, 72, color=STROKE, w=2)
    P.oval(56, 16, 18, 18, fill=a, line=None)                                      # 노드(액센트)
    P.oval(16, 38, 18, 18, fill=NAVY, line=None)
    P.oval(16, 64, 18, 18, fill=NAVY, line=None)

def g_lock(P):
    a = P.acc
    P.box(30, 14, 36, 30, fill=None, line=STROKE, lw=4, shape=M.BLOCK_ARC, rot=180)  # 고리
    P.box(24, 40, 48, 40, fill=NAVY, line=None, shape=M.ROUNDED_RECTANGLE)         # 몸체
    P.oval(42, 50, 12, 12, fill=a, line=None)                                      # 열쇠구멍(액센트)
    P.box(46, 58, 4, 12, fill=a, line=None)

def g_key(P):
    a = P.acc
    P.box(16, 32, 30, 30, fill=None, line=a, lw=4.5, shape=M.DONUT)                # 머리(액센트)
    P.line(44, 47, 80, 47, color=NAVY, w=4.5)                                      # 샤프트
    P.box(66, 47, 4, 14, fill=NAVY, line=None)                                     # 이빨
    P.box(74, 47, 4, 12, fill=NAVY, line=None)

def g_star(P):
    a = P.acc
    P.box(18, 18, 60, 60, fill=a, line=STROKE, lw=1.5, shape=M.STAR_5_POINT)       # 별(액센트)

def g_heart(P):
    a = P.acc
    P.box(18, 22, 60, 56, fill=a, line=None, shape=M.HEART)                        # 하트(액센트)

def g_bookmark(P):
    a = P.acc
    P.box(30, 12, 36, 62, fill=a, line=None)                                       # 리본(액센트)
    P.box(30, 60, 36, 22, fill=WHITE, line=None, shape=M.ISOSCELES_TRIANGLE)       # 하단 V노치

def g_bell(P):
    a = P.acc
    P.oval(44, 8, 8, 8, fill=NAVY, line=None)                                      # 꼭지
    P.oval(28, 16, 40, 40, fill=a, line=None)                                      # 종 상단(액센트)
    P.box(28, 36, 40, 26, fill=a, line=None)                                       # 종 몸
    P.box(22, 58, 52, 9, fill=NAVY, line=None, shape=M.ROUNDED_RECTANGLE)          # 림
    P.oval(43, 66, 10, 10, fill=NAVY, line=None)                                   # 추

def g_user(P):
    a = P.acc
    P.oval(37, 16, 22, 22, fill=NAVY, line=None)                                   # 머리
    P.box(28, 44, 40, 36, fill=a, line=None, shape=M.ROUND_2_SAME_RECTANGLE)       # 어깨(액센트)

def g_group(P):
    a = P.acc
    P.oval(18, 22, 18, 18, fill=NAVY6, line=None)                                  # 뒤 사람
    P.box(12, 44, 30, 28, fill=NAVY6, line=None, shape=M.ROUND_2_SAME_RECTANGLE)
    P.oval(50, 22, 18, 18, fill=a, line=None)                                      # 앞 사람(액센트)
    P.box(46, 44, 30, 28, fill=a, line=None, shape=M.ROUND_2_SAME_RECTANGLE)

def g_building(P):
    a = P.acc
    P.box(26, 16, 44, 64, fill=WHITE, line=STROKE)                                 # 본체
    for ry in (24, 38, 52):
        for rx in (32, 44, 56):
            P.box(rx, ry, 8, 8, fill=a, line=None)                                 # 창문(액센트)
    P.box(42, 66, 12, 14, fill=NAVY, line=None)                                    # 출입구

def g_factory(P):
    a = P.acc
    P.box(18, 40, 20, 20, fill=NAVY, line=None, shape=M.ISOSCELES_TRIANGLE)        # 톱니지붕
    P.box(38, 40, 20, 20, fill=NAVY, line=None, shape=M.ISOSCELES_TRIANGLE)
    P.box(18, 44, 60, 36, fill=NAVY6, line=None)                                   # 본체
    P.box(62, 20, 10, 30, fill=NAVY, line=None)                                    # 굴뚝
    P.box(30, 60, 12, 20, fill=a, line=None)                                       # 문(액센트)

def g_globe(P):
    a = P.acc
    P.oval(20, 20, 56, 56, fill=WHITE, line=STROKE, lw=2.2)                        # 지구
    P.oval(38, 20, 20, 56, fill=None, line=a, lw=2)                                # 경선(액센트)
    P.line(22, 48, 74, 48, color=a, w=2)                                           # 적도
    P.line(28, 34, 68, 34, color=NAVY6, w=1.5)
    P.line(28, 62, 68, 62, color=NAVY6, w=1.5)

def g_compass(P):
    a = P.acc
    P.box(20, 20, 56, 56, fill=WHITE, line=STROKE, lw=2.5, shape=M.DONUT)          # 링
    P.box(43, 28, 10, 22, fill=a, line=None, shape=M.ISOSCELES_TRIANGLE)           # 북침(액센트)
    P.box(43, 46, 10, 22, fill=NAVY6, line=None, shape=M.ISOSCELES_TRIANGLE, rot=180)  # 남침
    P.oval(43, 43, 10, 10, fill=NAVY, line=None)                                   # 중심

def g_timer(P):
    a = P.acc
    P.box(43, 8, 10, 8, fill=NAVY, line=None)                                      # 상단 버튼
    P.box(22, 20, 52, 56, fill=WHITE, line=STROKE, lw=2.5, shape=M.DONUT)          # 링
    P.line(48, 48, 48, 32, color=NAVY, w=3.5)                                      # 시침
    P.line(48, 48, 62, 54, color=a, w=2.5)                                         # 분침(액센트)
    P.oval(44, 44, 8, 8, fill=NAVY, line=None)

def g_tablet(P):
    a = P.acc
    P.box(24, 12, 48, 72, fill=WHITE, line=STROKE, shape=M.ROUNDED_RECTANGLE)      # 본체
    P.box(29, 20, 38, 52, fill=a, line=None)                                       # 화면(액센트)
    P.oval(45, 75, 6, 6, fill=NAVY, line=None)                                     # 홈버튼

def g_laptop(P):
    a = P.acc
    P.box(24, 18, 48, 34, fill=WHITE, line=STROKE)                                 # 화면 프레임
    P.box(28, 22, 40, 26, fill=a, line=None)                                       # 스크린(액센트)
    P.box(16, 52, 64, 10, fill=NAVY, line=None, shape=M.TRAPEZOID)                 # 받침

def g_camera(P):
    a = P.acc
    P.box(30, 22, 18, 8, fill=NAVY, line=None)                                     # 뷰파인더 돌출
    P.box(18, 28, 60, 42, fill=NAVY6, line=None, shape=M.ROUNDED_RECTANGLE)        # 바디
    P.oval(37, 37, 24, 24, fill=WHITE, line=STROKE, lw=2)                          # 렌즈
    P.oval(43, 43, 12, 12, fill=a, line=None)                                      # 렌즈 중심(액센트)
    P.oval(64, 33, 6, 6, fill=a, line=None)                                        # 플래시(액센트)

def g_filter(P):
    a = P.acc
    P.box(18, 20, 60, 24, fill=a, line=None, shape=M.TRAPEZOID, rot=180)           # 깔때기(액센트)
    P.box(40, 42, 16, 20, fill=NAVY6, line=None)                                   # 관
    P.box(40, 60, 16, 16, fill=NAVY, line=None, shape=M.TRAPEZOID)                 # 하단 방울

def g_refresh(P):
    a = P.acc
    P.box(22, 22, 52, 52, fill=a, line=None, shape=M.CIRCULAR_ARROW)               # 순환 화살표(액센트)

def g_plus(P):
    a = P.acc
    P.oval(18, 18, 60, 60, fill=WHITE, line=STROKE, lw=2.2)                        # 원
    P.box(34, 34, 28, 28, fill=a, line=None, shape=M.MATH_PLUS)                    # 플러스(액센트)

def g_battery(P):
    a = P.acc
    P.box(20, 34, 52, 28, fill=WHITE, line=STROKE, shape=M.ROUNDED_RECTANGLE)      # 몸체
    P.box(72, 42, 6, 12, fill=NAVY, line=None)                                     # 단자
    P.box(24, 38, 30, 20, fill=a, line=None)                                       # 충전량(액센트)

def g_map(P):
    a = P.acc
    P.box(18, 20, 60, 56, fill=WHITE, line=STROKE)                                 # 지도면
    P.line(38, 20, 38, 76, color=NAVY6, w=1.5)                                     # 접힘선
    P.line(58, 20, 58, 76, color=NAVY6, w=1.5)
    P.oval(44, 34, 16, 16, fill=a, line=None)                                      # 핀 머리(액센트)
    P.box(48, 46, 8, 14, fill=a, line=None, shape=M.ISOSCELES_TRIANGLE, rot=180)   # 핀 꼬리

def g_home(P):
    a = P.acc
    P.box(18, 12, 60, 34, fill=a, line=None, shape=M.ISOSCELES_TRIANGLE)           # 지붕(액센트)
    P.box(28, 40, 40, 40, fill=WHITE, line=STROKE)                                 # 본체
    P.box(42, 56, 12, 24, fill=NAVY, line=None)                                    # 문


# =========================================================
# 글리프 레지스트리 (glyph_key, 한글라벨, fn, tags)
# =========================================================
GLYPHS = [
    ("folder",    "폴더",     g_folder,    ["폴더", "분류", "자료"]),
    ("file",      "파일",     g_file,      ["파일", "문서", "페이지"]),
    ("piechart",  "파이차트", g_piechart,  ["차트", "비중", "통계"]),
    ("linechart", "선그래프", g_linechart, ["차트", "추세", "성과"]),
    ("arrowdown", "감소",     g_arrowdown, ["감소", "하락", "하향"]),
    ("checklist", "체크리스트", g_checklist, ["체크리스트", "점검", "항목"]),
    ("xcircle",   "취소",     g_xcircle,   ["취소", "불가", "제외"]),
    ("warning",   "경고",     g_warning,   ["경고", "주의", "위험"]),
    ("info",      "정보",     g_info,      ["정보", "안내", "설명"]),
    ("question",  "문의",     g_question,  ["문의", "질문", "FAQ"]),
    ("server",    "서버",     g_server,    ["서버", "인프라", "호스팅"]),
    ("upload",    "업로드",   g_upload,    ["업로드", "전송", "등록"]),
    ("download",  "다운로드", g_download,  ["다운로드", "내려받기", "저장"]),
    ("share",     "공유",     g_share,     ["공유", "확산", "배포"]),
    ("lock",      "잠금",     g_lock,      ["잠금", "보안", "비밀"]),
    ("key",       "인증",     g_key,       ["인증", "권한", "열쇠"]),
    ("star",      "즐겨찾기", g_star,      ["즐겨찾기", "우수", "추천"]),
    ("heart",     "관심",     g_heart,     ["관심", "선호", "만족"]),
    ("bookmark",  "북마크",   g_bookmark,  ["북마크", "저장", "표시"]),
    ("bell",      "알림",     g_bell,      ["알림", "공지", "이벤트"]),
    ("user",      "사용자",   g_user,      ["사용자", "회원", "개인"]),
    ("group",     "그룹",     g_group,     ["그룹", "팀", "커뮤니티"]),
    ("building",  "기관",     g_building,  ["기관", "기업", "건물"]),
    ("factory",   "생산",     g_factory,   ["생산", "제조", "공장"]),
    ("globe",     "글로벌",   g_globe,     ["글로벌", "해외", "세계"]),
    ("compass",   "방향",     g_compass,   ["방향", "전략", "탐색"]),
    ("timer",     "타이머",   g_timer,     ["타이머", "측정", "속도"]),
    ("tablet",    "태블릿",   g_tablet,    ["태블릿", "기기", "콘텐츠"]),
    ("laptop",    "노트북",   g_laptop,    ["노트북", "업무", "온라인"]),
    ("camera",    "카메라",   g_camera,    ["카메라", "촬영", "미디어"]),
    ("filter",    "필터",     g_filter,    ["필터", "선별", "정제"]),
    ("refresh",   "새로고침", g_refresh,   ["새로고침", "갱신", "순환"]),
    ("plus",      "추가",     g_plus,      ["추가", "신규", "확대"]),
    ("battery",   "배터리",   g_battery,   ["배터리", "전력", "지속"]),
    ("map",       "지도",     g_map,       ["지도", "위치", "권역"]),
    ("home",      "홈",       g_home,      ["홈", "메인", "기본"]),
]
assert len(GLYPHS) == 36, len(GLYPHS)


# =========================================================
# 시트 배치
# =========================================================
def title_block(slide, text, sub):
    c.add_text(slide, 0.55, 0.42, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    c.add_text(slide, 0.55, 0.92, 12.2, 0.3, sub, size=11,
               color=R("muted_text"), align=PP_ALIGN.LEFT)


def place_icon(slide, asset_id, label, fn, acc, col, row):
    S = 1.25
    left, top = 0.75, 1.55
    colw, rowh = (13.333 - 1.5) / 4.0, 2.75
    cellx = left + col * colw
    ix = cellx + (colw - S) / 2.0
    iy = top + row * rowh + 0.30
    # 개별 ID 캡션 — 라이브러리 라벨(그룹 밖)
    c.add_text(slide, ix, iy - 0.30, S, 0.24, asset_id, size=9, bold=True,
               color=CC["gray_500"], align=PP_ALIGN.CENTER)
    P = Painter(slide, ix, iy, S, acc)
    fn(P)
    lbl = c.add_text(slide, cellx, iy + S + 0.06, colw, 0.32, label, size=12, bold=True,
                     color=R("body_text"), align=PP_ALIGN.CENTER)
    c.group_asset(slide, P.shapes + [lbl], asset_id)


# =========================================================
# 빌드 — glyph-major 넘버링, 60개/파일 3분할
# =========================================================
# 전체 아이템 목록 생성: ICN-021 .. ICN-200
items = []  # (asset_id, glyph_key, label, fn, tags, acc_name, acc_rgb, acc_kr)
n = 21
for gi, (gkey, label, fn, tags) in enumerate(GLYPHS):
    for (cname, crgb, ckr) in COLORS:
        aid = "ICN-%03d" % n
        items.append((aid, gkey, label, fn, tags, cname, crgb, ckr))
        n += 1
assert len(items) == 180, len(items)
assert items[-1][0] == "ICN-200", items[-1][0]

GROUPS = [
    ("g1", items[0:60]),     # ICN-021 ~ ICN-080
    ("g2", items[60:120]),   # ICN-081 ~ ICN-140
    ("g3", items[120:180]),  # ICN-141 ~ ICN-200
]

entries = []
saved_files = []

for gname, group in GROUPS:
    file_rel = "decks/11_icons/ICN_bulk_%s_v1.pptx" % gname
    p = c.new_deck()
    # 8개/슬라이드
    chunks = [group[i:i + 8] for i in range(0, len(group), 8)]
    for si, chunk in enumerate(chunks):
        s = c.blank_slide(p)
        rng = "%s ~ %s" % (chunk[0][0], chunk[-1][0])
        c.id_caption(s, rng)
        title_block(s, "아이콘 대량세트 · ICN [%s] (%d/%d)" % (gname, si + 1, len(chunks)),
                    "네이티브 도형 조합 · 2색(네이비+포인트) · 편집 가능 · 36글리프×5색")
        for k, it in enumerate(chunk):
            aid, gkey, label, fn, tags, cname, crgb, ckr = it
            place_icon(s, aid, label, fn, crgb, col=k % 4, row=k // 4)
        slide_idx = si + 1
        # 엔트리 기록
        for it in chunk:
            aid, gkey, label, fn, tags, cname, crgb, ckr = it
            full_tags = ["아이콘"] + tags + ["플랫", "2색", ckr]
            entries.append(c.entry(
                aid, "ICN", "아이콘 · %s (%s)" % (label, ckr), file_rel, slide_idx,
                full_tags,
                {"glyph": gkey, "colors": 2, "accent": cname},
                {"label": label},
                ["color", "label", "size"],
                recommended_use=["항목 아이콘", "불릿 강조", "섹션 상징"],
            ))
    out = c.save_deck(p, file_rel)
    saved_files.append(out)

frag = c.write_fragment("ICN_bulk", entries)

print("=== ICN BULK DONE ===")
for f in saved_files:
    print("SAVED:", f)
print("FRAGMENT:", frag)
print("ENTRIES:", len(entries))
print("RANGE:", entries[0]["id"], "~", entries[-1]["id"])
