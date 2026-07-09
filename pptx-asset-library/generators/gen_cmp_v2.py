# -*- coding: utf-8 -*-
"""CMP 카테고리 비교/매트릭스 변형팩 v2 — CMP-011~CMP-018.
v1과 다른 형태: 하비볼 비교표 / 3단계 요금·등급 카드 / 장단점 2열 /
가중치 스코어카드 / 기능 매트릭스 그리드 / SWOT 2x2 / 의사결정 매트릭스 / 순위 비교표."""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

R = c.role
CC = c.C

# ---------- 로컬 헬퍼 ----------
def add_table(slide, x, y, w, h, nrows, ncols, col_widths=None, row_h=None):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    return gf, tbl

def set_cell(cell, text, size=11, bold=False, color=None, fill=None,
             align=PP_ALIGN.CENTER, font=None):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill is not None else R("row_base")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(5)
    cell.margin_top = cell.margin_bottom = Pt(3)
    tf = cell.text_frame; tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))

def title_block(slide, text, sub=None):
    c.add_text(slide, 0.55, 0.52, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.02, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)

# =========================================================
# CMP-011 하비볼 비교표 (●◕◐◔○ 채움 정도)
# =========================================================
def harvey_table(slide, asset_id, title, sub, cols, rows):
    """rows: (항목, [glyph…]). glyph는 (기호,색) 튜플."""
    title_block(slide, title, sub)
    ncol = len(cols)
    x, y, w = 1.1, 1.7, 11.1
    first_w = 4.1
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    nrows = len(rows) + 1
    rh = 0.66
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, ncol,
                        col_widths=col_widths, row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    for j, hd in enumerate(cols):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (item, cells) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), item, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j, (sym, col) in enumerate(cells, start=1):
            set_cell(tbl.cell(i, j), sym, size=17, bold=True, color=col, fill=base)

# =========================================================
# CMP-012 3단계 요금/등급 카드 (v1과 다른 세로 강조 레이아웃)
# =========================================================
def tier_price_cards(slide, asset_id, title, sub, tiers):
    """tiers: dict {name, badge, price, unit, feats[], color, recommended}."""
    title_block(slide, title, sub)
    shapes = []
    n = len(tiers)
    x_start, total_w, gap = 1.1, 11.1, 0.4
    cw = (total_w - gap * (n - 1)) / n
    base_y, base_h = 1.95, 4.55
    for i, t in enumerate(tiers):
        rec = t.get("recommended")
        cx = x_start + i * (cw + gap)
        cy = base_y - (0.22 if rec else 0.0)
        ch = base_h + (0.44 if rec else 0.0)
        line_col = R("accent_point") if rec else R("border")
        line_w = 2.5 if rec else 1.0
        card = c.add_box(slide, cx, cy, cw, ch, fill=R("row_base"), line=line_col,
                         line_w=line_w, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(card)
        # 등급 뱃지(작은 pill)
        bw = 1.9
        badge = c.add_box(slide, cx + (cw - bw) / 2, cy + 0.22, bw, 0.38, fill=t["color"],
                          line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.set_shape_text(badge, t["badge"], size=10, bold=True, color=R("header_text"))
        shapes.append(badge)
        # 플랜명
        nm = c.add_text(slide, cx + 0.15, cy + 0.72, cw - 0.3, 0.4, t["name"], size=15,
                        bold=True, color=R("header_fill"), align=PP_ALIGN.CENTER, font=c.FONT_H)
        shapes.append(nm)
        # 대형 가격/등급 숫자
        pr = c.add_text(slide, cx + 0.15, cy + 1.14, cw - 0.3, 0.7, t["price"], size=34,
                        bold=True, color=t["color"], align=PP_ALIGN.CENTER, font=c.FONT_H)
        shapes.append(pr)
        un = c.add_text(slide, cx + 0.15, cy + 1.86, cw - 0.3, 0.28, t["unit"], size=10,
                        color=R("muted_text"), align=PP_ALIGN.CENTER)
        shapes.append(un)
        # 구분선
        div = c.add_box(slide, cx + 0.4, cy + 2.25, cw - 0.8, 0.02, fill=R("border"), line=None)
        shapes.append(div)
        # 기능 목록 (✓)
        fy = cy + 2.42
        for feat in t["feats"]:
            chk = c.add_text(slide, cx + 0.32, fy, 0.32, 0.34, "✓", size=12, bold=True,
                             color=t["color"], align=PP_ALIGN.CENTER)
            ft = c.add_text(slide, cx + 0.62, fy, cw - 0.8, 0.34, feat, size=10.5,
                            color=R("body_text"), align=PP_ALIGN.LEFT)
            shapes.extend([chk, ft]); fy += 0.42
        if rec:
            rb = 1.5
            ribbon = c.add_box(slide, cx + cw - rb - 0.12, cy - 0.02, rb, 0.4, fill=R("warn"),
                               line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            c.set_shape_text(ribbon, "★ 추천", size=10, bold=True, color=R("header_text"))
            shapes.append(ribbon)
    c.group_asset(slide, shapes, asset_id)

# =========================================================
# CMP-013 장단점 2열 (Pros / Cons 색 구분)
# =========================================================
def pros_cons(slide, asset_id, title, sub, pros, cons):
    title_block(slide, title, sub)
    shapes = []
    y0, ch, cw, gap = 1.85, 4.5, 5.35, 0.4
    x_p = 1.1
    x_c = x_p + cw + gap
    cfg = [(x_p, "장점 (Pros)", "＋", R("accent_secondary"), CC["gray_050"], pros),
           (x_c, "단점 (Cons)", "－", R("warn"), CC["beige_100"], cons)]
    for cx, head, mark, accent, bg, items in cfg:
        card = c.add_box(slide, cx, y0, cw, ch, fill=bg, line=accent, line_w=1.5,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(card)
        hd = c.add_box(slide, cx, y0, cw, 0.8, fill=accent, line=None,
                       shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        c.set_shape_text(hd, head, size=15, bold=True, color=R("header_text"))
        shapes.append(hd)
        fy = y0 + 1.05
        for it in items:
            badge = c.add_box(slide, cx + 0.3, fy + 0.04, 0.34, 0.34, fill=accent, line=None,
                              shape=MSO_SHAPE.OVAL)
            c.set_shape_text(badge, mark, size=13, bold=True, color=R("header_text"))
            tx = c.add_text(slide, cx + 0.78, fy - 0.02, cw - 1.05, 0.5, it, size=11.5,
                            color=R("body_text"), align=PP_ALIGN.LEFT)
            shapes.extend([badge, tx]); fy += 0.72
    c.group_asset(slide, shapes, asset_id)

# =========================================================
# CMP-014 가중치 스코어카드 (항목·가중·점수·환산 + 합계)
# =========================================================
def scorecard_table(slide, asset_id, title, sub, rows, total):
    """rows: (항목, 가중, 점수, 환산). total: (가중합, 환산합)."""
    title_block(slide, title, sub)
    x, y, w = 1.6, 1.7, 10.1
    nrows = len(rows) + 2
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 4,
                        col_widths=[4.3, 1.9, 1.9, 2.0], row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    for j, hd in enumerate(["평가 항목", "가중치", "점수(/100)", "환산점수"]):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (item, wgt, sc, conv) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), item, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        set_cell(tbl.cell(i, 1), wgt, size=11, fill=base, color=R("muted_text"))
        set_cell(tbl.cell(i, 2), sc, size=11, fill=base, color=R("body_text"))
        set_cell(tbl.cell(i, 3), conv, size=11, bold=True, fill=CC["gray_050"], color=R("accent_primary"))
    li = nrows - 1
    set_cell(tbl.cell(li, 0), "종합 환산점수", size=12, bold=True, color=R("header_text"),
             fill=CC["navy_600"], align=PP_ALIGN.LEFT)
    set_cell(tbl.cell(li, 1), total[0], size=12, bold=True, color=R("header_text"), fill=CC["navy_600"])
    set_cell(tbl.cell(li, 2), "—", size=12, bold=True, color=R("header_text"), fill=CC["navy_600"])
    set_cell(tbl.cell(li, 3), total[1], size=13, bold=True, color=R("accent_point"), fill=CC["navy_600"])

# =========================================================
# CMP-015 기능 매트릭스 그리드 (항목×옵션 ✓/✗)
# =========================================================
def feature_matrix(slide, asset_id, title, sub, cols, rows):
    """rows: (기능, [state…]). state: 'y'|'n'|'p'."""
    title_block(slide, title, sub)
    ncol = len(cols)
    x, y, w = 1.6, 1.7, 10.1
    first_w = 4.3
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    nrows = len(rows) + 1
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, ncol,
                        col_widths=col_widths, row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    glyph = {"y": ("✓", R("accent_secondary")), "n": ("✗", R("warn")), "p": ("△", R("muted_text"))}
    for j, hd in enumerate(cols):
        fill = R("accent_primary") if j == 1 else R("header_fill")
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (feat, states) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), feat, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j, st in enumerate(states, start=1):
            sym, col = glyph[st]
            cfill = CC["gray_050"] if j == 1 else base
            set_cell(tbl.cell(i, j), sym, size=15, bold=True, color=col, fill=cfill)

# =========================================================
# CMP-016 SWOT 2x2 (4분면 라벨 + 불릿 내용)
# =========================================================
def swot_2x2(slide, asset_id, title, sub, quads):
    """quads: 4 dict {letter, label, items[], color} 순서 S,W,O,T (TL,TR,BL,BR)."""
    title_block(slide, title, sub)
    shapes = []
    x0, x1, y0, y1 = 1.35, 11.95, 1.75, 6.45
    gap = 0.3
    bw = (x1 - x0 - gap) / 2
    bh = (y1 - y0 - gap) / 2
    coords = [(x0, y0), (x0 + bw + gap, y0), (x0, y0 + bh + gap), (x0 + bw + gap, y0 + bh + gap)]
    for (bx, by), q in zip(coords, quads):
        panel = c.add_box(slide, bx, by, bw, bh, fill=CC["gray_050"], line=q["color"],
                          line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(panel)
        # 코너 라벨 원
        badge = c.add_box(slide, bx + 0.16, by + 0.16, 0.6, 0.6, fill=q["color"], line=None,
                          shape=MSO_SHAPE.OVAL)
        c.set_shape_text(badge, q["letter"], size=20, bold=True, color=R("header_text"), font=c.FONT_H)
        shapes.append(badge)
        lab = c.add_text(slide, bx + 0.86, by + 0.22, bw - 1.0, 0.5, q["label"], size=14,
                         bold=True, color=q["color"], align=PP_ALIGN.LEFT, font=c.FONT_H)
        shapes.append(lab)
        fy = by + 0.92
        for it in q["items"]:
            dot = c.add_box(slide, bx + 0.28, fy + 0.11, 0.11, 0.11, fill=q["color"], line=None,
                            shape=MSO_SHAPE.OVAL)
            tx = c.add_text(slide, bx + 0.52, fy - 0.02, bw - 0.75, 0.36, it, size=10.5,
                            color=R("body_text"), align=PP_ALIGN.LEFT)
            shapes.extend([dot, tx]); fy += 0.42
    c.group_asset(slide, shapes, asset_id)

# =========================================================
# CMP-017 의사결정 매트릭스 (옵션×기준 점수표, 최고안 강조)
# =========================================================
def decision_matrix(slide, asset_id, title, sub, criteria, options, best_idx):
    """criteria: 기준 헤더 리스트(마지막=종합). options: (옵션명, [점수…])."""
    title_block(slide, title, sub)
    cols = ["대안"] + criteria
    ncol = len(cols)
    x, y, w = 1.1, 1.7, 11.1
    first_w = 3.0
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    nrows = len(options) + 1
    rh = 0.66
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, ncol,
                        col_widths=col_widths, row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    last = ncol - 1
    for j, hd in enumerate(cols):
        fill = R("accent_primary") if j == last else R("header_fill")
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (name, scores) in enumerate(options, start=1):
        is_best = (i - 1) == best_idx
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), name, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j, sv in enumerate(scores, start=1):
            if j == last:
                cfill = R("accent_primary") if is_best else CC["gray_050"]
                col = R("header_text") if is_best else R("accent_primary")
                set_cell(tbl.cell(i, j), sv, size=13, bold=True, color=col, fill=cfill)
            else:
                set_cell(tbl.cell(i, j), sv, size=11, fill=base, color=R("body_text"))

# =========================================================
# CMP-018 순위 비교표 (순위 뱃지)
# =========================================================
def ranking_table(slide, asset_id, title, sub, cols, rows):
    """rows: (순위int, 대상, 점수, 비고). 상위 3위 색 뱃지."""
    title_block(slide, title, sub)
    ncol = len(cols)
    x, y, w = 1.6, 1.7, 10.1
    col_widths = [1.4, 4.5, 2.1, 2.1]
    nrows = len(rows) + 1
    rh = 0.62
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, ncol,
                        col_widths=col_widths, row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    rank_fill = {1: R("accent_primary"), 2: R("accent_point"), 3: R("accent_secondary")}
    for j, hd in enumerate(cols):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER)
    for i, (rank, tgt, score, note) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        if rank in rank_fill:
            set_cell(tbl.cell(i, 0), "%d위" % rank, size=13, bold=True,
                     color=R("header_text"), fill=rank_fill[rank])
        else:
            set_cell(tbl.cell(i, 0), "%d위" % rank, size=11, bold=True,
                     color=R("muted_text"), fill=base)
        set_cell(tbl.cell(i, 1), tgt, size=11, bold=(rank <= 3), fill=base, align=PP_ALIGN.LEFT)
        set_cell(tbl.cell(i, 2), score, size=11, bold=True, fill=base, color=R("accent_primary"))
        set_cell(tbl.cell(i, 3), note, size=10, fill=base, color=R("muted_text"))

# =========================================================
# 빌드
# =========================================================
entries = []
def E(*a, **k):
    entries.append(c.entry(*a, **k))

FG = "decks/04_compare/CMP_matrix-grid_v2.pptx"
FS = "decks/04_compare/CMP_scorecard_v2.pptx"
FP = "decks/04_compare/CMP_panel-card_v2.pptx"

# 하비볼 기호
HB_FULL = ("●", R("accent_secondary"))
HB_HIGH = ("◕", R("accent_secondary"))
HB_HALF = ("◐", R("accent_primary"))
HB_LOW = ("◔", R("muted_text"))
HB_NONE = ("○", R("muted_text"))

# ---- FILE G: 매트릭스/그리드 (CMP-011, 015, 017) ----
pg = c.new_deck()

s = c.blank_slide(pg); c.id_caption(s, "CMP-011")
harvey_table(s, "CMP-011", "하비볼 비교표 (충족도 우열)", "●(우수) ◕ ◐ ◔ ○(미흡) 채움 정도로 항목별 우열 표기",
             ["평가 항목", "자사", "경쟁사 A", "경쟁사 B"],
             [("글로벌 대응 역량", [HB_FULL, HB_HALF, HB_LOW]),
              ("협업 프로세스", [HB_FULL, HB_HIGH, HB_HALF]),
              ("AI 활용 성숙도", [HB_HIGH, HB_LOW, HB_NONE]),
              ("품질 관리 체계", [HB_FULL, HB_HALF, HB_HALF]),
              ("가격 경쟁력", [HB_HALF, HB_FULL, HB_LOW]),
              ("납기 대응력", [HB_FULL, HB_HALF, HB_HIGH])])

s = c.blank_slide(pg); c.id_caption(s, "CMP-015")
feature_matrix(s, "CMP-015", "기능 매트릭스 그리드 (제공 여부)", "기능 항목별 옵션 제공 여부를 ✓ / △ / ✗ 로 표기",
               ["기능 항목", "기본", "표준", "프리미엄"],
               [("실시간 대시보드", ["y", "y", "y"]),
                ("AI 자동화 엔진", ["n", "y", "y"]),
                ("글로벌 다국어 지원", ["n", "p", "y"]),
                ("전담 지원 매니저", ["n", "n", "y"]),
                ("맞춤 개발/연동", ["n", "p", "y"]),
                ("고급 리포트/분석", ["p", "y", "y"])])

s = c.blank_slide(pg); c.id_caption(s, "CMP-017")
decision_matrix(s, "CMP-017", "의사결정 매트릭스 (대안 평가)", "기준별 점수와 종합점수로 최적 대안 선정 (최고안 강조)",
                ["기술성", "경제성", "실현성", "확장성", "종합점수"],
                [("대안 A (자체구축)", ["9", "7", "8", "9", "33"]),
                 ("대안 B (협업제휴)", ["8", "9", "7", "8", "32"]),
                 ("대안 C (외주위탁)", ["7", "8", "9", "7", "31"])],
                best_idx=0)

fg = c.save_deck(pg, FG)

# ---- FILE S: 스코어카드/순위 (CMP-014, 018) ----
ps = c.new_deck()

s = c.blank_slide(ps); c.id_caption(s, "CMP-014")
scorecard_table(s, "CMP-014", "가중치 스코어카드 (환산점수)", "항목별 가중치 × 점수 = 환산점수, 종합 환산점수 합산",
                [("글로벌 역량", "30%", "90", "27.0"),
                 ("협업 역량", "20%", "88", "17.6"),
                 ("AI 활용도", "20%", "85", "17.0"),
                 ("사업 이해도", "15%", "92", "13.8"),
                 ("실행 가능성", "15%", "86", "12.9")],
                ("100%", "88.3"))

s = c.blank_slide(ps); c.id_caption(s, "CMP-018")
ranking_table(s, "CMP-018", "순위 비교표 (종합 순위)", "종합 점수 기준 순위 뱃지 — 상위 3위 색상 강조",
              ["순위", "평가 대상", "종합 점수", "비고"],
              [(1, "자사 컨소시엄", "88.3", "선정 우선"),
               (2, "경쟁사 A", "82.1", "차순위"),
               (3, "경쟁사 B", "78.5", "예비"),
               (4, "경쟁사 C", "71.2", "탈락"),
               (5, "경쟁사 D", "65.9", "탈락")])

fs = c.save_deck(ps, FS)

# ---- FILE P: 패널/카드 (CMP-012, 013, 016) ----
pp = c.new_deck()

s = c.blank_slide(pp); c.id_caption(s, "CMP-012")
tier_price_cards(s, "CMP-012", "3단계 요금/등급 카드 (플랜 비교)", "등급 뱃지 · 대형 가격 · ✓ 기능 목록 · 추천 강조 (v1과 다른 세로 강조형)",
                 [{"name": "Basic", "badge": "기본형", "price": "100", "unit": "만원 / 월",
                   "color": R("muted_text"), "feats": ["핵심 기능 제공", "표준 기술 지원", "기본 리포트"]},
                  {"name": "Standard", "badge": "표준형", "price": "250", "unit": "만원 / 월",
                   "color": R("accent_primary"), "recommended": True,
                   "feats": ["전체 기능 제공", "AI 활용 지원", "글로벌 대응", "우선 지원"]},
                  {"name": "Enterprise", "badge": "확장형", "price": "맞춤", "unit": "별도 견적",
                   "color": R("sub_header"), "feats": ["맞춤 개발 포함", "전담 협업팀", "무제한 확장"]}])

s = c.blank_slide(pp); c.id_caption(s, "CMP-013")
pros_cons(s, "CMP-013", "장단점 2열 (Pros / Cons)", "장점·단점을 색으로 구분한 2열 대비 패널",
          ["글로벌 표준 대응 체계 확보", "AI 기반 품질 자동화", "통합 협업 플랫폼 운영", "데이터 기반 의사결정"],
          ["초기 도입 비용 부담", "내부 역량 확보 기간 필요", "레거시 시스템 연동 복잡", "변화 관리 리스크"])

s = c.blank_slide(pp); c.id_caption(s, "CMP-016")
swot_2x2(s, "CMP-016", "SWOT 2x2 분석", "강점·약점·기회·위협 4분면에 요인을 불릿으로 정리",
         [{"letter": "S", "label": "강점 (Strength)", "color": R("accent_primary"),
           "items": ["글로벌 네트워크 보유", "검증된 기술 역량", "풍부한 수행 실적"]},
          {"letter": "W", "label": "약점 (Weakness)", "color": R("warn"),
           "items": ["AI 인력 부족", "브랜드 인지도 열위", "자금 조달 제약"]},
          {"letter": "O", "label": "기회 (Opportunity)", "color": R("accent_secondary"),
           "items": ["해외 시장 확대", "정부 지원 정책", "디지털 전환 수요"]},
          {"letter": "T", "label": "위협 (Threat)", "color": CC["navy_600"],
           "items": ["경쟁 심화", "환율 변동성", "규제 강화 흐름"]}])

fp = c.save_deck(pp, FP)

# =========================================================
# 매니페스트 엔트리
# =========================================================
E("CMP-011", "CMP", "하비볼 비교표 (충족도 우열)", FG, 1,
  ["비교표", "하비볼", "우열비교", "충족도", "경쟁분석"],
  {"cols": 4, "rows": 6, "glyphs": ["●", "◕", "◐", "◔", "○"]},
  {"rows": ["글로벌 대응 역량", "협업 프로세스", "AI 활용 성숙도", "품질 관리 체계", "가격 경쟁력", "납기 대응력"],
   "cols": ["평가 항목", "자사", "경쟁사 A", "경쟁사 B"]},
  ["rows", "cell_glyphs", "glyph_color", "headers"],
  recommended_use=["경쟁력 비교", "충족도 평가", "제안 우위"])

E("CMP-012", "CMP", "3단계 요금/등급 카드 (플랜 비교)", FP, 1,
  ["비교표", "요금카드", "등급카드", "플랜비교", "추천강조"],
  {"columns": 3, "recommended": "Standard", "price_emphasis": True},
  {"rows": ["Basic", "Standard", "Enterprise"], "cols": ["기본형", "표준형", "확장형"]},
  ["tier_names", "badges", "prices", "units", "feats", "tier_color", "recommended_badge"],
  recommended_use=["요금제 비교", "등급 제안", "가격 정책"])

E("CMP-013", "CMP", "장단점 2열 (Pros / Cons)", FP, 2,
  ["비교표", "장단점", "proscons", "2열대비", "색구분"],
  {"columns": 2, "pros_items": 4, "cons_items": 4},
  {"quadrants": ["pros", "cons"], "axis": ["state"]},
  ["pros_items", "cons_items", "pros_color", "cons_color"],
  recommended_use=["장단점 정리", "의사결정 근거", "리스크 제시"])

E("CMP-014", "CMP", "가중치 스코어카드 (환산점수)", FS, 1,
  ["평가표", "스코어카드", "가중치", "환산점수", "정량평가"],
  {"cols": 4, "rows": 5, "total_row": True, "computed_col": "환산"},
  {"rows": ["글로벌 역량", "협업 역량", "AI 활용도", "사업 이해도", "실행 가능성"],
   "cols": ["평가 항목", "가중치", "점수", "환산점수"]},
  ["rows", "weights", "scores", "converted", "total"],
  recommended_use=["정량 평가", "심사 대응", "가중 채점"])

E("CMP-015", "CMP", "기능 매트릭스 그리드 (제공 여부)", FG, 2,
  ["비교표", "기능매트릭스", "그리드", "제공여부", "체크"],
  {"cols": 4, "rows": 6, "states": ["y", "p", "n"], "symbols": ["✓", "△", "✗"]},
  {"rows": ["실시간 대시보드", "AI 자동화 엔진", "글로벌 다국어 지원", "전담 지원 매니저", "맞춤 개발/연동", "고급 리포트/분석"],
   "cols": ["기능 항목", "기본", "표준", "프리미엄"]},
  ["rows", "cell_states", "state_symbols", "headers"],
  recommended_use=["기능 비교", "옵션 제공표", "제안 범위"])

E("CMP-016", "CMP", "SWOT 2x2 분석", FP, 3,
  ["매트릭스", "2x2", "SWOT", "사분면", "전략분석"],
  {"quadrants": 4, "items_per_quad": 3, "layout": "SWOT"},
  {"quadrants": ["S", "W", "O", "T"], "axis": ["internal", "external"]},
  ["quad_letters", "quad_labels", "quad_items", "quad_color"],
  recommended_use=["SWOT 분석", "전략 진단", "환경 분석"])

E("CMP-017", "CMP", "의사결정 매트릭스 (대안 평가)", FG, 3,
  ["매트릭스", "의사결정", "대안평가", "점수표", "최적안"],
  {"criteria": 4, "options": 3, "highlight_col": "종합점수", "best_row": True},
  {"rows": ["대안 A", "대안 B", "대안 C"],
   "cols": ["대안", "기술성", "경제성", "실현성", "확장성", "종합점수"]},
  ["options", "criteria", "scores", "best_idx", "highlight_col"],
  recommended_use=["대안 비교", "의사결정 근거", "선정 논리"])

E("CMP-018", "CMP", "순위 비교표 (종합 순위)", FS, 2,
  ["비교표", "순위표", "랭킹", "순위뱃지", "종합점수"],
  {"cols": 4, "rows": 5, "rank_badge": True, "top_n_highlight": 3},
  {"rows": ["1위", "2위", "3위", "4위", "5위"],
   "cols": ["순위", "평가 대상", "종합 점수", "비고"]},
  ["rows", "rank", "target", "score", "note", "rank_fill"],
  recommended_use=["순위 제시", "종합 평가", "심사 결과"])

frag = c.write_fragment("CMP_v2", entries)

print("FILE_G", fg)
print("FILE_S", fs)
print("FILE_P", fp)
print("FRAG", frag)
print("ENTRIES", len(entries))
