# -*- coding: utf-8 -*-
"""
06_org (ORG) 대량 확충 배치 생성기 — ORG-015 ~ ORG-200 (+186)
구조패밀리 × 색 스킴 × 노드수 스윕. 파라미터화 1-슬라이드-1-에셋.

규칙:
- 색은 c.role / c.C 만 사용(navy_600/navy_900 raw = c.C[]). 한글은 c 헬퍼(set_kfont) 경유.
- 다중도형 에셋 = c.group_asset(slide, shapes, aid) (그룹명이 앵커). id_caption/slide_title 은 그룹 제외.
- 단일 표 에셋 = graphicFrame 에 c.name_asset 만 (그룹 미사용).
- 페이지이미지·SmartArt 금지. 커넥터는 c.connector(add_connector) 기반.
- 파일은 패밀리별 분산: decks/06_org/ORG_bulk_<family>_v1.pptx (각 ≤25 슬라이드).
- 마지막에 c.write_fragment('ORG_bulk', entries) 1회.

함정 대응:
- cell.merge() 는 None 반환 → 병합 후 원점셀에 텍스트.
- MSO_SHAPE 는 사전 조회 상수화. 리터럴 중복 금지(라벨/색 풀 상수화).
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

R = c.role
CC = c.C
RR = MSO_SHAPE.ROUNDED_RECTANGLE
RECT = MSO_SHAPE.RECTANGLE
OVAL = MSO_SHAPE.OVAL
UPARROW = MSO_SHAPE.UP_ARROW

DIR = "decks/06_org/"
EMU = 914400.0

# ─────────────────────────────────────────────────────────────
# 색 스킴 스윕 (header 색 + node 색 팔레트). role/C 만.
# ─────────────────────────────────────────────────────────────
SCHEMES = [
    {"key": "navy",   "hdr": R("header_fill"),
     "nodes": [R("accent_primary"), R("accent_secondary"), R("sub_header"), R("accent_point"), CC["navy_600"]]},
    {"key": "blue",   "hdr": R("accent_primary"),
     "nodes": [CC["navy_600"], R("accent_secondary"), R("sub_header"), R("accent_point"), CC["navy_800"]]},
    {"key": "purple", "hdr": R("sub_header"),
     "nodes": [R("accent_primary"), R("accent_point"), R("accent_secondary"), CC["navy_600"], R("warn")]},
    {"key": "ink",    "hdr": CC["navy_900"],
     "nodes": [R("accent_primary"), R("accent_secondary"), R("accent_point"), R("sub_header"), CC["navy_600"]]},
    {"key": "teal",   "hdr": R("accent_secondary"),
     "nodes": [R("accent_primary"), R("sub_header"), R("accent_point"), CC["navy_600"], CC["navy_800"]]},
    {"key": "steel",  "hdr": CC["navy_600"],
     "nodes": [R("accent_primary"), R("accent_secondary"), R("sub_header"), R("accent_point"), R("warn")]},
]
NS = len(SCHEMES)
def sc_at(i):    return SCHEMES[i % NS]
def ncol(sc, i): return sc["nodes"][i % len(sc["nodes"])]

# ─────────────────────────────────────────────────────────────
# 라벨 풀 (더미 제안서 텍스트) — 리터럴 중복 방지용 상수
# ─────────────────────────────────────────────────────────────
OWNER   = "발주기관\n(콘텐츠진흥원)"
PM      = "총괄PM"
TEAMS   = ["기획팀\n전략·기획", "개발팀\n플랫폼 구현", "현지화팀\n번역·문화화",
           "품질(QA)팀\n검수·안정화", "운영팀\n서비스 운영", "데이터팀\n분석·리포팅"]
FUNCS   = ["기획본부", "개발본부", "현지화본부", "품질본부", "운영본부"]
PROJS   = ["사업 A\n주력시장", "사업 B\n신흥시장", "사업 C\n전략시장", "사업 D\n협력시장"]
RRTITLE = ["기획·총괄", "기술지원", "현지화", "품질관리", "운영지원"]
RRDUTY  = [["사업 총괄 관리", "일정·리스크", "발주처 협의"],
           ["플랫폼 구축", "API 연동", "인프라 운영"],
           ["다국어 번역", "문화 현지화", "현지 검수"],
           ["QA 테스트", "산출물 검수", "안정화 지원"],
           ["운영 모니터링", "장애 대응", "성과 관리"]]
GOV3    = ["운영위원회\n(의사결정·정책)", "실무협의체\n(조정·검수·소통)"]
GOVEXEC = ["분석·기획", "개발·구축", "현지화·운영", "품질·검수"]
ACTORS  = ["발주처\n(발주·관리)", "수행사\n(총괄 수행)", "협력사\n(분야 협력)", "운영사\n(운영·이관)"]
EDGEPAIR = [("보고", "검수"), ("위탁", "납품"), ("협의", "조정")]
PARTNERS = ["글로벌 OTT", "현지 배급사", "제작 스튜디오", "마케팅 대행",
            "번역 전문", "법무·자문", "결제 PG사", "데이터 분석"]
MEMBERS  = [("주관사 ㈜가", "플랫폼 개발"), ("참여사 ㈜나", "현지화·운영"),
            ("참여사 ㈜다", "마케팅·배급"), ("참여사 ㈜라", "품질·검수")]
SHARES   = [["100%"], ["55%", "45%"], ["45%", "30%", "25%"], ["40%", "25%", "20%", "15%"]]
COMMS    = [("기술 분과위", ["시스템·인프라 심의", "보안·품질 기준", "기술 자문"]),
            ("콘텐츠 분과위", ["콘텐츠 기획 심의", "현지화 방향", "저작권 검토"]),
            ("운영 분과위", ["사업 운영 점검", "예산·일정 조정", "성과 평가"]),
            ("평가 분과위", ["성과 지표 심의", "품질 평가", "개선 권고"])]
TIERS    = [("1계층 · 핵심 파트너", ["글로벌 OTT", "메이저 배급사"]),
            ("2계층 · 전략 파트너", ["제작 스튜디오", "마케팅 대행", "번역 전문"]),
            ("3계층 · 협력 파트너", ["법무·자문", "결제 PG사", "데이터 분석"])]
ESCLV    = [("담당자 대응", "4시간 이내 · 현장 1차 조치"),
            ("팀장 대응", "4시간 초과 · 자원 재배치"),
            ("총괄PM 대응", "24시간 초과 · 범위·일정 조정"),
            ("운영위원회 대응", "중대 이슈 · 계약·정책 결정"),
            ("이사회 대응", "최상위 · 전략·투자 판단")]
RACI_ROLES = ["총괄PM", "기획팀", "개발팀", "현지화팀", "품질(QA)"]
RACI_TASKS = ["사업 착수·킥오프", "요구사항 정의", "플랫폼 설계·개발", "콘텐츠 현지화",
              "통합 QA 테스트", "발주처 보고·검수", "안정화·운영 이관", "성과 정산·종료"]
RACI_CODES = ["A", "R", "C", "I"]
RACI_MARK  = {"A": (R("warn"), True), "R": (R("accent_primary"), True),
              "C": (R("accent_secondary"), False), "I": (R("muted_text"), False)}
RACI_LEGEND = "R 실무 담당 · A 최종 책임 · C 자문 · I 결과 공유"

# ─────────────────────────────────────────────────────────────
# ID 카운터
# ─────────────────────────────────────────────────────────────
_idc = [15]
def NID():
    a = "ORG-%03d" % _idc[0]; _idc[0] += 1; return a

entries = []
def E(aid, name, file_rel, slide_idx, tags, params, bindings, editable, rec):
    entries.append(c.entry(aid, "ORG", name, file_rel, slide_idx,
                           tags, params, bindings, editable, recommended_use=rec))

# ─────────────────────────────────────────────────────────────
# 저수준 드로잉 헬퍼 (좌표 inch)
# ─────────────────────────────────────────────────────────────
def cxf(sp):  return sp.left / EMU + sp.width / EMU / 2
def cyf(sp):  return sp.top / EMU + sp.height / EMU / 2
def topf(sp): return sp.top / EMU
def botf(sp): return sp.top / EMU + sp.height / EMU
def rgtf(sp): return sp.left / EMU + sp.width / EMU
def lftf(sp): return sp.left / EMU

def node(s, x, y, w, h, text, fill, txt=None, size=12, bold=True, line=None,
         shape=RR, radius=0.10):
    sp = c.add_box(s, x, y, w, h, fill=fill, line=line, line_w=1.0, shape=shape)
    if shape == RR:
        sp.adjustments[0] = radius
    c.set_shape_text(sp, text, size=size, bold=bold,
                     color=txt or R("header_text"), align=PP_ALIGN.CENTER, font=c.FONT_H)
    return sp

def conn(s, x1, y1, x2, y2, color=None, w=1.5, dashed=False, head=False, tail=False):
    cn = c.connector(s, x1, y1, x2, y2, color=color or R("muted_text"), w=w)
    ln = cn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn

def tree(s, parent, children, color=None, bus_gap=0.42, w=1.6):
    col = color or R("muted_text")
    out = []
    py = botf(parent); busy = py + bus_gap; pcx = cxf(parent)
    xs = [cxf(ch) for ch in children]
    out.append(conn(s, pcx, py, pcx, busy, color=col, w=w))
    out.append(conn(s, min(xs + [pcx]), busy, max(xs + [pcx]), busy, color=col, w=w))
    for ch in children:
        out.append(conn(s, cxf(ch), busy, cxf(ch), topf(ch), color=col, w=w))
    return out

def hslide(prs, aid, title, sub):
    s = c.blank_slide(prs)
    c.id_caption(s, aid + " · " + title)
    c.add_text(s, 0.55, 0.5, 12.2, 0.5, title, size=17, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    c.add_text(s, 0.55, 1.0, 12.2, 0.32, sub, size=11,
               color=R("muted_text"), align=PP_ALIGN.LEFT)
    return s

def row_x(n, total_w=12.4, x_left=0.467, min_w=1.7, max_w=2.7, gap=0.3):
    tw = (total_w - gap * (n - 1)) / n
    tw = max(min_w, min(max_w, tw))
    span = tw * n + gap * (n - 1)
    x0 = (13.333 - span) / 2
    return tw, [x0 + i * (tw + gap) for i in range(n)]

# ── 표 헬퍼 ──
def add_table(s, x, y, w, h, nrows, ncols, cw=None, rh=None):
    gf = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    if cw:
        for j, v in enumerate(cw): tbl.columns[j].width = Inches(v)
    if rh:
        for i, v in enumerate(rh): tbl.rows[i].height = Inches(v)
    return gf, tbl

def cell(cl, text, size=11, bold=False, color=None, fill=None,
         align=PP_ALIGN.CENTER, font=None):
    cl.fill.solid()
    cl.fill.fore_color.rgb = fill if fill is not None else R("row_base")
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_left = cl.margin_right = Pt(5)
    cl.margin_top = cl.margin_bottom = Pt(3)
    tf = cl.text_frame; tf.word_wrap = True
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))

# ═════════════════════════════════════════════════════════════
# 패밀리 빌더 — 각자 prs 생성 → count 슬라이드 → 파일 저장
# ═════════════════════════════════════════════════════════════
def save(prs, fname):
    rel = DIR + fname
    c.save_deck(prs, rel)
    return rel

# 1) 위계 트리 (발주-PM-팀 n)
def fam_hierarchy(count, lo=2, hi=5):
    rel = DIR + "ORG_bulk_hierarchy_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "위계 조직도 · 팀 %d" % n,
                   "발주기관 → 총괄PM → 실행팀 %d (%s)" % (n, sc["key"]))
        own = node(s, 5.17, 1.35, 3.0, 0.72, OWNER, sc["hdr"], size=12.5)
        pm  = node(s, 5.17, 2.85, 3.0, 0.66, PM, ncol(sc, 0), size=14)
        tw, xs = row_x(n)
        teams = [node(s, xs[k], 4.55, tw, 1.0, TEAMS[k % len(TEAMS)], ncol(sc, k + 1), size=11.5)
                 for k in range(n)]
        cn = [conn(s, cxf(own), botf(own), cxf(pm), topf(pm), color=R("muted_text"), w=1.75)]
        cn += tree(s, pm, teams)
        c.group_asset(s, cn + [own, pm] + teams, aid)
        labels = [TEAMS[k % len(TEAMS)].split("\n")[0] for k in range(n)]
        E(aid, "위계 조직도 (발주-PM-팀%d)" % n, rel, i + 1,
          ["조직도", "위계", "트리", "수행조직", sc["key"]],
          {"levels": 3, "teams": n, "orient": "vertical", "scheme": sc["key"]},
          {"nodes": [{"role": "발주기관", "level": 1}, {"role": "총괄PM", "level": 2}] +
                    [{"role": t, "level": 3} for t in labels],
           "edges": [["발주기관", "총괄PM"]] + [["총괄PM", t] for t in labels]},
          ["node-text", "add-node", "color", "connector"],
          ["수행조직", "위계 조직", "역할분담"])
    return save(prs, "ORG_bulk_hierarchy_v1.pptx"), count

# 2) 매트릭스 조직 (기능 f × 사업 p)
def fam_matrix(count):
    rel = DIR + "ORG_bulk_matrix_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        f = 3 + (i % 2); p = 2 + (i % 3); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "매트릭스 조직 · %d×%d" % (f, p),
                   "기능축 %d × 사업축 %d 이중 보고 체계 (%s)" % (f, p, sc["key"]))
        box = []; cn = []; dot = []
        pm = node(s, 5.4, 1.35, 2.6, 0.66, "총괄 PMO", sc["hdr"], size=13); box.append(pm)
        pw, pxs = row_x(p, total_w=8.0, x_left=3.3, min_w=1.7, max_w=2.2)
        # 사업 헤더 좌표는 노드 중심 x
        pcx = [pxs[k] + pw / 2 for k in range(p)]
        pboxes = [node(s, pxs[k], 2.35, pw, 0.6, PROJS[k % len(PROJS)], ncol(sc, k), size=10.5)
                  for k in range(p)]
        box += pboxes
        fys = [3.35 + k * 0.82 for k in range(f)]
        fboxes = [node(s, 0.55, fys[k] - 0.26, 2.0, 0.52, FUNCS[k % len(FUNCS)], ncol(sc, k + 1), size=11)
                  for k in range(f)]
        box += fboxes
        cn += tree(s, pm, pboxes, bus_gap=0.28, w=1.5)
        spine = 2.55
        cn.append(conn(s, lftf(pm), cyf(pm), spine, cyf(pm), color=R("muted_text"), w=1.5))
        cn.append(conn(s, spine, cyf(pm), spine, fys[-1], color=R("muted_text"), w=1.5))
        for k in range(p):
            cn.append(conn(s, pcx[k], 2.95, pcx[k], fys[-1] + 0.35, color=R("border"), w=1.2, dashed=True))
        for k in range(f):
            cn.append(conn(s, spine, fys[k], max(pcx) + 0.4, fys[k], color=R("border"), w=1.2, dashed=True))
        for k in range(p):
            for m in range(f):
                d = c.add_box(s, pcx[k] - 0.15, fys[m] - 0.15, 0.30, 0.30,
                              fill=R("accent_primary"), line=CC["white"], line_w=1.0, shape=OVAL)
                dot.append(d)
        c.group_asset(s, cn + box + dot, aid)
        E(aid, "매트릭스 조직 (기능%d×사업%d)" % (f, p), rel, i + 1,
          ["조직도", "매트릭스", "이중보고", "기능별", sc["key"]],
          {"functions": f, "projects": p, "orient": "matrix-grid", "scheme": sc["key"]},
          {"pmo": "총괄 PMO",
           "functions": [FUNCS[k % len(FUNCS)] for k in range(f)],
           "projects": [PROJS[k % len(PROJS)].split("\n")[0] for k in range(p)],
           "cells": f * p},
          ["node-text", "add-function", "add-project", "color", "connector"],
          ["매트릭스 조직", "이중 보고 체계", "수행조직"])
    return save(prs, "ORG_bulk_matrix_v1.pptx"), count

# 3) R&R 카드 (총괄 + 역할카드 n)
def fam_rrcard(count, lo=2, hi=5):
    rel = DIR + "ORG_bulk_rrcard_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "R&R 역할·책임 카드 · %d" % n,
                   "총괄PM + 역할별 책임(R&R) 카드 %d (%s)" % (n, sc["key"]))
        head = node(s, 5.17, 1.35, 3.0, 0.64, PM, sc["hdr"], size=13)
        cw, xs = row_x(n, min_w=2.1, max_w=2.7)
        shp = [head]; cn = []
        cards = []
        for k in range(n):
            ac = ncol(sc, k)
            card = c.add_box(s, xs[k], 2.9, cw, 1.95, fill=R("panel_bg"),
                             line=R("border"), line_w=0.8, shape=RR); card.adjustments[0] = 0.05
            band = c.add_box(s, xs[k], 2.9, cw, 0.44, fill=ac, line=None, shape=RECT)
            c.set_shape_text(band, RRTITLE[k % len(RRTITLE)], size=12, bold=True,
                             color=R("header_text"), align=PP_ALIGN.CENTER, font=c.FONT_H)
            duties = RRDUTY[k % len(RRDUTY)]
            body = c.add_box(s, xs[k] + 0.12, 3.44, cw - 0.24, 1.35, fill=None, line=None, shape=RECT)
            c.set_shape_text(body, "\n".join("· " + d for d in duties), size=10.5, bold=False,
                             color=R("body_text"), align=PP_ALIGN.LEFT, font=c.FONT_B)
            body.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            for pp in body.text_frame.paragraphs: pp.alignment = PP_ALIGN.LEFT
            shp += [card, band, body]; cards.append(card)
        busy = 2.55
        cn.append(conn(s, cxf(head), botf(head), cxf(head), busy, color=R("muted_text"), w=1.5))
        ccx = [cxf(cd) for cd in cards]
        cn.append(conn(s, min(ccx), busy, max(ccx), busy, color=R("muted_text"), w=1.5))
        for cd in cards:
            cn.append(conn(s, cxf(cd), busy, cxf(cd), topf(cd), color=R("muted_text"), w=1.5))
        c.group_asset(s, cn + shp, aid)
        E(aid, "R&R 역할·책임 카드 (총괄+%d)" % n, rel, i + 1,
          ["조직도", "R&R", "역할책임", "역할카드", sc["key"]],
          {"levels": 2, "cards": n, "style": "role-card", "scheme": sc["key"]},
          {"nodes": [{"role": "총괄PM", "level": 1}] +
                    [{"role": RRTITLE[k % len(RRTITLE)], "duties": RRDUTY[k % len(RRDUTY)], "level": 2}
                     for k in range(n)],
           "edges": [["총괄PM", RRTITLE[k % len(RRTITLE)]] for k in range(n)]},
          ["node-text", "add-card", "duty-text", "color", "connector"],
          ["역할분담", "R&R 정의", "책임 명세"])
    return save(prs, "ORG_bulk_rrcard_v1.pptx"), count

# 4) 거버넌스 3층 (운영위-실무-실행 n)
def fam_governance(count, lo=2, hi=4):
    rel = DIR + "ORG_bulk_governance_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "거버넌스 체계 3층 · 실행 %d" % n,
                   "운영위원회 → 실무협의체 → 실행조직 %d (%s)" % (n, sc["key"]))
        gov = node(s, 4.17, 1.35, 5.0, 0.74, GOV3[0], sc["hdr"], size=13)
        ops = node(s, 4.17, 3.05, 5.0, 0.74, GOV3[1], ncol(sc, 0), size=13)
        ew, xs = row_x(n, min_w=2.2, max_w=3.0)
        ex = [node(s, xs[k], 4.85, ew, 0.9, "실행조직\n" + GOVEXEC[k % len(GOVEXEC)], ncol(sc, k + 1), size=12)
              for k in range(n)]
        cn = [conn(s, cxf(gov), botf(gov), cxf(ops), topf(ops), color=R("muted_text"), w=1.75)]
        cn += tree(s, ops, ex)
        c.group_asset(s, cn + [gov, ops] + ex, aid)
        E(aid, "거버넌스 체계 3층 (실행%d)" % n, rel, i + 1,
          ["조직도", "거버넌스", "체계도", "3층", sc["key"]],
          {"levels": 3, "exec": n, "tiers": ["운영위", "실무협의체", "실행조직"], "scheme": sc["key"]},
          {"nodes": [{"role": "운영위원회", "level": 1}, {"role": "실무협의체", "level": 2}] +
                    [{"role": "실행조직-" + GOVEXEC[k % len(GOVEXEC)], "level": 3} for k in range(n)],
           "edges": [["운영위원회", "실무협의체"]] +
                    [["실무협의체", "실행조직-" + GOVEXEC[k % len(GOVEXEC)]] for k in range(n)]},
          ["node-text", "add-node", "color", "connector"],
          ["거버넌스", "추진체계", "운영 구조"])
    return save(prs, "ORG_bulk_governance_v1.pptx"), count

# 5) 추진체계 n주체 양방향
def fam_pushsystem(count, lo=3, hi=4):
    rel = DIR + "ORG_bulk_pushsystem_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "추진체계도 · %d주체 양방향" % n,
                   "주체 %d 간 보고·검수·위탁·납품 양방향 흐름 (%s)" % (n, sc["key"]))
        aw = 2.5 if n >= 4 else 3.0
        span = 12.4; gap = (span - aw * n) / (n - 1)
        xs = [0.467 + k * (aw + gap) for k in range(n)]
        ay = 3.2; ah = 1.2
        boxes = [node(s, xs[k], ay, aw, ah, ACTORS[k % len(ACTORS)],
                      sc["hdr"] if k == 0 else ncol(sc, k), size=12.5) for k in range(n)]
        cn = []; txts = []
        for k in range(n - 1):
            a, b = boxes[k], boxes[k + 1]
            cn.append(conn(s, rgtf(a), cyf(a), lftf(b), cyf(b),
                           color=R("muted_text"), w=1.75, head=True, tail=True))
            top_l, bot_l = EDGEPAIR[k % len(EDGEPAIR)]
            t1 = c.add_text(s, rgtf(a), cyf(a) - 0.5, lftf(b) - rgtf(a), 0.32, top_l,
                            size=10.5, bold=True, color=R("accent_primary"), align=PP_ALIGN.CENTER)
            t2 = c.add_text(s, rgtf(a), cyf(a) + 0.18, lftf(b) - rgtf(a), 0.32, bot_l,
                            size=10.5, bold=True, color=R("warn"), align=PP_ALIGN.CENTER)
            txts += [t1, t2]
        c.group_asset(s, cn + boxes + txts, aid)
        acts = [ACTORS[k % len(ACTORS)].split("\n")[0] for k in range(n)]
        E(aid, "추진체계도 (%d주체 양방향)" % n, rel, i + 1,
          ["추진체계", "체계도", "양방향", "%d주체" % n, sc["key"]],
          {"actors": n, "orient": "horizontal", "bidirectional": True, "scheme": sc["key"]},
          {"nodes": [{"role": a} for a in acts],
           "edges": [[acts[k], acts[k + 1]] for k in range(n - 1)]},
          ["node-text", "add-node", "edge-label", "color", "connector"],
          ["추진체계", "역할분담", "협업 구조"])
    return save(prs, "ORG_bulk_pushsystem_v1.pptx"), count

# 6) 자문위 포함형 (위계 n팀 + 자문 점선)
def fam_advisory(count, lo=2, hi=4):
    rel = DIR + "ORG_bulk_advisory_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "자문위 포함 조직도 · 팀 %d" % n,
                   "발주-PM-팀%d 본조직 + 자문위원회 점선 연결 (%s)" % (n, sc["key"]))
        own = node(s, 2.7, 1.35, 3.0, 0.68, OWNER, sc["hdr"], size=12)
        pm  = node(s, 2.7, 2.85, 3.0, 0.64, PM, ncol(sc, 0), size=13.5)
        tw, xs = row_x(n, total_w=7.4, x_left=0.5, min_w=1.9, max_w=2.5)
        teams = [node(s, xs[k], 4.55, tw, 0.95, TEAMS[k % len(TEAMS)], ncol(sc, k + 1), size=11)
                 for k in range(n)]
        adv = node(s, 9.35, 2.55, 3.15, 1.05, "자문위원회\n(전문가 자문·검토)", R("warn"), size=12)
        cn = [conn(s, cxf(own), botf(own), cxf(pm), topf(pm), color=R("muted_text"), w=1.75)]
        cn += tree(s, pm, teams)
        cn.append(conn(s, rgtf(pm), cyf(pm), lftf(adv), cyf(adv), color=R("warn"), w=1.75, dashed=True))
        adv_t = c.add_text(s, rgtf(pm) + 0.05, cyf(pm) - 0.4, lftf(adv) - rgtf(pm) - 0.1, 0.3,
                           "자문", size=10.5, bold=True, color=R("warn"), align=PP_ALIGN.CENTER)
        c.group_asset(s, cn + [own, pm] + teams + [adv, adv_t], aid)
        labels = [TEAMS[k % len(TEAMS)].split("\n")[0] for k in range(n)]
        E(aid, "자문위 포함 조직도 (팀%d+자문)" % n, rel, i + 1,
          ["조직도", "자문위원", "위계", "점선연결", sc["key"]],
          {"levels": 3, "teams": n, "advisory": True, "advisory_link": "dashed", "scheme": sc["key"]},
          {"nodes": [{"role": "발주기관", "level": 1}, {"role": "총괄PM", "level": 2}] +
                    [{"role": t, "level": 3} for t in labels] +
                    [{"role": "자문위원회", "level": 2, "type": "advisory"}],
           "edges": [["발주기관", "총괄PM"]] + [["총괄PM", t] for t in labels] +
                    [["총괄PM", "자문위원회"]]},
          ["node-text", "add-node", "color", "connector"],
          ["수행조직", "자문 구조", "거버넌스"])
    return save(prs, "ORG_bulk_advisory_v1.pptx"), count

# 7) 협업 네트워크 방사 (허브 + 파트너 n)
def fam_network(count, lo=4, hi=8):
    rel = DIR + "ORG_bulk_network_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "협업 네트워크 · 파트너 %d" % n,
                   "중앙 허브 + 방사형 파트너 %d 협력 네트워크 (%s)" % (n, sc["key"]))
        hcx, hcy, hr = 6.667, 4.35, 0.98
        cn = []; pboxes = []
        pw, ph = 2.05, 0.82
        for k in range(n):
            ang = math.radians(-90 + k * (360.0 / n))
            pxc = hcx + hr * 3.15 * math.cos(ang) * 1.32
            pyc = hcy + hr * 3.15 * math.sin(ang) * 0.74
            cn.append(conn(s, hcx, hcy, pxc, pyc, color=R("border"), w=1.4))
            b = node(s, pxc - pw / 2, pyc - ph / 2, pw, ph, PARTNERS[k % len(PARTNERS)],
                     ncol(sc, k), size=11)
            pboxes.append(b)
        hub = node(s, hcx - hr, hcy - hr * 0.72, hr * 2, hr * 1.44, "통합\n플랫폼",
                   sc["hdr"], size=13, shape=OVAL)
        c.group_asset(s, cn + pboxes + [hub], aid)
        labels = [PARTNERS[k % len(PARTNERS)] for k in range(n)]
        E(aid, "협업 네트워크 (허브+파트너%d)" % n, rel, i + 1,
          ["조직도", "네트워크", "허브", "방사형", sc["key"]],
          {"hub": 1, "partners": n, "orient": "radial", "scheme": sc["key"]},
          {"nodes": [{"role": "통합 플랫폼", "type": "hub"}] +
                    [{"role": p, "type": "partner"} for p in labels],
           "edges": [["통합 플랫폼", p] for p in labels]},
          ["node-text", "add-partner", "color", "connector"],
          ["협업 구조", "파트너 네트워크", "생태계"])
    return save(prs, "ORG_bulk_network_v1.pptx"), count

# 8) 컨소시엄 구성 (주관 + 참여 n)
def fam_consortium(count, lo=2, hi=4):
    rel = DIR + "ORG_bulk_consortium_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "컨소시엄 구성도 · 참여 %d" % n,
                   "주관사 + 참여사 %d 역할·분담율 구성 (%s)" % (n, sc["key"]))
        lead = node(s, 4.17, 1.35, 5.0, 0.9, "주관사\n(사업 총괄·계약 주체)", sc["hdr"], size=13)
        cw, xs = row_x(n, min_w=2.4, max_w=3.1)
        shares = SHARES[n - 1] if n - 1 < len(SHARES) else SHARES[-1]
        cards = []; shp = [lead]
        for k in range(n):
            nm, duty = MEMBERS[(k + 1) % len(MEMBERS)] if n > 1 else MEMBERS[0]
            ac = ncol(sc, k)
            card = c.add_box(s, xs[k], 3.55, cw, 1.6, fill=R("panel_bg"),
                             line=R("border"), line_w=0.8, shape=RR); card.adjustments[0] = 0.06
            band = c.add_box(s, xs[k], 3.55, cw, 0.5, fill=ac, line=None, shape=RECT)
            c.set_shape_text(band, nm, size=12, bold=True, color=R("header_text"),
                             align=PP_ALIGN.CENTER, font=c.FONT_H)
            info = c.add_box(s, xs[k], 4.1, cw, 1.0, fill=None, line=None, shape=RECT)
            c.set_shape_text(info, duty + "\n분담 " + shares[k % len(shares)], size=11.5, bold=True,
                             color=R("body_text"), align=PP_ALIGN.CENTER, font=c.FONT_H)
            shp += [card, band, info]; cards.append(card)
        busy = 3.2
        cn = [conn(s, cxf(lead), botf(lead), cxf(lead), busy, color=R("muted_text"), w=1.75)]
        ccx = [cxf(cd) for cd in cards]
        cn.append(conn(s, min(ccx), busy, max(ccx), busy, color=R("muted_text"), w=1.75))
        for cd in cards:
            cn.append(conn(s, cxf(cd), busy, cxf(cd), topf(cd), color=R("muted_text"), w=1.75))
        c.group_asset(s, cn + shp, aid)
        mem = [MEMBERS[(k + 1) % len(MEMBERS)][0] for k in range(n)]
        E(aid, "컨소시엄 구성도 (주관+참여%d)" % n, rel, i + 1,
          ["조직도", "컨소시엄", "분담", "참여사", sc["key"]],
          {"levels": 2, "members": n, "show_share": True, "scheme": sc["key"]},
          {"nodes": [{"role": "주관사", "level": 1}] +
                    [{"role": mem[k], "share": shares[k % len(shares)], "level": 2} for k in range(n)],
           "edges": [["주관사", m] for m in mem]},
          ["node-text", "add-member", "share", "color", "connector"],
          ["컨소시엄", "역할분담", "수행조직"])
    return save(prs, "ORG_bulk_consortium_v1.pptx"), count

# 9) RACI 표 (역할 r × 업무 t)
def fam_raci(count):
    rel = DIR + "ORG_bulk_raci_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        r = 3 + (i % 3); t = 5 + (i % 4); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "RACI 책임 배분표 · %d업무×%d역할" % (t, r),
                   "핵심 업무 %d × 역할 %d 책임(R/A/C/I) 정의 (%s)" % (t, r, sc["key"]))
        roles = RACI_ROLES[:r]; tasks = RACI_TASKS[:t]
        ncols = 1 + r; nrows = 1 + t + 1
        cw = [2.6] + [min(2.0, 9.2 / r)] * r
        rh = [0.5] * nrows
        gf, tbl = add_table(s, 0.7, 1.6, sum(cw), min(sum(rh), 5.4), nrows, ncols, cw, rh)
        c.name_asset(gf, aid)
        cell(tbl.cell(0, 0), "핵심 업무", size=12, bold=True, color=R("header_text"),
             fill=sc["hdr"], align=PP_ALIGN.LEFT)
        for j, rn in enumerate(roles):
            cell(tbl.cell(0, j + 1), rn, size=11, bold=True, color=R("header_text"), fill=sc["hdr"])
        for a in range(t):
            base = R("row_stripe") if a % 2 else R("row_base")
            cell(tbl.cell(a + 1, 0), tasks[a], size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
            for j in range(r):
                code = "A" if j == 0 else RACI_CODES[(a + j) % len(RACI_CODES)]
                col, bd = RACI_MARK[code]
                cell(tbl.cell(a + 1, j + 1), code, size=13, bold=bd, color=col, fill=base)
        li = nrows - 1
        tbl.cell(li, 0).merge(tbl.cell(li, ncols - 1))
        cell(tbl.cell(li, 0), RACI_LEGEND, size=10, bold=True, color=R("muted_text"), fill=R("panel_bg"))
        E(aid, "RACI 책임 배분표 (%d업무×%d역할)" % (t, r), rel, i + 1,
          ["RACI", "책임배분", "역할분담", "매트릭스", sc["key"]],
          {"tasks": t, "roles": r, "codes": RACI_CODES, "scheme": sc["key"]},
          {"roles": roles,
           "activities": [{"task": tasks[a],
                           "assign": {roles[j]: ("A" if j == 0 else RACI_CODES[(a + j) % len(RACI_CODES)])
                                      for j in range(r)}} for a in range(t)],
           "legend": {"R": "실무 담당", "A": "최종 책임", "C": "자문", "I": "결과 공유"}},
          ["cell-text", "add-row", "add-col", "marker", "color"],
          ["RACI", "책임배분(R&R)", "역할 책임 정의"])
    return save(prs, "ORG_bulk_raci_v1.pptx"), count

# 10) 위원회 구조 (위원장 + 분과위 n)
def fam_committee(count, lo=2, hi=4):
    rel = DIR + "ORG_bulk_committee_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "추진위원회 구조 · 분과 %d" % n,
                   "위원장 · 사무국(간사) · 분과위원회 %d 운영 (%s)" % (n, sc["key"]))
        chair = node(s, 5.4, 1.35, 2.55, 0.7, "추진위원장", sc["hdr"], size=14)
        sec = node(s, 9.9, 1.4, 2.1, 0.6, "사무국(간사)", R("muted_text"), size=11.5)
        cw, xs = row_x(n, min_w=2.6, max_w=3.4)
        subs = []; panels = []; shp = [chair, sec]
        for k in range(n):
            nm, duties = COMMS[k % len(COMMS)]
            b = node(s, xs[k], 3.0, cw, 0.66, nm, ncol(sc, k), size=12.5)
            panel = c.add_box(s, xs[k], 3.85, cw, 1.8, fill=R("panel_bg"),
                              line=R("border"), line_w=0.8, shape=RR); panel.adjustments[0] = 0.05
            c.set_shape_text(panel, "\n".join("· " + d for d in duties), size=10.5, bold=False,
                             color=R("body_text"), align=PP_ALIGN.LEFT, font=c.FONT_B)
            panel.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            for pp in panel.text_frame.paragraphs: pp.alignment = PP_ALIGN.LEFT
            subs.append(b); panels.append(panel); shp += [b, panel]
        cn = [conn(s, rgtf(chair), cyf(chair), lftf(sec), cyf(sec), color=R("muted_text"), w=1.4, dashed=True)]
        cn += tree(s, chair, subs, bus_gap=0.5)
        for b, panel in zip(subs, panels):
            cn.append(conn(s, cxf(b), botf(b), cxf(b), topf(panel), color=R("border"), w=1.2))
        c.group_asset(s, cn + shp, aid)
        E(aid, "위원회 구조 (위원장-분과위%d)" % n, rel, i + 1,
          ["위원회", "거버넌스", "분과위", "심의", sc["key"]],
          {"chair": 1, "secretariat": 1, "subcommittees": n, "scheme": sc["key"]},
          {"chair": "추진위원장", "secretariat": "사무국(간사)",
           "subcommittees": [{"name": COMMS[k % len(COMMS)][0], "duties": COMMS[k % len(COMMS)][1]}
                             for k in range(n)]},
          ["node-text", "add-subcommittee", "color", "connector"],
          ["위원회 구조", "거버넌스", "심의 체계"])
    return save(prs, "ORG_bulk_committee_v1.pptx"), count

# 11) 파트너 생태계 (허브 + 계층 n)
def fam_ecosystem(count, lo=2, hi=3):
    rel = DIR + "ORG_bulk_ecosystem_v1.pptx"; prs = c.new_deck()
    layout = [(4.85, 0.95, 3.6, 1.5), (0.6, 5.05, 4.15, 1.85), (8.6, 5.05, 4.15, 1.85)]
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "파트너 생태계 · 계층 %d" % n,
                   "중앙 플랫폼 + 계층별 파트너군 %d 협력 구조 (%s)" % (n, sc["key"]))
        cn = []; shp = []
        hub = node(s, 5.47, 3.15, 2.4, 1.35, "통합\n플랫폼", sc["hdr"], size=14, shape=OVAL)
        centers = []
        for k in range(n):
            header, chips = TIERS[k % len(TIERS)]
            px, py, pw, ph = layout[k % len(layout)]
            band_c = ncol(sc, k)
            panel = c.add_box(s, px, py, pw, ph, fill=R("panel_bg"), line=R("border"), line_w=0.9, shape=RR)
            panel.adjustments[0] = 0.06
            band = c.add_box(s, px, py, pw, 0.42, fill=band_c, line=None, shape=RECT)
            c.set_shape_text(band, header, size=11, bold=True, color=R("header_text"),
                             align=PP_ALIGN.CENTER, font=c.FONT_H)
            shp += [panel, band]
            m = len(chips); chw = (pw - 0.3 - 0.15 * (m - 1)) / m
            for j, ch in enumerate(chips):
                chx = px + 0.15 + j * (chw + 0.15)
                chip = node(s, chx, py + 0.56, chw, ph - 0.72, ch, CC["white"],
                            txt=R("body_text"), size=10, bold=True, line=band_c, radius=0.12)
                shp.append(chip)
            centers.append((px + pw / 2, py + ph / 2))
        for pcx, pcy in centers:
            cn.append(conn(s, cxf(hub), cyf(hub), pcx, pcy, color=R("accent_primary"), w=1.6))
        c.group_asset(s, cn + shp + [hub], aid)
        E(aid, "파트너 생태계 (허브+계층%d)" % n, rel, i + 1,
          ["생태계", "파트너", "네트워크", "계층", sc["key"]],
          {"hub": 1, "tiers": n, "orient": "ecosystem", "scheme": sc["key"]},
          {"hub": "통합 플랫폼",
           "tiers": [{"tier": TIERS[k % len(TIERS)][0], "partners": TIERS[k % len(TIERS)][1]}
                     for k in range(n)]},
          ["node-text", "add-partner", "add-tier", "color", "connector"],
          ["파트너 생태계", "협업 네트워크", "협력 구조"])
    return save(prs, "ORG_bulk_ecosystem_v1.pptx"), count

# 12) 에스컬레이션 체계 (단계 n 상향)
def fam_escalation(count, lo=3, hi=5):
    rel = DIR + "ORG_bulk_escalation_v1.pptx"; prs = c.new_deck()
    for i in range(count):
        n = lo + (i % (hi - lo + 1)); sc = sc_at(i); aid = NID()
        s = hslide(prs, aid, "이슈 에스컬레이션 체계 · %d단계" % n,
                   "미해결 이슈의 단계별 대응 주체 상향 %d (%s)" % (n, sc["key"]))
        lb_x, lb_w, lb_h = 3.3, 5.3, 0.72
        y_bottom, y_gap = 6.05, (min(3.6, 0.95 * (n - 1)))
        step = y_gap / max(1, n - 1) if n > 1 else 0
        # 아래(1단계)→위(n단계) 배치
        ys = [y_bottom - k * (lb_h + max(0.32, 0.42)) for k in range(n)]
        boxes = []; anns = []; cn = []
        for k in range(n):
            lab, sla = ESCLV[k % len(ESCLV)]
            fc = R("warn") if k == n - 1 else ncol(sc, k)
            b = node(s, lb_x, ys[k], lb_w, lb_h, "%d단계 · %s" % (k + 1, lab), fc, size=12.5)
            ann = c.add_box(s, lb_x + lb_w + 0.35, ys[k] + 0.05, 3.55, lb_h - 0.1,
                            fill=R("panel_bg"), line=R("border"), line_w=0.8, shape=RR)
            ann.adjustments[0] = 0.08
            c.set_shape_text(ann, sla, size=10, bold=False, color=R("muted_text"),
                             align=PP_ALIGN.LEFT, font=c.FONT_B)
            boxes.append(b); anns.append(ann)
        for lo_b, up_b in zip(boxes[:-1], boxes[1:]):
            cn.append(conn(s, cxf(lo_b), topf(lo_b), cxf(up_b), botf(up_b),
                           color=R("warn"), w=2.2, head=True))
        big = c.add_box(s, 1.1, min(ys), 1.35, (max(ys) + lb_h) - min(ys), fill=sc["hdr"],
                        line=None, shape=UPARROW)
        c.set_shape_text(big, "대응\n상향", size=12, bold=True, color=R("header_text"),
                         align=PP_ALIGN.CENTER, font=c.FONT_H)
        c.group_asset(s, cn + boxes + anns + [big], aid)
        E(aid, "에스컬레이션 체계 (%d단계 상향)" % n, rel, i + 1,
          ["에스컬레이션", "이슈대응", "단계", "상향", sc["key"]],
          {"levels": n, "direction": "upward", "scheme": sc["key"]},
          {"levels": [{"step": k + 1, "owner": ESCLV[k % len(ESCLV)][0], "sla": ESCLV[k % len(ESCLV)][1]}
                      for k in range(n)]},
          ["node-text", "add-level", "sla-text", "color", "connector"],
          ["에스컬레이션", "이슈 대응 체계", "리스크 관리"])
    return save(prs, "ORG_bulk_escalation_v1.pptx"), count

# ═════════════════════════════════════════════════════════════
# 실행 — 총 186 (16,15,16,15,15,15,16,15,15,16,16,16)
# ═════════════════════════════════════════════════════════════
PLAN = [
    (fam_hierarchy,   16),
    (fam_matrix,      15),
    (fam_rrcard,      16),
    (fam_governance,  15),
    (fam_pushsystem,  15),
    (fam_advisory,    15),
    (fam_network,     16),
    (fam_consortium,  15),
    (fam_raci,        15),
    (fam_committee,   16),
    (fam_ecosystem,   16),
    (fam_escalation,  16),
]

if __name__ == "__main__":
    files = []
    for fn, cnt in PLAN:
        rel, made = fn(cnt)
        files.append((rel, made))
        print("SAVED:", rel, "->", made, "slides")
    frag = c.write_fragment("ORG_bulk", entries)
    print("FRAGMENT:", frag)
    print("ENTRIES:", len(entries), "range:", entries[0]["id"], "..", entries[-1]["id"])
    assert len(entries) == 186, "expected 186, got %d" % len(entries)
    assert entries[-1]["id"] == "ORG-200", "last id %s" % entries[-1]["id"]
