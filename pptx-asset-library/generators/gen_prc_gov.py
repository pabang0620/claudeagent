# -*- coding: utf-8 -*-
"""PRC gov 트랙 — 정부보고서(정답지) 캔버스·팔레트 전용 프로세스 흐름도 6종.

정답지 재현 목표: GROUP 조합 수작업 프로세스 흐름도 스타일(9장 커버 대표 세트).
네이티브 SmartArt/차트 절대 금지 — 전부 도형(autoshape)+커넥터 GROUP 조합.

캔버스: gov 전용 11.93 x 8.50in = EMU 10905360 x 7772400 (design-tokens.json 최상위와
독립된 gov 트랙 캔버스 — new_deck_gov()가 정확한 EMU 값으로 슬라이드 크기를 지정).
색·폰트: design-tokens.json["gov_theme"] 참조만(매직 헥스 금지). common.set_kfont로
a:latin/a:ea/a:cs 3계열 동시 설정(한글 안전).

ID: PRC-201 ~ PRC-206 (기존 PRC-001~200 다음 free 대역, manifest.json 확인 완료).
파일: decks/03_process/PRC_gov_flow_v1.pptx (슬라이드 1~6 = 에셋 1~6개, 1슬라이드=1에셋)
매니페스트 조각: _incoming/manifest_PRC_gov.json (기존 manifest_PRC*.json 미변경)
"""
import sys, math
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DECK = 'decks/03_process'
FILE_REL = DECK + '/PRC_gov_flow_v1.pptx'

# ---------------------------------------------------------------------------
# gov 캔버스 (design-tokens.json 최상위 slide.width_in/height_in과 별개, additive)
# ---------------------------------------------------------------------------
GOV_W = Emu(10905360)
GOV_H = Emu(7772400)
SW = 11.93  # in — 레이아웃 산술용 근사치(실제 슬라이드 경계는 GOV_W/GOV_H EMU가 정확히 지정)
SH = 8.50


def new_deck_gov():
    prs = Presentation()
    prs.slide_width = GOV_W
    prs.slide_height = GOV_H
    return prs


def blank_slide_gov(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# gov 팔레트 (design-tokens.json["gov_theme"] 참조만, 매직 헥스 금지)
# ---------------------------------------------------------------------------
GOV = c.TOKENS['gov_theme']
GC = {k: RGBColor.from_string(v) for k, v in GOV['color'].items()}


def R(name):
    return GC[GOV['role'][name]]


WHITE = GC['bg_primary']  # FFFFFF — gov 팔레트엔 별도 white 키가 없어 bg_primary로 대체
F_BODY = GOV['font']['body']['typeface']
F_LABEL = GOV['font']['label']['typeface']
F_SUB = GOV['font']['subheading']['typeface']
F_CAP = GOV['font']['caption']['typeface']

RADIUS_PCT = GOV['geom']['round_radius_pct']
BORDER_PT = GOV['geom']['border_pt']
LINE_ACCENT_PT = GOV['geom']['line_accent_pt']

# 단계 차등 램프 — gov_theme 내 "도형용" 색만 사용(표헤더 전용색 8F99AF/76829E는 배제,
# provenance 근거: 그 두 색은 정답지에서 100% <a:tbl> 내부에서만 관찰됨)
RAMP_GOV = [GC['accent_navy'], GC['shape_fill_gray'], GC['category_purple'], GC['line_accent']]


def tone(i):
    return RAMP_GOV[i % len(RAMP_GOV)]


def gtext(slide, x, y, w, h, text, size=None, bold=False, color=None,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=None):
    return c.add_text(slide, x, y, w, h, text, size=size or GOV['size_pt']['body'],
                       bold=bold, color=color or R('body_text'), align=align,
                       anchor=anchor, font=font or F_BODY)


def gshape_text(sp, text, size=None, bold=True, color=None, align=PP_ALIGN.CENTER, font=None):
    return c.set_shape_text(sp, text, size=size or GOV['size_pt']['body'], bold=bold,
                             color=color or WHITE, align=align, font=font or F_LABEL)


def arrowize(cn, head=False, tail=True):
    ln = cn.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def gcaption(slide, asset_id):
    """id_caption과 동일 역할이나 gov 캔버스 좌표계 사용 — 그룹 밖에 두어야 함(원칙 1)."""
    c.add_text(slide, 0.12, 0.06, 3.6, 0.28, asset_id, size=9, bold=True,
               color=GC['shape_fill_gray'], align=PP_ALIGN.LEFT, font=F_CAP)


# =============================================================================
# PRC-201 단계 박스+화살표 수평흐름 (5단계, box+arrow horizontal)
# =============================================================================
def build_box_arrow_horizontal(prs):
    s = blank_slide_gov(prs); gcaption(s, 'PRC-201'); S = []
    steps = ['과제 발굴', '현황 조사', '계획 수립', '사업 추진', '성과 관리']
    n = len(steps)
    bw, bh = 1.95, 1.15
    gap = 0.22
    tot = n * bw + (n - 1) * gap
    x0 = (SW - tot) / 2
    y = (SH - bh) / 2 - 0.3
    x = x0
    for i, lab in enumerate(steps):
        x = x0 + i * (bw + gap)
        box = c.add_box(s, x, y, bw, bh, fill=tone(i), line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        box.adjustments[0] = RADIUS_PCT
        gshape_text(box, lab, size=13)
        S.append(box)
        if i < n - 1:
            ar = c.add_box(s, x + bw + 0.02, y + bh / 2 - 0.18, gap - 0.04, 0.36,
                           fill=R('line'), line=None, shape=MSO_SHAPE.RIGHT_ARROW)
            S.append(ar)
    S.append(gtext(s, x0, y + bh + 0.25, tot, 0.4, '단계별 세부 추진계획은 하위 실행계획에 따라 조정',
                   size=GOV['size_pt']['caption'], color=GC['shape_fill_gray'], font=F_CAP))
    c.group_asset(s, S, 'PRC-201')
    return s


# =============================================================================
# PRC-202 세로 스텝(계단형) — 4단계 상승 계단(우상향 스텝래더)
# =============================================================================
def build_vertical_stairs(prs):
    s = blank_slide_gov(prs); gcaption(s, 'PRC-202'); S = []
    steps = [('1단계', '기반 조성'), ('2단계', '시범 운영'), ('3단계', '본격 확산'), ('4단계', '고도화·정착')]
    n = len(steps)
    pw, ph = 2.55, 0.95
    x0, y0 = 0.9, 5.55
    dx, dy = 2.35, -1.15
    pos = []
    for i, (no, lab) in enumerate(steps):
        px = x0 + i * dx; py = y0 + i * dy
        plat = c.add_box(s, px, py, pw, ph, fill=tone(i), line=None, shape=MSO_SHAPE.RECTANGLE)
        gshape_text(plat, no + '\n' + lab, size=12)
        S.append(plat)
        pos.append((px, py))
    for i in range(n - 1):
        (px, py) = pos[i]; (nx, ny) = pos[i + 1]
        cn = c.connector(s, px + pw, py + ph / 2, nx + pw * 0.15, ny + ph, color=R('line'), w=LINE_ACCENT_PT)
        arrowize(cn)
        S.append(cn)
    (lx, ly) = pos[-1]
    S.append(gtext(s, lx - 0.3, ly - 0.5, pw + 0.6, 0.4, '단계별 목표 달성 시 다음 단계 전환',
                   size=GOV['size_pt']['caption'], color=GC['shape_fill_gray'], font=F_CAP))
    c.group_asset(s, S, 'PRC-202')
    return s


# =============================================================================
# PRC-203 chevron(갈매기) 흐름 — 5단계 CHEVRON 화살 연속
# =============================================================================
def build_chevron(prs):
    s = blank_slide_gov(prs); gcaption(s, 'PRC-203'); S = []
    steps = ['접수', '검토', '심사', '승인', '집행']
    n = len(steps)
    cw, ch = 1.98, 1.2
    overlap = 0.35
    tot = cw + (n - 1) * (cw - overlap)
    x0 = (SW - tot) / 2
    y = (SH - ch) / 2 - 0.2
    for i, lab in enumerate(steps):
        x = x0 + i * (cw - overlap)
        chv = c.add_box(s, x, y, cw, ch, fill=tone(i), line=WHITE, line_w=BORDER_PT,
                        shape=MSO_SHAPE.CHEVRON)
        gshape_text(chv, lab, size=14)
        S.append(chv)
    S.append(gtext(s, x0, y + ch + 0.25, tot, 0.4, '민원 처리 5단계 표준 절차',
                   size=GOV['size_pt']['caption'], color=GC['shape_fill_gray'], font=F_CAP))
    c.group_asset(s, S, 'PRC-203')
    return s


# =============================================================================
# PRC-204 순환(cycle) — 4단계 원형 순환(PDCA류)
# =============================================================================
def build_cycle(prs):
    s = blank_slide_gov(prs); gcaption(s, 'PRC-204'); S = []
    cx, cy = SW / 2, SH / 2 + 0.1
    radius = 2.35; node_d = 1.7
    labels = ['계획(Plan)', '실행(Do)', '점검(Check)', '개선(Act)']
    n = len(labels)
    pos = []
    for i in range(n):
        ang = math.radians(-90 + i * (360 / n))
        pos.append((cx + radius * math.cos(ang), cy + radius * math.sin(ang)))
    core = c.add_box(s, cx - 1.05, cy - 0.75, 2.1, 1.5, fill=R('panel_bg'),
                     line=R('shape_fill'), line_w=BORDER_PT, shape=MSO_SHAPE.OVAL)
    c.set_shape_text(core, '지속 개선\n순환관리', size=12, bold=True, color=R('body_text'), font=F_SUB)
    S.append(core)
    for i in range(n):
        x1, y1 = pos[i]; x2, y2 = pos[(i + 1) % n]
        dxx, dyy = x2 - x1, y2 - y1
        d = math.hypot(dxx, dyy); ux, uy = dxx / d, dyy / d
        off = node_d / 2 + 0.06
        cn = c.connector(s, x1 + ux * off, y1 + uy * off, x2 - ux * off, y2 - uy * off,
                         color=R('line'), w=LINE_ACCENT_PT)
        arrowize(cn)
        S.append(cn)
    for i, (nx, ny) in enumerate(pos):
        nd = c.add_box(s, nx - node_d / 2, ny - node_d / 2, node_d, node_d, fill=tone(i),
                       line=WHITE, line_w=BORDER_PT * 2, shape=MSO_SHAPE.OVAL)
        gshape_text(nd, labels[i], size=12)
        S.append(nd)
    c.group_asset(s, S, 'PRC-204')
    return s


# =============================================================================
# PRC-205 허브-스포크 — 중앙 허브 + 방사 5스포크
# =============================================================================
def build_hub_spoke(prs):
    s = blank_slide_gov(prs); gcaption(s, 'PRC-205'); S = []
    cx, cy = SW / 2, SH / 2 + 0.1
    hub_d = 2.0; spoke_d = 1.5; radius = 2.55
    spokes = ['정책 기획', '예산 편성', '사업 집행', '성과 평가', '환류·개선']
    n = len(spokes)
    pos = []
    for i in range(n):
        ang = math.radians(-90 + i * (360 / n))
        pos.append((cx + radius * math.cos(ang), cy + radius * math.sin(ang)))
    for (sx, sy) in pos:
        cn = c.connector(s, cx, cy, sx, sy, color=R('shape_fill'), w=BORDER_PT * 1.5)
        S.append(cn)
    hub = c.add_box(s, cx - hub_d / 2, cy - hub_d / 2, hub_d, hub_d, fill=R('accent_primary'),
                    line=WHITE, line_w=BORDER_PT * 2, shape=MSO_SHAPE.OVAL)
    gshape_text(hub, '통합 추진\n체계', size=13)
    S.append(hub)
    for i, (sx, sy) in enumerate(pos):
        sp = c.add_box(s, sx - spoke_d / 2, sy - spoke_d / 2, spoke_d, spoke_d,
                       fill=tone(1 + i % 3), line=WHITE, line_w=BORDER_PT * 1.5, shape=MSO_SHAPE.OVAL)
        gshape_text(sp, spokes[i], size=11)
        S.append(sp)
    c.group_asset(s, S, 'PRC-205')
    return s


# =============================================================================
# PRC-206 게이트형 파이프라인 — 박스+판단(마름모) 게이트+반려경로
# =============================================================================
def build_gated_pipeline(prs):
    s = blank_slide_gov(prs); gcaption(s, 'PRC-206'); S = []
    ymid = SH / 2 - 0.2
    bw, bh = 1.75, 1.0
    dw = 1.45
    aw = 0.55
    x = 0.85

    def box(label, tn):
        nonlocal x
        sp = c.add_box(s, x, ymid - bh / 2, bw, bh, fill=tone(tn), line=None,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        sp.adjustments[0] = RADIUS_PCT
        gshape_text(sp, label, size=12)
        S.append(sp); x += bw
        return sp

    def diamond(label, tn):
        nonlocal x
        sp = c.add_box(s, x, ymid - dw / 2, dw, dw, fill=tone(tn), line=None, shape=MSO_SHAPE.DIAMOND)
        gshape_text(sp, label, size=10.5)
        S.append(sp); x += dw
        return sp

    def arrow():
        nonlocal x
        sp = c.add_box(s, x + 0.03, ymid - 0.2, aw, 0.4, fill=R('line'), line=None,
                       shape=MSO_SHAPE.RIGHT_ARROW)
        S.append(sp); x += aw + 0.06

    box('신청 접수', 0); arrow()
    diamond('요건\n검토', 2); arrow()
    box('사업 수행', 1); arrow()
    d2 = diamond('실적\n검증', 2); arrow()
    box('정산·종료', 0)

    d2cx = Emu(d2.left).inches + Emu(d2.width).inches / 2
    reject_top = ymid + dw / 2 + 0.05
    da = c.add_box(s, d2cx - 0.2, reject_top, 0.4, 0.5, fill=R('warn'), line=None,
                   shape=MSO_SHAPE.DOWN_ARROW)
    S.append(da)
    rb = c.add_box(s, d2cx - bw / 2, reject_top + 0.55, bw, 0.8, fill=R('warn'), line=None,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rb.adjustments[0] = RADIUS_PCT
    gshape_text(rb, '보완 요청', size=11)
    S.append(rb)
    c.group_asset(s, S, 'PRC-206')
    return s


# =============================================================================
# 매니페스트 조각
# =============================================================================
def manifest():
    E = c.entry
    ed_full = ['step-text', 'step-desc', 'step-count', 'color']
    return [
        E('PRC-201', 'PRC', '[gov] 단계 박스+화살표 수평흐름 5단', FILE_REL, 1,
          ['프로세스', 'gov', '수평흐름', '박스화살표', '5단계'],
          {'count': 5, 'style': 'box-arrow-horizontal', 'canvas': 'gov'},
          {'steps': [{'label': l} for l in ['과제 발굴', '현황 조사', '계획 수립', '사업 추진', '성과 관리']],
           'count': 5, 'direction': 'horizontal'}, ed_full, master='gov',
          recommended_use=['추진절차', '업무흐름', '단계흐름']),
        E('PRC-202', 'PRC', '[gov] 세로 스텝 계단형 4단(상승)', FILE_REL, 2,
          ['프로세스', 'gov', '계단형', '세로스텝', '상승', '4단계'],
          {'count': 4, 'style': 'vertical-stairs-ascending', 'canvas': 'gov'},
          {'steps': [{'label': l} for l in ['기반 조성', '시범 운영', '본격 확산', '고도화·정착']],
           'count': 4, 'direction': 'up-right'}, ed_full, master='gov',
          recommended_use=['단계적확산', '추진로드맵', '점진고도화']),
        E('PRC-203', 'PRC', '[gov] 갈매기(chevron) 흐름 5단', FILE_REL, 3,
          ['프로세스', 'gov', 'chevron', '갈매기', '5단계', '절차'],
          {'count': 5, 'style': 'chevron', 'canvas': 'gov'},
          {'steps': [{'label': l} for l in ['접수', '검토', '심사', '승인', '집행']],
           'count': 5, 'direction': 'horizontal'}, ed_full, master='gov',
          recommended_use=['민원처리절차', '심사절차', '표준프로세스']),
        E('PRC-204', 'PRC', '[gov] 순환(cycle) 4단계 PDCA형', FILE_REL, 4,
          ['프로세스', 'gov', '순환', '사이클', 'PDCA', '4단계'],
          {'count': 4, 'style': 'circular', 'canvas': 'gov'},
          {'steps': [{'label': l} for l in ['계획(Plan)', '실행(Do)', '점검(Check)', '개선(Act)']],
           'count': 4, 'direction': 'cyclic'}, ed_full, master='gov',
          recommended_use=['환류체계', '지속개선', '운영사이클']),
        E('PRC-205', 'PRC', '[gov] 허브-스포크 5방사', FILE_REL, 5,
          ['프로세스', 'gov', '허브앤스포크', '방사형', '중앙집중', '5개'],
          {'count': 5, 'style': 'hub-spoke', 'canvas': 'gov'},
          {'hub': '통합 추진 체계',
           'spokes': [{'label': l} for l in ['정책 기획', '예산 편성', '사업 집행', '성과 평가', '환류·개선']],
           'count': 5, 'direction': 'radial'}, ['hub-text', 'spoke-text', 'spoke-count', 'color'],
          master='gov', recommended_use=['추진체계', '통합운영', '연계구조']),
        E('PRC-206', 'PRC', '[gov] 게이트형 파이프라인 (판단+반려경로)', FILE_REL, 6,
          ['프로세스', 'gov', '파이프라인', '게이트', '판단', '반려'],
          {'count': 5, 'style': 'gated-pipeline', 'canvas': 'gov'},
          {'steps': [{'label': '신청 접수'}, {'label': '요건 검토', 'type': 'gate'},
                     {'label': '사업 수행'}, {'label': '실적 검증', 'type': 'gate'},
                     {'label': '정산·종료'}], 'count': 5, 'direction': 'horizontal'}, ed_full,
          master='gov', recommended_use=['품질게이트', '검수절차', '승인프로세스']),
    ]


if __name__ == '__main__':
    prs = new_deck_gov()
    build_box_arrow_horizontal(prs)
    build_vertical_stairs(prs)
    build_chevron(prs)
    build_cycle(prs)
    build_hub_spoke(prs)
    build_gated_pipeline(prs)
    out = c.save_deck(prs, FILE_REL)
    frag = c.write_fragment('PRC_gov', manifest())
    print('OK', out)
    print('FRAGMENT', frag)
