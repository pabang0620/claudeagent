# -*- coding: utf-8 -*-
"""CMP 카테고리(비교표/평가표/매트릭스) 에셋 생성 — CMP-001~CMP-010."""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

R = c.role
CC = c.C

# ---------- 로컬 테이블 헬퍼 ----------
def add_table(slide, x, y, w, h, nrows, ncols, col_widths=None, row_h=None):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    return gf, tbl

def set_cell(cell, text, size=11, bold=False, color=None, fill=None,
             align=PP_ALIGN.CENTER, font=None):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = R("row_base")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(5)
    cell.margin_top = cell.margin_bottom = Pt(3)
    tf = cell.text_frame; tf.word_wrap = True
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))

def title_block(slide, text, sub=None):
    c.add_text(slide, 0.55, 0.52, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.02, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)

# =========================================================
# 매트릭스 계열
# =========================================================
def matrix_2x2(slide, asset_id, title, sub, axis_x, axis_y, quads, palette):
    """quads: 4 dict {title, desc}. palette: list of 4 fill colors (q1 TL,q2 TR,q3 BL,q4 BR)."""
    title_block(slide, title, sub)
    x0, x1 = 2.95, 11.35
    y0, y1 = 1.55, 6.15
    gap = 0.30
    bw = (x1 - x0 - gap) / 2
    bh = (y1 - y0 - gap) / 2
    cx = x0 + bw + gap / 2
    cy = y0 + bh + gap / 2
    # 최상위 앵커(투명 컨테이너)
    cont = c.add_box(slide, x0 - 0.05, y0 - 0.05, (x1 - x0) + 0.1, (y1 - y0) + 0.1, fill=None, line=None)
    c.name_asset(cont, asset_id)
    coords = [(x0, y0), (cx + gap / 2, y0), (x0, cy + gap / 2), (cx + gap / 2, cy + gap / 2)]
    for (bx, by), q, fill in zip(coords, quads, palette):
        box = c.add_box(slide, bx, by, bw, bh, fill=fill, line=R("border"), line_w=1.0,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = q["title"]
        c.set_kfont(r, c.FONT_H, 15, True, R("header_text"))
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = q["desc"]
        c.set_kfont(r2, c.FONT_B, 10, False, R("header_text"))
    # 축: 십자선
    c.connector(slide, x0, cy, x1, cy, color=R("muted_text"), w=1.5)
    c.connector(slide, cx, y0, cx, y1, color=R("muted_text"), w=1.5)
    # X축 라벨(하단)
    c.add_text(slide, x0, y1 + 0.12, x1 - x0, 0.35, "%s  →" % axis_x, size=12, bold=True,
               color=R("accent_primary"), align=PP_ALIGN.CENTER)
    # Y축 라벨(좌측 회전)
    yl = c.add_text(slide, 0.1, cy - 0.18, y1 - y0, 0.36, "%s  →" % axis_y, size=12, bold=True,
                    color=R("accent_primary"), align=PP_ALIGN.CENTER)
    yl.left = Inches((x0 - 0.45) - (y1 - y0) / 2)
    yl.top = Inches(cy - 0.18)
    yl.rotation = 270

def priority_matrix(slide, asset_id, title, sub, axis_x, axis_y, quad_titles, items):
    """items: list of (label, quadrant_index 0..3). 4분면에 칩 배치."""
    title_block(slide, title, sub)
    x0, x1 = 2.95, 11.35
    y0, y1 = 1.55, 6.15
    gap = 0.28
    bw = (x1 - x0 - gap) / 2
    bh = (y1 - y0 - gap) / 2
    cx = x0 + bw + gap / 2
    cy = y0 + bh + gap / 2
    cont = c.add_box(slide, x0 - 0.05, y0 - 0.05, (x1 - x0) + 0.1, (y1 - y0) + 0.1, fill=None, line=None)
    c.name_asset(cont, asset_id)
    coords = [(x0, y0), (cx + gap / 2, y0), (x0, cy + gap / 2), (cx + gap / 2, cy + gap / 2)]
    tints = [CC["gray_050"], CC["beige_100"], CC["gray_100"], CC["gray_050"]]
    qcolors = [R("warn"), R("accent_primary"), R("accent_secondary"), R("muted_text")]
    for (bx, by), qt, bg in zip(coords, quad_titles, tints):
        c.add_box(slide, bx, by, bw, bh, fill=bg, line=R("border"), line_w=1.0)
        c.add_text(slide, bx + 0.12, by + 0.08, bw - 0.24, 0.32, qt, size=12, bold=True,
                   color=R("header_fill"), align=PP_ALIGN.LEFT)
    # 칩 배치: 사분면별 세로 스택
    stack = {0: 0, 1: 0, 2: 0, 3: 0}
    for label, qi in items:
        bx, by = coords[qi]
        idx = stack[qi]; stack[qi] += 1
        chip_y = by + 0.5 + idx * 0.62
        chip = c.add_box(slide, bx + 0.25, chip_y, bw - 0.5, 0.5, fill=R("row_base"),
                         line=qcolors[qi], line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.set_shape_text(chip, label, size=11, bold=True, color=qcolors[qi])
    c.connector(slide, x0, cy, x1, cy, color=R("muted_text"), w=1.5)
    c.connector(slide, cx, y0, cx, y1, color=R("muted_text"), w=1.5)
    c.add_text(slide, x0, y1 + 0.12, x1 - x0, 0.35, "%s  →" % axis_x, size=12, bold=True,
               color=R("accent_primary"), align=PP_ALIGN.CENTER)
    yl = c.add_text(slide, 0.1, cy - 0.18, y1 - y0, 0.36, "%s  →" % axis_y, size=12, bold=True,
                    color=R("accent_primary"), align=PP_ALIGN.CENTER)
    yl.left = Inches((x0 - 0.45) - (y1 - y0) / 2)
    yl.top = Inches(cy - 0.18)
    yl.rotation = 270

# =========================================================
# 평가표 계열
# =========================================================
def eval_table_3col(slide, asset_id, title, sub, headers, rows):
    """rows: list of (항목, 현행(text,color), 개선(text,color))."""
    title_block(slide, title, sub)
    x, y = 1.1, 1.55
    w = 11.1
    nrows = len(rows) + 1
    rh = 0.62
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 3,
                        col_widths=[4.7, 3.2, 3.2], row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    for j, hd in enumerate(headers):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"))
    for i, (item, cur, imp) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), item, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        ct, ccolor = cur
        set_cell(tbl.cell(i, 1), ct, size=11, bold=True, color=ccolor, fill=base)
        it, icolor = imp
        set_cell(tbl.cell(i, 2), it, size=11, bold=True, color=icolor, fill=base)

def weighted_table(slide, asset_id, title, sub, rows, total):
    """항목|배점|평가 가중치 평가표. rows: (항목, 배점, 평가). total: (배점합, 평가합)."""
    title_block(slide, title, sub)
    x, y = 1.6, 1.55
    w = 10.1
    nrows = len(rows) + 2
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 3,
                        col_widths=[5.9, 2.1, 2.1], row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    for j, hd in enumerate(["평가 항목", "배점", "평가 점수"]):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"))
    for i, (item, pts, score) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), item, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        set_cell(tbl.cell(i, 1), pts, size=11, fill=base, color=R("muted_text"))
        set_cell(tbl.cell(i, 2), score, size=11, bold=True, fill=base, color=R("accent_primary"))
    li = nrows - 1
    set_cell(tbl.cell(li, 0), "합계", size=12, bold=True, color=R("header_text"),
             fill=CC["navy_600"], align=PP_ALIGN.LEFT)
    set_cell(tbl.cell(li, 1), total[0], size=12, bold=True, color=R("header_text"), fill=CC["navy_600"])
    set_cell(tbl.cell(li, 2), total[1], size=12, bold=True, color=R("header_text"), fill=CC["navy_600"])

def checklist_table(slide, asset_id, title, sub, rows):
    """rows: (항목, checked_bool, note). 체크박스 도형 오버레이."""
    title_block(slide, title, sub)
    x, y = 1.6, 1.55
    w = 10.1
    nrows = len(rows) + 1
    rh = 0.62
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 3,
                        col_widths=[1.3, 6.5, 2.3], row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    for j, hd in enumerate(["확인", "점검 항목", "상태"]):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"))
    for i, (item, checked, note) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), "", fill=base)
        set_cell(tbl.cell(i, 1), item, size=11, bold=False, fill=base, align=PP_ALIGN.LEFT)
        ncolor = R("accent_secondary") if checked else R("muted_text")
        set_cell(tbl.cell(i, 2), note, size=10, bold=True, color=ncolor, fill=base)
        # 체크박스 도형 오버레이 (확인 칼럼 중앙)
        cbx = x + (1.3 - 0.32) / 2
        cby = y + rh * i + (rh - 0.32) / 2
        cb_fill = R("accent_secondary") if checked else R("row_base")
        cb_line = R("accent_secondary") if checked else R("border")
        cb = c.add_box(slide, cbx, cby, 0.32, 0.32, fill=cb_fill, line=cb_line, line_w=1.5,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        if checked:
            c.set_shape_text(cb, "✓", size=13, bold=True, color=R("header_text"))

# =========================================================
# 우열 비교표 (자사 강조)
# =========================================================
def versus_table(slide, asset_id, title, sub, cols, rows, own_col=1):
    """cols: 헤더 리스트. rows: (항목, [셀…]) 셀은 (symbol, color) or text. own_col 강조."""
    title_block(slide, title, sub)
    ncol = len(cols)
    x = 1.1
    w = 11.1
    first_w = 4.4
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    nrows = len(rows) + 1
    rh = 0.6
    gf, tbl = add_table(slide, x, 1.55, w, rh * nrows, nrows, ncol,
                        col_widths=col_widths, row_h=[rh] * nrows)
    c.name_asset(gf, asset_id)
    for j, hd in enumerate(cols):
        fill = R("accent_primary") if j == own_col else R("header_fill")
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (item, cells) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), item, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j, cell in enumerate(cells, start=1):
            if j == own_col:
                cfill = CC["gray_050"]
            else:
                cfill = base
            if isinstance(cell, tuple):
                sym, col = cell
                set_cell(tbl.cell(i, j), sym, size=14, bold=True, color=col, fill=cfill)
            else:
                set_cell(tbl.cell(i, j), cell, size=11, fill=cfill, color=R("muted_text"))
    # 자사열 강조 테두리
    ax = x + sum(col_widths[:own_col])
    frame = c.add_box(slide, ax, 1.55, col_widths[own_col], rh * nrows, fill=None,
                      line=R("accent_primary"), line_w=2.25)

# =========================================================
# Before / After 대비 카드
# =========================================================
def before_after(slide, asset_id, title, sub, before, after):
    title_block(slide, title, sub)
    y0 = 1.65
    ch = 4.6
    cw = 5.35
    gap_mid = 1.1
    x_b = 1.1
    x_a = x_b + cw + gap_mid
    cont = c.add_box(slide, x_b, y0, x_a + cw - x_b, ch, fill=None, line=None)
    c.name_asset(cont, asset_id)
    # Before 카드 (muted)
    cb = c.add_box(slide, x_b, y0, cw, ch, fill=CC["gray_100"], line=R("border"), line_w=1.0,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    hb = c.add_box(slide, x_b, y0, cw, 0.75, fill=R("muted_text"), line=None,
                   shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    c.set_shape_text(hb, before["head"], size=15, bold=True, color=R("header_text"))
    _bullet_list(slide, x_b + 0.35, y0 + 1.0, cw - 0.7, before["items"], R("muted_text"))
    # After 카드 (accent)
    ca = c.add_box(slide, x_a, y0, cw, ch, fill=CC["gray_050"], line=R("accent_primary"), line_w=1.5,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ha = c.add_box(slide, x_a, y0, cw, 0.75, fill=R("accent_primary"), line=None,
                   shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    c.set_shape_text(ha, after["head"], size=15, bold=True, color=R("header_text"))
    _bullet_list(slide, x_a + 0.35, y0 + 1.0, cw - 0.7, after["items"], R("accent_secondary"))
    # 중앙 화살표
    ar = c.add_box(slide, x_b + cw + 0.15, y0 + ch / 2 - 0.4, gap_mid - 0.3, 0.8,
                   fill=R("accent_point"), line=None, shape=MSO_SHAPE.RIGHT_ARROW)
    c.set_shape_text(ar, "개선", size=11, bold=True, color=R("header_text"))

def _bullet_list(slide, x, y, w, items, dot_color):
    for i, it in enumerate(items):
        yy = y + i * 0.62
        c.add_box(slide, x, yy + 0.12, 0.12, 0.12, fill=dot_color, line=None, shape=MSO_SHAPE.OVAL)
        c.add_text(slide, x + 0.28, yy - 0.02, w - 0.28, 0.5, it, size=11,
                   color=R("body_text"), align=PP_ALIGN.LEFT)

# =========================================================
# 선택지 비교 3열 카드 (플랜 A/B/C, 추천 뱃지)
# =========================================================
def plan_cards(slide, asset_id, title, sub, plans):
    """plans: list of dict {name, tier, items, color, recommended}."""
    title_block(slide, title, sub)
    y0 = 1.7
    ch = 4.75
    n = len(plans)
    x_start = 1.1
    total_w = 11.1
    gap = 0.35
    cw = (total_w - gap * (n - 1)) / n
    cont = c.add_box(slide, x_start, y0 - 0.1, total_w, ch + 0.2, fill=None, line=None)
    c.name_asset(cont, asset_id)
    for i, pl in enumerate(plans):
        cx = x_start + i * (cw + gap)
        rec = pl.get("recommended")
        line_col = R("accent_point") if rec else R("border")
        line_w = 2.5 if rec else 1.0
        c.add_box(slide, cx, y0, cw, ch, fill=CC["gray_050"], line=line_col, line_w=line_w,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        hd = c.add_box(slide, cx, y0, cw, 1.05, fill=pl["color"], line=None,
                       shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        tf = hd.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = pl["name"]; c.set_kfont(r, c.FONT_H, 16, True, R("header_text"))
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = pl["tier"]; c.set_kfont(r2, c.FONT_B, 10, False, R("header_text"))
        _bullet_list(slide, cx + 0.3, y0 + 1.35, cw - 0.6, pl["items"], pl["color"])
        if rec:
            bw_ = 1.5
            bd = c.add_box(slide, cx + cw - bw_ - 0.15, y0 + 0.12, bw_, 0.42, fill=R("warn"),
                           line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            c.set_shape_text(bd, "★ 추천", size=10, bold=True, color=R("header_text"))

# =========================================================
# 빌드
# =========================================================
entries = []

def E(*a, **k):
    entries.append(c.entry(*a, **k))

# ---- FILE 1: 매트릭스 (CMP-001, 002, 009) ----
p1 = c.new_deck()

s = c.blank_slide(p1); c.id_caption(s, "CMP-001")
matrix_2x2(s, "CMP-001", "2x2 전략 포지셔닝 매트릭스", "축 라벨과 4개 사분면 박스를 자유 편집",
           "시장 매력도", "실행 역량",
           [{"title": "집중 육성", "desc": "글로벌 역량 우위 영역"},
            {"title": "선택 확대", "desc": "협업 역량 강화 영역"},
            {"title": "현상 유지", "desc": "기존 사업 안정화"},
            {"title": "재검토", "desc": "AI 활용도 낮음"}],
           [R("accent_primary"), R("accent_secondary"), R("muted_text"), CC["navy_600"]])

s = c.blank_slide(p1); c.id_caption(s, "CMP-002")
matrix_2x2(s, "CMP-002", "2x2 우선순위 매트릭스 (색 변형)", "네이비·틸 코퍼레이트 색 변형",
           "기대 효과", "투자 규모",
           [{"title": "핵심 과제", "desc": "즉시 착수 대상"},
            {"title": "전략 과제", "desc": "단계적 확대"},
            {"title": "보완 과제", "desc": "리스크 관리"},
            {"title": "관망 과제", "desc": "우선순위 후순위"}],
           [CC["navy_800"], R("accent_secondary"), R("sub_header"), R("muted_text")])

s = c.blank_slide(p1); c.id_caption(s, "CMP-009")
priority_matrix(s, "CMP-009", "우선순위 매트릭스 (중요도 x 긴급도)", "과제를 4분면에 칩으로 배치",
                "긴급도", "중요도",
                ["즉시 실행 (중요·긴급)", "계획 수립 (중요)", "위임 (긴급)", "보류 (후순위)"],
                [("글로벌 진출", 0), ("AI 활용도", 0), ("협업 역량", 1),
                 ("품질 관리", 1), ("현지화", 2), ("문서 정비", 3)])

f1 = c.save_deck(p1, "decks/04_compare/CMP_matrix_v1.pptx")

# ---- FILE 2: 평가표 (CMP-003, 006, 007) ----
p2 = c.new_deck()

s = c.blank_slide(p2); c.id_caption(s, "CMP-003")
eval_table_3col(s, "CMP-003", "3열 평가표 (현행 vs 개선)", "항목별 현행·개선 상태를 기호로 표기",
                ["평가 항목", "현행", "개선"],
                [("글로벌 역량", ("✕ 미흡", R("warn")), ("○ 확보", R("accent_secondary"))),
                 ("협업 역량", ("△ 보통", R("muted_text")), ("○ 우수", R("accent_secondary"))),
                 ("AI 활용도", ("✕ 없음", R("warn")), ("○ 적용", R("accent_secondary"))),
                 ("품질 관리", ("△ 부분", R("muted_text")), ("○ 전면", R("accent_secondary"))),
                 ("대응 속도", ("✕ 지연", R("warn")), ("○ 즉시", R("accent_secondary")))])

s = c.blank_slide(p2); c.id_caption(s, "CMP-006")
weighted_table(s, "CMP-006", "가중치 평가표 (배점 · 평가)", "항목별 배점과 획득 점수 합산",
               [("글로벌 역량", "30", "27"),
                ("협업 역량", "20", "18"),
                ("AI 활용도", "20", "16"),
                ("사업 이해도", "15", "14"),
                ("실행 가능성", "15", "13")],
               ("100", "88"))

s = c.blank_slide(p2); c.id_caption(s, "CMP-007")
checklist_table(s, "CMP-007", "체크리스트 표 (점검 항목)", "체크박스 도형으로 완료 여부 표시",
                [("사업 목표 정합성 검토", True, "완료"),
                 ("글로벌 역량 자료 확보", True, "완료"),
                 ("협업 파트너 확정", False, "진행중"),
                 ("AI 활용 방안 수립", True, "완료"),
                 ("리스크 대응 계획", False, "예정"),
                 ("최종 산출물 검수", False, "예정")])

f2 = c.save_deck(p2, "decks/04_compare/CMP_eval-table_v1.pptx")

# ---- FILE 3: 우열 비교표 + 카드 (CMP-004, 005, 008, 010) ----
p3 = c.new_deck()

O = ("○", R("accent_secondary"))   # 우수
T = ("△", R("muted_text"))         # 보통
X = ("✕", R("warn"))               # 미흡

s = c.blank_slide(p3); c.id_caption(s, "CMP-004")
versus_table(s, "CMP-004", "우열 비교표 (자사 vs 경쟁사)", "자사 열을 강조, ○△✕ 기호로 우열 표기",
             ["평가 항목", "자사", "경쟁사"],
             [("글로벌 역량", [O, T]),
              ("협업 역량", [O, X]),
              ("AI 활용도", [O, X]),
              ("품질 관리", [O, T]),
              ("가격 경쟁력", [T, O]),
              ("납기 대응", [O, T])],
             own_col=1)

s = c.blank_slide(p3); c.id_caption(s, "CMP-005")
versus_table(s, "CMP-005", "우열 비교표 (3사 비교 변형)", "자사·경쟁사 A·B 3열 비교, 자사 강조",
             ["평가 항목", "자사", "경쟁사 A", "경쟁사 B"],
             [("글로벌 역량", [O, T, X]),
              ("협업 역량", [O, T, T]),
              ("AI 활용도", [O, X, T]),
              ("품질 관리", [O, O, T]),
              ("가격 경쟁력", [T, O, X]),
              ("납기 대응", [O, T, T])],
             own_col=1)

s = c.blank_slide(p3); c.id_caption(s, "CMP-008")
before_after(s, "CMP-008", "Before / After 대비 카드", "개선 전·후를 2열 카드로 대비",
             {"head": "Before (현행)",
              "items": ["글로벌 대응 체계 부재", "수작업 중심 품질 관리", "협업 프로세스 분산", "AI 미적용"]},
             {"head": "After (개선)",
              "items": ["글로벌 표준 대응 체계", "AI 기반 품질 자동화", "통합 협업 플랫폼", "데이터 기반 의사결정"]})

s = c.blank_slide(p3); c.id_caption(s, "CMP-010")
plan_cards(s, "CMP-010", "선택지 비교 카드 (플랜 A/B/C)", "헤더 색 구분 · 추천 뱃지 강조",
           [{"name": "플랜 A", "tier": "기본형", "color": R("muted_text"),
             "items": ["핵심 기능 제공", "표준 지원", "기본 리포트"]},
            {"name": "플랜 B", "tier": "표준형", "color": R("accent_primary"), "recommended": True,
             "items": ["전체 기능 제공", "AI 활용 지원", "글로벌 대응", "우선 지원"]},
            {"name": "플랜 C", "tier": "확장형", "color": R("sub_header"),
             "items": ["맞춤 개발 포함", "전담 협업팀", "무제한 확장"]}])

f3 = c.save_deck(p3, "decks/04_compare/CMP_versus-card_v1.pptx")

# =========================================================
# 매니페스트 엔트리
# =========================================================
TAGS_BASE = ["비교표", "매트릭스", "2x2", "평가표", "우열비교"]
F1 = "decks/04_compare/CMP_matrix_v1.pptx"
F2 = "decks/04_compare/CMP_eval-table_v1.pptx"
F3 = "decks/04_compare/CMP_versus-card_v1.pptx"

E("CMP-001", "compare", "2x2 전략 포지셔닝 매트릭스", F1, 1,
  ["비교표", "매트릭스", "2x2", "사분면", "포지셔닝"],
  {"variant": "A", "quadrants": 4, "axis": 2},
  {"quadrants": ["q1", "q2", "q3", "q4"], "axis": ["x", "y"]},
  ["quad_titles", "quad_desc", "axis_labels", "quad_fill"],
  recommended_use=["경쟁력 비교", "전략 포지셔닝", "우선순위"])

E("CMP-002", "compare", "2x2 우선순위 매트릭스 (색 변형)", F1, 2,
  ["비교표", "매트릭스", "2x2", "사분면", "색변형"],
  {"variant": "B", "quadrants": 4, "axis": 2},
  {"quadrants": ["q1", "q2", "q3", "q4"], "axis": ["x", "y"]},
  ["quad_titles", "quad_desc", "axis_labels", "quad_fill"],
  recommended_use=["경쟁력 비교", "투자 우선순위", "포지셔닝"])

E("CMP-003", "compare", "3열 평가표 (현행 vs 개선)", F2, 1,
  ["평가표", "비교표", "현행개선", "기호평가"],
  {"cols": 3, "rows": 5, "symbols": ["○", "△", "✕"]},
  {"rows": ["글로벌 역량", "협업 역량", "AI 활용도", "품질 관리", "대응 속도"],
   "cols": ["평가 항목", "현행", "개선"]},
  ["rows", "cell_symbols", "cell_color", "headers"],
  recommended_use=["평가기준", "개선 효과 제시", "현황 진단"])

E("CMP-004", "compare", "우열 비교표 (자사 vs 경쟁사)", F3, 1,
  ["비교표", "우열비교", "자사강조", "경쟁분석"],
  {"cols": 3, "rows": 6, "highlight_col": "자사", "symbols": ["○", "△", "✕"]},
  {"rows": ["글로벌 역량", "협업 역량", "AI 활용도", "품질 관리", "가격 경쟁력", "납기 대응"],
   "cols": ["항목", "자사", "경쟁"]},
  ["rows", "cell_symbols", "highlight_col", "headers"],
  recommended_use=["경쟁력 비교", "차별성 강조", "제안 우위"])

E("CMP-005", "compare", "우열 비교표 (3사 비교 변형)", F3, 2,
  ["비교표", "우열비교", "자사강조", "3사비교"],
  {"cols": 4, "rows": 6, "highlight_col": "자사", "symbols": ["○", "△", "✕"]},
  {"rows": ["글로벌 역량", "협업 역량", "AI 활용도", "품질 관리", "가격 경쟁력", "납기 대응"],
   "cols": ["항목", "자사", "경쟁사 A", "경쟁사 B"]},
  ["rows", "cell_symbols", "highlight_col", "headers"],
  recommended_use=["경쟁력 비교", "다자 비교", "제안 우위"])

E("CMP-006", "compare", "가중치 평가표 (배점 · 평가)", F2, 2,
  ["평가표", "가중치", "배점", "점수합산"],
  {"cols": 3, "rows": 5, "total_row": True},
  {"rows": ["글로벌 역량", "협업 역량", "AI 활용도", "사업 이해도", "실행 가능성"],
   "cols": ["항목", "배점", "평가"]},
  ["rows", "points", "scores", "total"],
  recommended_use=["평가기준", "정량 평가", "심사 대응"])

E("CMP-007", "compare", "체크리스트 표 (점검 항목)", F2, 3,
  ["평가표", "체크리스트", "체크박스", "점검"],
  {"cols": 3, "rows": 6, "checkbox_shape": True},
  {"rows": ["사업 목표 정합성 검토", "글로벌 역량 자료 확보", "협업 파트너 확정",
            "AI 활용 방안 수립", "리스크 대응 계획", "최종 산출물 검수"],
   "cols": ["확인", "점검 항목", "상태"]},
  ["rows", "checked_state", "note", "checkbox_fill"],
  recommended_use=["평가기준", "진행 점검", "완료 관리"])

E("CMP-008", "compare", "Before / After 대비 카드", F3, 3,
  ["비교표", "전후비교", "대비카드", "개선효과"],
  {"columns": 2, "before_items": 4, "after_items": 4},
  {"quadrants": ["before", "after"], "axis": ["state"]},
  ["before_head", "before_items", "after_head", "after_items", "card_color"],
  recommended_use=["개선 효과 제시", "전후 대비", "성과 강조"])

E("CMP-009", "compare", "우선순위 매트릭스 (중요도 x 긴급도)", F1, 3,
  ["매트릭스", "2x2", "우선순위", "중요도긴급도"],
  {"quadrants": 4, "axis": 2, "items": 6},
  {"quadrants": ["q1", "q2", "q3", "q4"], "axis": ["x", "y"]},
  ["quad_titles", "axis_labels", "items", "item_placement"],
  recommended_use=["우선순위", "과제 분류", "실행 계획"])

E("CMP-010", "compare", "선택지 비교 카드 (플랜 A/B/C)", F3, 4,
  ["비교표", "선택지비교", "플랜카드", "추천뱃지"],
  {"columns": 3, "recommended": "플랜 B", "items_per_card": 4},
  {"rows": ["플랜 A", "플랜 B", "플랜 C"], "cols": ["기본형", "표준형", "확장형"]},
  ["plan_names", "tiers", "items", "header_color", "recommended_badge"],
  recommended_use=["선택지 비교", "제안 옵션", "가격 정책"])

frag = c.write_fragment("CMP", entries)

print("FILE1", f1)
print("FILE2", f2)
print("FILE3", f3)
print("FRAG", frag)
print("ENTRIES", len(entries))
