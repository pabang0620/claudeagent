# -*- coding: utf-8 -*-
"""
차트 스타일(CHT) 대량 확충 생성기 — CHT-011 ~ CHT-200 (+190).
python-pptx 네이티브 차트(add_chart)만 사용. 파라미터화 스윕:
  차트타입군 × 색팔레트 오프셋 × 카테고리수(3~7) × 계열수(1~3).
콤보(막대+선)는 python-pptx 미지원 → AREA/2계열로 대체. RADAR는 RADAR_MARKERS 시도 후
실패 시 COLUMN_CLUSTERED 대체. 리터럴은 풀/템플릿으로 파라미터화(중복 최소화).
파일당 ≤20 슬라이드로 분산 저장. 매니페스트 조각 1회 기록(c.write_fragment).
"""
import sys, random, itertools
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# 계열 색 팔레트 (스윕 대상) — 토큰만 사용
PAL = [c.C["navy_800"], c.C["blue_500"], c.C["teal_500"],
       c.C["cyan_500"], c.C["purple_600"], c.C["gray_500"]]

LINE_TYPES = {"LINE", "LINE_MARKERS", "RADAR", "RADAR_MARKERS"}
PIE_TYPES = {"PIE", "DOUGHNUT"}

CX, CY, CW, CH = 0.6, 1.35, 12.1, 5.7

# ---------------- 데이터 풀 (더미, 제안서 맥락) ----------------
POOLS = {
    "genre":      ["게임", "음악", "방송", "영화", "캐릭터", "애니메이션", "만화"],
    "region":     ["아시아", "북미", "유럽", "중남미", "중동", "아프리카", "오세아니아"],
    "year":       ["2019", "2020", "2021", "2022", "2023", "2024", "2025"],
    "country":    ["일본", "미국", "중국", "대만", "베트남", "태국", "인도네시아"],
    "channel":    ["온라인", "오프라인", "모바일", "콘솔", "플랫폼", "오픈마켓", "자사몰"],
    "capability": ["기획", "제작", "유통", "마케팅", "현지화", "투자유치", "인프라"],
    "quarter":    ["1분기", "2분기", "3분기", "4분기"],
    "program":    ["R&D", "수출지원", "인력양성", "인프라구축", "글로벌마케팅", "네트워킹", "컨설팅"],
    "platform":   ["넷플릭스", "유튜브", "틱톡", "스포티파이", "아마존", "디즈니플러스"],
    "stage":      ["도입", "성장", "성숙", "확산", "고도화"],
}
SUBJECT = {
    "genre": "장르별", "region": "권역별", "year": "연도별", "country": "국가별",
    "channel": "채널별", "capability": "역량", "quarter": "분기별", "program": "사업별",
    "platform": "플랫폼별", "stage": "단계별",
}
METRICS = ["수출액", "매출", "점유율", "성과지표", "투자액", "이용자수"]
SUFFIX = ["현황", "추이", "비교", "분석", "전망", "구성"]
MUNIT = {"수출액": "백만$", "매출": "억원", "점유율": "%", "성과지표": "점",
         "투자액": "억원", "이용자수": "만명"}

# 계열명 세트 (다계열 스윕)
SNS = {
    "years":   ["2023", "2024", "2025"],
    "segment": ["대기업", "중소기업", "벤처"],
    "plan":    ["계획", "실적", "목표"],
    "compare": ["자사", "경쟁사", "업계평균"],
    "region3": ["아시아", "북미", "유럽"],
}
SNS_KEYS = list(SNS.keys())


def kfont(font, name=None, size=None, bold=None, color=None):
    name = name or c.FONT_B
    if size is not None: font.size = Pt(size)
    if bold is not None: font.bold = bold
    if color is not None: font.color.rgb = color
    font.name = name
    rPr = font._element
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def set_title(chart, text):
    chart.has_title = True
    tf = chart.chart_title.text_frame
    tf.text = text
    c.set_kfont(tf.paragraphs[0].runs[0], c.FONT_H, 15, True, c.role("body_text"))


def set_legend(chart, position=XL_LEGEND_POSITION.BOTTOM, show=True):
    chart.has_legend = show
    if show:
        chart.legend.position = position
        chart.legend.include_in_layout = False
        kfont(chart.legend.font, c.FONT_B, 11, False, c.role("muted_text"))


def style_axes(chart):
    for ax_name in ("category_axis", "value_axis"):
        try:
            ax = getattr(chart, ax_name)
            kfont(ax.tick_labels.font, c.FONT_B, 10, False, c.role("muted_text"))
        except Exception:
            pass


def color_series(chart, ctype, off=0):
    plot = chart.plots[0]
    if ctype in PIE_TYPES:
        ser = plot.series[0]
        for j, pt in enumerate(ser.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = PAL[(off + j) % len(PAL)]
    elif ctype in LINE_TYPES:
        for i, ser in enumerate(plot.series):
            ser.format.line.color.rgb = PAL[(off + i) % len(PAL)]
            ser.format.line.width = Pt(2.5)
    else:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = PAL[(off + i) % len(PAL)]


def add_labels(chart, ctype):
    plot = chart.plots[0]
    try:
        if ctype in PIE_TYPES:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_percentage = True; dl.show_value = False
            dl.number_format = "0%"; dl.number_format_is_linked = False
            dl.position = XL_LABEL_POSITION.OUTSIDE_END if ctype == "PIE" else XL_LABEL_POSITION.CENTER
            kfont(dl.font, c.FONT_B, 10, True, c.role("body_text"))
        else:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = "#,##0"; dl.number_format_is_linked = False
            kfont(dl.font, c.FONT_B, 9, False, c.role("body_text"))
    except Exception:
        pass


def make_chart(slide, asset_id, ctype, title, categories, series, pal_off=0,
               legend_pos=XL_LEGEND_POSITION.BOTTOM, show_legend=True, labels=True):
    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(s["name"], s["values"])
    xl = getattr(XL_CHART_TYPE, ctype)
    gf = slide.shapes.add_chart(xl, Inches(CX), Inches(CY), Inches(CW), Inches(CH), cd)
    chart = gf.chart
    kfont(chart.font, c.FONT_B, 11, False, c.role("body_text"))
    set_title(chart, title)
    set_legend(chart, legend_pos, show_legend)
    style_axes(chart)
    color_series(chart, ctype, pal_off)
    if labels:
        add_labels(chart, ctype)
    c.name_asset(gf, asset_id)
    return gf


def meta(asset_id, name, file_rel, slide_idx, ctype, categories, series,
         tags, params, recommended, fallback_note=None):
    bindings = {
        "categories": list(categories),
        "series": [{"name": s["name"], "values": list(s["values"])} for s in series],
        "chart_type": ctype,
    }
    if fallback_note:
        bindings["_note"] = fallback_note
    return c.entry(asset_id, "CHT", name, file_rel, slide_idx,
                   tags=tags, params=params, bindings=bindings,
                   editable=["data", "series-color", "title", "legend"],
                   recommended_use=recommended)


def vals(n, base, step, spread, seed):
    r = random.Random(seed)
    out = []; v = float(base)
    for _ in range(n):
        out.append(int(max(1, v + r.randint(-spread, spread))))
        v += step
    return tuple(out)


# ---------------- 스펙 스윕 ----------------
# family -> (target, [(ctype, nseries), ...], (cat_lo, cat_hi), pool_keys, series_style)
FAMILY_PLAN = [
    ("column", 50,
     [("COLUMN_CLUSTERED", 1), ("COLUMN_CLUSTERED", 2), ("COLUMN_CLUSTERED", 3),
      ("COLUMN_STACKED", 3), ("COLUMN_STACKED_100", 3), ("COLUMN_STACKED", 2)],
     (3, 7), ["genre", "region", "country", "channel", "program", "capability", "platform"]),
    ("bar", 30,
     [("BAR_CLUSTERED", 1), ("BAR_CLUSTERED", 2), ("BAR_CLUSTERED", 3),
      ("BAR_STACKED", 3), ("BAR_STACKED", 2)],
     (3, 7), ["country", "genre", "program", "platform", "channel", "region"]),
    ("line", 30,
     [("LINE", 1), ("LINE", 2), ("LINE_MARKERS", 2), ("LINE_MARKERS", 3), ("LINE", 3)],
     (4, 7), ["year", "quarter"]),
    ("area", 25,
     [("AREA", 1), ("AREA_STACKED", 2), ("AREA_STACKED", 3), ("AREA", 2)],
     (4, 7), ["year", "quarter"]),
    ("pie", 30,
     [("PIE", 1), ("DOUGHNUT", 1)],
     (3, 6), ["genre", "region", "channel", "program", "platform", "stage", "capability"]),
    ("radar", 25,
     [("RADAR_MARKERS", 2), ("RADAR", 2), ("RADAR_MARKERS", 3), ("RADAR", 3)],
     (4, 7), ["capability", "stage"]),
]


def build_series(family, ns, nc, gseed, metric):
    """계열 리스트 생성. 단일계열은 metric, 다계열은 SNS 세트에서 슬라이스."""
    if ns == 1:
        unit = MUNIT.get(metric, "")
        nm = "%s(%s)" % (metric, unit) if unit else metric
        return [{"name": nm, "values": vals(nc, 800 + (gseed % 5) * 300, 180, 260, gseed)}]
    style = SNS_KEYS[gseed % len(SNS_KEYS)]
    names = SNS[style][:ns]
    out = []
    for k, nm in enumerate(names):
        base = 400 + k * 260 + (gseed % 4) * 120
        out.append({"name": nm, "values": vals(nc, base, 130, 190, gseed * 7 + k)})
    return out


def build_specs():
    specs = []
    gseed = 0
    for family, target, variants, (lo, hi), pkeys in FAMILY_PLAN:
        made = 0
        i = 0
        while made < target:
            ctype, ns = variants[i % len(variants)]
            pool_key = pkeys[i % len(pkeys)]
            pool = POOLS[pool_key]
            nc = lo + (i % (hi - lo + 1))
            nc = max(3, min(nc, len(pool)))
            cats = pool[:nc]
            metric = METRICS[gseed % len(METRICS)]
            series = build_series(family, ns, nc, gseed, metric)
            subj = SUBJECT[pool_key]
            suf = SUFFIX[gseed % len(SUFFIX)]
            title = ("%s %s %s" % (subj, metric, suf)).strip()
            tags = ["차트", family, ctype.lower(), "%d계열" % ns, "%d항목" % nc, subj]
            params = {"type": ctype, "series": ns, "categories": nc,
                      "palette_offset": gseed % len(PAL), "pool": pool_key}
            rec = REC.get(family, ["데이터 시각화"])
            specs.append(dict(family=family, ctype=ctype, ns=ns, nc=nc, cats=cats,
                              series=series, title=title, tags=tags, params=params,
                              rec=rec, pal_off=gseed % len(PAL)))
            made += 1; i += 1; gseed += 1
    return specs


REC = {
    "column": ["실적 비교", "항목별 성과", "연도 비교"],
    "bar": ["순위 비교", "국가별 실적", "랭킹"],
    "line": ["시계열 추세", "성장 추이", "추세 비교"],
    "area": ["누적 추세", "총량 변화", "성장 곡선"],
    "pie": ["구성비", "점유율", "비중"],
    "radar": ["역량 진단", "다축 비교", "gap 분석"],
}

FAMILY_KO = {"column": "세로막대", "bar": "가로막대", "line": "꺾은선",
             "area": "영역", "pie": "원형/도넛", "radar": "레이더"}


def main():
    specs = build_specs()
    assert len(specs) == 190, "spec count=%d (기대 190)" % len(specs)

    entries = []
    files_written = []
    fallbacks = []
    n = 11  # CHT-011 시작

    # family 별로 ≤20 슬라이드 분산
    idx = 0
    while idx < len(specs):
        family = specs[idx]["family"]
        # 같은 family 연속 구간 수집
        group = []
        while idx < len(specs) and specs[idx]["family"] == family:
            group.append(specs[idx]); idx += 1
        # 20개씩 파일 분할
        for fi, start in enumerate(range(0, len(group), 20), start=1):
            chunk = group[start:start + 20]
            frel = "decks/08_charts/CHT_bulk_%s_v%d.pptx" % (family, fi)
            prs = c.new_deck()
            for slide_idx, sp in enumerate(chunk, start=1):
                aid = "CHT-%03d" % n
                s = c.blank_slide(prs); c.id_caption(s, aid)
                ctype = sp["ctype"]; note = None
                legend = not (ctype in PIE_TYPES)
                legend_pos = XL_LEGEND_POSITION.RIGHT if ctype in PIE_TYPES else XL_LEGEND_POSITION.BOTTOM
                show_legend = True
                lbl = ctype not in ("RADAR", "RADAR_MARKERS") and (sp["nc"] * sp["ns"] <= 21)
                use_type = ctype
                try:
                    make_chart(s, aid, ctype, sp["title"], sp["cats"], sp["series"],
                               pal_off=sp["pal_off"], legend_pos=legend_pos,
                               show_legend=show_legend, labels=lbl)
                except Exception as e:
                    # RADAR 등 실패 → COLUMN_CLUSTERED 대체
                    for sh in list(s.shapes):
                        if getattr(sh, "has_chart", False):
                            sh._element.getparent().remove(sh._element)
                    use_type = "COLUMN_CLUSTERED"
                    note = "%s 생성 실패(%s) → COLUMN_CLUSTERED 대체" % (ctype, type(e).__name__)
                    fallbacks.append((aid, ctype, type(e).__name__))
                    make_chart(s, aid, use_type, sp["title"], sp["cats"], sp["series"],
                               pal_off=sp["pal_off"])
                name = "%s %s (%d계열·%d항목)" % (FAMILY_KO[family], use_type, sp["ns"], sp["nc"])
                params = dict(sp["params"]); params["type"] = use_type
                entries.append(meta(aid, name, frel, slide_idx, use_type,
                                    sp["cats"], sp["series"], sp["tags"], params,
                                    sp["rec"], fallback_note=note))
                n += 1
            out = c.save_deck(prs, frel)
            files_written.append((out, len(chunk)))

    frag = c.write_fragment("CHT_bulk", entries)

    print("ENTRIES:", len(entries))
    print("REACHED:", "CHT-%03d" % (n - 1))
    print("FRAG:", frag)
    for f, cnt in files_written:
        print("FILE:", f, "slides=%d" % cnt)
    print("FALLBACKS:", fallbacks if fallbacks else "none")


if __name__ == "__main__":
    main()
