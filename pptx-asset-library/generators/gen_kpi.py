# -*- coding: utf-8 -*-
"""
KPI 카드/숫자 강조 카테고리 생성기 (KPI-001 ~ KPI-010)
공공기관 제안서 성과지표/실적 강조용. 네이티브 도형+텍스트만 사용.
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RR = MSO_SHAPE.ROUNDED_RECTANGLE
RADIUS = c.TOKENS["geom"]["round_radius"]

# ---------- 공용 헬퍼 ----------
def group_last(slide, n, asset_id):
    """슬라이드에 마지막으로 추가된 n개 도형을 그룹화하고 asset 이름 부여."""
    shapes = list(slide.shapes)
    targets = shapes[-n:]
    # python-pptx group: build grpSp via shapes.add_group_shape 은 없음 → 이름만 대표도형에.
    # 대표(첫 배경 도형)에 asset 이름을 부여한다.
    return targets

def stat_number(slide, x, y, w, num, unit="", num_color=None, unit_color=None,
                num_size=40, unit_size=16, align=PP_ALIGN.CENTER):
    """큰 숫자 + 단위 텍스트박스 (숫자 강조)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(num)
    c.set_kfont(r, c.FONT_H, num_size, True, num_color or c.role("accent_primary"))
    if unit:
        ru = p.add_run(); ru.text = " " + unit
        c.set_kfont(ru, c.FONT_H, unit_size, True, unit_color or c.role("muted_text"))
    return tb


# ============================================================
# FILE 1: KPI_cards_v1.pptx  (카드 세트류)
# ============================================================
prs1 = c.new_deck()
entries = []

CARD3 = [
    {"value": "1,000", "unit": "건", "label": "검증완료 기업정보"},
    {"value": "40",    "unit": "개사", "label": "참여기업"},
    {"value": "42",    "unit": "만$",  "label": "수출상담액"},
]

# ---- KPI-001: 네이비 채움 3카드 ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-001")
c.add_text(s, 0.15, 0.5, 6, 0.4, "성과지표 (네이비 채움 카드)", size=12, bold=True,
           color=c.role("muted_text"))
n0 = len(s.shapes._spTree)
gap = 0.4; cw = (13.333 - 1.0 - gap*2) / 3; x0 = 0.5; y = 1.5; ch = 2.6
for i, cd in enumerate(CARD3):
    x = x0 + i*(cw+gap)
    bg = c.add_box(s, x, y, cw, ch, fill=c.role("header_fill"), shape=RR)
    bg.adjustments[0] = RADIUS
    c.add_box(s, x+0.35, y+0.5, 0.6, 0.06, fill=c.role("accent_point"))
    stat_number(s, x, y+0.75, cw, cd["value"], cd["unit"],
                num_color=c.C["white"], unit_color=c.C["gray_300"], num_size=44)
    c.add_text(s, x, y+1.85, cw, 0.5, cd["label"], size=13, bold=False,
               color=c.C["gray_100"], align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[2], "KPI-001")  # 첫 카드 배경을 대표 앵커로

# ---- KPI-002: 흰바탕 + 네이비 테두리 3카드 ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-002")
c.add_text(s, 0.15, 0.5, 6, 0.4, "성과지표 (흰바탕 테두리 카드)", size=12, bold=True,
           color=c.role("muted_text"))
for i, cd in enumerate(CARD3):
    x = x0 + i*(cw+gap)
    bg = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"),
                   line=c.role("header_fill"), line_w=1.5, shape=RR)
    bg.adjustments[0] = RADIUS
    stat_number(s, x, y+0.55, cw, cd["value"], cd["unit"],
                num_color=c.role("accent_primary"), unit_color=c.role("muted_text"), num_size=44)
    c.connector(s, x+0.7, y+1.75, x+cw-0.7, y+1.75, color=c.role("border"), w=1.0)
    c.add_text(s, x, y+1.85, cw, 0.5, cd["label"], size=13,
               color=c.role("body_text"), align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[2], "KPI-002")

# ---- KPI-003: 상단 컬러바 카드 3개 ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-003")
c.add_text(s, 0.15, 0.5, 6, 0.4, "성과지표 (상단 컬러바 카드)", size=12, bold=True,
           color=c.role("muted_text"))
barcolors = [c.role("accent_primary"), c.role("accent_secondary"), c.role("accent_point")]
for i, cd in enumerate(CARD3):
    x = x0 + i*(cw+gap)
    bg = c.add_box(s, x, y, cw, ch, fill=c.role("panel_bg"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = RADIUS
    top = c.add_box(s, x, y, cw, 0.28, fill=barcolors[i], shape=RR)
    top.adjustments[0] = RADIUS
    stat_number(s, x, y+0.7, cw, cd["value"], cd["unit"],
                num_color=c.role("header_fill"), unit_color=c.role("muted_text"), num_size=44)
    c.add_text(s, x, y+1.9, cw, 0.5, cd["label"], size=13,
               color=c.role("body_text"), align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[2], "KPI-003")

# ---- KPI-004: 아이콘 원형 placeholder + 숫자 카드 3개 ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-004")
c.add_text(s, 0.15, 0.5, 6, 0.4, "성과지표 (아이콘 + 숫자 카드)", size=12, bold=True,
           color=c.role("muted_text"))
ICON3 = [
    {"value": "4.6",  "unit": "점", "label": "이용자 만족도"},
    {"value": "128",  "unit": "건", "label": "누적 컨설팅"},
    {"value": "95",   "unit": "%",  "label": "목표 달성률"},
]
icc = [c.role("accent_primary"), c.role("accent_secondary"), c.role("sub_header")]
ch4 = 2.9
for i, cd in enumerate(ICON3):
    x = x0 + i*(cw+gap)
    bg = c.add_box(s, x, y, cw, ch4, fill=c.role("row_base"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = RADIUS
    # 아이콘 자리 = 원형 placeholder
    icd = 0.9
    ic = c.add_box(s, x+cw/2-icd/2, y+0.35, icd, icd, fill=icc[i], shape=MSO_SHAPE.OVAL)
    c.set_shape_text(ic, "ICON", size=9, bold=True, color=c.C["white"])
    stat_number(s, x, y+1.45, cw, cd["value"], cd["unit"],
                num_color=c.role("header_fill"), unit_color=c.role("muted_text"), num_size=40)
    c.add_text(s, x, y+2.25, cw, 0.5, cd["label"], size=13,
               color=c.role("body_text"), align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[2], "KPI-004")

# ---- KPI-006: 가로 롱 카드 (라벨 좌 + 숫자 우) 3행 ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-006")
c.add_text(s, 0.15, 0.5, 8, 0.4, "성과지표 (가로 롱 카드)", size=12, bold=True,
           color=c.role("muted_text"))
LONG = [
    {"label": "검증완료 기업정보", "value": "1,000", "unit": "건"},
    {"label": "참여기업",         "value": "40",    "unit": "개사"},
    {"label": "수출상담액",       "value": "42",    "unit": "만$"},
]
lw = 12.0; lh = 1.1; lx = 0.65; ly0 = 1.6; lgap = 0.35
for i, cd in enumerate(LONG):
    yy = ly0 + i*(lh+lgap)
    bg = c.add_box(s, lx, yy, lw, lh, fill=c.role("panel_bg"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = 0.12
    # 좌측 컬러 액센트 바
    c.add_box(s, lx, yy, 0.14, lh, fill=c.role("accent_primary"))
    c.add_text(s, lx+0.45, yy, 6.5, lh, cd["label"], size=16, bold=True,
               color=c.role("body_text"), align=PP_ALIGN.LEFT)
    stat_number(s, lx+lw-4.3, yy, 4.0, cd["value"], cd["unit"],
                num_color=c.role("accent_primary"), unit_color=c.role("muted_text"),
                num_size=32, align=PP_ALIGN.RIGHT)
c.name_asset(s.shapes[2], "KPI-006")

# ---- KPI-009: 4분할 KPI 대시보드 카드 ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-009")
c.add_text(s, 0.15, 0.5, 8, 0.4, "성과지표 (4분할 대시보드)", size=12, bold=True,
           color=c.role("muted_text"))
QUAD = [
    {"value": "1,000", "unit": "건",   "label": "검증완료 기업정보"},
    {"value": "40",    "unit": "개사", "label": "참여기업"},
    {"value": "42",    "unit": "만$",  "label": "수출상담액"},
    {"value": "4.6",   "unit": "점",   "label": "이용자 만족도"},
]
qcol = [c.role("accent_primary"), c.role("accent_secondary"),
        c.role("accent_point"), c.role("sub_header")]
gx = 0.6; gy = 1.5; gw = 12.13; gh = 5.2; qgap = 0.3
qw = (gw - qgap) / 2; qh = (gh - qgap) / 2
# 외곽 프레임
frame = c.add_box(s, gx-0.05, gy-0.05, gw+0.1, gh+0.1, fill=c.role("row_base"),
                  line=c.role("border"), line_w=1.0, shape=RR)
frame.adjustments[0] = 0.04
for i, cd in enumerate(QUAD):
    r_, cc_ = divmod(i, 2)
    x = gx + cc_*(qw+qgap); yy = gy + r_*(qh+qgap)
    bg = c.add_box(s, x, yy, qw, qh, fill=c.role("panel_bg"), shape=RR)
    bg.adjustments[0] = 0.06
    c.add_box(s, x, yy, 0.14, qh, fill=qcol[i])
    stat_number(s, x+0.3, yy+0.35, qw-0.6, cd["value"], cd["unit"],
                num_color=c.role("header_fill"), unit_color=c.role("muted_text"),
                num_size=40, align=PP_ALIGN.LEFT)
    c.add_text(s, x+0.35, yy+qh-0.75, qw-0.6, 0.5, cd["label"], size=14,
               color=c.role("body_text"), align=PP_ALIGN.LEFT)
c.name_asset(s.shapes[2], "KPI-009")

out1 = c.save_deck(prs1, "decks/02_kpi/KPI_cards_v1.pptx")


# ============================================================
# FILE 2: KPI_metric-highlight_v1.pptx  (증감/게이지/히어로)
# ============================================================
prs2 = c.new_deck()

# ---- KPI-005: 증감 화살표(▲▼) 전년대비 % 카드 ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-005")
c.add_text(s, 0.15, 0.5, 8, 0.4, "성과지표 (전년대비 증감)", size=12, bold=True,
           color=c.role("muted_text"))
DELTA = [
    {"value": "1,000", "unit": "건",  "label": "검증완료 기업정보", "delta": "+18.5%", "up": True},
    {"value": "42",    "unit": "만$", "label": "수출상담액",       "delta": "+27.0%", "up": True},
    {"value": "3",     "unit": "건",  "label": "처리 지연",         "delta": "-12.4%", "up": False},
]
for i, cd in enumerate(DELTA):
    x = x0 + i*(cw+gap)
    bg = c.add_box(s, x, y, cw, 2.8, fill=c.role("row_base"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = RADIUS
    c.add_text(s, x, y+0.3, cw, 0.4, cd["label"], size=13, bold=True,
               color=c.role("muted_text"), align=PP_ALIGN.CENTER)
    stat_number(s, x, y+0.8, cw, cd["value"], cd["unit"],
                num_color=c.role("header_fill"), unit_color=c.role("muted_text"), num_size=40)
    # 증감 pill
    up = cd["up"]
    dcol = c.role("accent_secondary") if up else c.role("warn")
    arrow = "▲" if up else "▼"  # ▲ ▼
    pill = c.add_box(s, x+cw/2-1.0, y+1.95, 2.0, 0.5, fill=None,
                     line=dcol, line_w=1.2, shape=RR)
    pill.adjustments[0] = 0.5
    c.set_shape_text(pill, "%s %s  전년比" % (arrow, cd["delta"]), size=12, bold=True, color=dcol)
c.name_asset(s.shapes[2], "KPI-005")

# ---- KPI-007: 그룹 라벨 + 숫자 묶음 (+기호 연결) ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-007")
c.add_text(s, 0.15, 0.5, 10, 0.4, "핵심 가치 지표 (라벨 + 숫자 묶음)", size=12, bold=True,
           color=c.role("muted_text"))
GRP = [
    {"label": "신뢰성", "value": "1,000", "unit": "건"},
    {"label": "확장성", "value": "10",    "unit": "종"},
    {"label": "활용성", "value": "20",    "unit": "개"},
]
gcw = 3.4; ggap = 0.9; gtot = gcw*3 + ggap*2
gx0 = (13.333 - gtot)/2; gy = 2.4; gch = 2.4
for i, cd in enumerate(GRP):
    x = gx0 + i*(gcw+ggap)
    bg = c.add_box(s, x, gy, gcw, gch, fill=c.role("panel_bg"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = 0.1
    # 그룹 라벨 배지
    badge = c.add_box(s, x+gcw/2-0.9, gy+0.3, 1.8, 0.5, fill=c.role("header_fill"), shape=RR)
    badge.adjustments[0] = 0.5
    c.set_shape_text(badge, cd["label"], size=14, bold=True, color=c.C["white"])
    stat_number(s, x, gy+1.05, gcw, cd["value"], cd["unit"],
                num_color=c.role("accent_primary"), unit_color=c.role("muted_text"), num_size=42)
    # + 연결 기호
    if i < 2:
        c.add_text(s, x+gcw+ggap/2-0.35, gy+0.6, 0.7, gch-1.0, "+", size=40, bold=True,
                   color=c.role("accent_point"), align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[2], "KPI-007")

# ---- KPI-008: 도넛/원형 게이지 느낌 비율 강조 (네이티브 도형) ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-008")
c.add_text(s, 0.15, 0.5, 10, 0.4, "달성률 게이지 (도형 기반 비율)", size=12, bold=True,
           color=c.role("muted_text"))
GAUGE = [
    {"pct": 95, "label": "목표 달성률", "col": "accent_primary"},
    {"pct": 78, "label": "만족도 지수", "col": "accent_secondary"},
    {"pct": 64, "label": "수출 전환율", "col": "sub_header"},
]
gd = 2.6; gy = 1.9; gtot2 = gd*3 + 0.8*2; gx0 = (13.333-gtot2)/2
for i, cd in enumerate(GAUGE):
    x = gx0 + i*(gd+0.8)
    frac = cd["pct"]/100.0
    accent = c.role(cd["col"])
    # 뒤 링(회색 전체 원)
    c.add_box(s, x, gy, gd, gd, fill=c.role("border"), shape=MSO_SHAPE.OVAL)
    # 채움 PIE (top -90도 시작)
    start = -90.0; end = -90.0 + 360.0*frac
    pie = c.add_box(s, x, gy, gd, gd, fill=accent, shape=MSO_SHAPE.PIE)
    pie.adjustments[0] = start
    pie.adjustments[1] = end
    # 중앙 흰 원 → 도넛 효과
    inr = gd*0.58
    c.add_box(s, x+(gd-inr)/2, gy+(gd-inr)/2, inr, inr, fill=c.role("row_base"),
              shape=MSO_SHAPE.OVAL)
    # 중앙 % 텍스트
    stat_number(s, x, gy+gd/2-0.45, gd, str(cd["pct"]), "%",
                num_color=c.role("header_fill"), unit_color=accent, num_size=34)
    c.add_text(s, x, gy+gd+0.1, gd, 0.5, cd["label"], size=14, bold=True,
               color=c.role("body_text"), align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[2], "KPI-008")

# ---- KPI-010: 강조 단일 대형 숫자 (히어로 지표) ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-010")
hero = c.add_box(s, 0.6, 1.4, 12.13, 5.3, fill=c.role("header_fill"), shape=RR)
hero.adjustments[0] = 0.04
c.add_box(s, 6.0, 2.0, 1.33, 0.09, fill=c.role("accent_point"))
c.add_text(s, 0.6, 2.3, 12.13, 0.6, "누적 수출상담액", size=20, bold=True,
           color=c.C["gray_100"], align=PP_ALIGN.CENTER)
# 히어로 대형 숫자
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.13), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "42"
c.set_kfont(r, c.FONT_H, 120, True, c.C["white"])
ru = p.add_run(); ru.text = " 만$"
c.set_kfont(ru, c.FONT_H, 44, True, c.role("accent_point"))
c.add_text(s, 0.6, 5.3, 12.13, 0.6, "2026년 사업 목표 대비 108% 조기 달성", size=15,
           color=c.C["gray_300"], align=PP_ALIGN.CENTER)
c.name_asset(s.shapes[1], "KPI-010")

out2 = c.save_deck(prs2, "decks/02_kpi/KPI_metric-highlight_v1.pptx")

f1 = "decks/02_kpi/KPI_cards_v1.pptx"
f2 = "decks/02_kpi/KPI_metric-highlight_v1.pptx"

# ============================================================
# 메타데이터
# ============================================================
def E(aid, name, file_rel, slide_idx, cards, count, tags, editable, params, ruse):
    return c.entry(
        asset_id=aid, category="KPI", name=name, file_rel=file_rel, slide_idx=slide_idx,
        tags=tags, params=params,
        bindings={"cards": cards, "count": count},
        editable=editable, recommended_use=ruse)

entries = [
    E("KPI-001", "네이비 채움 3카드", f1, 1,
      CARD3, 3,
      ["KPI","성과지표","3카드","네이비","숫자강조"],
      ["value","unit","label","card-color"],
      {"cards": 3, "style": "filled-navy", "num_size": 44},
      ["성과목표","실적강조","기대효과"]),
    E("KPI-002", "흰바탕 테두리 3카드", f1, 2,
      CARD3, 3,
      ["KPI","성과지표","3카드","화이트","테두리"],
      ["value","unit","label","border-color"],
      {"cards": 3, "style": "outline-navy", "num_size": 44},
      ["성과목표","실적강조","기대효과"]),
    E("KPI-003", "상단 컬러바 3카드", f1, 3,
      CARD3, 3,
      ["KPI","성과지표","3카드","컬러바","포인트컬러"],
      ["value","unit","label","bar-color"],
      {"cards": 3, "style": "topbar", "num_size": 44},
      ["성과목표","실적강조","기대효과"]),
    E("KPI-004", "아이콘+숫자 3카드", f1, 4,
      [{"value": d["value"], "unit": d["unit"], "label": d["label"]} for d in ICON3], 3,
      ["KPI","성과지표","3카드","아이콘","원형플레이스홀더"],
      ["value","unit","label","icon","icon-color"],
      {"cards": 3, "style": "icon-top", "icon_placeholder": True, "num_size": 40},
      ["성과목표","실적강조","서비스지표"]),
    E("KPI-006", "가로 롱 카드 3행", f1, 5,
      [{"value": d["value"], "unit": d["unit"], "label": d["label"]} for d in LONG], 3,
      ["KPI","성과지표","가로카드","라벨좌숫자우","리스트"],
      ["value","unit","label","accent-color"],
      {"rows": 3, "style": "horizontal-long", "num_size": 32},
      ["성과목표","실적강조","기대효과"]),
    E("KPI-009", "4분할 KPI 대시보드", f1, 6,
      [{"value": d["value"], "unit": d["unit"], "label": d["label"]} for d in QUAD], 4,
      ["KPI","성과지표","4분할","대시보드","종합지표"],
      ["value","unit","label","quadrant-color"],
      {"cards": 4, "style": "quad-dashboard", "num_size": 40},
      ["성과목표","실적강조","종합성과"]),
    E("KPI-005", "전년대비 증감 3카드", f2, 1,
      [{"value": d["value"], "unit": d["unit"], "label": d["label"], "delta": d["delta"], "up": d["up"]} for d in DELTA], 3,
      ["KPI","성과지표","3카드","증감화살표","전년대비"],
      ["value","unit","label","delta","direction","delta-color"],
      {"cards": 3, "style": "delta-arrow", "num_size": 40},
      ["성과목표","실적강조","증감비교"]),
    E("KPI-007", "그룹 라벨+숫자 묶음(+연결)", f2, 2,
      [{"value": d["value"], "unit": d["unit"], "label": d["label"]} for d in GRP], 3,
      ["KPI","성과지표","핵심가치","묶음","플러스연결"],
      ["value","unit","label","badge-color"],
      {"cards": 3, "style": "grouped-plus", "connector": "+", "num_size": 42},
      ["핵심가치","기대효과","차별성"]),
    E("KPI-008", "도넛 게이지 비율 3종", f2, 3,
      [{"value": str(d["pct"]), "unit": "%", "label": d["label"]} for d in GAUGE], 3,
      ["KPI","성과지표","게이지","도넛","비율강조"],
      ["value","label","gauge-color","pct"],
      {"gauges": 3, "style": "donut-gauge", "shape": "native-pie", "num_size": 34},
      ["달성률","목표대비","성과강조"]),
    E("KPI-010", "히어로 단일 대형 숫자", f2, 4,
      [{"value": "42", "unit": "만$", "label": "누적 수출상담액"}], 1,
      ["KPI","성과지표","히어로","대형숫자","단일강조"],
      ["value","unit","label","caption","hero-bg"],
      {"cards": 1, "style": "hero-single", "num_size": 120},
      ["대표성과","핵심지표","임팩트"]),
]

frag = c.write_fragment("KPI", entries)
print("OUT1:", out1)
print("OUT2:", out2)
print("FRAG:", frag)
print("ENTRIES:", len(entries))
