# -*- coding: utf-8 -*-
"""PRC v3 — 복합 분기/병합 프로세스 (외부연동 루프 + 회원유형 인지 + 부가 산출물 라벨).

PRC-019(분기/병합 흐름, 1->2갈래->합류, 4노드 단순형)보다 실전 시나리오에 더 가까운
"심화형" 분기/병합 패턴 1종 추가.

리비전 2 (플로우 구조 자체 수정 — 노드 배치가 아니라 로직 순서가 실제 서비스와
달랐다는 피드백 반영): 처음 만든 버전은 "필수정보입력/역량진단 이후 자리에 있던
판단 다이아몬드"를 "회원유형 분기점"으로 잘못 해석했으나, 실제 서비스는
  1) 로그인 직후 회원유형(기업/개인)이 "인지"만 되고(실제 분기 아님, 단일 트렁크 유지)
  2) 필수정보입력/역량진단(자가진단 설문)은 기업·개인 공통 단일 경로
  3) "결과 산출" 단계에서만 실제로 처리 로직이 물리적으로 갈라짐
     (기업: 나이스(NICE·DNB) API 연계 -> 신용등급 조회 -> 자가진단 점수 결합,
      개인: 신용등급 결합 없이 자가진단 점수만)
  4) 두 경로가 다시 합류해 "완료" 단일 노드로 종료
  순서다. 이를 반영해 다이아몬드(실제 분기/합류 도형)는 "결과 산출"과
  "결과 통합" 두 곳에만 쓰고, "회원 구분 인지"는 일반 사각형(단일 트렁크,
  화살표 하나만 들고 남) + 보라색으로만 구분해 "인지이지 분기가 아님"을
  시각적으로 못박는다.

파일: decks/03_process/PRC_branch-merge-loop_v3.pptx (PRC-207, 1슬라이드 1에셋)
다중 도형 -> group_asset 필수. 페이지 이미지·SmartArt 금지. 색은 c.role/c.C만 사용.
"""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

R = c.role
CC = c.C
DECK = 'decks/03_process'

# 쿨블루 '같은 계열 명도' 램프 (dark -> light), 전부 c.C 원색 (raw 헥스 금지)
RAMP = [CC['navy_900'], CC['navy_800'], CC['navy_600'], CC['blue_500'], CC['cyan_500']]


def tone(i):
    return RAMP[min(i, len(RAMP) - 1)]


def arrowize(cn, head=False, tail=True):
    """커넥터 라인에 삼각 화살촉 부여."""
    ln = cn.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def dashize(cn, val='dash'):
    """커넥터 라인을 점선으로 (부가/보조 흐름 구분용)."""
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': val}))
    return cn


# =========================================================
# PRC-207 복합 분기/병합 흐름 (세로형 — 로그인부터 역량진단까지는 완전 공통
# 단일 트렁크, "결과 산출" 단계에서만 실제로 갈라졌다가 "완료"로 합류)
# =========================================================
def build_branch_merge_loop():
    prs = c.new_deck(); s = c.blank_slide(prs); c.id_caption(s, 'PRC-207'); S = []
    CX = 6.667  # 트렁크(중앙 축) 중심 x — 슬라이드 중앙
    biz_cx, ind_cx = CX - 2.1, CX + 2.1   # 좌: 기업 경로 열 / 우: 개인 경로 열

    def center_box(cx, y, w, h, text, tn, shape=MSO_SHAPE.RECTANGLE, size=12, bold=True,
                   text_color=None):
        sp = c.add_box(s, cx - w / 2, y, w, h, fill=tone(tn), line=None, shape=shape)
        c.set_shape_text(sp, text, size=size, bold=bold, color=text_color or R('header_text'))
        S.append(sp)
        return sp

    def title_sub_box(cx, y, w, h, title, subtitle, fill_color, shape=MSO_SHAPE.RECTANGLE,
                      title_size=12, sub_size=8, text_color=None):
        sp = c.add_box(s, cx - w / 2, y, w, h, fill=fill_color, line=None, shape=shape)
        tf = sp.text_frame; tf.word_wrap = True
        tc = text_color or R('header_text')
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = title
        c.set_kfont(r0, size=title_size, bold=True, color=tc)
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = subtitle
        c.set_kfont(r1, size=sub_size, bold=False, color=tc)
        S.append(sp)
        return sp

    def arrow(x1, y1, x2, y2, w=1.6, color=None):
        cn = c.connector(s, x1, y1, x2, y2, color=color or R('accent_primary'), w=w)
        arrowize(cn)
        S.append(cn)
        return cn

    # ---- 공통 트렁크 1: 로그인 ----
    login_y, login_h = 0.70, 0.42
    login = center_box(CX, login_y, 1.8, login_h, '로그인', 0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    # ---- 공통 트렁크 2: 회원 구분 인지 — 실제 분기 아님(단일 트렁크 유지).
    #      다이아몬드가 아닌 사각형 + 보라색으로만 구분해 "인지"임을 표시 ----
    recog_y, recog_h = login_y + login_h + 0.12, 0.48
    recog = title_sub_box(CX, recog_y, 2.6, recog_h, '회원 구분 인지',
                          '(기업회원 · 개인회원 자동 인식)', R('sub_header'),
                          title_size=11, sub_size=8)

    # ---- 공통 트렁크 3: 필수 정보 입력 (헤더밴드+불릿) + 좌우 곁가지 라벨 2개 ----
    info_y = recog_y + recog_h + 0.12
    info_w, info_h = 3.0, 0.80
    info_x = CX - info_w / 2
    head_h = 0.28
    info_head = c.add_box(s, info_x, info_y, info_w, head_h, fill=tone(1), line=None)
    c.set_shape_text(info_head, '필수 정보 입력', size=11.5, bold=True, color=R('header_text'))
    S.append(info_head)
    info_body = c.add_box(s, info_x, info_y + head_h, info_w, info_h - head_h,
                          fill=R('panel_bg'), line=R('border'), line_w=0.75)
    tf = info_body.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_top = Pt(2)
    bullets = ['관심 콘텐츠 및 국가(각 3개)', '최근 3년간 수출액', '주요 수출국 Top3']
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = '· ' + b
        c.set_kfont(run, size=8, bold=False, color=R('muted_text'))
    S.append(info_body)

    # 곁가지 라벨 — 필수정보입력 카드 좌/우에 작은 태그로 대칭 부착(폭을 넓히지
    # 않도록 카드 아래로 쌓지 않고 좌우 각 1개씩만 배치, 점선으로 본 흐름과 구분)
    info_mid_y = info_y + info_h / 2
    note_w, note_h = 1.5, 0.4
    note_gap = 0.15
    note1_x = info_x - note_gap - note_w
    note1 = c.add_box(s, note1_x, info_mid_y - note_h / 2, note_w, note_h, fill=R('panel_bg'),
                      line=R('border'), line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.set_shape_text(note1, '사업자등록번호와 연동', size=8.5, bold=False, color=R('muted_text'))
    S.append(note1)
    note2_x = info_x + info_w + note_gap
    note2 = c.add_box(s, note2_x, info_mid_y - note_h / 2, note_w, note_h, fill=R('panel_bg'),
                      line=R('border'), line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.set_shape_text(note2, '국가 및 콘텐츠 보고서 압축 다운', size=8.5, bold=False, color=R('muted_text'))
    S.append(note2)
    ln1 = dashize(c.connector(s, info_x, info_mid_y, note1_x + note_w, info_mid_y,
                              color=R('border'), w=1.0))
    S.append(ln1)
    ln2 = dashize(c.connector(s, info_x + info_w, info_mid_y, note2_x, info_mid_y,
                              color=R('border'), w=1.0))
    S.append(ln2)

    # ---- 공통 트렁크 4: 역량 진단(자가진단 설문) ----
    sc_y, sc_h = info_y + info_h + 0.12, 0.40
    selfcheck = center_box(CX, sc_y, 2.4, sc_h, '역량 진단\n(자가진단 설문)', 1, size=10.5)

    # ---- 실제 분기점: 결과 산출(다이아몬드) — 이 슬라이드에서 유일하게 물리적으로
    #      갈라지는 지점. 로그인~역량진단까지는 전부 공통 단일 경로였음을 강조 ----
    dec_y = sc_y + sc_h + 0.14
    dec_w, dec_h = 1.9, 0.55
    decision = c.add_box(s, CX - dec_w / 2, dec_y, dec_w, dec_h, fill=R('accent_point'),
                         line=None, shape=MSO_SHAPE.DIAMOND)
    c.set_shape_text(decision, '결과 산출', size=10.5, bold=True, color=R('header_text'))
    S.append(decision)
    dec_left = (CX - dec_w / 2, dec_y + dec_h / 2)
    dec_right = (CX + dec_w / 2, dec_y + dec_h / 2)

    # ---- 기업 경로(결과산출 이후에만 존재): 나이스(NICE·DNB) 연계 -> 기업 신용등급
    #      조회 -> 자가진단 점수 결합 -> 기업 회원 결과. 전부 기업 열 아래로 수직 적재 ----
    nice_y, nice_h = dec_y + dec_h + 0.14, 0.38
    nice = center_box(biz_cx, nice_y, 1.95, nice_h, '나이스(NICE·DNB) 연계', 0,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE, size=10, text_color=R('header_text'))
    nice.fill.solid(); nice.fill.fore_color.rgb = R('accent_secondary')

    credit_y, credit_h = nice_y + nice_h + 0.10, 0.38
    credit = center_box(biz_cx, credit_y, 1.95, credit_h, '기업 신용등급 조회', 3, size=10)

    combine_y, combine_h = credit_y + credit_h + 0.10, 0.38
    combine = center_box(biz_cx, combine_y, 1.95, combine_h, '자가진단 점수 결합', 3, size=10)

    # ---- 개인/기업 결과 (동일 행에서 정렬 — 합류 전 마지막 단계) ----
    res_y, res_h = combine_y + combine_h + 0.12, 0.48
    biz_res = center_box(biz_cx, res_y, 1.95, res_h, '기업 회원 결과', 2, size=10.5)
    ind_res = center_box(ind_cx, res_y, 1.95, res_h, '개인 회원 결과', 2, size=10.5)

    # ---- 합류: 결과 통합(다이아몬드) — 결과산출에서 갈라진 두 경로가 여기서 재합류 ----
    merge_y = res_y + res_h + 0.14
    merge_w, merge_h = 1.7, 0.50
    merge = c.add_box(s, CX - merge_w / 2, merge_y, merge_w, merge_h, fill=R('accent_point'),
                      line=None, shape=MSO_SHAPE.DIAMOND)
    c.set_shape_text(merge, '결과 통합', size=10, bold=True, color=R('header_text'))
    S.append(merge)
    merge_left = (CX - merge_w / 2, merge_y + merge_h / 2)
    merge_right = (CX + merge_w / 2, merge_y + merge_h / 2)

    # ---- 최종: 완료 (지원 사업 추천 등은 완료 화면 내 하위 요소이므로 부제로만 표기) ----
    final_y, final_h = merge_y + merge_h + 0.14, 0.55
    final = title_sub_box(CX, final_y, 2.6, final_h, '완료', '(지원 사업 추천 등 결과 제공)',
                          tone(0), shape=MSO_SHAPE.ROUNDED_RECTANGLE, title_size=14, sub_size=8.5)

    # =========================================================
    # 커넥터(화살표) — 로그인부터 결과산출 다이아몬드까지는 전부 수직 단일선.
    # 실제로 두 갈래로 나뉘는 화살표는 결과산출 다이아몬드 좌/우 꼭짓점에서
    # 뻗어나가는 것 단 한 쌍뿐이다.
    # =========================================================
    arrow(CX, login_y + login_h, CX, recog_y)                      # 로그인 -> 회원구분인지
    arrow(CX, recog_y + recog_h, CX, info_y)                       # 회원구분인지 -> 필수정보입력
    arrow(CX, info_y + info_h, CX, sc_y)                           # 필수정보입력 -> 역량진단
    arrow(CX, sc_y + sc_h, CX, dec_y)                              # 역량진단 -> 결과산출(분기점)

    # 결과산출에서 처음이자 유일하게 두 갈래로 갈라짐
    arrow(dec_left[0], dec_left[1], biz_cx, nice_y, color=R('accent_secondary'))
    arrow(dec_right[0], dec_right[1], ind_cx, res_y, color=R('accent_secondary'))
    biz_lbl = c.add_text(s, biz_cx - 1.05, (dec_y + dec_h / 2 + nice_y) / 2 - 0.14, 1.0, 0.28,
                         'API 연계', size=8, bold=True, color=R('accent_secondary'),
                         align=PP_ALIGN.RIGHT)
    S.append(biz_lbl)
    ind_lbl = c.add_text(s, ind_cx + 0.05, (dec_y + dec_h / 2 + res_y) / 2 - 0.14, 1.7, 0.28,
                         '자가진단 점수만', size=8, bold=True, color=R('accent_secondary'),
                         align=PP_ALIGN.LEFT)
    S.append(ind_lbl)

    # 기업 경로 내부 — 전부 수직 단일선(공통 경로가 아니라 기업 열 전용)
    arrow(biz_cx, nice_y + nice_h, biz_cx, credit_y, w=1.4)         # 나이스연계 -> 신용등급조회
    arrow(biz_cx, credit_y + credit_h, biz_cx, combine_y, w=1.4)    # 신용등급조회 -> 점수결합
    arrow(biz_cx, combine_y + combine_h, biz_cx, res_y, w=1.4)      # 점수결합 -> 기업결과

    # 결과 -> 합류(결과통합) — 갈라졌던 두 경로가 다시 모이는 지점
    arrow(biz_cx, res_y + res_h, merge_left[0], merge_left[1], color=R('accent_secondary'))
    arrow(ind_cx, res_y + res_h, merge_right[0], merge_right[1], color=R('accent_secondary'))
    # 합류 -> 완료(공통 단일 경로로 복귀)
    arrow(CX, merge_y + merge_h, CX, final_y)

    c.group_asset(s, S, 'PRC-207')
    return c.save_deck(prs, DECK + '/PRC_branch-merge-loop_v3.pptx')


def build_manifest_entries():
    E = c.entry
    entries = [
        E('PRC-207', 'PRC', '복합 분기/병합 흐름 (세로형 · 공통트렁크 + 결과산출 단일분기점 + 외부연동)',
          DECK + '/PRC_branch-merge-loop_v3.pptx', 1,
          ['프로세스', '분기', '병합', '루프백', '외부연동', '단일분기점', '판단다이아몬드', '세로형', '수직흐름'],
          {'nodes': 14, 'style': 'branch-merge-loop-vertical', 'orientation': 'vertical',
           'branches': 2, 'asymmetric': True, 'branch_point_count': 1,
           'note': '분기(다이아몬드)는 결과산출/결과통합 두 곳뿐 — 회원구분 인지는 사각형 단일 트렁크로 처리, 실제 분기 아님'},
          {
              'start': '로그인',
              'recognition_title': '회원 구분 인지',
              'recognition_sub': '(기업회원 · 개인회원 자동 인식)',
              'info_title': '필수 정보 입력',
              'info_bullets': ['관심 콘텐츠 및 국가(각 3개)', '최근 3년간 수출액', '주요 수출국 Top3'],
              'info_side_note_1': '사업자등록번호와 연동',
              'info_side_note_2': '국가 및 콘텐츠 보고서 압축 다운',
              'selfcheck': '역량 진단(자가진단 설문)',
              'decision_label': '결과 산출',
              'branch_a_external': '나이스(NICE·DNB) 연계',
              'branch_a_external_step': '기업 신용등급 조회',
              'branch_a_combine': '자가진단 점수 결합',
              'branch_a_result': '기업 회원 결과',
              'branch_b_result': '개인 회원 결과',
              'merge_label': '결과 통합',
              'final_title': '완료',
              'final_sub': '(지원 사업 추천 등 결과 제공)',
          },
          ['start', 'recognition_title', 'recognition_sub', 'info_title', 'info_bullets[3]',
           'info_side_note_1', 'info_side_note_2', 'selfcheck', 'decision_label',
           'branch_a_external', 'branch_a_external_step', 'branch_a_combine', 'branch_a_result',
           'branch_b_result', 'merge_label', 'final_title', 'final_sub'],
          recommended_use=['공통 절차 후 단일 분기점만 있는 진단·심사 플로우',
                            '회원유형 인지(비분기) + 결과산출 단계 실제분기',
                            '외부기관(신용평가 등) 연동 프로세스', '비대칭 분기(한쪽만 외부호출)']),
    ]
    return entries


if __name__ == '__main__':
    out = build_branch_merge_loop()
    print('SAVED', out)
    entries = build_manifest_entries()
    frag = c.write_fragment('PRC', entries)
    print('FRAGMENT', frag)
