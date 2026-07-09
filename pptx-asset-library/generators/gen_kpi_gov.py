# -*- coding: utf-8 -*-
"""KPI gov 트랙 — 정부보고서(정답지) 캔버스·팔레트 전용 KPI 콜아웃 6종.

정답지 재현 목표: GROUP 조합 수작업 KPI콜아웃(slide40,41)·카드형 반복 레이아웃 커버.
네이티브 SmartArt/차트 절대 금지 — 전부 도형(autoshape)+텍스트 GROUP 조합.
숫자강조는 gov_theme.font.number_emphasis(페이퍼로지/Paperlogy) 전용 사용.

캔버스: gov 전용 11.93 x 8.50in = EMU 10905360 x 7772400 — PRC gov와 동일 EMU.
색·폰트: design-tokens.json["gov_theme"] 참조만(매직 헥스 금지).

ID: KPI-201 ~ KPI-206 (기존 KPI-001~200 다음 free 대역, manifest.json 확인 완료).
파일: decks/02_kpi/KPI_gov_callout_v1.pptx (슬라이드 1~6 = 에셋 1~6개, 1슬라이드=1에셋)
매니페스트 조각: _incoming/manifest_KPI_gov.json (기존 manifest_KPI*.json 미변경)
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DECK = 'decks/02_kpi'
FILE_REL = DECK + '/KPI_gov_callout_v1.pptx'

RR = MSO_SHAPE.ROUNDED_RECTANGLE
OVAL = MSO_SHAPE.OVAL
PIE = MSO_SHAPE.PIE

# ---------------------------------------------------------------------------
# gov 캔버스 (PRC gov와 동일 EMU — additive, 기존 표준 13.333x7.5 캔버스와 무관)
# ---------------------------------------------------------------------------
GOV_W = Emu(10905360)
GOV_H = Emu(7772400)
SW = 11.93
SH = 8.50


def new_deck_gov():
    prs = Presentation()
    prs.slide_width = GOV_W
    prs.slide_height = GOV_H
    return prs


def blank_slide_gov(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# gov 팔레트·폰트 (design-tokens.json["gov_theme"] 참조만)
# ---------------------------------------------------------------------------
GOV = c.TOKENS['gov_theme']
GC = {k: RGBColor.from_string(v) for k, v in GOV['color'].items()}


def R(name):
    return GC[GOV['role'][name]]


WHITE = GC['bg_primary']
F_BODY = GOV['font']['body']['typeface']
F_LABEL = GOV['font']['label']['typeface']
F_SUB = GOV['font']['subheading']['typeface']
F_CAP = GOV['font']['caption']['typeface']
F_NUM = GOV['font']['number_emphasis']['typeface']  # 페이퍼로지(Paperlogy) — 숫자강조 전용

RADIUS_PCT = GOV['geom']['round_radius_pct']
BORDER_PT = GOV['geom']['border_pt']

# 강조색 스윕 — gov_theme 도형용 색만(표헤더 전용 8F99AF/76829E 배제)
ACC = [GC['accent_navy'], GC['category_purple'], GC['line_accent'], GC['shape_fill_gray']]


def gcaption(slide, asset_id):
    c.add_text(slide, 0.12, 0.06, 3.6, 0.28, asset_id, size=9, bold=True,
               color=GC['shape_fill_gray'], align=PP_ALIGN.LEFT, font=F_CAP)


def gtext(slide, x, y, w, h, text, size=None, bold=False, color=None,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=None):
    return c.add_text(slide, x, y, w, h, text, size=size or GOV['size_pt']['body'],
                       bold=bold, color=color or R('body_text'), align=align,
                       anchor=anchor, font=font or F_BODY)


def numtext(slide, x, y, w, h, num, unit='', num_color=None, unit_color=None,
            num_size=40, unit_size=16, align=PP_ALIGN.CENTER):
    """gov 숫자강조 콜아웃 — 큰 숫자(Paperlogy)+단위(gov body)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(num)
    c.set_kfont(r, F_NUM, num_size, True, num_color or R('accent_primary'))
    if unit:
        ru = p.add_run(); ru.text = ' ' + unit
        c.set_kfont(ru, F_BODY, unit_size, True, unit_color or R('body_text'))
    return tb


# =============================================================================
# KPI-201 배지형 지표 (둥근 배지 안 숫자 4, Paperlogy)
# =============================================================================
def build_badge(prs):
    s = blank_slide_gov(prs); gcaption(s, 'KPI-201'); A = []
    BADGE = [
        {'value': '98%', 'label': '정시 집행률', 'col': 0},
        {'value': '4.8', 'label': '주민 만족도', 'col': 1},
        {'value': '120', 'label': '참여 기관수', 'col': 2},
        {'value': '32', 'label': '연계 지자체', 'col': 3},
    ]
    bd = 1.8; by = 2.5; bgap = 0.55
    tot = bd * 4 + bgap * 3; bx0 = (SW - tot) / 2
    for i, cd in enumerate(BADGE):
        x = bx0 + i * (bd + bgap); accent = ACC[cd['col']]
        ring = c.add_box(s, x - 0.09, by - 0.09, bd + 0.18, bd + 0.18, fill=None,
                         line=accent, line_w=BORDER_PT * 2, shape=OVAL)
        A.append(ring)
        badge = c.add_box(s, x, by, bd, bd, fill=accent, shape=OVAL)
        c.set_shape_text(badge, cd['value'], size=26, bold=True, color=WHITE, font=F_NUM)
        A.append(badge)
        A.append(gtext(s, x - 0.2, by + bd + 0.15, bd + 0.4, 0.5, cd['label'], size=13,
                       bold=True, color=R('body_text'), font=F_LABEL))
    c.group_asset(s, A, 'KPI-201')
    return s


# =============================================================================
# KPI-202 후프(ring/hoop) 프로그레스 3종 — 얇은 도넛 링 + % 중앙(Paperlogy)
# =============================================================================
def build_hoop(prs):
    s = blank_slide_gov(prs); gcaption(s, 'KPI-202'); A = []
    RINGS = [
        {'pct': 95, 'label': '사업 목표 달성률', 'col': 0},
        {'pct': 82, 'label': '주민 만족도', 'col': 1},
        {'pct': 68, 'label': '예산 집행률', 'col': 2},
    ]
    gd = 2.4; gy = 2.4; ggap = 0.8
    tot = gd * 3 + ggap * 2; gx0 = (SW - tot) / 2
    for i, cd in enumerate(RINGS):
        x = gx0 + i * (gd + ggap); frac = cd['pct'] / 100.0
        accent = ACC[cd['col']]
        A.append(c.add_box(s, x, gy, gd, gd, fill=R('panel_bg'), shape=OVAL))
        pie = c.add_box(s, x, gy, gd, gd, fill=accent, shape=PIE)
        pie.adjustments[0] = -90.0; pie.adjustments[1] = -90.0 + 360.0 * frac
        A.append(pie)
        inr = gd * 0.74
        A.append(c.add_box(s, x + (gd - inr) / 2, gy + (gd - inr) / 2, inr, inr,
                           fill=WHITE, shape=OVAL))
        A.append(numtext(s, x, gy + gd / 2 - 0.45, gd, 0.9, str(cd['pct']), '%',
                         num_color=R('body_text'), unit_color=accent, num_size=28, unit_size=13))
        A.append(gtext(s, x, gy + gd + 0.12, gd, 0.5, cd['label'], size=13, bold=True,
                       color=R('body_text'), font=F_LABEL))
    c.group_asset(s, A, 'KPI-202')
    return s


# =============================================================================
# KPI-203 hero — 단일 대형 숫자강조 콜아웃 (전면 히어로)
# =============================================================================
def build_hero(prs):
    s = blank_slide_gov(prs); gcaption(s, 'KPI-203'); A = []
    panel = c.add_box(s, 1.4, 1.7, SW - 2.8, 4.6, fill=R('panel_bg'), shape=RR)
    panel.adjustments[0] = RADIUS_PCT
    A.append(panel)
    A.append(gtext(s, 1.7, 2.05, SW - 3.4, 0.55, '연간 핵심 성과지표', size=GOV['size_pt']['subheading'],
                   bold=True, color=R('accent_primary'), align=PP_ALIGN.LEFT, font=F_SUB))
    A.append(numtext(s, 1.7, 2.75, SW - 3.4, 2.2, '12,500', '건',
                     num_color=R('accent_primary'), unit_color=R('body_text'),
                     num_size=96, unit_size=28))
    A.append(gtext(s, 1.7, 5.05, SW - 3.4, 0.55, '누적 민원 처리건수 — 전년 대비 28% 증가',
                   size=14, bold=True, color=R('body_text'), align=PP_ALIGN.LEFT, font=F_LABEL))
    bar_w = (SW - 3.4) * 0.86
    A.append(c.add_box(s, 1.7, 5.75, SW - 3.4, 0.14, fill=GC['bg_panel']))
    A.append(c.add_box(s, 1.7, 5.75, bar_w, 0.14, fill=R('line')))
    c.group_asset(s, A, 'KPI-203')
    return s


# =============================================================================
# KPI-204 quad — 2x2 그리드 4지표
# =============================================================================
def build_quad(prs):
    s = blank_slide_gov(prs); gcaption(s, 'KPI-204'); A = []
    QUAD = [
        {'value': '1,000', 'unit': '건', 'label': '검증완료 사업 DB', 'col': 0},
        {'value': '40', 'unit': '개소', 'label': '참여 지역기관', 'col': 1},
        {'value': '42', 'unit': '억', 'label': '연계 예산 규모', 'col': 2},
        {'value': '95', 'unit': '%', 'label': '사업 목표 달성률', 'col': 3},
    ]
    cw, ch = 5.15, 2.15
    gx, gy = 0.9, 1.55
    ggap = 0.35
    for i, cd in enumerate(QUAD):
        row, col = divmod(i, 2)
        x = gx + col * (cw + ggap); y = gy + row * (ch + ggap)
        accent = ACC[cd['col']]
        bg = c.add_box(s, x, y, cw, ch, fill=WHITE, line=GC['bg_panel'], line_w=BORDER_PT * 1.5, shape=RR)
        bg.adjustments[0] = RADIUS_PCT; A.append(bg)
        A.append(c.add_box(s, x, y, 0.16, ch, fill=accent))
        A.append(numtext(s, x + 0.45, y + 0.25, cw - 0.7, 1.1, cd['value'], cd['unit'],
                         num_color=accent, unit_color=R('body_text'), num_size=40, unit_size=16,
                         align=PP_ALIGN.LEFT))
        A.append(gtext(s, x + 0.45, y + ch - 0.65, cw - 0.7, 0.5, cd['label'], size=14, bold=True,
                       color=R('body_text'), align=PP_ALIGN.LEFT, font=F_LABEL))
    c.group_asset(s, A, 'KPI-204')
    return s


# =============================================================================
# KPI-205 트렌드 콜아웃 3종 (증감표식 + 미니 막대 스파크라인)
# =============================================================================
def build_trend(prs):
    s = blank_slide_gov(prs); gcaption(s, 'KPI-205'); A = []
    TREND = [
        {'value': '1,000', 'unit': '건', 'label': '민원 처리건수', 'delta': '+18%', 'col': 0},
        {'value': '42', 'unit': '억', 'label': '예산 집행액', 'delta': '+27%', 'col': 2},
        {'value': '4.6', 'unit': '점', 'label': '주민 만족도', 'delta': '+5%', 'col': 1},
    ]
    cw = (SW - 1.0 - 0.4 * 2) / 3; x0 = 0.5; y = 1.9; ch = 3.9
    heights = [0.5, 0.72, 0.62, 0.95, 1.25]
    for i, cd in enumerate(TREND):
        x = x0 + i * (cw + 0.4); accent = ACC[cd['col']]
        bg = c.add_box(s, x, y, cw, ch, fill=WHITE, line=GC['bg_panel'], line_w=BORDER_PT * 1.5, shape=RR)
        bg.adjustments[0] = RADIUS_PCT; A.append(bg)
        A.append(gtext(s, x + 0.25, y + 0.3, cw - 0.5, 0.4, cd['label'], size=14, bold=True,
                       color=R('body_text'), align=PP_ALIGN.LEFT, font=F_LABEL))
        A.append(numtext(s, x + 0.25, y + 0.72, cw - 0.5, 0.9, cd['value'], cd['unit'],
                         num_color=R('accent_primary'), unit_color=R('body_text'),
                         num_size=34, unit_size=14, align=PP_ALIGN.LEFT))
        A.append(gtext(s, x + 0.25, y + 1.6, cw - 0.5, 0.35, '▲ %s 전년比' % cd['delta'],
                       size=12, bold=True, color=accent, align=PP_ALIGN.LEFT, font=F_LABEL))
        nb = len(heights); bw = 0.34; bgp = 0.16
        tot = nb * bw + (nb - 1) * bgp; bx = x + (cw - tot) / 2
        base = y + ch - 0.4
        for j, hh in enumerate(heights):
            col = accent if j == nb - 1 else GC['bg_panel']
            A.append(c.add_box(s, bx + j * (bw + bgp), base - hh, bw, hh, fill=col))
    c.group_asset(s, A, 'KPI-205')
    return s


# =============================================================================
# KPI-206 카드 로우 — 가로 4카드 균등 반복 레이아웃
# =============================================================================
def build_card_row(prs):
    s = blank_slide_gov(prs); gcaption(s, 'KPI-206'); A = []
    CARDS = [
        {'value': '1,000', 'unit': '건', 'label': '검증 사업 DB', 'col': 0},
        {'value': '40', 'unit': '개소', 'label': '참여기관', 'col': 1},
        {'value': '42', 'unit': '억', 'label': '연계 예산', 'col': 2},
        {'value': '4.6', 'unit': '점', 'label': '만족도', 'col': 3},
    ]
    n = len(CARDS)
    cw = 2.55; ch = 2.7; ggap = 0.3
    tot = cw * n + ggap * (n - 1); x0 = (SW - tot) / 2; y = 2.55
    for i, cd in enumerate(CARDS):
        x = x0 + i * (cw + ggap); accent = ACC[cd['col']]
        bg = c.add_box(s, x, y, cw, ch, fill=WHITE, line=accent, line_w=BORDER_PT * 1.5, shape=RR)
        bg.adjustments[0] = RADIUS_PCT; A.append(bg)
        A.append(c.add_box(s, x, y, cw, 0.5, fill=accent, shape=RR))
        A.append(gtext(s, x, y, cw, 0.5, '지표 %d' % (i + 1), size=13, bold=True,
                       color=WHITE, font=F_LABEL))
        A.append(numtext(s, x, y + 0.75, cw, 1.1, cd['value'], cd['unit'],
                         num_color=accent, unit_color=R('body_text'), num_size=30, unit_size=13))
        A.append(gtext(s, x + 0.1, y + ch - 0.6, cw - 0.2, 0.5, cd['label'], size=13, bold=True,
                       color=R('body_text'), font=F_LABEL))
    c.group_asset(s, A, 'KPI-206')
    return s


# =============================================================================
# 매니페스트 조각
# =============================================================================
def E(aid, name, slide_idx, cards, count, tags, editable, params, ruse):
    return c.entry(asset_id=aid, category='KPI', name=name, file_rel=FILE_REL,
                   slide_idx=slide_idx, tags=tags, params=params,
                   bindings={'cards': cards, 'count': count},
                   editable=editable, recommended_use=ruse, master='gov')


def manifest():
    return [
        E('KPI-201', '[gov] 배지형 지표 4종(Paperlogy)', 1,
          [{'value': '98%', 'label': '정시 집행률'}, {'value': '4.8', 'label': '주민 만족도'},
           {'value': '120', 'label': '참여 기관수'}, {'value': '32', 'label': '연계 지자체'}], 4,
          ['KPI', 'gov', '배지', '원형', '핵심수치', 'Paperlogy'],
          ['value', 'label', 'badge-color'],
          {'badges': 4, 'style': 'circle-badge', 'canvas': 'gov', 'num_font': 'paperlogy'},
          ['핵심지표', '실적강조', '한눈요약']),
        E('KPI-202', '[gov] 후프(ring) 프로그레스 3종', 2,
          [{'value': '95', 'unit': '%', 'label': '사업 목표 달성률'},
           {'value': '82', 'unit': '%', 'label': '주민 만족도'},
           {'value': '68', 'unit': '%', 'label': '예산 집행률'}], 3,
          ['KPI', 'gov', '후프', '링', '도넛', '달성률'],
          ['value', 'label', 'ring-color', 'pct'],
          {'rings': 3, 'style': 'ring-progress', 'shape': 'native-pie', 'canvas': 'gov'},
          ['달성률', '목표대비', '진행률']),
        E('KPI-203', '[gov] hero 대형 숫자강조 콜아웃', 3,
          [{'value': '12,500', 'unit': '건', 'label': '누적 민원 처리건수'}], 1,
          ['KPI', 'gov', 'hero', '대형숫자', '단일강조', 'Paperlogy'],
          ['value', 'unit', 'label', 'subtitle', 'progress-bar'],
          {'style': 'hero-callout', 'canvas': 'gov', 'num_font': 'paperlogy', 'num_size': 96},
          ['핵심성과대표', '표지강조', '메인지표']),
        E('KPI-204', '[gov] quad 2x2 그리드 4지표', 4,
          [{'value': '1,000', 'unit': '건', 'label': '검증완료 사업 DB'},
           {'value': '40', 'unit': '개소', 'label': '참여 지역기관'},
           {'value': '42', 'unit': '억', 'label': '연계 예산 규모'},
           {'value': '95', 'unit': '%', 'label': '사업 목표 달성률'}], 4,
          ['KPI', 'gov', 'quad', '2x2', '그리드', '4지표'],
          ['value', 'unit', 'label', 'accent-color'],
          {'cells': 4, 'style': 'quad-grid', 'canvas': 'gov'},
          ['종합성과', '핵심지표', '한눈요약']),
        E('KPI-205', '[gov] 트렌드 콜아웃 3종 (증감+스파크라인)', 5,
          [{'value': '1,000', 'unit': '건', 'label': '민원 처리건수', 'delta': '+18%'},
           {'value': '42', 'unit': '억', 'label': '예산 집행액', 'delta': '+27%'},
           {'value': '4.6', 'unit': '점', 'label': '주민 만족도', 'delta': '+5%'}], 3,
          ['KPI', 'gov', '트렌드', '스파크라인', '증감', '콜아웃'],
          ['value', 'unit', 'label', 'delta', 'trend-color'],
          {'cards': 3, 'style': 'trend-sparkbar', 'bars': 5, 'canvas': 'gov'},
          ['성과추이', '증감비교', '성장강조']),
        E('KPI-206', '[gov] 카드 로우 가로 4카드', 6,
          [{'value': '1,000', 'unit': '건', 'label': '검증 사업 DB'},
           {'value': '40', 'unit': '개소', 'label': '참여기관'},
           {'value': '42', 'unit': '억', 'label': '연계 예산'},
           {'value': '4.6', 'unit': '점', 'label': '만족도'}], 4,
          ['KPI', 'gov', '카드로우', '가로반복', '카드형'],
          ['value', 'unit', 'label', 'accent-color'],
          {'cards': 4, 'style': 'card-row', 'canvas': 'gov'},
          ['성과요약', '지표나열', '반복레이아웃']),
    ]


if __name__ == '__main__':
    prs = new_deck_gov()
    build_badge(prs)
    build_hoop(prs)
    build_hero(prs)
    build_quad(prs)
    build_trend(prs)
    build_card_row(prs)
    out = c.save_deck(prs, FILE_REL)
    frag = c.write_fragment('KPI_gov', manifest())
    print('OK', out)
    print('FRAGMENT', frag)
