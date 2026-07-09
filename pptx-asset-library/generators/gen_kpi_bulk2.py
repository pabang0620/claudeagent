# -*- coding: utf-8 -*-
"""
KPI 대량 확충 생성기 2차 (KPI-201 ~ KPI-400, +200)
gen_kpi_bulk.py 의 14개 구조 패밀리를 재사용하되 파라미터 공간을 확대:
- accent 6색 전체 스윕
- 카드수 2~6 (기존 2~5 → 6 추가)
- 크기변형 소/중/대 (폰트·높이·지름 스케일)
- 숫자포맷 다양화 (정수/%/억/만$/명/점/건/소수)
- 라벨 풀 새 주제군 3종 (R&D·디지털 / 제조·품질 / 서비스·고객)
- 완전 동일 슬라이드 금지 (accent·count·size·theme 다중 스트라이드 조합)

규칙 준수: 색은 c.role/c.C 명명키만, 한글은 c 헬퍼, 다중도형은 c.group_asset 자체완결
그룹, 좌상단 c.id_caption(그룹 밖). 페이지이미지·SmartArt 금지, 네이티브 도형+텍스트만.
파일당 ≤25슬라이드, 패밀리별 1파일 분산. 마지막에 c.write_fragment('KPI_bulk2') 1회.
"""
import sys, math
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RR = MSO_SHAPE.ROUNDED_RECTANGLE
OVAL = MSO_SHAPE.OVAL
PIE = MSO_SHAPE.PIE
RADIUS = c.TOKENS["geom"]["round_radius"]
SW = 13.333
CENTER, LEFT, RIGHT = PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT
WHITE = c.C["white"]
GRAY_100 = c.C["gray_100"]
GRAY_300 = c.C["gray_300"]
UP_MARK = "▲"  # ▲

# accent 6색 전체 (c.C 명명키만; raw hex 금지)
ANAMES = ["blue_500", "teal_500", "cyan_500", "purple_600", "navy_600", "navy_800"]
ACCENTS = [c.C[n] for n in ANAMES]

# ---------- 라벨 풀 새 주제군 (숫자포맷 다양화: 억/만$/명/%/점/건/소수/정수) ----------
THEMES = [
    {
        "title": "R&D·디지털 성과",
        "data": [
            {"value": "128",    "unit": "억",  "label": "연구개발 투자액", "pct": 91, "prev": "104",   "cur": "128",    "delta": "+23%"},
            {"value": "3.2",    "unit": "만$", "label": "1인당 생산성",   "pct": 77, "prev": "2.7",   "cur": "3.2",    "delta": "+19%"},
            {"value": "12,500", "unit": "명",  "label": "플랫폼 가입자",   "pct": 86, "prev": "9,800", "cur": "12,500", "delta": "+28%"},
            {"value": "99.2",   "unit": "%",       "label": "시스템 가동률",   "pct": 99, "prev": "97.5",  "cur": "99.2",   "delta": "+2%"},
            {"value": "7.8",    "unit": "점",  "label": "ESG 평가등급",   "pct": 78, "prev": "6.9",   "cur": "7.8",    "delta": "+13%"},
            {"value": "340",    "unit": "건",  "label": "특허 출원 건수",   "pct": 84, "prev": "268",   "cur": "340",    "delta": "+27%"},
        ],
    },
    {
        "title": "제조·품질 성과",
        "data": [
            {"value": "1.2",  "unit": "만개", "label": "월 생산량",   "pct": 88, "prev": "0.95", "cur": "1.2",  "delta": "+26%"},
            {"value": "99.6", "unit": "%",            "label": "공정 수율",   "pct": 96, "prev": "98.1", "cur": "99.6", "delta": "+2%"},
            {"value": "18",   "unit": "억",       "label": "원가 절감액", "pct": 72, "prev": "12",   "cur": "18",   "delta": "+50%"},
            {"value": "0.8",  "unit": "%",            "label": "불량률",         "pct": 92, "prev": "1.4",  "cur": "0.8",  "delta": "-43%"},
            {"value": "4.7",  "unit": "점",       "label": "품질 만족도", "pct": 94, "prev": "4.3",  "cur": "4.7",  "delta": "+9%"},
            {"value": "52",   "unit": "개사", "label": "협력사 네트워크", "pct": 65, "prev": "38",   "cur": "52",   "delta": "+37%"},
        ],
    },
    {
        "title": "서비스·고객 성과",
        "data": [
            {"value": "2.4",   "unit": "만명", "label": "월 활성사용자", "pct": 83, "prev": "1.8",   "cur": "2.4",   "delta": "+33%"},
            {"value": "4.8",   "unit": "점",       "label": "앱 평점",       "pct": 96, "prev": "4.4",   "cur": "4.8",   "delta": "+9%"},
            {"value": "87",    "unit": "%",            "label": "재구매율",     "pct": 87, "prev": "79",    "cur": "87",    "delta": "+10%"},
            {"value": "15",    "unit": "억",       "label": "월 결제액",     "pct": 80, "prev": "11",    "cur": "15",    "delta": "+36%"},
            {"value": "320",   "unit": "만$",      "label": "해외 매출",     "pct": 74, "prev": "240",   "cur": "320",   "delta": "+33%"},
            {"value": "96",    "unit": "%",            "label": "응답 처리율",   "pct": 96, "prev": "90",    "cur": "96",    "delta": "+7%"},
        ],
    },
]

# ---------- 크기변형 소/중/대 ----------
SIZES = [("소", 0.86), ("중", 1.0), ("대", 1.16)]

# 카드수별 기본 사이즈 맵 (n=6 추가)
NSZ = {2: 46, 3: 40, 4: 34, 5: 28, 6: 24}
LSZ = {2: 16, 3: 15, 4: 14, 5: 12, 6: 11}
RINGNSZ = {2: 34, 3: 30, 4: 26, 5: 22, 6: 19}
GD = {2: 2.6, 3: 2.3, 4: 1.95, 5: 1.6, 6: 1.35}
GGAP = {2: 1.0, 3: 0.8, 4: 0.6, 5: 0.45, 6: 0.35}
BADGEGD = {2: 2.0, 3: 1.8, 4: 1.5, 5: 1.25, 6: 1.05}
BADGESZ = {2: 34, 3: 30, 4: 24, 5: 20, 6: 17}


def fz(base, sc):
    """폰트 크기 스케일 (하한 8pt)."""
    return max(8, int(round(base * sc)))


def sn(slide, x, y, w, num, unit="", num_color=None, unit_color=None,
       num_size=40, unit_size=16, align=CENTER, h=0.9):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(num)
    c.set_kfont(r, c.FONT_H, num_size, True, num_color or c.role("accent_primary"))
    if unit:
        ru = p.add_run(); ru.text = " " + unit
        c.set_kfont(ru, c.FONT_H, unit_size, True, unit_color or c.role("muted_text"))
    return tb


def title(slide, text):
    c.add_text(slide, 0.15, 0.5, 12.5, 0.4, text, size=12, bold=True,
               color=c.role("muted_text"))


def hx(n, w, gap):
    tot = n * w + (n - 1) * gap
    x0 = (SW - tot) / 2
    return [x0 + i * (w + gap) for i in range(n)]


# ---------- 메타(bindings.cards) 빌더 ----------
def meta_default(D): return [{"value": d["value"], "unit": d["unit"], "label": d["label"]} for d in D]
def meta_pct(D):     return [{"value": str(d["pct"]), "unit": "%", "label": d["label"]} for d in D]
def meta_yoy(D):     return [{"value": d["cur"], "unit": d["unit"], "label": d["label"], "prev": d["prev"]} for d in D]
def meta_trend(D):   return [{"value": d["value"], "unit": d["unit"], "label": d["label"], "delta": d["delta"]} for d in D]


# ============================================================
# 패밀리 렌더 — signature: (s, accent, n, D, sc, gtitle) -> (shapes, meta)
# D = n개로 슬라이스된 데이터, sc = 크기 스케일, gtitle = 그룹 타이틀
# ============================================================
def render_filled(s, accent, n, D, sc, gt):
    A = []; gap = 0.4; cw = (11.5 - (n - 1) * gap) / n; ch = 2.8 * (1 + (sc - 1) * 0.5); y = 2.1
    for x, cd in zip(hx(n, cw, gap), D):
        card = c.add_box(s, x, y, cw, ch, fill=accent, shape=RR); card.adjustments[0] = RADIUS; A.append(card)
        A.append(sn(s, x, y + 0.55, cw, cd["value"], cd["unit"], num_color=WHITE, unit_color=GRAY_100,
                    num_size=fz(NSZ[n], sc), unit_size=fz(15, sc)))
        A.append(c.add_text(s, x + 0.1, y + ch - 0.85, cw - 0.2, 0.7, cd["label"], size=fz(LSZ[n], sc), bold=True,
                            color=GRAY_100, align=CENTER))
    return A, meta_default(D)


def render_border(s, accent, n, D, sc, gt):
    A = []; gap = 0.4; cw = (11.5 - (n - 1) * gap) / n; ch = 2.8 * (1 + (sc - 1) * 0.5); y = 2.1
    for x, cd in zip(hx(n, cw, gap), D):
        card = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"), line=accent, line_w=2.0, shape=RR)
        card.adjustments[0] = RADIUS; A.append(card)
        A.append(sn(s, x, y + 0.55, cw, cd["value"], cd["unit"], num_color=accent, unit_color=c.role("muted_text"),
                    num_size=fz(NSZ[n], sc), unit_size=fz(15, sc)))
        A.append(c.add_text(s, x + 0.1, y + ch - 0.85, cw - 0.2, 0.7, cd["label"], size=fz(LSZ[n], sc), bold=True,
                            color=c.role("body_text"), align=CENTER))
    return A, meta_default(D)


def render_topbar(s, accent, n, D, sc, gt):
    A = []; gap = 0.4; cw = (11.5 - (n - 1) * gap) / n; ch = 2.8 * (1 + (sc - 1) * 0.5); y = 2.1
    for x, cd in zip(hx(n, cw, gap), D):
        card = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"), line=c.role("border"), line_w=1.0, shape=RR)
        card.adjustments[0] = RADIUS; A.append(card)
        bar = c.add_box(s, x, y, cw, 0.62, fill=accent, shape=RR); bar.adjustments[0] = RADIUS; A.append(bar)
        A.append(c.add_text(s, x, y, cw, 0.62, cd["label"], size=fz(LSZ[n], sc), bold=True, color=WHITE, align=CENTER))
        A.append(sn(s, x, y + 1.05, cw, cd["value"], cd["unit"], num_color=accent, unit_color=c.role("muted_text"),
                    num_size=fz(NSZ[n], sc), unit_size=fz(15, sc), h=1.2))
    return A, meta_default(D)


def render_iconnum(s, accent, n, D, sc, gt):
    A = []; gap = 0.4; cw = (11.5 - (n - 1) * gap) / n; ch = 3.0 * (1 + (sc - 1) * 0.4); y = 2.0
    for x, cd in zip(hx(n, cw, gap), D):
        card = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"), line=c.role("border"), line_w=1.2, shape=RR)
        card.adjustments[0] = RADIUS; A.append(card)
        icd = min(1.0, cw * 0.42)
        ic = c.add_box(s, x + (cw - icd) / 2, y + 0.3, icd, icd, fill=accent, shape=OVAL)
        c.set_shape_text(ic, "ICON", size=9, bold=True, color=WHITE); A.append(ic)
        A.append(sn(s, x, y + 0.35 + icd, cw, cd["value"], cd["unit"], num_color=accent,
                    unit_color=c.role("muted_text"), num_size=fz(min(40, NSZ[n] + 2), sc), unit_size=fz(14, sc), h=0.85))
        A.append(c.add_text(s, x + 0.1, y + ch - 0.65, cw - 0.2, 0.55, cd["label"], size=fz(LSZ[n], sc), bold=True,
                            color=c.role("body_text"), align=CENTER))
    return A, meta_default(D)


def render_yoy(s, accent, n, D, sc, gt):
    A = []; gap = 0.4; cw = (11.5 - (n - 1) * gap) / n; ch = 3.6 * (1 + (sc - 1) * 0.35); y = 1.7
    up = c.role("accent_secondary")
    for x, cd in zip(hx(n, cw, gap), D):
        card = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"), line=c.role("border"), line_w=1.2, shape=RR)
        card.adjustments[0] = RADIUS; A.append(card)
        bar = c.add_box(s, x, y, cw, 0.52, fill=accent, shape=RR); bar.adjustments[0] = RADIUS; A.append(bar)
        A.append(c.add_text(s, x, y, cw, 0.52, cd["label"], size=fz(LSZ[n], sc), bold=True, color=WHITE, align=CENTER))
        A.append(c.add_text(s, x + 0.1, y + 0.72, cw - 0.2, 0.35, "전년 %s%s" % (cd["prev"], cd["unit"]),
                            size=fz(12, sc), bold=True, color=c.role("muted_text"), align=CENTER))
        A.append(c.add_text(s, x, y + 1.12, cw, 0.4, UP_MARK, size=fz(18, sc), bold=True, color=up, align=CENTER))
        A.append(sn(s, x, y + 1.6, cw, cd["cur"], cd["unit"], num_color=c.role("header_fill"), unit_color=accent,
                    num_size=fz(NSZ[n], sc), unit_size=fz(15, sc)))
        A.append(c.add_text(s, x, y + ch - 0.6, cw, 0.45, "%s 전년比" % cd["delta"], size=fz(13, sc), bold=True,
                            color=up, align=CENTER))
    return A, meta_yoy(D)


def render_ring(s, accent, n, D, sc, gt):
    A = []; gd = GD[n] * (1 + (sc - 1) * 0.5); gy = 2.2; ggap = GGAP[n]
    for x, cd in zip(hx(n, gd, ggap), D):
        frac = cd["pct"] / 100.0
        A.append(c.add_box(s, x, gy, gd, gd, fill=c.role("border"), shape=OVAL))
        pie = c.add_box(s, x, gy, gd, gd, fill=accent, shape=PIE)
        pie.adjustments[0] = -90.0; pie.adjustments[1] = -90.0 + 360.0 * frac; A.append(pie)
        inr = gd * 0.74
        A.append(c.add_box(s, x + (gd - inr) / 2, gy + (gd - inr) / 2, inr, inr, fill=c.role("row_base"), shape=OVAL))
        A.append(sn(s, x, gy + gd / 2 - 0.4, gd, str(cd["pct"]), "%", num_color=c.role("header_fill"),
                    unit_color=accent, num_size=fz(RINGNSZ[n], sc), unit_size=fz(13, sc)))
        A.append(c.add_text(s, x - 0.2, gy + gd + 0.1, gd + 0.4, 0.5, cd["label"], size=fz(13, sc), bold=True,
                            color=c.role("body_text"), align=CENTER))
    return A, meta_pct(D)


def render_gauge(s, accent, n, D, sc, gt):
    A = []; gd = GD[n] * (1 + (sc - 1) * 0.5); gy = 2.4; ggap = GGAP[n]
    for x, cd in zip(hx(n, gd, ggap), D):
        frac = cd["pct"] / 100.0; cyl = gy + gd / 2
        A.append(c.add_box(s, x, gy, gd, gd, fill=c.role("border"), shape=OVAL))
        pie = c.add_box(s, x, gy, gd, gd, fill=accent, shape=PIE)
        pie.adjustments[0] = 180.0; pie.adjustments[1] = 180.0 + 180.0 * frac; A.append(pie)
        inr = gd * 0.58
        A.append(c.add_box(s, x + (gd - inr) / 2, gy + (gd - inr) / 2, inr, inr, fill=c.role("row_base"), shape=OVAL))
        A.append(c.add_box(s, x - 0.05, cyl, gd + 0.1, gd / 2 + 0.15, fill=c.role("row_base")))  # 하단 반원 클립
        A.append(sn(s, x, cyl - 0.68, gd, str(cd["pct"]), "%", num_color=c.role("header_fill"),
                    unit_color=accent, num_size=fz(RINGNSZ[n], sc), unit_size=fz(13, sc)))
        A.append(c.add_text(s, x - 0.2, cyl + 0.12, gd + 0.4, 0.5, cd["label"], size=fz(13, sc), bold=True,
                            color=c.role("body_text"), align=CENTER))
    return A, meta_pct(D)


def render_badge(s, accent, n, D, sc, gt):
    A = []; bd = BADGEGD[n] * (1 + (sc - 1) * 0.5); by = 2.4; bgap = GGAP[n]
    for x, cd in zip(hx(n, bd, bgap), D):
        A.append(c.add_box(s, x - 0.09, by - 0.09, bd + 0.18, bd + 0.18, fill=None, line=accent, line_w=1.5, shape=OVAL))
        badge = c.add_box(s, x, by, bd, bd, fill=accent, shape=OVAL)
        c.set_shape_text(badge, cd["value"], size=fz(BADGESZ[n], sc), bold=True, color=WHITE); A.append(badge)
        A.append(c.add_text(s, x - 0.2, by + bd + 0.15, bd + 0.4, 0.5, cd["label"], size=fz(13, sc), bold=True,
                            color=c.role("body_text"), align=CENTER))
    return A, meta_default(D)


def render_hero(s, accent, n, D, sc, gt):
    A = []; y = 1.8; H = 3.9; lw = 5.4; lx = (SW - 11.0) / 2
    left = c.add_box(s, lx, y, lw, H, fill=accent, shape=RR); left.adjustments[0] = RADIUS; A.append(left)
    d0 = D[0]
    A.append(sn(s, lx + 0.2, y + 0.7, lw - 0.4, d0["value"], d0["unit"], num_color=WHITE, unit_color=GRAY_100,
                num_size=fz(72, sc), unit_size=fz(24, sc), h=1.6))
    A.append(c.add_text(s, lx + 0.3, y + H - 1.3, lw - 0.6, 1.0, d0["label"], size=fz(18, sc), bold=True,
                        color=GRAY_100, align=CENTER))
    rest = D[1:]; rx = lx + lw + 0.4; rw = 11.0 - lw - 0.4; m = len(rest)
    if m > 0:
        rgap = 0.25; rh = (H - (m - 1) * rgap) / m
        for j, cd in enumerate(rest):
            yy = y + j * (rh + rgap)
            card = c.add_box(s, rx, yy, rw, rh, fill=c.role("row_base"), line=c.role("border"), line_w=1.2, shape=RR)
            card.adjustments[0] = 0.08; A.append(card)
            A.append(c.add_box(s, rx, yy, 0.14, rh, fill=accent))
            A.append(sn(s, rx + 0.4, yy, rw * 0.45, cd["value"], cd["unit"], num_color=accent,
                        unit_color=c.role("muted_text"), num_size=fz(24, sc), unit_size=fz(12, sc), align=LEFT, h=rh))
            A.append(c.add_text(s, rx + rw * 0.45 + 0.5, yy, rw * 0.55 - 0.6, rh, cd["label"], size=fz(13, sc), bold=True,
                                color=c.role("body_text"), align=LEFT))
    return A, meta_default(D)


def render_strip(s, accent, n, D, sc, gt):
    A = []; rw = 11.0; rx = (SW - rw) / 2; top = 1.7; bottom = 6.1; H = bottom - top
    rgap = 0.22; rh = (H - (n - 1) * rgap) / n
    numsz = {2: 34, 3: 30, 4: 26, 5: 22, 6: 19}[n]
    for j, cd in enumerate(D):
        yy = top + j * (rh + rgap)
        card = c.add_box(s, rx, yy, rw, rh, fill=c.role("panel_bg"), line=c.role("border"), line_w=1.0, shape=RR)
        card.adjustments[0] = min(0.14, RADIUS * 1.5); A.append(card)
        A.append(c.add_box(s, rx, yy, 0.15, rh, fill=accent))
        A.append(c.add_text(s, rx + 0.5, yy, rw * 0.6, rh, cd["label"], size=fz(16, sc), bold=True,
                            color=c.role("body_text"), align=LEFT))
        A.append(sn(s, rx + rw * 0.6, yy, rw * 0.4 - 0.35, cd["value"], cd["unit"], num_color=accent,
                    unit_color=c.role("muted_text"), num_size=fz(numsz, sc), unit_size=fz(13, sc),
                    align=RIGHT, h=rh))
    return A, meta_default(D)


def render_trend(s, accent, n, D, sc, gt):
    A = []; gap = 0.4; cw = (11.5 - (n - 1) * gap) / n; ch = 3.3 * (1 + (sc - 1) * 0.4); y = 1.8
    heights = [0.5, 0.72, 0.62, 0.95, 1.25]
    numsz = {2: 40, 3: 36, 4: 30, 5: 26, 6: 22}[n]
    for x, cd in zip(hx(n, cw, gap), D):
        bg = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"), line=c.role("border"), line_w=1.0, shape=RR)
        bg.adjustments[0] = RADIUS; A.append(bg)
        A.append(c.add_text(s, x + 0.2, y + 0.25, cw - 0.4, 0.4, cd["label"], size=fz(LSZ[n], sc), bold=True,
                            color=c.role("muted_text"), align=LEFT))
        A.append(sn(s, x + 0.2, y + 0.65, cw - 0.4, cd["value"], cd["unit"], num_color=c.role("header_fill"),
                    unit_color=c.role("muted_text"), num_size=fz(numsz, sc), unit_size=fz(14, sc), align=LEFT))
        A.append(c.add_text(s, x + 0.2, y + 1.5, cw - 0.4, 0.35, "%s %s 전년比" % (UP_MARK, cd["delta"]),
                            size=fz(12, sc), bold=True, color=c.role("accent_secondary"), align=LEFT))
        nb = len(heights); area = cw - 0.5; bw = area / (1.5 * nb - 0.5); bgp = bw * 0.5
        tot = nb * bw + (nb - 1) * bgp; bx = x + (cw - tot) / 2; base = y + ch - 0.35
        for jj, hh in enumerate(heights):
            col = accent if jj == nb - 1 else GRAY_300
            A.append(c.add_box(s, bx + jj * (bw + bgp), base - hh, bw, hh, fill=col))
    return A, meta_trend(D)


def render_quad(s, accent, n, D, sc, gt):
    A = []; cols = 2; rows = math.ceil(n / 2); pw = 9.0; px = (SW - pw) / 2; py = 1.9
    cellgap = 0.3; cw = (pw - cellgap * (cols - 1)) / cols; chh = 1.5 * (1 + (sc - 1) * 0.3)
    for idx, cd in enumerate(D):
        r = idx // cols; col = idx % cols
        x = px + col * (cw + cellgap); yy = py + r * (chh + cellgap)
        card = c.add_box(s, x, yy, cw, chh, fill=c.role("row_base"), line=c.role("border"), line_w=1.2, shape=RR)
        card.adjustments[0] = 0.08; A.append(card)
        A.append(c.add_box(s, x, yy, cw, 0.12, fill=accent))
        A.append(sn(s, x + 0.3, yy + 0.28, cw - 0.6, cd["value"], cd["unit"], num_color=accent,
                    unit_color=c.role("muted_text"), num_size=fz(32, sc), unit_size=fz(14, sc), align=LEFT, h=0.8))
        A.append(c.add_text(s, x + 0.3, yy + chh - 0.55, cw - 0.6, 0.45, cd["label"], size=fz(13, sc), bold=True,
                            color=c.role("muted_text"), align=LEFT))
    return A, meta_default(D)


def render_grouplabel(s, accent, n, D, sc, gt):
    A = []; cw = 11.0; cx = (SW - cw) / 2; y = 2.0
    hb = c.add_box(s, cx, y, cw, 0.7, fill=accent, shape=RR); hb.adjustments[0] = 0.12; A.append(hb)
    A.append(c.add_text(s, cx + 0.3, y, cw - 0.6, 0.7, gt, size=fz(16, sc), bold=True, color=WHITE, align=LEFT))
    ch = 2.6 * (1 + (sc - 1) * 0.4); cy = y + 0.85
    cont = c.add_box(s, cx, cy, cw, ch, fill=c.role("row_base"), line=c.role("border"), line_w=1.2, shape=RR)
    cont.adjustments[0] = 0.06; A.append(cont)
    seg = cw / n
    numsz = {2: 40, 3: 34, 4: 30, 5: 26, 6: 22}[n]
    for i, cd in enumerate(D):
        sx = cx + i * seg
        if i > 0:
            A.append(c.add_box(s, sx, cy + 0.4, 0.02, ch - 0.8, fill=c.role("border")))
        A.append(sn(s, sx, cy + 0.55, seg, cd["value"], cd["unit"], num_color=accent, unit_color=c.role("muted_text"),
                    num_size=fz(numsz, sc), unit_size=fz(14, sc)))
        A.append(c.add_text(s, sx, cy + ch - 0.8, seg, 0.5, cd["label"], size=fz(13, sc), bold=True,
                            color=c.role("body_text"), align=CENTER))
    return A, meta_default(D)


def render_band(s, accent, n, D, sc, gt):
    A = []; bw = 12.0; bh = 2.0 * (1 + (sc - 1) * 0.4); bx = (SW - bw) / 2; byy = 2.7
    band = c.add_box(s, bx, byy, bw, bh, fill=c.role("header_fill"), shape=RR); band.adjustments[0] = 0.06; A.append(band)
    A.append(c.add_box(s, bx, byy, bw, 0.14, fill=accent))
    seg = bw / n
    numsz = {2: 40, 3: 34, 4: 30, 5: 26, 6: 22}[n]
    for i, cd in enumerate(D):
        sx = bx + i * seg
        if i > 0:
            A.append(c.add_box(s, sx, byy + 0.4, 0.02, bh - 0.8, fill=accent))
        A.append(sn(s, sx, byy + 0.45, seg, cd["value"], cd["unit"], num_color=WHITE, unit_color=GRAY_300,
                    num_size=fz(numsz, sc), unit_size=fz(15, sc)))
        A.append(c.add_text(s, sx, byy + bh - 0.6, seg, 0.45, cd["label"], size=fz(13, sc), bold=True,
                            color=GRAY_100, align=CENTER))
    return A, meta_default(D)


# ============================================================
# 패밀리 설정: (slug, kr, style, render, tags_extra, editable_extra, recommended_use)
# ============================================================
FAMILIES = [
    ("filledcards", "채움 카드 세트", "filled-card", render_filled,
     ["채움카드", "솔리드", "카드세트"], [], ["대표성과", "실적강조", "핵심지표"]),
    ("bordercards", "테두리 카드", "outline-card", render_border,
     ["테두리카드", "아웃라인", "카드세트"], [], ["핵심지표", "실적강조", "성과요약"]),
    ("topbar", "상단 컴러바 카드", "topbar-card", render_topbar,
     ["상단컴러바", "헤더바", "카드세트"], [], ["항목별성과", "핵심지표", "실적강조"]),
    ("iconnum", "아이콘원+숫자 카드", "icon-number-card", render_iconnum,
     ["아이콘", "원형아이콘", "숫자카드"], ["icon", "icon-color"], ["서비스지표", "성과목표", "실적강조"]),
    ("yoy", "전년대비 증감 카드", "yoy-change-card", render_yoy,
     ["전년대비", "증감", "성장강조"], ["prev", "cur", "delta"], ["증감비교", "전년대비", "성장강조"]),
    ("ring", "링 프로그레스", "ring-progress", render_ring,
     ["링", "도넯", "프로그레스", "비율강조"], ["pct"], ["달성률", "목표대비", "진행률"]),
    ("gauge", "반원 게이지", "semicircle-gauge", render_gauge,
     ["게이지", "반원", "semicircle", "달성률"], ["pct"], ["달성률", "목표대비", "가동률"]),
    ("badge", "배지형 지표", "circle-badge", render_badge,
     ["배지", "원형", "핵심수치"], [], ["핵심지표", "실적강조", "한눈요약"]),
    ("hero", "히어로 대형숫자", "hero-bignumber", render_hero,
     ["히어로", "대형숫자", "빅넘버", "대표지표"], [], ["대표성과", "핵심수치", "메인지표"]),
    ("strip", "가로 스트립", "label-left-number-right", render_strip,
     ["가로스트립", "라벨좌", "숫자우", "리스트"], [], ["성과목표", "실적강조", "지표리스트"]),
    ("trend", "트렌드 카드", "trend-sparkbar", render_trend,
     ["트렌드", "스파크라인", "미니막대", "증감"], ["delta"], ["성과추이", "증감비교", "성장강조"]),
    ("quad", "4분할 대시보드", "quadrant-dashboard", render_quad,
     ["대시보드", "4분할", "그리드", "요약패널"], [], ["종합성과", "대시보드", "한눈요약"]),
    ("grouplabel", "그룹라벨+숫자묶음", "group-label-cluster", render_grouplabel,
     ["그룹라벨", "숫자묶음", "구분선", "요약"], [], ["종합성과", "핵심지표", "한눈요약"]),
    ("band", "통계 밴드", "stat-band", render_band,
     ["통계밴드", "가로띄", "구분선", "지표"], [], ["종합성과", "핵심지표", "한눈요약"]),
]

CARD_COUNTS = [2, 3, 4, 5, 6]
START_ID = 201
END_ID = 400  # inclusive → 200개

# 200개를 14패밀리로 분산: 앞 4패밀리 15개, 나머지 10패밀리 14개 = 60+140 = 200 (파일당 ≤25)
PER_FAMILY = [15, 15, 15, 15] + [14] * 10  # 합계 200


def main():
    assert sum(PER_FAMILY) == (END_ID - START_ID + 1), sum(PER_FAMILY)
    allentries = []
    files = []
    gid = START_ID
    for f, (slug, kr, style, render, tags_x, edit_x, ruse) in enumerate(FAMILIES):
        prs = c.new_deck()
        file_rel = "decks/02_kpi/KPI_bulk2_%s_v1.pptx" % slug
        cnt = PER_FAMILY[f]
        for k in range(cnt):
            g = gid - START_ID  # 전역 인덱스 (다중 스트라이드로 조합 다양화)
            n = CARD_COUNTS[g % 5]
            ci = g % 6
            accent = ACCENTS[ci]; aname = ANAMES[ci]
            si = (g // 2) % 3; sname, sc = SIZES[si]
            ti = (g // 3) % len(THEMES); theme = THEMES[ti]
            D = theme["data"][:n]; gt = theme["title"]
            aid = "KPI-%03d" % gid; gid += 1
            s = c.blank_slide(prs); c.id_caption(s, aid)
            title(s, "%s · %d카드 · %s · %s크기 · %s" % (kr, n, aname, sname, theme["title"]))
            A, meta = render(s, accent, n, D, sc, gt)
            c.group_asset(s, A, aid)
            allentries.append(c.entry(
                asset_id=aid, category="KPI", name="%s %d종 (%s·%s)" % (kr, n, aname, sname),
                file_rel=file_rel, slide_idx=k + 1,
                tags=["KPI", "성과지표", kr, "크기-" + sname, theme["title"]] + tags_x,
                params={"style": style, "count": n, "accent": aname, "family": slug,
                        "size": sname, "scale": sc, "theme": theme["title"], "num_size": fz(NSZ[n], sc)},
                bindings={"cards": meta, "count": n},
                editable=["value", "unit", "label", "accent-color"] + edit_x,
                recommended_use=ruse))
        out = c.save_deck(prs, file_rel); files.append(out)
    frag = c.write_fragment("KPI_bulk2", allentries)
    print("FILES:")
    for p in files:
        print("  ", p)
    print("FRAG:", frag)
    print("ENTRIES:", len(allentries))
    print("ID_RANGE:", allentries[0]["id"], "->", allentries[-1]["id"])


if __name__ == "__main__":
    main()
