"""
후처리 그룹화 — 1에셋=1슬라이드 덱의 각 슬라이드 콘텐츠 도형을 asset:<ID> 그룹으로 묶는다.
조합(addElement 1콜) 대응. 다중에셋 슬라이드(HDR/ICN)는 스킵(별도 재생성 필요).
멱등: 이미 그룹된 슬라이드는 건너뜀. 사용: python3 generators/regroup.py
"""
import sys, os, json, collections
sys.path.insert(0, "/home/pabang/myapp/.claude/pptx-asset-library/generators/lib")
import common as c
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = c.ROOT

def content_shapes(slide, asset_id):
    out = []
    for sh in slide.shapes:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip() == asset_id:
                continue  # id_caption 제외
        except Exception:
            pass
        out.append(sh)
    return out

def main():
    man = json.load(open(os.path.join(ROOT, "manifest.json"), encoding="utf-8"))
    # (file, slide) -> [ids]
    groups = collections.defaultdict(list)
    for a in man["assets"]:
        groups[(a["file"], a["slide"])].append(a["id"])

    stats = collections.Counter()
    changed_files = {}
    for (frel, sidx), ids in sorted(groups.items()):
        if len(ids) != 1:
            stats["skip_multi"] += 1
            continue
        aid = ids[0]
        path = os.path.join(ROOT, frel)
        prs = changed_files.get(path) or Presentation(path)
        changed_files[path] = prs
        slide = prs.slides[sidx-1]
        # 이미 그룹됨?
        grp_named = any(el.get("name") == "asset:%s" % aid
                        for el in slide.shapes._spTree.iter(qn("p:cNvPr")))
        top_grp = [sh for sh in slide.shapes if sh.shape_type == 6 and sh.name == "asset:%s" % aid]
        if top_grp:
            stats["already_group"] += 1; continue
        cs = content_shapes(slide, aid)
        if len(cs) <= 1:
            if cs:
                cs[0].name = "asset:%s" % aid  # 단일 객체 → 이름만 보장
            stats["single_obj"] += 1; continue
        # 기존 앵커명 중복 제거(그룹만 앵커 보유하도록)
        for sh in cs:
            if sh.name == "asset:%s" % aid:
                sh.name = "%s_bg" % aid
        c.group_asset(slide, cs, aid)
        stats["grouped"] += 1

    for path, prs in changed_files.items():
        prs.save(path)

    print("REGROUP:", dict(stats))
    print("touched files:", len(changed_files))

if __name__ == "__main__":
    main()
