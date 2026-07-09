# -*- coding: utf-8 -*-
"""TML 카테고리(타임라인/로드맵) 에셋 생성 — TML-001~TML-008.

파일 분산:
  TML_timeline_v1.pptx  : TML-001 수평타임라인 / TML-002 간트형 / TML-005 마일스톤플래그 / TML-006 3트랙
  TML_roadmap_v1.pptx   : TML-003 단계별로드맵 / TML-004 분기별 / TML-007 연차별계단 / TML-008 추진일정요약

네이티브 도형(사각/원/화살표/오각/계단)+커넥터만 사용. SmartArt·이미지 금지.
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

R = c.role
CC = c.C

# 코퍼레이트 팔레트 순환(네이비/블루/틸/시안/그레이)
PALETTE = [R("accent_primary"), R("accent_secondary"), R("accent_point"),
           CC["navy_600"], R("sub_header")]


def title_block(slide, text, sub=None):
    c.add_text(slide, 0.55, 0.52, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.02, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)


def label_card(slide, x, y, w, h, when, label, color):
    """상/하 라벨: when 뱃지 + 라벨 텍스트."""
    badge = c.add_box(slide, x + (w - 1.0) / 2, y, 1.0, 0.34, fill=color, line=None,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.set_shape_text(badge, when, size=10, bold=True, color=R("header_text"))
    c.add_text(slide, x, y + 0.40, w, h - 0.40, label, size=11, bold=True,
               color=R("body_text"), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# =========================================================
# TML-001 수평 타임라인 (가로선 + 마일스톤 원 5개 + 상하 라벨 교차)
# =========================================================
def horizontal_timeline(slide, asset_id, title, sub, milestones):
    """milestones: list of dict {label, when}."""
    title_block(slide, title, sub)
    x0, x1 = 1.35, 11.98
    ymid = 4.0
    cont = c.add_box(slide, x0 - 0.15, 1.55, (x1 - x0) + 0.3, 4.9, fill=None, line=None)
    c.name_asset(cont, asset_id)
    # 기준선
    c.connector(slide, x0, ymid, x1, ymid, color=R("accent_primary"), w=3.0)
    # 좌우 화살촉(진행 방향)
    tip = c.add_box(slide, x1 - 0.02, ymid - 0.16, 0.34, 0.32, fill=R("accent_primary"),
                    line=None, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
    tip.rotation = 90
    n = len(milestones)
    dia = 0.46
    cw = 2.15
    for i, m in enumerate(milestones):
        cx = x0 + (x1 - x0) * (i / (n - 1)) if n > 1 else (x0 + x1) / 2
        color = PALETTE[i % len(PALETTE)]
        above = (i % 2 == 0)
        # 노드에서 라벨카드로 연결하는 리더선
        if above:
            c.connector(slide, cx, ymid - dia / 2, cx, ymid - 0.95, color=R("border"), w=1.0)
            label_card(slide, cx - cw / 2, ymid - 1.95, cw, 1.0, m["when"], m["label"], color)
        else:
            c.connector(slide, cx, ymid + dia / 2, cx, ymid + 0.95, color=R("border"), w=1.0)
            label_card(slide, cx - cw / 2, ymid + 0.98, cw, 1.0, m["when"], m["label"], color)
        # 마일스톤 원(맨 위 도포 순서 유지)
        node = c.add_box(slide, cx - dia / 2, ymid - dia / 2, dia, dia, fill=color,
                         line=CC["white"], line_w=2.5, shape=MSO_SHAPE.OVAL)
        c.set_shape_text(node, str(i + 1), size=13, bold=True, color=R("header_text"))


# =========================================================
# TML-002 간트차트형 (작업 4행 × 기간 막대, 도형 기반)
# =========================================================
def gantt_chart(slide, asset_id, title, sub, periods, tasks):
    """periods: 헤더 라벨 리스트. tasks: list of (name, start_idx, span, color_idx)."""
    title_block(slide, title, sub)
    lx, lw = 1.1, 2.7          # 작업명 열
    gx0 = lx + lw              # 그리드 시작 X
    gx1 = 12.2
    gw = gx1 - gx0
    np_ = len(periods)
    colw = gw / np_
    hy = 1.75                  # 헤더 Y
    hh = 0.5
    ry0 = hy + hh              # 첫 작업 행 Y
    rh = 0.82
    nrows = len(tasks)
    grid_h = hh + rh * nrows
    cont = c.add_box(slide, lx - 0.05, hy - 0.05, (gx1 - lx) + 0.1, grid_h + 0.1, fill=None, line=None)
    c.name_asset(cont, asset_id)
    # 헤더 배경(작업 열 + 기간 열)
    c.add_box(slide, lx, hy, lw, hh, fill=R("header_fill"), line=CC["white"], line_w=1.0)
    c.add_text(slide, lx, hy, lw, hh, "작업 / 기간", size=11, bold=True,
               color=R("header_text"), align=PP_ALIGN.CENTER)
    for j, pl in enumerate(periods):
        px = gx0 + j * colw
        c.add_box(slide, px, hy, colw, hh, fill=CC["navy_600"], line=CC["white"], line_w=1.0)
        c.add_text(slide, px, hy, colw, hh, pl, size=10, bold=True,
                   color=R("header_text"), align=PP_ALIGN.CENTER)
    # 세로 그리드 커넥터
    for j in range(np_ + 1):
        gxj = gx0 + j * colw
        c.connector(slide, gxj, ry0, gxj, ry0 + rh * nrows, color=CC["gray_300"], w=0.75)
    # 작업 행 + 막대
    for i, (name, start, span, ci) in enumerate(tasks):
        ry = ry0 + i * rh
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        c.add_box(slide, lx, ry, lw, rh, fill=base, line=CC["gray_300"], line_w=0.75)
        c.add_text(slide, lx + 0.12, ry, lw - 0.2, rh, name, size=11, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT)
        c.add_box(slide, gx0, ry, gw, rh, fill=base, line=None)
        bx = gx0 + start * colw + 0.06
        bw = span * colw - 0.12
        bar = c.add_box(slide, bx, ry + rh * 0.22, bw, rh * 0.56, fill=PALETTE[ci % len(PALETTE)],
                        line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.set_shape_text(bar, "%d단계" % (i + 1), size=9, bold=True, color=R("header_text"))


# =========================================================
# TML-005 마일스톤 플래그 타임라인 (수평선 + 깃발 도형)
# =========================================================
def flag_timeline(slide, asset_id, title, sub, flags):
    """flags: list of dict {label, when}. 깃발(폴+페넌트) 도형."""
    title_block(slide, title, sub)
    x0, x1 = 1.4, 11.95
    baseline = 5.7
    cont = c.add_box(slide, x0 - 0.15, 1.6, (x1 - x0) + 0.3, 4.6, fill=None, line=None)
    c.name_asset(cont, asset_id)
    c.connector(slide, x0, baseline, x1, baseline, color=CC["navy_600"], w=3.0)
    n = len(flags)
    for i, fl in enumerate(flags):
        cx = x0 + (x1 - x0) * ((i + 0.5) / n)
        color = PALETTE[i % len(PALETTE)]
        pole_top = baseline - 2.5
        # 기준점 원
        c.add_box(slide, cx - 0.11, baseline - 0.11, 0.22, 0.22, fill=color, line=CC["white"],
                  line_w=1.5, shape=MSO_SHAPE.OVAL)
        # 깃대(폴)
        c.connector(slide, cx, baseline, cx, pole_top, color=R("muted_text"), w=2.0)
        # 페넌트(오각 깃발)
        flag = c.add_box(slide, cx, pole_top, 1.75, 0.7, fill=color, line=None,
                         shape=MSO_SHAPE.PENTAGON)
        c.set_shape_text(flag, fl["when"], size=11, bold=True, color=R("header_text"))
        # 라벨(깃발 아래)
        c.add_text(slide, cx - 0.15, pole_top + 0.78, 2.2, 0.9, fl["label"], size=11, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


# =========================================================
# TML-006 로드맵 3트랙 (트랙별 가로 레인 + 활동 박스)
# =========================================================
def three_track(slide, asset_id, title, sub, phases, tracks):
    """phases: 상단 단계 라벨. tracks: list of (track_name, color_idx, [활동 문자열…])."""
    title_block(slide, title, sub)
    lx, lw = 1.1, 2.0          # 트랙명 열
    gx0 = lx + lw
    gx1 = 12.2
    gw = gx1 - gx0
    npc = len(phases)
    colw = gw / npc
    hy, hh = 1.7, 0.5
    ry0 = hy + hh + 0.08
    lane_h = 1.35
    gap = 0.18
    cont = c.add_box(slide, lx - 0.05, hy - 0.05, (gx1 - lx) + 0.1,
                     (ry0 - hy) + lane_h * len(tracks) + gap * (len(tracks) - 1) + 0.1,
                     fill=None, line=None)
    c.name_asset(cont, asset_id)
    # 단계 헤더
    c.add_box(slide, lx, hy, lw, hh, fill=R("header_fill"), line=CC["white"], line_w=1.0)
    c.add_text(slide, lx, hy, lw, hh, "트랙 / 단계", size=10, bold=True,
               color=R("header_text"), align=PP_ALIGN.CENTER)
    for j, ph in enumerate(phases):
        px = gx0 + j * colw
        c.add_box(slide, px, hy, colw, hh, fill=CC["navy_600"], line=CC["white"], line_w=1.0)
        c.add_text(slide, px, hy, colw, hh, ph, size=10, bold=True,
                   color=R("header_text"), align=PP_ALIGN.CENTER)
    # 단계 구분 세로 커넥터
    lanes_bottom = ry0 + lane_h * len(tracks) + gap * (len(tracks) - 1)
    for j in range(1, npc):
        gxj = gx0 + j * colw
        c.connector(slide, gxj, ry0, gxj, lanes_bottom, color=CC["gray_300"], w=0.75)
    # 트랙 레인
    for ti, (tname, ci, acts) in enumerate(tracks):
        ly = ry0 + ti * (lane_h + gap)
        color = PALETTE[ci % len(PALETTE)]
        # 트랙명 박스
        c.add_box(slide, lx, ly, lw, lane_h, fill=CC["gray_050"], line=color, line_w=1.5,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.add_text(slide, lx + 0.1, ly, lw - 0.2, lane_h, tname, size=12, bold=True,
                   color=color, align=PP_ALIGN.CENTER)
        # 레인 배경 + 가로 기준선
        c.add_box(slide, gx0, ly, gw, lane_h, fill=CC["gray_100"], line=None)
        c.connector(slide, gx0, ly + lane_h / 2, gx0 + gw, ly + lane_h / 2, color=color, w=1.25)
        # 활동 박스(단계별)
        for j, act in enumerate(acts):
            if not act:
                continue
            ax = gx0 + j * colw + 0.1
            box = c.add_box(slide, ax, ly + 0.22, colw - 0.2, lane_h - 0.44,
                            fill=R("row_base"), line=color, line_w=1.25,
                            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            c.set_shape_text(box, act, size=9, bold=False, color=R("body_text"))


# =========================================================
# TML-003 단계별 로드맵 (1~4단계 화살표 밴드 + 각 단계 산출물 불릿)
# =========================================================
def phase_roadmap(slide, asset_id, title, sub, phases):
    """phases: list of dict {no, name, items, color_idx}."""
    title_block(slide, title, sub)
    x0 = 1.1
    total_w = 11.1
    y_band = 1.7
    band_h = 1.1
    n = len(phases)
    gap = 0.12
    bw = (total_w - gap * (n - 1)) / n
    cont = c.add_box(slide, x0 - 0.05, y_band - 0.05, total_w + 0.1, 4.7, fill=None, line=None)
    c.name_asset(cont, asset_id)
    # 단계 진행 기준선(밴드 하단)
    c.connector(slide, x0, y_band + band_h + 0.1, x0 + total_w, y_band + band_h + 0.1,
                color=R("accent_primary"), w=1.5)
    for i, ph in enumerate(phases):
        cx = x0 + i * (bw + gap)
        color = PALETTE[ph["color_idx"] % len(PALETTE)]
        # 화살표 밴드(마지막은 사각, 나머지 CHEVRON)
        shape = MSO_SHAPE.CHEVRON if i < n - 1 else MSO_SHAPE.PENTAGON
        band = c.add_box(slide, cx, y_band, bw, band_h, fill=color, line=None, shape=shape)
        tf = band.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "%s단계" % ph["no"]; c.set_kfont(r, c.FONT_H, 15, True, R("header_text"))
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = ph["name"]; c.set_kfont(r2, c.FONT_B, 10, False, R("header_text"))
        # 산출물 카드(밴드 아래)
        card_y = y_band + band_h + 0.25
        card = c.add_box(slide, cx, card_y, bw, 3.05, fill=CC["gray_050"], line=color, line_w=1.25,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.add_text(slide, cx + 0.08, card_y + 0.08, bw - 0.16, 0.32, "주요 산출물", size=10,
                   bold=True, color=color, align=PP_ALIGN.CENTER)
        for k, it in enumerate(ph["items"]):
            iy = card_y + 0.5 + k * 0.56
            c.add_box(slide, cx + 0.18, iy + 0.1, 0.1, 0.1, fill=color, line=None, shape=MSO_SHAPE.OVAL)
            c.add_text(slide, cx + 0.38, iy - 0.04, bw - 0.5, 0.5, it, size=10,
                       color=R("body_text"), align=PP_ALIGN.LEFT)


# =========================================================
# TML-004 분기별 로드맵 (Q1~Q4 컬럼 + 활동 카드)
# =========================================================
def quarter_roadmap(slide, asset_id, title, sub, quarters):
    """quarters: list of dict {q, theme, cards, color_idx}."""
    title_block(slide, title, sub)
    x0 = 1.1
    total_w = 11.1
    y0 = 1.7
    n = len(quarters)
    gap = 0.25
    cw = (total_w - gap * (n - 1)) / n
    col_h = 4.75
    cont = c.add_box(slide, x0 - 0.05, y0 - 0.05, total_w + 0.1, col_h + 0.1, fill=None, line=None)
    c.name_asset(cont, asset_id)
    # 컬럼 하단 기준선
    c.connector(slide, x0, y0 + col_h + 0.05, x0 + total_w, y0 + col_h + 0.05,
                color=R("accent_primary"), w=1.5)
    for i, q in enumerate(quarters):
        cx = x0 + i * (cw + gap)
        color = PALETTE[q["color_idx"] % len(PALETTE)]
        # 컬럼 헤더
        hd = c.add_box(slide, cx, y0, cw, 0.95, fill=color, line=None,
                       shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        tf = hd.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = q["q"]; c.set_kfont(r, c.FONT_H, 16, True, R("header_text"))
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = q["theme"]; c.set_kfont(r2, c.FONT_B, 9, False, R("header_text"))
        # 컬럼 바디
        c.add_box(slide, cx, y0 + 0.95, cw, col_h - 0.95, fill=CC["gray_050"],
                  line=R("border"), line_w=1.0)
        # 활동 카드 스택
        for k, card in enumerate(q["cards"]):
            ky = y0 + 1.15 + k * 0.92
            cb = c.add_box(slide, cx + 0.18, ky, cw - 0.36, 0.74, fill=R("row_base"),
                           line=color, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            c.set_shape_text(cb, card, size=10, bold=False, color=R("body_text"))


# =========================================================
# TML-007 연차별 발전 로드맵 (계단 상승형)
# =========================================================
def stair_roadmap(slide, asset_id, title, sub, steps):
    """steps: list of dict {year, focus, items, color_idx}. 계단식 상승 배치."""
    title_block(slide, title, sub)
    x0 = 1.15
    total_w = 11.15
    n = len(steps)
    gap = 0.22
    bw = (total_w - gap * (n - 1)) / n
    base_bottom = 6.35        # 최하단 기준선
    top_y = 1.75              # 최상단(마지막 계단) 상단
    rise = (base_bottom - (top_y + 1.4)) / (n - 1) if n > 1 else 0
    cont = c.add_box(slide, x0 - 0.05, top_y - 0.05, total_w + 0.1, (base_bottom - top_y) + 0.1,
                     fill=None, line=None)
    c.name_asset(cont, asset_id)
    prev = None
    for i, st in enumerate(steps):
        cx = x0 + i * (bw + gap)
        color = PALETTE[st["color_idx"] % len(PALETTE)]
        step_top = top_y + (n - 1 - i) * rise
        step_h = base_bottom - step_top
        # 계단 기둥(연한 배경)
        c.add_box(slide, cx, step_top, bw, step_h, fill=CC["gray_100"],
                  line=R("border"), line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 상단 연차 헤더 밴드
        hd = c.add_box(slide, cx, step_top, bw, 0.68, fill=color, line=None,
                       shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        c.set_shape_text(hd, st["year"], size=14, bold=True, color=R("header_text"))
        # 포커스 + 항목
        c.add_text(slide, cx + 0.1, step_top + 0.74, bw - 0.2, 0.36, st["focus"], size=11,
                   bold=True, color=color, align=PP_ALIGN.CENTER)
        for k, it in enumerate(st["items"]):
            iy = step_top + 1.16 + k * 0.5
            if iy + 0.4 > base_bottom:
                break
            c.add_box(slide, cx + 0.2, iy + 0.09, 0.09, 0.09, fill=color, line=None, shape=MSO_SHAPE.OVAL)
            c.add_text(slide, cx + 0.38, iy - 0.05, bw - 0.5, 0.44, it, size=9,
                       color=R("body_text"), align=PP_ALIGN.LEFT)
        # 상승 화살표 커넥터(이전 계단 → 현재 계단)
        if prev is not None:
            c.connector(slide, prev[0], prev[1], cx + 0.12, step_top, color=R("accent_point"), w=2.0)
        prev = (cx + bw - 0.12, step_top)


# =========================================================
# TML-008 추진일정 요약 (단계 + 주요활동 + 산출물 3열 매핑 표형)
# =========================================================
def schedule_summary(slide, asset_id, title, sub, rows):
    """rows: list of (단계, 주요활동, 산출물, color_idx)."""
    title_block(slide, title, sub)
    x, y = 1.1, 1.7
    w = 11.1
    col_w = [2.9, 5.1, 3.1]
    nrows = len(rows) + 1
    rh = 0.82
    gf = slide.shapes.add_table(nrows, 3, Inches(x), Inches(y), Inches(w), Inches(rh * nrows))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    for i in range(nrows):
        tbl.rows[i].height = Inches(rh)
    c.name_asset(gf, asset_id)

    def cell(cc, text, size=11, bold=False, color=None, fill=None, align=PP_ALIGN.CENTER):
        cc.fill.solid(); cc.fill.fore_color.rgb = fill if fill is not None else R("row_base")
        cc.vertical_anchor = MSO_ANCHOR.MIDDLE
        cc.margin_left = cc.margin_right = Pt(6); cc.margin_top = cc.margin_bottom = Pt(3)
        tf = cc.text_frame; tf.word_wrap = True
        for k, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run(); r.text = line
            c.set_kfont(r, c.FONT_B, size, bold, color or R("body_text"))

    for j, hd in enumerate(["추진 단계", "주요 활동", "산출물"]):
        cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=R("header_fill"))
    for i, (stage, act, deliv, ci) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        color = PALETTE[ci % len(PALETTE)]
        cell(tbl.cell(i, 0), stage, size=11, bold=True, color=color, fill=base)
        cell(tbl.cell(i, 1), act, size=10, fill=base, align=PP_ALIGN.LEFT, color=R("body_text"))
        cell(tbl.cell(i, 2), deliv, size=10, bold=True, fill=base, color=R("muted_text"))
    # 단계 열 강조 좌측 라인
    c.connector(slide, x, y, x, y + rh * nrows, color=R("accent_primary"), w=2.5)


# =========================================================
# 빌드
# =========================================================
entries = []


def E(*a, **k):
    entries.append(c.entry(*a, **k))


F1 = "decks/05_timeline/TML_timeline_v1.pptx"
F2 = "decks/05_timeline/TML_roadmap_v1.pptx"

# ---- FILE 1: 타임라인 계열 (TML-001, 002, 005, 006) ----
p1 = c.new_deck()

s = c.blank_slide(p1); c.id_caption(s, "TML-001")
horizontal_timeline(s, "TML-001", "수평 타임라인 (마일스톤 5개)",
                    "가로 기준선 위 마일스톤 원 5개, 라벨 상하 교차 배치",
                    [{"label": "착수·분석", "when": "M1"},
                     {"label": "설계", "when": "M2"},
                     {"label": "구축", "when": "M3"},
                     {"label": "시험·안정화", "when": "M4"},
                     {"label": "최종납품", "when": "M5"}])

s = c.blank_slide(p1); c.id_caption(s, "TML-002")
gantt_chart(s, "TML-002", "간트차트형 일정표 (작업 4행)",
            "작업별 기간 막대를 도형으로 표현 (차트 아님)",
            ["1분기", "2분기", "3분기", "4분기", "5분기", "6분기"],
            [("착수·분석", 0, 2, 0),
             ("설계", 1, 2, 1),
             ("구축", 2, 3, 2),
             ("시험·안정화", 4, 2, 3)])

s = c.blank_slide(p1); c.id_caption(s, "TML-005")
flag_timeline(s, "TML-005", "마일스톤 플래그 타임라인",
              "수평 기준선 위 깃발 도형으로 주요 마일스톤 표기",
              [{"label": "사업 착수", "when": "M1"},
               {"label": "중간보고", "when": "M3"},
               {"label": "시범 운영", "when": "M4"},
               {"label": "최종납품", "when": "M6"}])

s = c.blank_slide(p1); c.id_caption(s, "TML-006")
three_track(s, "TML-006", "로드맵 3트랙 (트랙별 레인)",
            "트랙별 가로 레인 위 단계별 활동 박스 배치",
            ["1단계", "2단계", "3단계", "4단계"],
            [("기획", 0, ["요구 분석", "과제 정의", "", ""]),
             ("구축", 1, ["", "아키텍처 설계", "시스템 개발", "통합 시험"]),
             ("운영", 2, ["", "", "시범 운영", "정식 운영·이관"])])

f1 = c.save_deck(p1, F1)

# ---- FILE 2: 로드맵 계열 (TML-003, 004, 007, 008) ----
p2 = c.new_deck()

s = c.blank_slide(p2); c.id_caption(s, "TML-003")
phase_roadmap(s, "TML-003", "단계별 로드맵 (4단계 화살표)",
              "1~4단계 화살표 밴드와 단계별 산출물 불릿",
              [{"no": 1, "name": "착수·분석", "color_idx": 0,
                "items": ["착수보고", "현황 분석서", "요구사항 정의서"]},
               {"no": 2, "name": "설계", "color_idx": 1,
                "items": ["기본 설계서", "상세 설계서", "화면 정의서"]},
               {"no": 3, "name": "구축", "color_idx": 2,
                "items": ["개발 산출물", "단위 시험 결과", "중간보고"]},
               {"no": 4, "name": "시험·안정화", "color_idx": 3,
                "items": ["통합 시험서", "안정화 보고", "최종납품"]}])

s = c.blank_slide(p2); c.id_caption(s, "TML-004")
quarter_roadmap(s, "TML-004", "분기별 로드맵 (Q1~Q4)",
                "분기 컬럼별 주제와 활동 카드 스택",
                [{"q": "Q1", "theme": "착수·분석", "color_idx": 0,
                  "cards": ["사업 착수", "현황 분석", "요구 정의"]},
                 {"q": "Q2", "theme": "설계", "color_idx": 1,
                  "cards": ["기본 설계", "상세 설계", "프로토타입"]},
                 {"q": "Q3", "theme": "구축", "color_idx": 2,
                  "cards": ["시스템 개발", "단위 시험", "중간보고"]},
                 {"q": "Q4", "theme": "안정화", "color_idx": 3,
                  "cards": ["통합 시험", "안정화", "최종납품"]}])

s = c.blank_slide(p2); c.id_caption(s, "TML-007")
stair_roadmap(s, "TML-007", "연차별 발전 로드맵 (계단 상승형)",
              "연차가 오를수록 계단이 상승하는 성장 로드맵",
              [{"year": "1년차", "focus": "기반 구축", "color_idx": 0,
                "items": ["체계 수립", "기반 시스템"]},
               {"year": "2년차", "focus": "확산·고도화", "color_idx": 1,
                "items": ["기능 확대", "고도화 개발"]},
               {"year": "3년차", "focus": "고도화·성과", "color_idx": 2,
                "items": ["전면 확산", "성과 창출"]}])

s = c.blank_slide(p2); c.id_caption(s, "TML-008")
schedule_summary(s, "TML-008", "추진일정 요약 (단계·활동·산출물)",
                 "추진 단계별 주요 활동과 산출물을 3열로 매핑",
                 [("1단계 착수·분석", "착수보고 · 현황 분석 · 요구사항 정의", "요구사항 정의서", 0),
                  ("2단계 설계", "기본·상세 설계 · 화면 정의", "설계서 일체", 1),
                  ("3단계 구축", "시스템 개발 · 단위 시험 · 중간보고", "개발 산출물 · 중간보고서", 2),
                  ("4단계 시험·안정화", "통합 시험 · 안정화 · 이관 교육", "시험 결과서 · 안정화 보고", 3),
                  ("5단계 최종납품", "최종 검수 · 완료보고 · 산출물 이관", "완료보고서 · 최종산출물", 4)])

f2 = c.save_deck(p2, F2)

# =========================================================
# 매니페스트 엔트리
# =========================================================
TAGS = ["타임라인", "로드맵", "마일스톤", "간트", "추진일정"]

E("TML-001", "timeline", "수평 타임라인 (마일스톤 5개)", F1, 1,
  ["타임라인", "수평", "마일스톤", "추진일정"],
  {"orientation": "horizontal", "count": 5, "label_alt": True},
  {"milestones": [{"label": "착수·분석", "when": "M1"}, {"label": "설계", "when": "M2"},
                  {"label": "구축", "when": "M3"}, {"label": "시험·안정화", "when": "M4"},
                  {"label": "최종납품", "when": "M5"}],
   "count": 5, "orientation": "horizontal"},
  ["milestone-text", "count", "color", "when-label"],
  recommended_use=["추진일정", "단계별 계획", "로드맵"])

E("TML-002", "timeline", "간트차트형 일정표 (작업 4행)", F1, 2,
  ["타임라인", "간트", "일정표", "추진일정", "막대"],
  {"tasks": 4, "periods": 6, "shape_based": True},
  {"periods": ["1분기", "2분기", "3분기", "4분기", "5분기", "6분기"],
   "tasks": [{"name": "착수·분석", "start": 0, "span": 2}, {"name": "설계", "start": 1, "span": 2},
             {"name": "구축", "start": 2, "span": 3}, {"name": "시험·안정화", "start": 4, "span": 2}],
   "count": 4, "orientation": "horizontal"},
  ["milestone-text", "count", "bar-length", "color"],
  recommended_use=["추진일정", "작업 일정", "로드맵"])

E("TML-003", "timeline", "단계별 로드맵 (4단계 화살표)", F2, 1,
  ["로드맵", "단계별", "화살표", "산출물", "추진일정"],
  {"phases": 4, "band": "chevron", "items_per_phase": 3},
  {"phases": [{"no": 1, "name": "착수·분석"}, {"no": 2, "name": "설계"},
              {"no": 3, "name": "구축"}, {"no": 4, "name": "시험·안정화"}],
   "count": 4, "orientation": "horizontal"},
  ["milestone-text", "count", "color", "bar-length"],
  recommended_use=["단계별 계획", "로드맵", "추진일정"])

E("TML-004", "timeline", "분기별 로드맵 (Q1~Q4)", F2, 2,
  ["로드맵", "분기별", "컬럼", "활동카드", "추진일정"],
  {"quarters": 4, "cards_per_q": 3},
  {"quarters": [{"q": "Q1", "theme": "착수·분석"}, {"q": "Q2", "theme": "설계"},
                {"q": "Q3", "theme": "구축"}, {"q": "Q4", "theme": "안정화"}],
   "count": 4, "orientation": "horizontal"},
  ["milestone-text", "count", "color", "bar-length"],
  recommended_use=["분기 계획", "로드맵", "추진일정"])

E("TML-005", "timeline", "마일스톤 플래그 타임라인", F1, 3,
  ["타임라인", "마일스톤", "플래그", "깃발", "추진일정"],
  {"orientation": "horizontal", "count": 4, "flag": True},
  {"milestones": [{"label": "사업 착수", "when": "M1"}, {"label": "중간보고", "when": "M3"},
                  {"label": "시범 운영", "when": "M4"}, {"label": "최종납품", "when": "M6"}],
   "count": 4, "orientation": "horizontal"},
  ["milestone-text", "count", "color", "when-label"],
  recommended_use=["추진일정", "주요 마일스톤", "로드맵"])

E("TML-006", "timeline", "로드맵 3트랙 (트랙별 레인)", F1, 4,
  ["로드맵", "트랙", "레인", "활동박스", "추진일정"],
  {"tracks": 3, "phases": 4},
  {"phases": ["1단계", "2단계", "3단계", "4단계"],
   "tracks": [{"name": "기획"}, {"name": "구축"}, {"name": "운영"}],
   "count": 3, "orientation": "horizontal"},
  ["milestone-text", "count", "color", "bar-length"],
  recommended_use=["로드맵", "추진체계", "단계별 계획"])

E("TML-007", "timeline", "연차별 발전 로드맵 (계단 상승형)", F2, 3,
  ["로드맵", "연차별", "계단", "성장", "추진일정"],
  {"steps": 3, "shape": "stair"},
  {"steps": [{"year": "1년차", "focus": "기반 구축"}, {"year": "2년차", "focus": "확산·고도화"},
             {"year": "3년차", "focus": "고도화·성과"}],
   "count": 3, "orientation": "horizontal"},
  ["milestone-text", "count", "color", "bar-length"],
  recommended_use=["중장기 계획", "로드맵", "발전 방향"])

E("TML-008", "timeline", "추진일정 요약 (단계·활동·산출물)", F2, 4,
  ["추진일정", "로드맵", "표형", "산출물", "매핑"],
  {"cols": 3, "rows": 5, "table": True},
  {"rows": [{"stage": "1단계 착수·분석", "deliverable": "요구사항 정의서"},
            {"stage": "2단계 설계", "deliverable": "설계서 일체"},
            {"stage": "3단계 구축", "deliverable": "개발 산출물"},
            {"stage": "4단계 시험·안정화", "deliverable": "시험 결과서"},
            {"stage": "5단계 최종납품", "deliverable": "완료보고서"}],
   "count": 5, "orientation": "table"},
  ["milestone-text", "count", "color", "bar-length"],
  recommended_use=["추진일정", "단계별 계획", "로드맵"])

frag = c.write_fragment("TML", entries)

print("FILE1", f1)
print("FILE2", f2)
print("FRAG", frag)
print("ENTRIES", len(entries))
