# -*- coding: utf-8 -*-
"""
차트 스타일(CHT) 카테고리 생성기 (CHT-001 ~ CHT-010)
python-pptx 네이티브 차트(add_chart)만 사용 — 페이지 이미지/도형 근사 금지.
조합 파이프라인이 bindings 스키마로 데이터만 교체(스타일 유지)할 수 있도록 실제
차트 오브젝트(GraphicFrame + c:chart 관계 + ppt/charts/ 파트)를 생성한다.
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# 계열 색 팔레트 (네이비/블루/틸/그레이/시안) — 하드코딩 색 금지, 토큰만 사용
PAL = [c.C["navy_800"], c.C["blue_500"], c.C["teal_500"], c.C["gray_500"],
       c.C["cyan_500"], c.C["navy_600"], c.C["gray_700"], c.C["purple_600"]]

FILL_TYPES = {"COLUMN_CLUSTERED", "COLUMN_STACKED", "COLUMN_STACKED_100",
              "BAR_CLUSTERED", "AREA", "AREA_STACKED"}
LINE_TYPES = {"LINE", "LINE_MARKERS", "RADAR", "RADAR_MARKERS"}
PIE_TYPES = {"PIE", "DOUGHNUT"}

# 차트 배치
CX, CY, CW, CH = 0.6, 1.35, 12.1, 5.7


def kfont(font, name=None, size=None, bold=None, color=None):
    """차트 Font 오브젝트에 한글 안전 폰트(latin/ea/cs) 지정."""
    name = name or c.FONT_B
    if size is not None:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if color is not None:
        font.color.rgb = color
    font.name = name
    rPr = font._element  # defRPr / rPr 요소
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def set_title(chart, text):
    chart.has_title = True
    tf = chart.chart_title.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    c.set_kfont(run, c.FONT_H, 15, True, c.role("body_text"))


def set_legend(chart, position=XL_LEGEND_POSITION.BOTTOM, show=True):
    chart.has_legend = show
    if show:
        chart.legend.position = position
        chart.legend.include_in_layout = False
        kfont(chart.legend.font, c.FONT_B, 11, False, c.role("muted_text"))


def style_axes(chart):
    for ax_name in ("category_axis", "value_axis"):
        try:
            ax = getattr(chart, ax_name)
            kfont(ax.tick_labels.font, c.FONT_B, 10, False, c.role("muted_text"))
        except Exception:
            pass


def color_series(chart, ctype):
    plot = chart.plots[0]
    if ctype in PIE_TYPES:
        ser = plot.series[0]
        for j, pt in enumerate(ser.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = PAL[j % len(PAL)]
    elif ctype in LINE_TYPES:
        for i, ser in enumerate(plot.series):
            ser.format.line.color.rgb = PAL[i % len(PAL)]
            ser.format.line.width = Pt(2.5)
    else:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = PAL[i % len(PAL)]


def add_labels(chart, ctype):
    plot = chart.plots[0]
    try:
        if ctype in PIE_TYPES:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_percentage = True
            dl.show_value = False
            dl.number_format = "0%"
            dl.number_format_is_linked = False
            dl.position = XL_LABEL_POSITION.OUTSIDE_END if ctype == "PIE" else XL_LABEL_POSITION.CENTER
            kfont(dl.font, c.FONT_B, 10, True, c.role("body_text"))
        else:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = "#,##0"
            dl.number_format_is_linked = False
            kfont(dl.font, c.FONT_B, 9, False, c.role("body_text"))
    except Exception:
        pass


def make_chart(slide, asset_id, ctype, title, categories, series,
               legend_pos=XL_LEGEND_POSITION.BOTTOM, show_legend=True,
               labels=True):
    """네이티브 차트 1개 생성 + 스타일 적용. GraphicFrame 반환."""
    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(s["name"], s["values"])
    xl = getattr(XL_CHART_TYPE, ctype)
    gf = slide.shapes.add_chart(xl, Inches(CX), Inches(CY), Inches(CW), Inches(CH), cd)
    chart = gf.chart
    kfont(chart.font, c.FONT_B, 11, False, c.role("body_text"))  # 차트 기본 폰트
    set_title(chart, title)
    set_legend(chart, legend_pos, show_legend)
    style_axes(chart)
    color_series(chart, ctype)
    if labels:
        add_labels(chart, ctype)
    c.name_asset(gf, asset_id)
    return gf


def meta(asset_id, name, file_rel, slide_idx, ctype, categories, series,
         tags, params, recommended, fallback_note=None):
    bindings = {
        "categories": categories,
        "series": [{"name": s["name"], "values": list(s["values"])} for s in series],
        "chart_type": ctype,
    }
    if fallback_note:
        bindings["_note"] = fallback_note
    return c.entry(
        asset_id, "CHT", name, file_rel, slide_idx,
        tags=tags, params=params, bindings=bindings,
        editable=["data", "series-color", "title", "legend"],
        recommended_use=recommended,
    )


entries = []

# ============================================================
# FILE 1: CHT_bar_v1.pptx  (막대류 5)
# ============================================================
prs1 = c.new_deck()
F1 = "decks/08_charts/CHT_bar_v1.pptx"

# ---- CHT-001: 세로 막대 (단일 계열) ----
s = c.blank_slide(prs1); c.id_caption(s, "CHT-001")
cats = ["게임", "음악", "방송", "영화", "캐릭터"]
ser = [{"name": "2024 수출액(백만$)", "values": (6820, 1210, 980, 640, 745)}]
make_chart(s, "CHT-001", "COLUMN_CLUSTERED", "장르별 콘텐츠 수출액", cats, ser,
           show_legend=False)
entries.append(meta("CHT-001", "세로 막대(단일 계열)", F1, 1, "COLUMN_CLUSTERED",
                    cats, ser, ["차트", "막대", "세로막대", "단일계열"],
                    {"type": "COLUMN_CLUSTERED", "series": 1},
                    ["실적 비교", "장르별 성과", "단일 지표"]))

# ---- CHT-002: 그룹 세로 막대 (2계열, 2023 vs 2024) ----
s = c.blank_slide(prs1); c.id_caption(s, "CHT-002")
cats = ["아시아", "북미", "유럽", "중남미", "중동"]
ser = [{"name": "2023", "values": (3120, 1580, 940, 210, 160)},
       {"name": "2024", "values": (3880, 1720, 1120, 260, 205)}]
make_chart(s, "CHT-002", "COLUMN_CLUSTERED", "권역별 수출 실적 (전년 대비)", cats, ser)
entries.append(meta("CHT-002", "그룹 세로 막대(2계열)", F1, 2, "COLUMN_CLUSTERED",
                    cats, ser, ["차트", "막대", "그룹막대", "비교", "2계열"],
                    {"type": "COLUMN_CLUSTERED", "series": 2},
                    ["연도 비교", "전년 대비", "권역별 비교"]))

# ---- CHT-003: 누적 세로 막대 ----
s = c.blank_slide(prs1); c.id_caption(s, "CHT-003")
cats = ["2022", "2023", "2024"]
ser = [{"name": "아시아", "values": (2800, 3120, 3880)},
       {"name": "북미", "values": (1400, 1580, 1720)},
       {"name": "유럽", "values": (820, 940, 1120)}]
make_chart(s, "CHT-003", "COLUMN_STACKED", "연도별 권역 수출 누적", cats, ser)
entries.append(meta("CHT-003", "누적 세로 막대", F1, 3, "COLUMN_STACKED",
                    cats, ser, ["차트", "막대", "누적막대", "구성"],
                    {"type": "COLUMN_STACKED", "series": 3},
                    ["누적 추세", "구성 합계", "연도별 성장"]))

# ---- CHT-004: 가로 막대 (순위형) ----
s = c.blank_slide(prs1); c.id_caption(s, "CHT-004")
cats = ["일본", "미국", "중국", "대만", "베트남"]
ser = [{"name": "수출액(백만$)", "values": (1980, 1720, 1540, 620, 480)}]
make_chart(s, "CHT-004", "BAR_CLUSTERED", "국가별 콘텐츠 수출 순위", cats, ser,
           show_legend=False)
entries.append(meta("CHT-004", "가로 막대(순위형)", F1, 4, "BAR_CLUSTERED",
                    cats, ser, ["차트", "막대", "가로막대", "순위"],
                    {"type": "BAR_CLUSTERED", "series": 1},
                    ["순위 비교", "국가별 실적", "랭킹"]))

# ---- CHT-008: 영역(AREA) — 콤보 대체(막대류 파일 5번째) ----
s = c.blank_slide(prs1); c.id_caption(s, "CHT-008")
cats = ["2020", "2021", "2022", "2023", "2024"]
ser = [{"name": "총 수출액(백만$)", "values": (7020, 8080, 9250, 10360, 11840)}]
make_chart(s, "CHT-008", "AREA", "콘텐츠 수출 총액 추세", cats, ser,
           show_legend=False)
entries.append(meta("CHT-008", "영역 차트(AREA·콤보 대체)", F1, 5, "AREA",
                    cats, ser, ["차트", "영역", "추세", "콤보대체"],
                    {"type": "AREA", "series": 1,
                     "note": "python-pptx 막대+선 콤보 미지원 → AREA로 안전 대체"},
                    ["누적 추세", "성장 곡선", "총량 변화"],
                    fallback_note="막대+선 콤보는 python-pptx 미지원 → AREA로 대체"))

f1_path = c.save_deck(prs1, F1)

# ============================================================
# FILE 2: CHT_line-pie_v1.pptx  (선/원/도넛/구성비/레이더 5)
# ============================================================
prs2 = c.new_deck()
F2 = "decks/08_charts/CHT_line-pie_v1.pptx"

# ---- CHT-005: 꺾은선 (추세, 2계열) ----
s = c.blank_slide(prs2); c.id_caption(s, "CHT-005")
cats = ["2020", "2021", "2022", "2023", "2024"]
ser = [{"name": "게임", "values": (4200, 4900, 5600, 6200, 6820)},
       {"name": "음악", "values": (560, 720, 900, 1050, 1210)}]
make_chart(s, "CHT-005", "LINE_MARKERS", "장르별 수출 추세", cats, ser)
entries.append(meta("CHT-005", "꺾은선(추세·2계열)", F2, 1, "LINE_MARKERS",
                    cats, ser, ["차트", "꺾은선", "라인", "추세", "2계열"],
                    {"type": "LINE_MARKERS", "series": 2},
                    ["시계열 추세", "성장 추이", "장르 비교"]))

# ---- CHT-006: 원형 (PIE) ----
s = c.blank_slide(prs2); c.id_caption(s, "CHT-006")
cats = ["게임", "음악", "방송", "영화", "기타"]
ser = [{"name": "장르 구성비", "values": (63, 11, 9, 6, 11)}]
make_chart(s, "CHT-006", "PIE", "장르별 수출 구성비", cats, ser,
           legend_pos=XL_LEGEND_POSITION.RIGHT)
entries.append(meta("CHT-006", "원형(PIE)", F2, 2, "PIE",
                    cats, ser, ["차트", "원형", "파이", "구성비"],
                    {"type": "PIE", "series": 1},
                    ["구성비", "점유율", "비중"]))

# ---- CHT-007: 도넛 (DOUGHNUT) ----
s = c.blank_slide(prs2); c.id_caption(s, "CHT-007")
cats = ["아시아", "북미", "유럽", "기타"]
ser = [{"name": "권역 구성비", "values": (56, 25, 14, 5)}]
make_chart(s, "CHT-007", "DOUGHNUT", "권역별 수출 구성비", cats, ser,
           legend_pos=XL_LEGEND_POSITION.RIGHT)
entries.append(meta("CHT-007", "도넛(DOUGHNUT)", F2, 3, "DOUGHNUT",
                    cats, ser, ["차트", "도넛", "구성비", "점유율"],
                    {"type": "DOUGHNUT", "series": 1},
                    ["구성비", "권역 점유율", "비중"]))

# ---- CHT-009: 100% 누적 막대 (구성비) ----
s = c.blank_slide(prs2); c.id_caption(s, "CHT-009")
cats = ["2022", "2023", "2024"]
ser = [{"name": "아시아", "values": (2800, 3120, 3880)},
       {"name": "북미", "values": (1400, 1580, 1720)},
       {"name": "유럽", "values": (820, 940, 1120)}]
make_chart(s, "CHT-009", "COLUMN_STACKED_100", "연도별 권역 구성비 (100%)", cats, ser)
entries.append(meta("CHT-009", "100% 누적 막대(구성비)", F2, 4, "COLUMN_STACKED_100",
                    cats, ser, ["차트", "막대", "구성비", "100%누적"],
                    {"type": "COLUMN_STACKED_100", "series": 3},
                    ["구성비 추세", "비중 변화", "권역 점유"]))

# ---- CHT-010: 레이더 (RADAR) — 미지원/불안정 시 막대 대체 ----
s = c.blank_slide(prs2); c.id_caption(s, "CHT-010")
cats = ["기획", "제작", "유통", "마케팅", "현지화"]
ser = [{"name": "현재 역량", "values": (3, 4, 2, 3, 2)},
       {"name": "목표 역량", "values": (5, 5, 4, 4, 5)}]
radar_type = "RADAR_MARKERS"
radar_note = None
try:
    make_chart(s, "CHT-010", radar_type, "수출 역량 진단 (레이더)", cats, ser,
               labels=False)
except Exception as e:
    # 레이더 실패 → 그룹 막대로 대체
    for sh in list(s.shapes):
        if sh.has_chart:
            sh._element.getparent().remove(sh._element)
    radar_type = "COLUMN_CLUSTERED"
    radar_note = "RADAR 생성 실패(%s) → COLUMN_CLUSTERED로 대체" % type(e).__name__
    make_chart(s, "CHT-010", radar_type, "수출 역량 진단 (막대 대체)", cats, ser)
entries.append(meta("CHT-010", "레이더(RADAR)", F2, 5, radar_type,
                    cats, ser, ["차트", "레이더", "역량진단", "다축비교"],
                    {"type": radar_type, "series": 2},
                    ["역량 진단", "다축 비교", "gap 분석"],
                    fallback_note=radar_note))

f2_path = c.save_deck(prs2, F2)

# ============================================================
# 매니페스트 조각 저장
# ============================================================
frag = c.write_fragment("CHT", entries)

print("F1:", f1_path)
print("F2:", f2_path)
print("FRAG:", frag)
print("ENTRIES:", len(entries))
print("RADAR_TYPE:", radar_type, "| note:", radar_note)
