# -*- coding: utf-8 -*-
"""
09_headers (HDR) 대량 확충 — HDR-025 ~ HDR-200 (+176).
파라미터화 배치 생성기: 구조패밀리 × accent색 × 변형 스윕.

원칙(공통 규칙 준수):
- 색은 c.C / c.role 만 사용(raw RGBColor 리터럴 금지). accent 스윕 6색.
- 각 에셋 = 자체완결 c.group_asset(<p:grpSp> asset:<ID>). id_caption은 그룹 밖(슬라이드당 1회).
- 페이지 전체 이미지·SmartArt·pic 금지. 전부 네이티브 shape/textbox.
- 빌더 함수는 gen_09_headers(v1)/gen_09_headers_v2(v2)의 검증된 구현을 그대로 복사.
- 파일은 패밀리별 분산: decks/09_headers/HDR_bulk_<family>_v1.pptx (각 ≤25슬라이드).
- 마지막에 c.write_fragment('HDR_bulk', entries) 1회.
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from collections import OrderedDict
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

R = c.role
RR = MSO_SHAPE.ROUNDED_RECTANGLE
RECT = MSO_SHAPE.RECTANGLE
OVAL = MSO_SHAPE.OVAL

# ═══════════════════════════════════════════════════════════════
# 빌더 함수 (v1/v2 검증 구현 복사) — 각 함수는 도형 리스트를 반환(그룹 대상)
# ═══════════════════════════════════════════════════════════════
def gradient_bar(slide, x, y, w, h, c1, c2):
    sp = c.add_box(slide, x, y, w, h, fill=c.C["white"], line=None, shape=RECT)
    sp.fill.gradient()
    gs = sp.fill.gradient_stops
    gs[0].color.rgb = c1; gs[0].position = 0.0
    gs[1].color.rgb = c2; gs[1].position = 1.0
    sp.fill.gradient_angle = 0.0
    return sp

def build_section_pill(slide, x, y, number, title, fill_color, w=4.6, h=0.6):
    shp = []
    pill_w = 0.6
    numbox = c.add_box(slide, x, y, pill_w, h, fill=fill_color, line=None, shape=RECT)
    c.set_shape_text(numbox, number, size=20, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(numbox)
    titlebox = c.add_box(slide, x + pill_w, y, w - pill_w, h,
                         fill=c.C["gray_100"], line=None, shape=RECT)
    c.set_shape_text(titlebox, title, size=15, bold=True, color=R("body_text"),
                     align=PP_ALIGN.LEFT, font=c.FONT_H)
    titlebox.text_frame.margin_left = Pt(10)
    shp.append(titlebox)
    return shp

def build_subtitle_bar(slide, x, y, marker, title, accent, w=4.6, h=0.42):
    shp = []
    mk = c.add_text(slide, x, y, 0.5, h, marker, size=14, bold=True,
                    color=accent, align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(mk)
    t = c.add_text(slide, x + 0.5, y, w - 0.5, h, title, size=14, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(t)
    ul = c.add_box(slide, x, y + h - 0.03, w, 0.05, fill=accent, line=None, shape=RECT)
    shp.append(ul)
    return shp

def build_check_bar(slide, x, y, text, fill_color, w=5.0, h=0.55):
    shp = []
    bar = c.add_box(slide, x, y, w, h, fill=fill_color, line=None, shape=RR)
    bar.adjustments[0] = 0.5
    shp.append(bar)
    ck = c.add_box(slide, x + 0.14, y + (h - 0.32) / 2, 0.32, 0.32,
                   fill=R("header_text"), line=None, shape=OVAL)
    c.set_shape_text(ck, "✓", size=13, bold=True, color=fill_color, align=PP_ALIGN.CENTER)
    shp.append(ck)
    t = c.add_text(slide, x + 0.58, y, w - 0.7, h, text, size=14, bold=True,
                   color=R("header_text"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(t)
    return shp

def build_tab_label(slide, x, y, text, fill_color, w=2.0, h=0.5):
    shp = []
    tab = c.add_box(slide, x, y, w, h, fill=fill_color, line=None,
                    shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    tab.adjustments[0] = 0.35
    c.set_shape_text(tab, text, size=13, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(tab)
    base = c.add_box(slide, x, y + h, w, 0.05, fill=fill_color, line=None, shape=RECT)
    shp.append(base)
    return shp

def build_callout(slide, x, y, title, body, accent, w=4.8, h=1.1):
    shp = []
    bg = c.add_box(slide, x, y, w, h, fill=c.C["gray_050"], line=c.C["gray_300"],
                   line_w=0.75, shape=RECT)
    shp.append(bg)
    bar = c.add_box(slide, x, y, 0.09, h, fill=accent, line=None, shape=RECT)
    shp.append(bar)
    t = c.add_text(slide, x + 0.28, y + 0.1, w - 0.4, 0.35, title, size=13, bold=True,
                   color=accent, align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(t)
    b = c.add_text(slide, x + 0.28, y + 0.45, w - 0.4, h - 0.55, body, size=11, bold=False,
                   color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    shp.append(b)
    return shp

def build_analysis_box(slide, x, y, title, bullets, bg, accent, w=5.4, h=2.2):
    shp = []
    box = c.add_box(slide, x, y, w, h, fill=bg, line=c.C["gray_300"], line_w=0.75, shape=RR)
    box.adjustments[0] = 0.04
    shp.append(box)
    tb = c.add_box(slide, x + 0.18, y + 0.16, w - 0.36, 0.42, fill=accent, line=None, shape=RECT)
    c.set_shape_text(tb, title, size=13, bold=True, color=R("header_text"),
                     align=PP_ALIGN.LEFT, font=c.FONT_H)
    tb.text_frame.margin_left = Pt(8)
    shp.append(tb)
    yy = y + 0.72
    for b in bullets:
        mk = c.add_text(slide, x + 0.28, yy, 0.3, 0.32, "•", size=12, bold=True,
                        color=accent, align=PP_ALIGN.LEFT)
        shp.append(mk)
        bt = c.add_text(slide, x + 0.52, yy, w - 0.7, 0.32, b, size=11,
                        color=R("body_text"), align=PP_ALIGN.LEFT)
        shp.append(bt)
        yy += 0.36
    return shp

def build_ribbon(slide, x, y, text, fill_color, w=1.7, h=0.72):
    shp = []
    top = c.add_box(slide, x + w / 2 - 0.5, y - 0.06, 1.0, 0.12, fill=fill_color,
                    line=None, shape=RECT)
    shp.append(top)
    rb = c.add_box(slide, x, y, w, h, fill=fill_color, line=None, shape=MSO_SHAPE.DOWN_RIBBON)
    c.set_shape_text(rb, text, size=14, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(rb)
    return shp

def build_pill_badge(slide, x, y, text, fill_color, w=1.3, h=0.44):
    shp = []
    dot = c.add_box(slide, x - 0.02, y + h / 2 - 0.06, 0.12, 0.12, fill=fill_color,
                    line=None, shape=OVAL)
    shp.append(dot)
    bgp = c.add_box(slide, x + 0.16, y, w, h, fill=fill_color, line=None, shape=RR)
    bgp.adjustments[0] = 0.5
    c.set_shape_text(bgp, text, size=13, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(bgp)
    return shp

def build_number_labels(slide, x, y, items, accent, w=5.6, h=0.6):
    shp = []
    step = w / len(items)
    for i, (num, txt) in enumerate(items):
        cx = x + i * step
        circ = c.add_box(slide, cx, y, 0.6, 0.6, fill=accent, line=None, shape=OVAL)
        c.set_shape_text(circ, num, size=16, bold=True, color=R("header_text"),
                         align=PP_ALIGN.CENTER, font=c.FONT_H)
        shp.append(circ)
        t = c.add_text(slide, cx + 0.66, y, step - 0.7, 0.6, txt, size=12, bold=True,
                       color=R("body_text"), align=PP_ALIGN.LEFT, font=c.FONT_H)
        shp.append(t)
    return shp

def build_agenda_rows(slide, x, y, items, accent, row_h=1.0, w=8.4):
    shp = []
    for i, (num, title, desc) in enumerate(items):
        ry = y + i * row_h
        circ = c.add_box(slide, x, ry + 0.06, 0.62, 0.62, fill=accent, line=None, shape=OVAL)
        c.set_shape_text(circ, num, size=18, bold=True, color=R("header_text"),
                         align=PP_ALIGN.CENTER, font=c.FONT_H)
        shp.append(circ)
        t = c.add_text(slide, x + 0.82, ry, w - 0.82, 0.42, title, size=15, bold=True,
                       color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM, font=c.FONT_H)
        shp.append(t)
        d = c.add_text(slide, x + 0.82, ry + 0.44, w - 0.82, 0.4, desc, size=11, bold=False,
                       color=R("muted_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        shp.append(d)
        if i < len(items) - 1:
            ln = c.add_box(slide, x + 0.82, ry + row_h - 0.06, w - 0.82, 0.012,
                           fill=c.C["gray_300"], line=None, shape=RECT)
            shp.append(ln)
    return shp

def build_quote_block(slide, x, y, quote, author, accent, w=6.6, h=1.9):
    shp = []
    panel = c.add_box(slide, x, y, w, h, fill=c.C["gray_050"], line=c.C["gray_300"],
                      line_w=0.75, shape=RR)
    panel.adjustments[0] = 0.045
    shp.append(panel)
    bar = c.add_box(slide, x, y + 0.12, 0.10, h - 0.24, fill=accent, line=None, shape=RECT)
    shp.append(bar)
    qm = c.add_text(slide, x + 0.24, y + 0.02, 1.0, 0.9, "“", size=54, bold=True,
                    color=accent, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=c.FONT_H)
    shp.append(qm)
    q = c.add_text(slide, x + 0.95, y + 0.24, w - 1.25, h - 0.85, quote, size=15, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=c.FONT_H)
    shp.append(q)
    au = c.add_text(slide, x + 0.95, y + h - 0.5, w - 1.25, 0.36, "— " + author, size=11,
                    bold=False, color=R("muted_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    shp.append(au)
    return shp

def build_stat_band(slide, x, y, number, unit, desc, fill_color, w=7.2, h=1.15):
    shp = []
    band = c.add_box(slide, x, y, w, h, fill=fill_color, line=None, shape=RR)
    band.adjustments[0] = 0.08
    shp.append(band)
    pt = c.add_box(slide, x + 0.02, y + 0.02, 0.12, h - 0.04, fill=R("accent_point"),
                   line=None, shape=RECT)
    shp.append(pt)
    numbox = c.add_text(slide, x + 0.32, y, 2.7, h, number, size=44, bold=True,
                        color=R("header_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=c.FONT_H)
    shp.append(numbox)
    ub = c.add_text(slide, x + 0.32, y + h - 0.42, 2.7, 0.34, unit, size=13, bold=True,
                    color=R("accent_point"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=c.FONT_H)
    shp.append(ub)
    db = c.add_text(slide, x + 3.1, y, w - 3.3, h, desc, size=14, bold=True,
                    color=R("header_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=c.FONT_H)
    shp.append(db)
    return shp

def build_dividers(slide, x, y, w, accent):
    shp = []
    grad = gradient_bar(slide, x, y, w, 0.06, accent, c.C["white"])
    shp.append(grad)
    dia = c.add_box(slide, x + w / 2 - 0.08, y - 0.05, 0.16, 0.16, fill=accent,
                    line=None, shape=MSO_SHAPE.DIAMOND)
    shp.append(dia)
    dy = y + 0.42
    n = 24
    gap = w / (n - 1)
    for i in range(n):
        dot = c.add_box(slide, x + i * gap - 0.035, dy, 0.07, 0.07, fill=c.C["gray_500"],
                        line=None, shape=OVAL)
        shp.append(dot)
    return shp

def build_ribbon_set(slide, x, y, items, w_each=1.55, h=0.56, gap=0.28):
    shp = []
    for i, (text, col) in enumerate(items):
        bx = x + i * (w_each + gap)
        tag = c.add_box(slide, bx, y, w_each, h, fill=col, line=None, shape=MSO_SHAPE.PENTAGON)
        c.set_shape_text(tag, text, size=13, bold=True, color=R("header_text"),
                         align=PP_ALIGN.CENTER, font=c.FONT_H)
        shp.append(tag)
    return shp

def build_tab_header(slide, x, y, tabs, active_idx, accent, tab_w=2.1, h=0.52, gap=0.12):
    shp = []
    total_w = len(tabs) * tab_w + (len(tabs) - 1) * gap
    for i, txt in enumerate(tabs):
        tx = x + i * (tab_w + gap)
        active = (i == active_idx)
        fill = accent if active else c.C["gray_100"]
        tab = c.add_box(slide, tx, y, tab_w, h, fill=fill,
                        line=None if active else c.C["gray_300"], line_w=0.75,
                        shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        tab.adjustments[0] = 0.28
        c.set_shape_text(tab, txt, size=12.5, bold=active,
                         color=R("header_text") if active else R("muted_text"),
                         align=PP_ALIGN.CENTER, font=c.FONT_H)
        shp.append(tab)
    base = c.add_box(slide, x, y + h, total_w, 0.05, fill=accent, line=None, shape=RECT)
    shp.append(base)
    return shp

def build_step_label(slide, x, y, step_no, kicker, title, accent, w=5.4):
    shp = []
    circ = c.add_box(slide, x, y, 0.86, 0.86, fill=accent, line=None, shape=OVAL)
    c.set_shape_text(circ, step_no, size=26, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(circ)
    k = c.add_text(slide, x + 1.06, y + 0.02, w - 1.06, 0.32, kicker, size=11, bold=True,
                   color=accent, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM, font=c.FONT_H)
    shp.append(k)
    t = c.add_text(slide, x + 1.06, y + 0.36, w - 1.06, 0.46, title, size=16, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=c.FONT_H)
    shp.append(t)
    return shp

def build_inline_kpi_header(slide, x, y, title, metric, metric_label, accent, w=8.6, h=0.9):
    shp = []
    t = c.add_text(slide, x, y, w * 0.58, h, title, size=17, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=c.FONT_H)
    shp.append(t)
    ul = c.add_box(slide, x + 0.02, y + h - 0.14, 2.2, 0.06, fill=accent, line=None, shape=RECT)
    shp.append(ul)
    mx = x + w * 0.58
    mval = c.add_text(slide, mx, y - 0.02, w * 0.42, 0.58, metric, size=30, bold=True,
                      color=accent, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM, font=c.FONT_H)
    shp.append(mval)
    mlab = c.add_text(slide, mx, y + 0.56, w * 0.42, 0.3, metric_label, size=11, bold=False,
                      color=R("muted_text"), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)
    shp.append(mlab)
    return shp

# ═══════════════════════════════════════════════════════════════
# 스윕 설정 — accent 6색 + 한글 라벨
# ═══════════════════════════════════════════════════════════════
COLORS = ["navy_800", "blue_500", "teal_500", "purple_600", "cyan_500", "red_500"]
KCOL = {"navy_800": "네이비", "blue_500": "블루", "teal_500": "틸",
        "purple_600": "퍼플", "cyan_500": "시안", "red_500": "레드"}

# ═══════════════════════════════════════════════════════════════
# 패밀리 정의 — render(slide,x,y,color_name,variant) -> (shapes, params, bindings)
# layout: cols / col_xs / y0 / pitch / per_slide
# variants: 색 외 파라미터(변형) 리스트; 각 variant dict는 'lab'(이름 접미) 포함 가능
# ═══════════════════════════════════════════════════════════════
def r_section(sl, x, y, col, v):
    shp = build_section_pill(sl, x, y, v["num"], v["title"], c.C[col])
    return shp, {"fill": col, "w_in": 4.6, "h_in": 0.6}, {"number": v["num"], "title": v["title"]}

def r_subtitle(sl, x, y, col, v):
    shp = build_subtitle_bar(sl, x, y, v["marker"], v["title"], c.C[col])
    return shp, {"accent": col, "w_in": 4.6, "h_in": 0.42}, {"marker": v["marker"], "title": v["title"]}

def r_check(sl, x, y, col, v):
    shp = build_check_bar(sl, x, y, v["text"], c.C[col])
    return shp, {"fill": col, "w_in": 5.0, "h_in": 0.55}, {"text": v["text"]}

def r_tablabel(sl, x, y, col, v):
    shp = build_tab_label(sl, x, y, v["text"], c.C[col])
    return shp, {"fill": col, "w_in": 2.0, "h_in": 0.5}, {"text": v["text"]}

def r_callout(sl, x, y, col, v):
    shp = build_callout(sl, x, y, v["title"], v["body"], c.C[col])
    return shp, {"accent": col, "w_in": 4.8, "h_in": 1.1}, {"title": v["title"], "body": v["body"]}

def r_analysis(sl, x, y, col, v):
    shp = build_analysis_box(sl, x, y, v["title"], v["bullets"], c.C[v["bg"]], c.C[col])
    return shp, {"bg": v["bg"], "accent": col, "w_in": 5.4, "h_in": 2.2}, {"title": v["title"], "bullets": v["bullets"]}

def r_ribbon(sl, x, y, col, v):
    shp = build_ribbon(sl, x, y, v["text"], c.C[col])
    return shp, {"fill": col, "w_in": 1.7, "h_in": 0.72}, {"text": v["text"]}

def r_pill(sl, x, y, col, v):
    shp = build_pill_badge(sl, x, y, v["text"], c.C[col])
    return shp, {"fill": col, "w_in": 1.3, "h_in": 0.44}, {"text": v["text"]}

def r_number(sl, x, y, col, v):
    shp = build_number_labels(sl, x, y, v["items"], c.C[col], w=v["w"])
    return shp, {"accent": col, "count": len(v["items"]), "w_in": v["w"]}, {"items": v["items"]}

def r_quote(sl, x, y, col, v):
    shp = build_quote_block(sl, x, y, v["quote"], v["author"], c.C[col], w=6.0)
    return shp, {"accent": col, "w_in": 6.0, "h_in": 1.9}, {"quote": v["quote"], "author": v["author"]}

def r_stat(sl, x, y, col, v):
    shp = build_stat_band(sl, x, y, v["number"], v["unit"], v["desc"], c.C[col])
    return shp, {"fill": col, "w_in": 7.2, "h_in": 1.15}, {"number": v["number"], "unit": v["unit"], "desc": v["desc"]}

def r_agenda(sl, x, y, col, v):
    shp = build_agenda_rows(sl, x, y, v["items"], c.C[col], w=5.8)
    return shp, {"accent": col, "count": len(v["items"]), "row_h_in": 1.0}, {"items": v["items"]}

def r_divider(sl, x, y, col, v):
    shp = build_dividers(sl, x, y, v["w"], c.C[col])
    return shp, {"accent": col, "w_in": v["w"], "variants": ["gradient", "dot"]}, {}

def r_step(sl, x, y, col, v):
    shp = build_step_label(sl, x, y, v["no"], v["kicker"], v["title"], c.C[col])
    return shp, {"accent": col, "w_in": 5.4}, {"step_no": v["no"], "kicker": v["kicker"], "title": v["title"]}

def r_kpi(sl, x, y, col, v):
    shp = build_inline_kpi_header(sl, x, y, v["title"], v["metric"], v["mlabel"], c.C[col])
    return shp, {"accent": col, "w_in": 8.6, "h_in": 0.9}, {"title": v["title"], "metric": v["metric"], "metric_label": v["mlabel"]}

def r_ribbonset(sl, x, y, col, v):
    items = [(t, c.C[col]) for t in v["texts"]]
    shp = build_ribbon_set(sl, x, y, items)
    return shp, {"accent": col, "count": len(v["texts"]), "w_each_in": 1.55}, {"texts": v["texts"]}

def r_tabheader(sl, x, y, col, v):
    shp = build_tab_header(sl, x, y, v["tabs"], 0, c.C[col])
    return shp, {"accent": col, "count": len(v["tabs"]), "active_idx": 0, "tab_w_in": 2.1}, {"tabs": v["tabs"], "active": v["tabs"][0]}

FAMILIES = [
    dict(key="section-pill", name="섹션 번호 pill", render=r_section, editable=["text", "color"],
         tags=["섹션헤더", "pill"], rec=["섹션 구분", "목차 헤더"],
         layout=dict(cols=1, col_xs=[0.9], y0=0.95, pitch=1.0, per_slide=6),
         variants=[dict(lab="사업추진 전략", num="1", title="사업추진 전략"),
                   dict(lab="추진 체계", num="2", title="추진 체계 및 조직"),
                   dict(lab="기대효과", num="3", title="기대효과 및 활용방안")]),
    dict(key="subtitle-bar", name="소제목 하이라이트 바", render=r_subtitle, editable=["text", "color"],
         tags=["소제목", "하이라이트바", "밑줄"], rec=["소제목", "단락 구분"],
         layout=dict(cols=1, col_xs=[0.9], y0=1.0, pitch=1.0, per_slide=6),
         variants=[dict(lab="가.", marker="가.", title="사업 추진 배경"),
                   dict(lab="나.", marker="나.", title="핵심 추진 전략")]),
    dict(key="check-bar", name="체크 헤드라인 바", render=r_check, editable=["text", "color"],
         tags=["체크바", "헤드라인", "라운드"], rec=["핵심 강조", "헤드라인"],
         layout=dict(cols=1, col_xs=[0.9], y0=1.0, pitch=1.0, per_slide=6),
         variants=[dict(lab="전략 수립", text="차별화된 해외 진출 전략 수립"),
                   dict(lab="실행력 확보", text="현지 파트너십 기반 실행력 확보")]),
    dict(key="tab-label", name="탭 스타일 라벨", render=r_tablabel, editable=["text", "color"],
         tags=["탭", "라벨", "섹션헤더"], rec=["탭 네비", "카테고리 라벨"],
         layout=dict(cols=3, col_xs=[1.0, 5.0, 9.0], y0=1.9, pitch=1.6, per_slide=6),
         variants=[dict(lab="핵심 제안", text="핵심 제안"),
                   dict(lab="추진 방향", text="추진 방향")]),
    dict(key="callout", name="강조 콜아웃 박스", render=r_callout, editable=["text", "color"],
         tags=["강조박스", "콜아웃", "세로바"], rec=["핵심 강조", "제안 요약"],
         layout=dict(cols=2, col_xs=[0.6, 7.0], y0=1.0, pitch=1.55, per_slide=6),
         variants=[dict(lab="핵심 제안", title="핵심 제안",
                        body="현지화 콘텐츠 재가공 및 채널별 최적 배급 전략으로\n해외 시장 진입 장벽을 최소화합니다."),
                   dict(lab="기대효과", title="기대효과",
                        body="K-콘텐츠 IP의 글로벌 인지도 제고와\n신규 해외 매출 파이프라인을 확보합니다.")]),
    dict(key="analysis-box", name="분석 요약 박스", render=r_analysis, editable=["text", "color"],
         tags=["분석박스", "요약", "불릿"], rec=["분석 요약", "현황 정리"],
         layout=dict(cols=2, col_xs=[0.6, 7.0], y0=0.95, pitch=2.45, per_slide=4),
         variants=[dict(lab="시장 환경 (베이지)", bg="beige_100", title="시장 환경 분석",
                        bullets=["글로벌 OTT 수요 확대로 K-콘텐츠 진출 기회 증가",
                                 "현지 언어·문화 장벽이 핵심 리스크로 상존",
                                 "권역별 배급 파트너 확보가 성패 좌우"]),
                   dict(lab="SWOT (그레이)", bg="gray_100", title="SWOT 요약",
                        bullets=["강점: 검증된 IP·제작 역량 보유",
                                 "약점: 해외 유통망·현지 네트워크 부족",
                                 "기회: 한류 확산 및 정책 지원 확대"])]),
    dict(key="ribbon", name="리본 배지", render=r_ribbon, editable=["text", "color"],
         tags=["배지", "리본", "추천"], rec=["핵심 강조", "우선순위 표시"],
         layout=dict(cols=3, col_xs=[1.2, 5.2, 9.2], y0=1.7, pitch=1.7, per_slide=6),
         variants=[dict(lab="추천", text="추천")]),
    dict(key="pill-badge", name="알약형 배지", render=r_pill, editable=["text", "color"],
         tags=["배지", "알약", "상태"], rec=["신규 표시", "상태 라벨"],
         layout=dict(cols=3, col_xs=[1.6, 5.6, 9.6], y0=2.2, pitch=1.4, per_slide=6),
         variants=[dict(lab="NEW", text="NEW")]),
    dict(key="number-labels", name="넘버링 라벨 세트", render=r_number, editable=["text", "color"],
         tags=["넘버링", "라벨", "원형", "세트"], rec=["단계 표시", "순서 나열"],
         layout=dict(cols=1, col_xs=[0.9], y0=1.0, pitch=1.15, per_slide=5),
         variants=[dict(lab="3단계", items=[["01", "기획"], ["02", "제작"], ["03", "배급"]], w=8.0),
                   dict(lab="4단계", items=[["01", "발굴"], ["02", "기획"], ["03", "제작"], ["04", "배급"]], w=9.6)]),
    dict(key="quote-block", name="인용 블록", render=r_quote, editable=["text", "color"],
         tags=["인용", "쿼트", "따옴표", "블록"], rec=["인용", "핵심 문구 강조"],
         layout=dict(cols=2, col_xs=[0.4, 6.9], y0=0.95, pitch=2.2, per_slide=4),
         variants=[dict(lab="문화적 맥락",
                        quote="현지 시청자의 문화적 맥락을 재해석한 콘텐츠만이\n글로벌 OTT에서 지속 가능한 경쟁력을 갖는다.",
                        author="한국콘텐츠진흥원, 2025 해외 콘텐츠시장 분석보고서"),
                   dict(lab="IP·배급 결합",
                        quote="검증된 IP와 현지 배급 네트워크의 결합이\n해외 진출 성패를 가르는 핵심 요인이다.",
                        author="글로벌 콘텐츠 진출 전략 세미나, 2025")]),
    dict(key="stat-band", name="통계 콜아웃 밴드", render=r_stat, editable=["text", "color"],
         tags=["통계", "콜아웃", "밴드", "지표"], rec=["핵심 지표", "성과 강조"],
         layout=dict(cols=1, col_xs=[0.9], y0=1.0, pitch=1.5, per_slide=4),
         variants=[dict(lab="성장률", number="127%", unit="전년 대비", desc="해외 콘텐츠 수출액 성장률"),
                   dict(lab="매출 규모", number="84.2억", unit="누적", desc="해외 라이선싱 매출 규모")]),
    dict(key="agenda-rows", name="아젠다 항목 행 세트", render=r_agenda, editable=["text", "color"],
         tags=["아젠다", "목차", "행세트", "번호원"], rec=["아젠다", "목차 본문"],
         layout=dict(cols=2, col_xs=[0.5, 6.9], y0=1.1, pitch=3.3, per_slide=2),
         variants=[dict(lab="3항목", items=[
             ["01", "사업 추진 배경 및 필요성", "글로벌 콘텐츠 시장 환경 변화와 진출 기회 분석"],
             ["02", "핵심 추진 전략 및 실행 방안", "권역별 현지화·배급 파트너십 기반 단계별 실행"],
             ["03", "기대효과 및 성과 활용 계획", "신규 해외 매출 파이프라인 확보와 IP 가치 제고"]])]),
    dict(key="dividers", name="구분선 변형", render=r_divider, editable=["color", "width"],
         tags=["구분선", "디바이더", "그라데이션", "도트"], rec=["섹션 구분", "단락 구분"],
         layout=dict(cols=1, col_xs=[1.2], y0=1.1, pitch=1.05, per_slide=6),
         variants=[dict(lab="와이드", w=8.8), dict(lab="컴팩트", w=6.0)]),
    dict(key="step-label", name="스텝 라벨", render=r_step, editable=["text", "color"],
         tags=["스텝", "단계", "STEP", "원형"], rec=["단계 표시", "프로세스 라벨"],
         layout=dict(cols=2, col_xs=[0.6, 7.0], y0=1.2, pitch=1.2, per_slide=6),
         variants=[dict(lab="STEP 01", no="01", kicker="STEP 01", title="시장 조사 및 현황 진단"),
                   dict(lab="STEP 02", no="02", kicker="STEP 02", title="현지화 전략 수립 및 실행")]),
    dict(key="inline-kpi", name="인라인 KPI 헤더", render=r_kpi, editable=["text", "color"],
         tags=["KPI", "헤더", "인라인", "지표"], rec=["섹션 헤더", "지표 병기 제목"],
         layout=dict(cols=1, col_xs=[1.0], y0=1.0, pitch=1.25, per_slide=5),
         variants=[dict(lab="성과 요약", title="해외사업 성과 요약", metric="84.2억 원", mlabel="누적 수출 매출"),
                   dict(lab="수출 실적", title="콘텐츠 수출 실적", metric="127%", mlabel="전년 대비 성장")]),
    dict(key="ribbon-set", name="리본 배지 세트", render=r_ribbonset, editable=["text", "color"],
         tags=["배지", "리본", "태그", "세트"], rec=["분류 태그", "상태 라벨"],
         layout=dict(cols=2, col_xs=[0.6, 7.0], y0=1.4, pitch=0.95, per_slide=6),
         variants=[dict(lab="핵심/전략/성과", texts=["핵심", "전략", "성과"])]),
    dict(key="tab-header", name="탭 헤더 세트", render=r_tabheader, editable=["text", "color"],
         tags=["탭", "헤더", "네비", "세트"], rec=["탭 네비게이션", "섹션 헤더"],
         layout=dict(cols=1, col_xs=[1.2], y0=1.1, pitch=1.15, per_slide=6),
         variants=[dict(lab="3탭", tabs=["추진 전략", "추진 체계", "기대효과"])]),
]

# ═══════════════════════════════════════════════════════════════
# 드라이버 — flat 스윕 → 176개 컷 → 패밀리별 파일/슬라이드 배치
# ═══════════════════════════════════════════════════════════════
TARGET = 176           # HDR-025 ~ HDR-200
START_ID = 25

# 1) flat 스윕: 패밀리 → 변형 → 색
flat = []
for fam in FAMILIES:
    for vi, variant in enumerate(fam["variants"]):
        for col in COLORS:
            flat.append((fam, vi, variant, col))
assert len(flat) >= TARGET, "스윕 조합 부족: %d < %d" % (len(flat), TARGET)
flat = flat[:TARGET]   # 정확히 176개로 컷 (초과분은 말미 패밀리에서 잘림)

# 2) 패밀리별 그룹(순서 보존)
groups = OrderedDict()
for item in flat:
    groups.setdefault(item[0]["key"], []).append(item)

# 3) 배치 + 엔트리
entries = []
files_written = []
aid = START_ID
for key, items in groups.items():
    fam = items[0][0]
    lay = fam["layout"]
    per = lay["per_slide"]
    cols = lay["cols"]
    file_rel = "decks/09_headers/HDR_bulk_%s_v1.pptx" % key
    prs = c.new_deck()
    slide = None
    slide_idx = 0
    for i, (f, vi, variant, col) in enumerate(items):
        slot = i % per
        if slot == 0:
            slide = c.blank_slide(prs)
            slide_idx += 1
        col_i = slot % cols
        row_i = slot // cols
        x = lay["col_xs"][col_i]
        y = lay["y0"] + row_i * lay["pitch"]
        asset_id = "HDR-%03d" % aid
        shapes, params, bindings = fam["render"](slide, x, y, col, variant)
        c.group_asset(slide, shapes, asset_id)
        if slot == 0:
            c.id_caption(slide, "%s~ · %s" % (asset_id, fam["name"]))
        kcol = KCOL[col]
        vlab = variant.get("lab", "")
        multi_v = len(fam["variants"]) > 1
        name = "%s (%s)" % (fam["name"], kcol) if not multi_v else "%s · %s (%s)" % (fam["name"], vlab, kcol)
        tags = list(fam["tags"]) + [kcol]
        entries.append(c.entry(asset_id, "HDR", name, file_rel, slide_idx,
                               tags, params, bindings, fam["editable"],
                               recommended_use=fam["rec"]))
        aid += 1
    out = c.save_deck(prs, file_rel)
    files_written.append((file_rel, slide_idx, len(items)))

# 4) 매니페스트 조각 1회
frag = c.write_fragment("HDR_bulk", entries)

print("ENTRIES:", len(entries))
print("ID RANGE:", entries[0]["id"], "~", entries[-1]["id"])
print("FRAGMENT:", frag)
print("FILES:")
for fr, ns, na in files_written:
    print("  %-52s slides=%d assets=%d" % (fr, ns, na))
