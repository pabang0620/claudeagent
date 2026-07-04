# -*- coding: utf-8 -*-
"""TBL 대량 확충 — TBL-023 ~ TBL-200 (표 에셋 178종).
파라미터화 배치 생성기: 구조패밀리 × 헤더색 × 열수/행수/밀도 스윕.

규칙:
- 색은 c.role()/c.C[...] 만 (하드코딩 금지). 헤더색 스윕 = navy_800/blue_500/teal_500/purple_600/navy_600/red_500.
- 한글은 c.set_kfont / c.add_text 경유.
- 단일 표: 표 graphicFrame.name='asset:TBL-0NNN'. 결합형(칩/카드): c.group_asset 그룹, id_caption은 그룹 밖.
- 페이지이미지·SmartArt 금지. 완전 동일 리터럴 중복 금지 (색/열수/행수/테마로 실질 차이).
- 파일당 ≤25 슬라이드. 1슬라이드 1에셋. 마지막에 c.write_fragment('TBL_bulk', entries) 1회.
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

R = c.role
CC = c.C
NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid
L = PP_ALIGN.LEFT; Rg = PP_ALIGN.RIGHT; Ct = PP_ALIGN.CENTER

# 헤더색 스윕 (name, RGBColor) — 모두 어두워 header_text(white)와 대비 확보
HEADER_COLORS = [
    ("navy_800", CC["navy_800"]),
    ("blue_500", CC["blue_500"]),
    ("teal_500", CC["teal_500"]),
    ("purple_600", CC["purple_600"]),
    ("navy_600", CC["navy_600"]),
    ("red_500", CC["red_500"]),
]

# ---------- 테이블 헬퍼 (gen_tbl_v2 계열과 동일 스타일) ----------
def add_table(slide, x, y, w, h, nrows, ncols, col_widths=None, row_h=None):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    tblPr = tbl._tbl.tblPr
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        sid = etree.SubElement(tblPr, qn('a:tableStyleId'))
    sid.text = NO_STYLE
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    return gf, tbl

def set_cell(cell, text, size=11, bold=False, color=None, fill=None,
             align=PP_ALIGN.CENTER, font=None, ml=5):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = R("row_base")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(ml)
    cell.margin_top = cell.margin_bottom = Pt(2)
    tf = cell.text_frame; tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))

def cell_borders(cell, spec):
    tcPr = cell._tc.get_or_add_tcPr()
    for edge, tag in (('B', 'a:lnB'), ('T', 'a:lnT'), ('R', 'a:lnR'), ('L', 'a:lnL')):
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)
        if edge not in spec:
            continue
        color, wpt = spec[edge]
        ln = etree.Element(qn(tag))
        ln.set('w', str(int(wpt * 12700))); ln.set('cap', 'flat')
        if color is None:
            etree.SubElement(ln, qn('a:noFill'))
        else:
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr')).set('val', str(color))
        tcPr.insert(0, ln)

def grid(tbl, color, wpt=0.75):
    for row in tbl.rows:
        for cell in row.cells:
            cell_borders(cell, {'L': (color, wpt), 'R': (color, wpt),
                                'T': (color, wpt), 'B': (color, wpt)})

def title_block(slide, text, sub=None):
    c.add_text(slide, 0.55, 0.5, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.0, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)

# ---------- 더미 콘텐츠 풀 (제안서 도메인) ----------
THEMES = [
    "콘텐츠 해외진출", "스마트 물류", "AI 통합 플랫폼", "지역 관광 활성화",
    "디지털 헬스케어", "그린 에너지", "핀테크 서비스", "에듀테크 플랫폼",
    "제조 디지털전환", "공공 데이터 개방", "메타버스 전시", "친환경 모빌리티",
    "농식품 수출", "문화예술 지원", "스마트시티", "바이오 소재",
]
LABEL_HEADERS = ["구분", "항목", "영역", "과업", "지표", "부문", "단계",
                 "분류", "요소", "기준", "대상", "범주"]
ATTR_HEADERS = ["주요 내용", "담당 부서", "추진 일정", "핵심 산출물", "목표치",
                "추진 실적", "가중치", "우선순위", "진행 상태", "비고",
                "책임자", "기대 효과", "핵심 지표", "리스크", "대응 방안"]
NUM_HEADERS = ["2024년", "2025년", "2026년(E)", "1분기", "2분기", "3분기",
               "4분기", "상반기", "하반기", "누적", "목표", "실적", "증감률"]
ROW_LABELS = [
    "사업 기획", "요구사항 분석", "정보구조 설계", "핵심 기능 개발", "현지화 지원",
    "통합 테스트", "품질 검증", "운영 이관", "성과 관리", "위험 관리",
    "고객 지원", "데이터 분석", "보안 점검", "인프라 구축", "콘텐츠 제작",
    "마케팅 홍보", "파트너 협력", "교육 훈련", "유지보수", "서비스 고도화",
    "표준화 정립", "글로벌 확장", "채널 다각화", "브랜드 강화", "재무 관리",
]
CONTENTS = [
    "요구사항 정의 및 범위 확정 후 단계별 실행",
    "글로벌 유통 채널 15개국 확장 및 현지화 대응",
    "핵심 기능 개발과 외부 시스템 연동 안정화",
    "품질 기준 수립 및 결함 조치 프로세스 운영",
    "운영 매뉴얼 정비와 담당자 교육 병행 추진",
    "데이터 기반 성과 측정 및 개선 과제 도출",
    "이해관계자 협의체 운영으로 리스크 선제 대응",
    "표준 프로세스 정립과 자동화 도입으로 효율화",
    "브랜드 인지도 제고 및 신규 파트너 발굴",
    "예산 집행 점검과 정산 체계 투명성 확보",
]
DEPTS = ["사업기획팀", "PMO", "개발팀", "디자인팀", "해외사업팀", "QA팀",
         "전략기획실", "운영지원팀", "데이터분석팀", "마케팅팀"]
STATUS_LABELS = [("완료", "ok"), ("진행중", "go"), ("지연", "warn"),
                 ("대기", "wait"), ("검토", "go"), ("보류", "wait")]
SCHED = ["M1~M2", "M3~M4", "M5~M6", "M7~M9", "M10~M11", "M12", "1분기", "2분기"]

def pick(pool, n, seed, step=1):
    """seed 기준 rotating slice — 에셋마다 다른 리터럴 조합."""
    off = seed % len(pool)
    return [pool[(off + i * step) % len(pool)] for i in range(n)]

def money(seed, i, j):
    base = 1000 + (seed * 37 + i * 53 + j * 29) % 9000
    return f"{base:,}"

def pct(seed, i, j):
    v = 55 + (seed * 7 + i * 13 + j * 5) % 45
    return f"{v}.{(seed + i + j) % 10}%"

def delta(seed, i, j):
    v = (seed * 3 + i * 7 + j * 11) % 40 - 15
    sign = "+" if v >= 0 else ""
    return f"{sign}{v}.{(i + j) % 10}%", v < 0

def qty(seed, i, j):
    return f"{20 + (seed * 5 + i * 17 + j * 7) % 480}"

def value_for(kind, seed, i, j):
    if kind == "money":
        return money(seed, i, j), Rg
    if kind == "pct":
        return pct(seed, i, j), Rg
    if kind == "qty":
        return qty(seed, i, j), Rg
    if kind == "text":
        return CONTENTS[(seed + i + j) % len(CONTENTS)], L
    return qty(seed, i, j), Rg

# =========================================================
# 패밀리 렌더러 — 각 함수는 anchor(gf 또는 grpSp) 를 named 처리하고 entry meta 리턴
# 공통 시그니처: fam(slide, aid, hcname, hcol, ncols, nrows, seed)
# =========================================================
def _grid_geom(ncols, nrows, label_w=None, dense=False):
    """열폭/행높이 계산 (표 폭 ~11.33in)."""
    total_w = 11.33
    lw = label_w if label_w else (2.6 if ncols <= 4 else 2.2)
    rest = (total_w - lw) / (ncols - 1) if ncols > 1 else total_w
    cw = [lw] + [rest] * (ncols - 1)
    rh = (0.34 if dense else 0.6)
    return cw, rh, total_w

def fam_striped(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 추진 현황표",
                f"줄무늬 헤더({hcname}) · {ncols}열 × {nrows}행")
    x, y = 1.0, 1.55
    cw, rh, w = _grid_geom(ncols, nrows)
    heads = [pick(LABEL_HEADERS, 1, seed)[0]] + pick(ATTR_HEADERS, ncols - 1, seed + 2)
    labels = pick(ROW_LABELS, nrows, seed, step=2)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[i - 1], size=11, bold=True, fill=base, align=L)
        for j in range(1, ncols):
            kind = "text" if (ncols <= 3 and j == 1) else ("pct" if j == ncols - 1 else "qty")
            val, al = value_for(kind, seed, i, j)
            set_cell(tbl.cell(i, j), val, size=11, fill=base, align=al)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "striped", "header", hcname, "줄무늬", "제안서"],
            {"cols": ncols, "rows": nrows, "header_color": hcname},
            ["header", "row_label", "cell_value"],
            ["header", "row_label", "cell_value"],
            ["현황요약", "실적표", "항목별 정리"])

def fam_borderless(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 개요 요약표",
                f"보더리스 미니멀 · {ncols}열 · 헤더 {hcname}")
    x, y = 1.0, 1.7
    cw, rh, w = _grid_geom(ncols, nrows, label_w=2.4)
    rh = 0.7
    heads = [pick(LABEL_HEADERS, 1, seed)[0]] + pick(ATTR_HEADERS, ncols - 1, seed + 3)
    labels = pick(ROW_LABELS, nrows, seed + 1, step=3)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=13, bold=True, color=hcol, fill=R("row_base"), align=L)
        cell_borders(tbl.cell(0, j), {'B': (hcol, 2.0), 'L': (None, .5), 'R': (None, .5), 'T': (None, .5)})
    for i in range(1, n):
        for j in range(ncols):
            if j == 0:
                v, al, col, bd = labels[i - 1], L, R("header_fill"), True
            else:
                kind = "text" if j == 1 else "qty"
                v, al = value_for(kind, seed, i, j); col, bd = R("body_text"), False
            set_cell(tbl.cell(i, j), v, size=11, bold=bd, color=col, fill=R("row_base"), align=al)
            hair = R("border") if i < nrows else None
            cell_borders(tbl.cell(i, j), {'B': (hair, 0.5) if hair else (None, .5),
                                          'L': (None, .5), 'R': (None, .5), 'T': (None, .5)})
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "borderless", "minimal", hcname, "미니멀"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "borderless": True},
            ["header", "row_label", "content"],
            ["header", "row_label", "content"],
            ["사업개요", "깔끔한 요약", "본문 삽입표"])

def fam_underline(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 세부 내역표",
                f"헤더 밑줄형 · {ncols}열 · {hcname} 액센트")
    x, y = 1.0, 1.7
    cw, rh, w = _grid_geom(ncols, nrows)
    heads = [pick(LABEL_HEADERS, 1, seed + 4)[0]] + pick(ATTR_HEADERS, ncols - 1, seed + 5)
    labels = pick(ROW_LABELS, nrows, seed + 6, step=2)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=hcol, fill=R("row_base"),
                 align=L if j == 0 else Ct)
        cell_borders(tbl.cell(0, j), {'B': (hcol, 1.75), 'L': (None, .5), 'R': (None, .5), 'T': (None, .5)})
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        for j in range(ncols):
            if j == 0:
                v, al, bd, col = labels[i - 1], L, True, R("body_text")
            else:
                kind = "money" if j == ncols - 1 else "qty"
                v, al = value_for(kind, seed, i, j); bd, col = False, R("body_text")
            set_cell(tbl.cell(i, j), v, size=11, bold=bd, color=col, fill=base, align=al)
            cell_borders(tbl.cell(i, j), {'B': (R("border"), 0.5),
                                          'L': (None, .5), 'R': (None, .5), 'T': (None, .5)})
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "underline", "header-underline", hcname, "밑줄"],
            {"cols": ncols, "rows": nrows, "header_color": hcname},
            ["header", "row_label", "cell_value"],
            ["header", "row_label", "cell_value"],
            ["세부내역", "정갈한 표", "본문 표"])

def fam_group2(slide, aid, hcname, hcol, ncols, nrows, seed):
    """2단 그룹헤더표 — 상단 그룹 span + 하단 서브헤더."""
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 실적 비교표",
                f"2단 그룹헤더 · 5열 · {hcname}")
    x, y = 1.0, 1.55
    cw = [2.53, 2.2, 2.2, 2.2, 2.2]; w = sum(cw); rh = 0.56
    n = nrows + 2  # 2 header rows
    gf, tbl = add_table(slide, x, y, w, rh * n, n, 5, col_widths=cw, row_h=[rh] * n)
    gA = pick(THEMES, 2, seed + 1)
    # 그룹 헤더(행0): 라벨 span(세로) + 그룹A(2) + 그룹B(2)
    tbl.cell(0, 0).merge(tbl.cell(1, 0))
    set_cell(tbl.cell(0, 0), pick(LABEL_HEADERS, 1, seed)[0], size=12, bold=True,
             color=R("header_text"), fill=hcol, align=Ct)
    tbl.cell(0, 1).merge(tbl.cell(0, 2))
    set_cell(tbl.cell(0, 1), f"{gA[0]} 부문", size=12, bold=True, color=R("header_text"), fill=hcol)
    tbl.cell(0, 3).merge(tbl.cell(0, 4))
    set_cell(tbl.cell(0, 3), f"{gA[1]} 부문", size=12, bold=True, color=R("header_text"), fill=hcol)
    subs = pick(NUM_HEADERS, 4, seed + 2)
    for j, hd in enumerate(subs, start=1):
        set_cell(tbl.cell(1, j), hd, size=11, bold=True, color=R("header_text"),
                 fill=CC["navy_600"] if hcname != "navy_600" else CC["navy_800"], align=Ct)
    labels = pick(ROW_LABELS, nrows, seed + 3, step=2)
    for r in range(nrows):
        i = r + 2
        base = R("row_stripe") if r % 2 == 1 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[r], size=11, bold=True, fill=base, align=L)
        for j in range(1, 5):
            val, al = value_for("qty" if j % 2 else "money", seed, i, j)
            set_cell(tbl.cell(i, j), val, size=11, fill=base, align=Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "group-header", "two-tier", hcname, "그룹헤더"],
            {"cols": 5, "rows": nrows, "header_color": hcname, "group_header": True},
            ["group_header", "sub_header", "row_label", "cell_value"],
            ["group_header", "sub_header", "row_label", "cell_value"],
            ["실적비교", "부문별 대비", "그룹 요약"])

def fam_leftlabel(slide, aid, hcname, hcol, ncols, nrows, seed):
    """좌측 라벨열 강조표 — 첫 열을 진한 컬러 패널로."""
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 항목별 정리표",
                f"좌측 라벨열 강조 · {ncols}열 · {hcname}")
    x, y = 1.0, 1.55
    cw, rh, w = _grid_geom(ncols, nrows, label_w=2.8)
    heads = [pick(LABEL_HEADERS, 1, seed + 7)[0]] + pick(ATTR_HEADERS, ncols - 1, seed + 8)
    labels = pick(ROW_LABELS, nrows, seed + 2, step=3)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[i - 1], size=11, bold=True, color=R("header_text"),
                 fill=hcol, align=L)
        for j in range(1, ncols):
            kind = "text" if (ncols <= 4 and j == 1) else ("pct" if j == ncols - 1 else "qty")
            v, al = value_for(kind, seed, i, j)
            set_cell(tbl.cell(i, j), v, size=11, fill=base, align=al)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "left-label", "row-header", hcname, "라벨강조"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "left_label": True},
            ["header", "row_label_emph", "cell_value"],
            ["header", "row_label", "cell_value"],
            ["항목정리", "속성표", "좌측 헤더표"])

def fam_totalrow(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 집계 요약표",
                f"합계행 강조 · {ncols}열 · {hcname}")
    x, y = 1.0, 1.55
    cw, rh, w = _grid_geom(ncols, nrows, label_w=2.9)
    heads = [pick(LABEL_HEADERS, 1, seed + 9)[0]] + pick(NUM_HEADERS, ncols - 1, seed + 1)
    labels = pick(ROW_LABELS, nrows, seed + 4, step=2)
    n = nrows + 2  # +헤더 +합계
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Rg)
    for i in range(1, nrows + 1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[i - 1], size=11, bold=True, fill=base, align=L)
        for j in range(1, ncols):
            v, al = value_for("money", seed, i, j)
            set_cell(tbl.cell(i, j), v, size=11, fill=base, align=Rg)
    li = n - 1
    set_cell(tbl.cell(li, 0), "합계", size=12, bold=True, color=R("header_text"),
             fill=CC["navy_800"], align=L)
    for j in range(1, ncols):
        set_cell(tbl.cell(li, j), money(seed, 99, j), size=12, bold=True,
                 color=R("header_text"), fill=CC["navy_800"], align=Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "total-row", "sum", hcname, "합계행"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "total_row": True},
            ["header", "row_label", "amount", "total_row"],
            ["header", "row_label", "amount", "total_row"],
            ["집계요약", "합계표", "금액 정산"])

def fam_signal(slide, aid, hcname, hcol, ncols, nrows, seed):
    """신호등 상태셀표 (teal/blue/red/gray)."""
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 과업 상태 점검표",
                f"신호등 상태셀 · {ncols}열 · {hcname}")
    x, y = 1.0, 1.55
    cw, rh, w = _grid_geom(ncols, nrows, label_w=3.6)
    heads = [pick(LABEL_HEADERS, 1, seed + 1)[0], "상태"] + pick(ATTR_HEADERS, ncols - 2, seed + 5)
    labels = pick(ROW_LABELS, nrows, seed + 5, step=2)
    signal = {"ok": R("accent_secondary"), "go": R("accent_primary"),
              "warn": R("warn"), "wait": R("muted_text")}
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        slab, kind = STATUS_LABELS[(seed + i) % len(STATUS_LABELS)]
        set_cell(tbl.cell(i, 0), labels[i - 1], size=11, bold=True, fill=base, align=L)
        set_cell(tbl.cell(i, 1), slab, size=11, bold=True, color=R("header_text"),
                 fill=signal[kind], align=Ct)
        for j in range(2, ncols):
            v = DEPTS[(seed + i + j) % len(DEPTS)] if j == 2 else CONTENTS[(seed + i + j) % len(CONTENTS)]
            set_cell(tbl.cell(i, j), v, size=11, fill=base, align=Ct if j == 2 else L,
                     color=R("muted_text") if j == 2 else R("body_text"))
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "status", "traffic-light", hcname, "신호등", "상태"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "status_cell": True},
            ["header", "item", "status_cell", "owner"],
            ["header", "item", "status_cell"],
            ["진행상태", "과업점검", "RAG 상태"])

def fam_check(slide, aid, hcname, hcol, ncols, nrows, seed):
    """체크(✓/△/✗) 글리프 상태표 — 단일 표."""
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 점검 체크표",
                f"체크 글리프(✓/△/✗) · {ncols}열 · {hcname}")
    x, y = 1.0, 1.55
    cw, rh, w = _grid_geom(ncols, nrows, label_w=4.2)
    heads = [pick(LABEL_HEADERS, 1, seed + 2)[0]] + pick(ATTR_HEADERS, ncols - 2, seed + 6) + ["점검"]
    labels = pick(ROW_LABELS, nrows, seed + 6, step=3)
    glyphs = [("✓", R("accent_secondary")), ("△", R("warn")), ("✗", R("muted_text"))]
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[i - 1], size=11, bold=True, fill=base, align=L)
        for j in range(1, ncols - 1):
            v = CONTENTS[(seed + i + j) % len(CONTENTS)] if j == 1 else DEPTS[(seed + i + j) % len(DEPTS)]
            set_cell(tbl.cell(i, j), v, size=11, fill=base, align=L if j == 1 else Ct,
                     color=R("body_text") if j == 1 else R("muted_text"))
        g, gc = glyphs[(seed + i) % 3]
        set_cell(tbl.cell(i, ncols - 1), g, size=14, bold=True, color=gc, fill=base, align=Ct)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "check", "checklist", hcname, "체크", "글리프"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "check_glyph": True},
            ["header", "item", "check_glyph"],
            ["header", "item", "check_glyph"],
            ["점검표", "완료여부", "체크리스트"])

def fam_kpichip(slide, aid, hcname, hcol, ncols, nrows, seed):
    """셀내 KPI칩표 — 값 + 증감칩 오버레이 (결합형 group_asset)."""
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 핵심지표(KPI) 표",
                f"셀내 KPI칩 · 증감칩 오버레이 · {hcname}")
    x, y, w = 1.0, 1.6, 11.33
    heads = [pick(LABEL_HEADERS, 1, seed + 3)[0], "현재 값", "목표", "전기 대비"]
    labels = pick(ROW_LABELS, nrows, seed + 7, step=2)
    rh = 0.66; cw = [4.03, 2.9, 2.1, 2.3]
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, 4, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[i - 1], size=12, bold=True, fill=base, align=L)
        set_cell(tbl.cell(i, 1), money(seed, i, 1), size=15, bold=True, color=hcol, fill=base, align=Ct)
        set_cell(tbl.cell(i, 2), money(seed, i, 2), size=11, color=R("muted_text"), fill=base, align=Ct)
        set_cell(tbl.cell(i, 3), "", fill=base)
    grid(tbl, R("border"), 0.75)
    chips = []
    chip_w, chip_h = 1.7, 0.42
    col3_x = x + sum(cw[:3])
    for i in range(1, n):
        dtxt, neg = delta(seed, i, 3)
        cx = col3_x + (cw[3] - chip_w) / 2
        cy = y + rh * i + (rh - chip_h) / 2
        col = R("warn") if neg else R("accent_secondary")
        chip = c.add_box(slide, cx, cy, chip_w, chip_h, fill=col, line=None,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        arrow = "▼" if neg else "▲"
        c.set_shape_text(chip, f"{arrow} {dtxt.lstrip('+-')}", size=11, bold=True, color=R("header_text"))
        chips.append(chip)
    c.group_asset(slide, [gf] + chips, aid)
    c.id_caption(slide, aid)
    return (["table", "kpi", "chip", "delta", hcname, "증감칩", "결합형"],
            {"cols": 4, "rows": nrows, "header_color": hcname, "combined": True},
            ["header", "kpi_value", "delta_chip"],
            ["header", "kpi_value", "delta_chip"],
            ["KPI 대시보드", "성과요약", "지표표"])

def fam_cardrow(slide, aid, hcname, hcol, ncols, nrows, seed):
    """카드행형표 — 행이 카드 (도형 결합형 group_asset)."""
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 성과 하이라이트 (카드형)",
                f"카드행 · 좌 라벨 / 우 내용 2톤 · {hcname}")
    x, w = 1.0, 11.33
    y0 = 1.65; card_h, gap = 0.86, 0.16; label_w = 3.0
    labels = pick(ROW_LABELS, nrows, seed + 8, step=3)
    shapes = []
    for i in range(nrows):
        cy = y0 + i * (card_h + gap)
        card = c.add_box(slide, x, cy, w, card_h, fill=CC["gray_050"], line=R("border"),
                         line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(card)
        panel = c.add_box(slide, x, cy, label_w, card_h, fill=hcol, line=None,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.set_shape_text(panel, labels[i], size=13, bold=True, color=R("header_text"))
        shapes.append(panel)
        metric = f"{qty(seed, i, 1)}건" if i % 2 else money(seed, i, 1)
        vb = c.add_text(slide, x + label_w + 0.2, cy + 0.06, 2.9, card_h - 0.12, metric,
                        size=17, bold=True, color=hcol, align=L, anchor=MSO_ANCHOR.MIDDLE, font=c.FONT_H)
        shapes.append(vb)
        desc = CONTENTS[(seed + i) % len(CONTENTS)]
        db = c.add_text(slide, x + label_w + 3.15, cy + 0.06, w - label_w - 3.35, card_h - 0.12,
                        desc, size=11, color=R("body_text"), align=L, anchor=MSO_ANCHOR.MIDDLE)
        shapes.append(db)
    c.group_asset(slide, shapes, aid)
    c.id_caption(slide, aid)
    return (["table", "card", "two-tone", hcname, "카드형", "도형결합"],
            {"rows": nrows, "header_color": hcname, "shape_based": True, "combined": True},
            ["label_panel", "metric", "detail"],
            ["label_panel", "metric", "detail"],
            ["성과 하이라이트", "카드형 요약", "강조 표"])

def fam_budget(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 사업비 집행표",
                f"예산·집행·잔액·집행률 · 합계 강조 · {hcname}")
    x, y = 1.0, 1.55
    cw = [3.33, 2.0, 2.0, 2.0, 2.0]; w = sum(cw); rh = 0.56
    heads = ["사업 항목", "예산액", "집행액", "잔액", "집행률"]
    labels = pick(ROW_LABELS, nrows, seed + 9, step=2)
    n = nrows + 2
    gf, tbl = add_table(slide, x, y, w, rh * n, n, 5, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct)
    tb = te = 0
    for i in range(1, nrows + 1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        bud = 40 + (seed * 11 + i * 23) % 200
        exe = int(bud * (0.55 + ((seed + i) % 40) / 100.0))
        rem = bud - exe
        tb += bud; te += exe
        set_cell(tbl.cell(i, 0), labels[i - 1], size=11, bold=True, fill=base, align=L)
        set_cell(tbl.cell(i, 1), f"{bud*1000:,}", size=11, fill=base, align=Rg, color=R("muted_text"))
        set_cell(tbl.cell(i, 2), f"{exe*1000:,}", size=11, bold=True, fill=base, align=Rg, color=R("accent_primary"))
        set_cell(tbl.cell(i, 3), f"{rem*1000:,}", size=11, fill=base, align=Rg)
        set_cell(tbl.cell(i, 4), f"{exe*100//bud}.{(seed+i)%10}%", size=11, bold=True, fill=base, align=Rg, color=hcol)
    li = n - 1
    tot = ["합계", f"{tb*1000:,}", f"{te*1000:,}", f"{(tb-te)*1000:,}", f"{te*100//tb}.0%"]
    for j, v in enumerate(tot):
        set_cell(tbl.cell(li, j), v, size=12, bold=True, color=R("header_text"),
                 fill=CC["navy_800"], align=L if j == 0 else Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "budget", "execution", hcname, "예산", "집행"],
            {"cols": 5, "rows": nrows, "header_color": hcname, "total_row": True},
            ["header", "row_label", "amount", "total_row"],
            ["header", "row_label", "amount", "total_row"],
            ["예산집행", "사업비 현황", "정산표"])

def fam_schedule(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 단계별 추진 일정표",
                f"단계·활동·산출물·일정 · {hcname}")
    x, y = 1.0, 1.55
    cw = [1.7, 5.13, 2.7, 1.8]; w = sum(cw); rh = 0.68
    heads = ["단계", "주요 활동", "산출물", "일정"]
    phase_cols = [CC["navy_800"], CC["blue_500"], CC["teal_500"], CC["cyan_500"],
                  CC["purple_600"], CC["navy_600"]]
    labels = pick(ROW_LABELS, nrows, seed + 3, step=2)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, 4, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=Ct if j != 1 else L)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        pcol = phase_cols[(seed + i) % len(phase_cols)]
        set_cell(tbl.cell(i, 0), f"{i}단계\n{labels[i-1][:2]}", size=11, bold=True,
                 color=R("header_text"), fill=pcol, align=Ct)
        set_cell(tbl.cell(i, 1), CONTENTS[(seed + i) % len(CONTENTS)], size=11, fill=base, align=L)
        set_cell(tbl.cell(i, 2), labels[i - 1] + " 산출물", size=11, bold=True, fill=base,
                 align=Ct, color=R("accent_primary"))
        set_cell(tbl.cell(i, 3), SCHED[(seed + i) % len(SCHED)], size=11, bold=True,
                 fill=base, align=Ct, color=hcol)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "schedule", "phase", hcname, "추진일정", "산출물"],
            {"cols": 4, "rows": nrows, "header_color": hcname},
            ["header", "phase", "activity", "deliverable"],
            ["header", "phase", "activity", "deliverable"],
            ["추진일정", "단계별 계획", "WBS 요약"])

def fam_compact(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 상세 실적표 (컴팩트)",
                f"고밀도 · {ncols}열 소폰트 · {hcname}")
    x, y, w = 0.7, 1.5, 11.93
    heads = [pick(LABEL_HEADERS, 1, seed + 5)[0]] + pick(NUM_HEADERS, ncols - 1, seed)
    labels = pick(ROW_LABELS, nrows, seed + 4, step=1)
    rh = 0.34
    cw = [2.13] + [(w - 2.13) / (ncols - 1)] * (ncols - 1)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=9, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Ct, ml=3)
    for i in range(1, n):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), labels[i - 1], size=9, bold=True, fill=base, align=L, ml=3)
        for j in range(1, ncols):
            last = j == ncols - 1
            col = R("accent_primary") if last else R("body_text")
            v = pct(seed, i, j) if last else qty(seed, i, j)
            set_cell(tbl.cell(i, j), v, size=9, bold=last, color=col, fill=base, align=Rg, ml=3)
    grid(tbl, R("border"), 0.5)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "compact", "dense", hcname, "고밀도", "다열"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "font_pt": 9},
            ["header", "cell_value"],
            ["header", "cell_value"],
            ["KPI 상세", "분기실적", "다지표 요약"])

def fam_numeric(slide, aid, hcname, hcol, ncols, nrows, seed):
    theme = THEMES[seed % len(THEMES)]
    title_block(slide, f"{theme} 요약 재무표",
                f"우측정렬 수치 · 감소 적색 · {ncols}열 · {hcname}")
    x, y = 1.0, 1.55
    cw, rh, w = _grid_geom(ncols, nrows, label_w=3.0)
    heads = ["계정 과목"] + pick(NUM_HEADERS, ncols - 2, seed + 1) + ["증감률"]
    labels = pick(ROW_LABELS, nrows, seed + 2, step=2)
    n = nrows + 1
    gf, tbl = add_table(slide, x, y, w, rh * n, n, ncols, col_widths=cw, row_h=[rh] * n)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=hcol, align=L if j == 0 else Rg)
    for i in range(1, n):
        emph = (i == n - 1)
        base = CC["gray_100"] if emph else (R("row_stripe") if i % 2 == 0 else R("row_base"))
        set_cell(tbl.cell(i, 0), labels[i - 1] + (" 총계" if emph else ""), size=11, bold=emph,
                 fill=base, align=L, color=R("header_fill") if emph else R("body_text"))
        for j in range(1, ncols - 1):
            set_cell(tbl.cell(i, j), money(seed, i, j), size=11, bold=emph, fill=base, align=Rg,
                     color=R("body_text") if j == 1 else R("muted_text"))
        dtxt, neg = delta(seed, i, ncols)
        dcol = R("warn") if neg else R("accent_secondary")
        set_cell(tbl.cell(i, ncols - 1), dtxt, size=11, bold=True, color=dcol, fill=base, align=Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, aid); c.id_caption(slide, aid)
    return (["table", "numeric", "right-align", hcname, "재무정렬", "수치"],
            {"cols": ncols, "rows": nrows, "header_color": hcname, "align": "right"},
            ["header", "account", "figure"],
            ["header", "account", "figure"],
            ["재무상태표", "회계 요약", "수치정렬"])

# =========================================================
# 스윕 계획 — 178개 스펙 생성 (TBL-023 ~ TBL-200)
# =========================================================
FAM_TITLE = {
    "striped": "줄무늬 헤더표", "borderless": "보더리스 미니멀표", "underline": "헤더밑줄표",
    "group2": "2단 그룹헤더표", "leftlabel": "좌측 라벨열강조표", "totalrow": "합계행 강조표",
    "signal": "신호등 상태셀표", "check": "체크(✓/△/✗)표", "kpichip": "셀내 KPI칩표",
    "cardrow": "카드행형표", "budget": "예산집행표", "schedule": "추진일정표",
    "compact": "고밀도 컴팩트표", "numeric": "우측정렬 수치표",
}
FAM_FN = {
    "striped": fam_striped, "borderless": fam_borderless, "underline": fam_underline,
    "group2": fam_group2, "leftlabel": fam_leftlabel, "totalrow": fam_totalrow,
    "signal": fam_signal, "check": fam_check, "kpichip": fam_kpichip,
    "cardrow": fam_cardrow, "budget": fam_budget, "schedule": fam_schedule,
    "compact": fam_compact, "numeric": fam_numeric,
}

def build_specs():
    """각 스펙: dict(fam, cols, rows). id/seed/file은 후처리에서 부여."""
    specs = []
    # striped 48: 6색 × 8(cols,rows) 조합
    striped_cr = [(3, 4), (3, 6), (4, 4), (4, 6), (5, 4), (5, 6), (6, 4), (6, 6)]
    for ci in range(6):
        for (cols, rows) in striped_cr:
            specs.append({"fam": "striped", "cols": cols, "rows": rows})
    # 나머지 패밀리: (count, [cols options], [rows options])
    plan = [
        ("borderless", 12, [3, 4, 5], [4, 5]),
        ("underline",  12, [3, 4, 5], [4, 5]),
        ("group2",     12, [5],       [4, 5, 6]),
        ("leftlabel",  12, [4, 5, 6], [4, 5]),
        ("totalrow",   12, [4, 5, 6], [4, 5]),
        ("signal",     12, [4, 5],    [4, 5, 6]),
        ("check",      10, [4, 5],    [4, 5]),
        ("kpichip",    10, [4],       [4, 5]),
        ("cardrow",    10, [1],       [3, 4]),
        ("budget",      8, [5],       [4, 5]),
        ("schedule",    8, [4],       [4, 5]),
        ("compact",     6, [6, 7, 8], [6, 7]),
        ("numeric",     6, [4, 5],    [5, 6]),
    ]
    for fam, cnt, cols_opts, rows_opts in plan:
        for k in range(cnt):
            cols = cols_opts[k % len(cols_opts)]
            rows = rows_opts[(k // len(cols_opts)) % len(rows_opts)]
            specs.append({"fam": fam, "cols": cols, "rows": rows})
    return specs

def build():
    specs = build_specs()
    assert len(specs) == 178, f"spec count={len(specs)} (expected 178)"
    # id/seed/color/file 부여
    for idx, sp in enumerate(specs):
        sp["id"] = f"TBL-{23 + idx:03d}"
        sp["seed"] = 23 + idx
        cname, ccol = HEADER_COLORS[idx % len(HEADER_COLORS)]
        sp["hcname"] = cname; sp["hcol"] = ccol

    # 파일 분할: 패밀리별, ≤25 슬라이드
    from collections import OrderedDict
    by_fam = OrderedDict()
    for sp in specs:
        by_fam.setdefault(sp["fam"], []).append(sp)

    entries = []
    saved_files = []
    for fam, items in by_fam.items():
        chunks = [items[i:i + 24] for i in range(0, len(items), 24)]
        for vi, chunk in enumerate(chunks, start=1):
            frel = f"decks/01_tables/TBL_bulk_{fam}_v{vi}.pptx"
            prs = c.new_deck()
            fn = FAM_FN[fam]
            for si, sp in enumerate(chunk, start=1):
                slide = c.blank_slide(prs)
                tags, params, bindings, editable, rec = fn(
                    slide, sp["id"], sp["hcname"], sp["hcol"], sp["cols"], sp["rows"], sp["seed"])
                entries.append(c.entry(
                    sp["id"], "TBL",
                    f"{THEMES[sp['seed'] % len(THEMES)]} {FAM_TITLE[fam]} ({sp['hcname']}·{sp['cols']}열)",
                    frel, si, tags, params, bindings, editable, recommended_use=rec))
            out = c.save_deck(prs, frel)
            saved_files.append(frel)
            print("SAVED:", out, f"({len(chunk)} slides)")

    entries.sort(key=lambda e: e["id"])
    frag = c.write_fragment("TBL_bulk", entries)
    print("FRAGMENT:", frag)
    print("ENTRIES:", len(entries))
    print("FILES:", len(saved_files))

if __name__ == "__main__":
    build()
