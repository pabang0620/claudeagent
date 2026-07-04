"""
제안서 에셋 라이브러리 — 공통 생성 유틸 (python-pptx)
모든 카테고리 생성 스크립트는 이 모듈을 import 해서 일관성을 유지한다.

핵심 규칙:
- 모든 에셋의 최상위 도형/그룹 이름은 "asset:<ID>" 로 지정 (pptx-automizer 타겟팅 앵커).
- 한글 텍스트는 set_kfont() 로 넣는다 (a:latin/a:ea/a:cs 동시 설정 — python-pptx 한글 버그 우회).
- 페이지 전체 이미지 금지. 모든 요소는 네이티브 shape/table/chart.
- 각 에셋 슬라이드 좌상단에 회색 ID 캡션을 단다(사람이 라이브러리에서 식별용).
"""
import json, os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
TOKENS = json.load(open(os.path.join(ROOT, "design-tokens.json"), encoding="utf-8"))

C = {k: RGBColor.from_string(v) for k, v in TOKENS["color"].items()}
def role(name):  # role -> RGBColor
    return C[TOKENS["role"][name]]
FONT_H = TOKENS["font"]["heading"]
FONT_B = TOKENS["font"]["body"]
SLIDE_W = Inches(TOKENS["meta"]["slide"]["width_in"])
SLIDE_H = Inches(TOKENS["meta"]["slide"]["height_in"])

def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank

def set_kfont(run, name=None, size=None, bold=None, color=None):
    """한글 안전 폰트 지정 — latin/ea/cs 3계열 모두 세팅."""
    name = name or FONT_B
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def add_text(slide, x, y, w, h, text, size=11, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=None, fill=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    for i, line in enumerate(str(text).split("\n")):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run(); r.text = line
        set_kfont(r, font or FONT_B, size, bold, color or role("body_text"))
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill; tb.line.fill.background()
    return tb

def add_box(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
            shape=MSO_SHAPE.RECTANGLE, shadow_off=True):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if shadow_off:
        sp.shadow.inherit = False
    return sp

def set_shape_text(sp, text, size=12, bold=False, color=None, align=PP_ALIGN.CENTER, font=None):
    tf = sp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for i, line in enumerate(str(text).split("\n")):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run(); r.text = line
        set_kfont(r, font or FONT_B, size, bold, color or role("body_text"))
    return sp

def connector(slide, x1, y1, x2, y2, color=None, w=1.5):
    from pptx.enum.shapes import MSO_CONNECTOR
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color or role("border"); cn.line.width = Pt(w)
    return cn

def name_asset(shape, asset_id):
    """pptx-automizer 앵커: 도형/그룹 이름을 asset:<ID> 로."""
    shape.name = "asset:%s" % asset_id
    return shape

def id_caption(slide, asset_id):
    add_text(slide, 0.15, 0.08, 4.0, 0.3, asset_id, size=9, bold=True,
             color=C["gray_500"], align=PP_ALIGN.LEFT)

def _normalize_table_styles(prs):
    """표스타일 GUID를 tableStyles.xml에 실존하는 No-Style로 통일.
    python-pptx 기본 표스타일 GUID는 tableStyles.xml에 정의가 없어, 한컴독스 등
    엄격한 렌더러가 상속 해석 중 크래시한다(우리 표는 수동 테두리라 스타일 불필요)."""
    for slide in prs.slides:
        for sid in slide.shapes._spTree.iter(qn("a:tableStyleId")):
            sid.text = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"


def save_deck(prs, rel_path):
    out = os.path.join(ROOT, rel_path)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    _normalize_table_styles(prs)
    prs.save(out)
    return out

def try_thumbnail(pptx_path):
    """LibreOffice 있으면 슬라이드별 png 생성, 없으면 None(보류)."""
    import shutil, subprocess, tempfile, glob
    if not shutil.which("soffice"):
        return None
    try:
        d = tempfile.mkdtemp()
        subprocess.run(["soffice","--headless","--convert-to","pdf","--outdir",d,pptx_path],
                       check=True, timeout=120, capture_output=True)
        pdf = glob.glob(os.path.join(d,"*.pdf"))[0]
        import fitz
        doc = fitz.open(pdf); outs=[]
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        for i,pg in enumerate(doc):
            p = os.path.join(ROOT,"thumbnails","%s_s%02d.png"%(base,i+1))
            pg.get_pixmap(dpi=96).save(p); outs.append(p)
        return outs
    except Exception:
        return None

def write_fragment(category, entries):
    """카테고리별 매니페스트 조각 저장 (Opus가 병합)."""
    d = os.path.join(ROOT, "_incoming"); os.makedirs(d, exist_ok=True)
    p = os.path.join(d, "manifest_%s.json" % category)
    json.dump(entries, open(p,"w",encoding="utf-8"), ensure_ascii=False, indent=2)
    return p

def group_asset(slide, shapes, asset_id):
    """여러 도형을 asset:<ID> 이름의 <p:grpSp>(그룹)로 묶는다.
    조합 파이프라인(pptx-automizer addElement)이 에셋당 1콜로 통째 복사하게 하는 핵심.
    chOff/chExt를 off/ext와 동일하게 둬서 자식 좌표를 1:1 유지(이동/스케일 없음).
    id_caption 등 라이브러리 라벨은 shapes에 포함하지 말 것.

    ※ 표/차트(graphicFrame)가 shapes에 있으면 그룹으로 묶지 않는다 — 한컴독스 등 엄격
      렌더러가 그룹 내 표스타일 상속을 해석하다 크래시하기 때문. 이 경우 graphicFrame을
      앵커로 삼고 나머지 도형은 top-level로 둔다."""
    gf = next((s for s in shapes
               if getattr(s, "has_table", False) or getattr(s, "has_chart", False)), None)
    if gf is not None:
        gf.name = "asset:%s" % asset_id
        return gf._element
    from lxml import etree
    spTree = slide.shapes._spTree
    xs=[s.left for s in shapes]; ys=[s.top for s in shapes]
    x2=[s.left+s.width for s in shapes]; y2=[s.top+s.height for s in shapes]
    ox,oy=min(xs),min(ys); ex,ey=max(x2)-ox,max(y2)-oy
    # 슬라이드 내 유일 shape id 확보
    ids=[int(e.get("id")) for e in spTree.iter(qn("p:cNvPr")) if e.get("id","").isdigit()]
    nid=(max(ids)+1) if ids else 100
    grp=('<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
         'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
         '<p:nvGrpSpPr><p:cNvPr id="%d" name="asset:%s"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
         '<p:grpSpPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/>'
         '<a:chOff x="%d" y="%d"/><a:chExt cx="%d" cy="%d"/></a:xfrm></p:grpSpPr>'
         '</p:grpSp>')%(nid,asset_id,ox,oy,ex,ey,ox,oy,ex,ey)
    grp_el=etree.fromstring(grp)
    for s in shapes:
        el=s._element; el.getparent().remove(el); grp_el.append(el)
    spTree.append(grp_el)
    return grp_el


def entry(asset_id, category, name, file_rel, slide_idx, tags, params,
          bindings, editable, source="generated:python-pptx", license="self",
          recommended_use=None, quality="approved"):
    """표준 매니페스트 엔트리 생성."""
    return {
        "id": asset_id, "category": category, "name": name,
        "file": file_rel, "slide": slide_idx, "anchor": "asset:%s" % asset_id,
        "tags": tags, "params": params, "bindings": bindings, "editable": editable,
        "source": source, "license": license,
        "recommended_use": recommended_use or [], "quality": quality,
        "thumbnail": None
    }
