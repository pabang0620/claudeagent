# -*- coding: utf-8 -*-
"""
09_headers (HDR) 카테고리 생성기 — 섹션 헤더 / 배지 / 라벨 / 강조 박스
HDR-001 ~ HDR-016. 소형 요소이므로 세트 배치(슬라이드당 다수).

[v1 재생성] 각 에셋은 이제 자체완결 그룹이다:
그 에셋을 구성하는 모든 도형(배경+번호/마커+제목+내용)을 리스트로 모아
c.group_asset(slide, shapes, ID)로 <p:grpSp>(asset:<ID>)로 묶는다.
→ 조합(addElement) 시 그룹 통째로 복사되어 내용이 딸려온다.
속 빈 배경 앵커(과거 결함) 제거. id_caption은 그룹에 넣지 않는다(좌상단 라벨).
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

R = c.role
RR = MSO_SHAPE.ROUNDED_RECTANGLE
RECT = MSO_SHAPE.RECTANGLE

entries = []

def E(asset_id, name, file_rel, slide_idx, tags, params, bindings, editable, rec):
    entries.append(c.entry(asset_id, "HDR", name, file_rel, slide_idx,
                           tags, params, bindings, editable, recommended_use=rec))

# ─────────────────────────────────────────────────────────────
# 에셋 빌더들 — 각 함수는 생성한 도형 리스트를 반환(그룹 대상)
# ─────────────────────────────────────────────────────────────
def build_section_pill(slide, x, y, number, title, fill_color, w=4.6, h=0.6):
    """섹션 번호 pill: 좌측 컬러 사각(번호) + 제목 하이라이트."""
    shp = []
    pill_w = 0.6
    # 좌측 컬러 번호 블록
    numbox = c.add_box(slide, x, y, pill_w, h, fill=fill_color, line=None, shape=RECT)
    c.set_shape_text(numbox, number, size=20, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(numbox)
    # 제목 하이라이트 바
    titlebox = c.add_box(slide, x + pill_w, y, w - pill_w, h,
                         fill=c.C["gray_100"], line=None, shape=RECT)
    c.set_shape_text(titlebox, title, size=15, bold=True, color=R("body_text"),
                     align=PP_ALIGN.LEFT, font=c.FONT_H)
    titlebox.text_frame.margin_left = Pt(10)
    shp.append(titlebox)
    return shp

def build_subtitle_bar(slide, x, y, marker, title, accent, w=4.6, h=0.42):
    """'가./나.' 소제목 하이라이트 바: 좌측 마커 + 밑줄 강조."""
    shp = []
    mk = c.add_text(slide, x, y, 0.5, h, marker, size=14, bold=True,
                    color=accent, align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(mk)
    t = c.add_text(slide, x + 0.5, y, w - 0.5, h, title, size=14, bold=True,
                   color=R("body_text"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(t)
    # 하단 강조 밑줄
    ul = c.add_box(slide, x, y + h - 0.03, w, 0.05, fill=accent, line=None, shape=RECT)
    shp.append(ul)
    return shp

def build_check_bar(slide, x, y, text, fill_color, w=5.0, h=0.55):
    """체크(✓) 헤드라인 바: 네이비 라운드 배경 + 흰 글자."""
    shp = []
    bar = c.add_box(slide, x, y, w, h, fill=fill_color, line=None, shape=RR)
    bar.adjustments[0] = 0.5
    shp.append(bar)
    # 체크 원
    ck = c.add_box(slide, x + 0.14, y + (h - 0.32) / 2, 0.32, 0.32,
                   fill=R("header_text"), line=None, shape=MSO_SHAPE.OVAL)
    c.set_shape_text(ck, "✓", size=13, bold=True, color=fill_color, align=PP_ALIGN.CENTER)
    shp.append(ck)
    t = c.add_text(slide, x + 0.58, y, w - 0.7, h, text, size=14, bold=True,
                   color=R("header_text"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(t)
    return shp

def build_tab_label(slide, x, y, text, fill_color, w=2.0, h=0.5):
    """탭 스타일 라벨: 상단 두 모서리 라운드 + 하단 베이스라인."""
    shp = []
    tab = c.add_box(slide, x, y, w, h, fill=fill_color, line=None,
                    shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    tab.adjustments[0] = 0.35
    c.set_shape_text(tab, text, size=13, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(tab)
    # 하단 베이스라인(자체완결 강조)
    base = c.add_box(slide, x, y + h, w, 0.05, fill=fill_color, line=None, shape=RECT)
    shp.append(base)
    return shp

def build_callout(slide, x, y, title, body, accent, w=4.8, h=1.1):
    """강조 콜아웃: 좌측 컬러 세로바 + 제목/본문."""
    shp = []
    bg = c.add_box(slide, x, y, w, h, fill=c.C["gray_050"], line=c.C["gray_300"],
                   line_w=0.75, shape=RECT)
    shp.append(bg)
    bar = c.add_box(slide, x, y, 0.09, h, fill=accent, line=None, shape=RECT)  # 세로바
    shp.append(bar)
    t = c.add_text(slide, x + 0.28, y + 0.1, w - 0.4, 0.35, title, size=13, bold=True,
                   color=accent, align=PP_ALIGN.LEFT, font=c.FONT_H)
    shp.append(t)
    b = c.add_text(slide, x + 0.28, y + 0.45, w - 0.4, h - 0.55, body, size=11, bold=False,
                   color=R("body_text"), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    shp.append(b)
    return shp

def build_analysis_box(slide, x, y, title, bullets, bg, accent, w=5.4, h=2.2):
    """분석 요약 박스: 제목 밴드 + 불릿."""
    shp = []
    box = c.add_box(slide, x, y, w, h, fill=bg, line=c.C["gray_300"], line_w=0.75, shape=RR)
    box.adjustments[0] = 0.04
    shp.append(box)
    # 제목 밴드
    tb = c.add_box(slide, x + 0.18, y + 0.16, w - 0.36, 0.42, fill=accent, line=None, shape=RECT)
    c.set_shape_text(tb, title, size=13, bold=True, color=R("header_text"),
                     align=PP_ALIGN.LEFT, font=c.FONT_H)
    tb.text_frame.margin_left = Pt(8)
    shp.append(tb)
    yy = y + 0.72
    for b in bullets:
        mk = c.add_text(slide, x + 0.28, yy, 0.3, 0.32, "•", size=12, bold=True,
                        color=accent, align=PP_ALIGN.LEFT)
        shp.append(mk)
        bt = c.add_text(slide, x + 0.52, yy, w - 0.7, 0.32, b, size=11,
                        color=R("body_text"), align=PP_ALIGN.LEFT)
        shp.append(bt)
        yy += 0.36
    return shp

def build_ribbon(slide, x, y, text, fill_color, w=1.7, h=0.72):
    """리본/배지: 하향 리본 도형 + 상단 걸이 바(자체완결)."""
    shp = []
    top = c.add_box(slide, x + w / 2 - 0.5, y - 0.06, 1.0, 0.12, fill=fill_color,
                    line=None, shape=RECT)
    shp.append(top)
    rb = c.add_box(slide, x, y, w, h, fill=fill_color, line=None, shape=MSO_SHAPE.DOWN_RIBBON)
    c.set_shape_text(rb, text, size=14, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(rb)
    return shp

def build_pill_badge(slide, x, y, text, fill_color, w=1.3, h=0.44):
    """알약형 배지(NEW 등) + 좌측 도트 마커(자체완결)."""
    shp = []
    dot = c.add_box(slide, x - 0.02, y + h / 2 - 0.06, 0.12, 0.12, fill=fill_color,
                    line=None, shape=MSO_SHAPE.OVAL)
    shp.append(dot)
    bgp = c.add_box(slide, x + 0.16, y, w, h, fill=fill_color, line=None, shape=RR)
    bgp.adjustments[0] = 0.5
    c.set_shape_text(bgp, text, size=13, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    shp.append(bgp)
    return shp

def build_number_labels(slide, x, y, items, accent, w=5.6, h=0.6):
    """넘버링 라벨 세트: 01/02/03 원형 + 텍스트 (가로 배치, 단일 에셋)."""
    shp = []
    step = w / len(items)
    for i, (num, txt) in enumerate(items):
        cx = x + i * step
        circ = c.add_box(slide, cx, y, 0.6, 0.6, fill=accent, line=None, shape=MSO_SHAPE.OVAL)
        c.set_shape_text(circ, num, size=16, bold=True, color=R("header_text"),
                         align=PP_ALIGN.CENTER, font=c.FONT_H)
        shp.append(circ)
        t = c.add_text(slide, cx + 0.66, y, step - 0.7, 0.6, txt, size=12, bold=True,
                       color=R("body_text"), align=PP_ALIGN.LEFT, font=c.FONT_H)
        shp.append(t)
    return shp

def build_note_box(slide, x, y, text, w=5.6, h=0.85):
    """인용/노트 박스: ※ 출처 표기용."""
    shp = []
    box = c.add_box(slide, x, y, w, h, fill=c.C["gray_100"], line=c.C["gray_300"],
                    line_w=0.75, shape=RECT)
    shp.append(box)
    bar = c.add_box(slide, x, y, 0.09, h, fill=c.C["gray_500"], line=None, shape=RECT)
    shp.append(bar)
    mk = c.add_text(slide, x + 0.24, y + 0.06, 0.4, h - 0.12, "※", size=13, bold=True,
                    color=c.C["gray_700"], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    shp.append(mk)
    t = c.add_text(slide, x + 0.56, y + 0.06, w - 0.7, h - 0.12, text, size=10,
                   color=c.C["gray_700"], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    shp.append(t)
    return shp

# ─────────────────────────────────────────────────────────────
# DECK 1: HDR_section-pill_v1  (섹션 헤더/소제목/체크바/탭)
# ─────────────────────────────────────────────────────────────
F1 = "decks/09_headers/HDR_section-pill_v1.pptx"
prs = c.new_deck()

# slide 1 — 섹션 번호 pill 3색
s = c.blank_slide(prs); c.id_caption(s, "HDR-001~003 · 섹션 번호 pill")
c.group_asset(s, build_section_pill(s, 0.9, 1.3, "1", "사업추진 전략", R("header_fill")), "HDR-001")
c.group_asset(s, build_section_pill(s, 0.9, 2.5, "2", "추진 체계 및 조직", R("accent_primary")), "HDR-002")
c.group_asset(s, build_section_pill(s, 0.9, 3.7, "3", "기대효과 및 활용방안", R("accent_secondary")), "HDR-003")
E("HDR-001", "섹션 번호 pill (네이비)", F1, 1, ["섹션헤더","pill","네이비"],
  {"fill":"navy_800","w_in":4.6,"h_in":0.6}, {"number":"1","title":"사업추진 전략"},
  ["text","color"], ["섹션 구분","목차 헤더"])
E("HDR-002", "섹션 번호 pill (블루)", F1, 1, ["섹션헤더","pill","블루"],
  {"fill":"blue_500","w_in":4.6,"h_in":0.6}, {"number":"2","title":"추진 체계 및 조직"},
  ["text","color"], ["섹션 구분","목차 헤더"])
E("HDR-003", "섹션 번호 pill (틸)", F1, 1, ["섹션헤더","pill","틸"],
  {"fill":"teal_500","w_in":4.6,"h_in":0.6}, {"number":"3","title":"기대효과 및 활용방안"},
  ["text","color"], ["섹션 구분","목차 헤더"])

# slide 2 — 소제목 하이라이트 바 2 + 탭 라벨 1
s = c.blank_slide(prs); c.id_caption(s, "HDR-004~005 소제목바 · HDR-016 탭")
c.group_asset(s, build_subtitle_bar(s, 0.9, 1.4, "가.", "사업 추진 배경", R("accent_primary")), "HDR-004")
c.group_asset(s, build_subtitle_bar(s, 0.9, 2.6, "나.", "핵심 추진 전략", R("sub_header")), "HDR-005")
c.group_asset(s, build_tab_label(s, 0.9, 4.0, "핵심 제안", R("header_fill")), "HDR-016")
# 시각 비교용(메타 미등록) — 그룹/에셋 아님, 단순 참고 도형
build_tab_label(s, 3.1, 4.0, "추진 방향", R("accent_primary"))
E("HDR-004", "소제목 하이라이트 바 (가.)", F1, 2, ["소제목","하이라이트바","밑줄"],
  {"accent":"blue_500","w_in":4.6,"h_in":0.42}, {"marker":"가.","title":"사업 추진 배경"},
  ["text","color"], ["소제목","단락 구분"])
E("HDR-005", "소제목 하이라이트 바 (나.)", F1, 2, ["소제목","하이라이트바","밑줄"],
  {"accent":"purple_600","w_in":4.6,"h_in":0.42}, {"marker":"나.","title":"핵심 추진 전략"},
  ["text","color"], ["소제목","단락 구분"])
E("HDR-016", "탭 스타일 라벨", F1, 2, ["탭","라벨","섹션헤더"],
  {"fill":"navy_800","w_in":2.0,"h_in":0.5}, {"text":"핵심 제안"},
  ["text","color"], ["탭 네비","카테고리 라벨"])

# slide 3 — 체크 헤드라인 바 2
s = c.blank_slide(prs); c.id_caption(s, "HDR-006~007 · 체크 헤드라인 바")
c.group_asset(s, build_check_bar(s, 0.9, 1.5, "차별화된 해외 진출 전략 수립", R("header_fill")), "HDR-006")
c.group_asset(s, build_check_bar(s, 0.9, 2.7, "현지 파트너십 기반 실행력 확보", R("accent_secondary")), "HDR-007")
E("HDR-006", "체크 헤드라인 바 (네이비)", F1, 3, ["체크바","헤드라인","라운드"],
  {"fill":"navy_800","w_in":5.0,"h_in":0.55}, {"text":"차별화된 해외 진출 전략 수립"},
  ["text","color"], ["핵심 강조","헤드라인"])
E("HDR-007", "체크 헤드라인 바 (틸)", F1, 3, ["체크바","헤드라인","라운드"],
  {"fill":"teal_500","w_in":5.0,"h_in":0.55}, {"text":"현지 파트너십 기반 실행력 확보"},
  ["text","color"], ["핵심 강조","헤드라인"])

c.save_deck(prs, F1)

# ─────────────────────────────────────────────────────────────
# DECK 2: HDR_callout-badge_v1  (콜아웃/배지/넘버링)
# ─────────────────────────────────────────────────────────────
F2 = "decks/09_headers/HDR_callout-badge_v1.pptx"
prs = c.new_deck()

# slide 1 — 콜아웃 박스 2
s = c.blank_slide(prs); c.id_caption(s, "HDR-008~009 · 강조 콜아웃 박스")
c.group_asset(s, build_callout(s, 0.9, 1.3, "핵심 제안",
              "현지화 콘텐츠 재가공 및 채널별 최적 배급 전략으로\n해외 시장 진입 장벽을 최소화합니다.", R("accent_primary")), "HDR-008")
c.group_asset(s, build_callout(s, 0.9, 2.9, "기대효과",
              "K-콘텐츠 IP의 글로벌 인지도 제고와\n신규 해외 매출 파이프라인 확보가 가능합니다.", R("warn")), "HDR-009")
E("HDR-008", "강조 콜아웃 박스 (블루)", F2, 1, ["강조박스","콜아웃","세로바"],
  {"accent":"blue_500","w_in":4.8,"h_in":1.1}, {"title":"핵심 제안","body":"현지화 콘텐츠 재가공..."},
  ["text","color"], ["핵심 강조","제안 요약"])
E("HDR-009", "강조 콜아웃 박스 (레드)", F2, 1, ["강조박스","콜아웃","세로바"],
  {"accent":"red_500","w_in":4.8,"h_in":1.1}, {"title":"기대효과","body":"K-콘텐츠 IP..."},
  ["text","color"], ["핵심 강조","기대효과"])

# slide 2 — 리본/배지
s = c.blank_slide(prs); c.id_caption(s, "HDR-012~013 · 리본/배지")
c.group_asset(s, build_ribbon(s, 1.0, 1.5, "추천", R("warn")), "HDR-012")
c.group_asset(s, build_pill_badge(s, 1.0, 3.0, "NEW", R("accent_point")), "HDR-013")
# 시각 비교용(메타 미등록)
build_ribbon(s, 3.0, 1.5, "핵심", R("header_fill"))
build_pill_badge(s, 2.9, 3.0, "필수", R("accent_secondary"))
E("HDR-012", "리본 배지 (추천/핵심)", F2, 2, ["배지","리본","추천"],
  {"fill":"red_500","w_in":1.7,"h_in":0.72}, {"text":"추천"},
  ["text","color"], ["핵심 강조","우선순위 표시"])
E("HDR-013", "알약형 배지 (NEW)", F2, 2, ["배지","알약","NEW"],
  {"fill":"cyan_500","w_in":1.3,"h_in":0.44}, {"text":"NEW"},
  ["text","color"], ["신규 표시","상태 라벨"])

# slide 3 — 넘버링 라벨 세트
s = c.blank_slide(prs); c.id_caption(s, "HDR-014 · 넘버링 라벨 세트")
c.group_asset(s, build_number_labels(s, 0.9, 2.0,
              [("01","기획"),("02","제작"),("03","배급")], R("accent_primary"), w=8.0), "HDR-014")
E("HDR-014", "넘버링 라벨 세트 (01/02/03)", F2, 3, ["넘버링","라벨","원형","세트"],
  {"accent":"blue_500","count":3,"w_in":8.0}, {"items":[["01","기획"],["02","제작"],["03","배급"]]},
  ["text","color"], ["단계 표시","순서 나열"])

c.save_deck(prs, F2)

# ─────────────────────────────────────────────────────────────
# DECK 3: HDR_analysis-box_v1  (분석 박스/노트)
# ─────────────────────────────────────────────────────────────
F3 = "decks/09_headers/HDR_analysis-box_v1.pptx"
prs = c.new_deck()

# slide 1 — 베이지 분석 박스
s = c.blank_slide(prs); c.id_caption(s, "HDR-010 · 베이지 분석 요약 박스")
c.group_asset(s, build_analysis_box(s, 0.9, 1.2, "시장 환경 분석",
              ["글로벌 OTT 수요 확대로 K-콘텐츠 진출 기회 증가",
               "현지 언어·문화 장벽이 핵심 리스크로 상존",
               "권역별 배급 파트너 확보가 성패 좌우"],
              c.C["beige_100"], R("sub_header")), "HDR-010")
E("HDR-010", "베이지 분석 요약 박스", F3, 1, ["분석박스","요약","베이지","불릿"],
  {"bg":"beige_100","accent":"purple_600","w_in":5.4,"h_in":2.2},
  {"title":"시장 환경 분석","bullets":["글로벌 OTT 수요 확대...","현지 언어·문화 장벽...","권역별 배급 파트너..."]},
  ["text","color"], ["분석 요약","현황 정리"])

# slide 2 — 그레이 분석 박스
s = c.blank_slide(prs); c.id_caption(s, "HDR-011 · 그레이 분석 요약 박스")
c.group_asset(s, build_analysis_box(s, 0.9, 1.2, "SWOT 요약",
              ["강점: 검증된 IP·제작 역량 보유",
               "약점: 해외 유통망·현지 네트워크 부족",
               "기회: 한류 확산 및 정책 지원 확대"],
              c.C["gray_100"], R("header_fill")), "HDR-011")
E("HDR-011", "그레이 분석 요약 박스", F3, 2, ["분석박스","요약","그레이","불릿"],
  {"bg":"gray_100","accent":"navy_800","w_in":5.4,"h_in":2.2},
  {"title":"SWOT 요약","bullets":["강점: 검증된 IP...","약점: 해외 유통망...","기회: 한류 확산..."]},
  ["text","color"], ["분석 요약","SWOT"])

# slide 3 — 인용/노트 박스
s = c.blank_slide(prs); c.id_caption(s, "HDR-015 · 인용/노트(출처) 박스")
c.group_asset(s, build_note_box(s, 0.9, 1.5,
              "출처: 한국콘텐츠진흥원, 「2025 해외 콘텐츠시장 분석보고서」, 2025."), "HDR-015")
E("HDR-015", "인용/노트 박스 (출처)", F3, 3, ["노트","인용","출처","각주"],
  {"w_in":5.6,"h_in":0.85}, {"text":"출처: 한국콘텐츠진흥원, 「2025 해외 콘텐츠시장 분석보고서」, 2025."},
  ["text"], ["출처 표기","각주"])

c.save_deck(prs, F3)

# ─────────────────────────────────────────────────────────────
frag = c.write_fragment("HDR", entries)
print("SAVED:", F1, F2, F3)
print("FRAGMENT:", frag)
print("ENTRIES:", len(entries))
for e in entries:
    print(" ", e["id"], e["name"])
