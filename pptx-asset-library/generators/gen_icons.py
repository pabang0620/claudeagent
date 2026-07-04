# -*- coding: utf-8 -*-
"""ICN 카테고리(플랫 아이콘형 도식) 에셋 생성 — ICN-001~ICN-020.

- 전부 네이티브 도형 조합 (외부 이미지/아이콘 파일 금지, 라이선스 self).
- 2색 원칙: 네이비(navy_800) 스트로크/솔리드 + 포인트1(cyan_500) 액센트. 화이트/그레이는 중립.
- 96x96 가상 그리드에 통일 스타일로 배치. 각 아이콘 = 자체완결 그룹(<p:grpSp> asset:ICN-0NN 앵커).
  아이콘 획 도형들(+하단 한글 라벨)을 전부 c.group_asset 으로 묶어 그룹 자체가 앵커가 된다.
  → 조합(addElement) 시 아이콘 획이 통째로 딸려옴(투명 컨테이너 rect 방식 폐기).
- 아이콘 시트: 슬라이드당 6~8개 배치(개별 ID 캡션은 그룹 밖 + 하단 한글 라벨은 그룹 안).
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

R = c.role
CC = c.C

# ---- 2색 팔레트 (네이비 + 포인트1) + 중립 ----
STROKE = CC["navy_800"]
NAVY = CC["navy_800"]
NAVY6 = CC["navy_600"]
ACC = CC["cyan_500"]      # 포인트 컬러 (유일)
WHITE = CC["white"]
GRAY = CC["gray_300"]
LW = 2.0                    # 통일 스트로크 두께(pt)


# ---------- 그리드 페인터 (96x96 가상 그리드 → 아이콘 박스 내 절대 인치) ----------
class Painter:
    def __init__(self, slide, ix, iy, S):
        self.s, self.ix, self.iy, self.S = slide, ix, iy, S
        self.shapes = []   # 이 아이콘이 그린 모든 도형(그룹 대상)

    def X(self, u): return self.ix + (u / 96.0) * self.S
    def Y(self, u): return self.iy + (u / 96.0) * self.S
    def L(self, u): return (u / 96.0) * self.S

    def box(self, x, y, w, h, fill=None, line=STROKE, lw=LW,
            shape=MSO_SHAPE.RECTANGLE, rot=None):
        sp = c.add_box(self.s, self.X(x), self.Y(y), self.L(w), self.L(h),
                       fill=fill, line=line, line_w=lw, shape=shape)
        if rot is not None:
            sp.rotation = rot
        self.shapes.append(sp)
        return sp

    def oval(self, x, y, w, h, fill=None, line=STROKE, lw=LW):
        # box()가 self.shapes 에 이미 append 하므로 여기서 중복 추가 금지
        return self.box(x, y, w, h, fill=fill, line=line, lw=lw, shape=MSO_SHAPE.OVAL)

    def line(self, x1, y1, x2, y2, color=STROKE, w=LW):
        cn = c.connector(self.s, self.X(x1), self.Y(y1), self.X(x2), self.Y(y2),
                         color=color, w=w)
        self.shapes.append(cn)
        return cn


# =========================================================
# 아이콘 드로잉 (각 함수: painter P 위에 도형 조합)
# =========================================================
def i_document(P):
    P.box(26, 12, 44, 72, fill=WHITE, line=STROKE)          # 문서 본체
    P.box(34, 22, 20, 5, fill=NAVY, line=None)              # 제목줄(네이비)
    for yy in (38, 50, 62):                                  # 본문줄(액센트)
        P.line(34, yy, 62, yy, color=ACC, w=2.5)

def i_barchart(P):
    P.line(24, 16, 24, 80, color=STROKE, w=2.5)             # y축
    P.line(24, 80, 82, 80, color=STROKE, w=2.5)             # x축
    P.box(32, 56, 12, 24, fill=NAVY, line=None)
    P.box(50, 40, 12, 40, fill=NAVY6, line=None)
    P.box(68, 26, 12, 54, fill=ACC, line=None)              # 최고값 = 액센트

def i_growth(P):
    P.box(24, 60, 12, 20, fill=NAVY6, line=None)            # 상승 막대
    P.box(40, 46, 12, 34, fill=NAVY6, line=None)
    P.box(56, 32, 12, 48, fill=NAVY, line=None)
    P.box(72, 20, 14, 60, fill=ACC, line=None,
          shape=MSO_SHAPE.UP_ARROW)                          # 성장 상향 화살표

def i_magnifier(P):
    P.box(22, 22, 42, 42, fill=WHITE, line=STROKE, lw=2.5,
          shape=MSO_SHAPE.DONUT)                             # 렌즈 링
    P.line(56, 56, 80, 80, color=ACC, w=5)                  # 손잡이 = 액센트

def i_checkcircle(P):
    P.oval(20, 20, 56, 56, fill=NAVY, line=None)            # 원
    P.line(33, 50, 45, 62, color=ACC, w=5)                  # 체크(액센트)
    P.line(45, 62, 66, 34, color=ACC, w=5)

def i_pin(P):
    P.box(36, 40, 24, 40, fill=NAVY, line=None,
          shape=MSO_SHAPE.ISOSCELES_TRIANGLE, rot=180)        # 하단 뾰족
    P.oval(28, 14, 40, 40, fill=NAVY, line=None)            # 핀 머리
    P.oval(40, 26, 16, 16, fill=ACC, line=None)            # 중심점(액센트)

def i_database(P):
    P.box(26, 18, 44, 60, fill=WHITE, line=STROKE, lw=2.2,
          shape=MSO_SHAPE.CAN)                               # 실린더
    P.line(26, 44, 70, 44, color=ACC, w=2)                  # 데이터 레이어(액센트)
    P.line(26, 60, 70, 60, color=ACC, w=2)

def i_network(P):
    P.line(48, 46, 52, 22, color=STROKE, w=2)               # 연결선(먼저=뒤)
    P.line(44, 52, 26, 68, color=STROKE, w=2)
    P.line(56, 52, 70, 68, color=STROKE, w=2)
    P.oval(40, 40, 16, 16, fill=NAVY, line=None)            # 중심 노드
    P.oval(44, 12, 16, 16, fill=ACC, line=None)            # 상단 노드(액센트)
    P.oval(16, 62, 16, 16, fill=NAVY, line=None)
    P.oval(64, 62, 16, 16, fill=NAVY, line=None)

def i_shield(P):
    P.box(28, 44, 44, 36, fill=NAVY, line=None,
          shape=MSO_SHAPE.ISOSCELES_TRIANGLE, rot=180)        # 하단 포인트
    P.box(28, 14, 44, 40, fill=NAVY, line=None,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)                 # 상단 몸체
    P.line(38, 40, 46, 50, color=ACC, w=4)                  # 체크(액센트)
    P.line(46, 50, 62, 28, color=ACC, w=4)

def i_target(P):
    P.box(20, 20, 56, 56, fill=WHITE, line=STROKE, lw=2.5,
          shape=MSO_SHAPE.DONUT)                             # 외곽 링
    P.box(34, 34, 28, 28, fill=WHITE, line=ACC, lw=3.5,
          shape=MSO_SHAPE.DONUT)                             # 중간 링(액센트)
    P.oval(42, 42, 12, 12, fill=NAVY, line=None)            # 중심

def i_monitor(P):
    P.box(18, 16, 60, 44, fill=WHITE, line=STROKE,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)                 # 화면 프레임
    P.box(24, 22, 48, 32, fill=ACC, line=None)             # 화면(액센트)
    P.box(44, 60, 8, 10, fill=NAVY, line=None)             # 목
    P.box(32, 70, 32, 6, fill=NAVY, line=None,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)                 # 받침

def i_mobile(P):
    P.box(32, 12, 32, 72, fill=WHITE, line=STROKE,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)                 # 본체
    P.box(37, 22, 22, 48, fill=ACC, line=None)             # 스크린(액센트)
    P.oval(45, 74, 6, 6, fill=NAVY, line=None)             # 홈버튼

def i_mail(P):
    P.box(18, 26, 60, 44, fill=WHITE, line=STROKE)          # 봉투
    P.line(18, 28, 48, 52, color=ACC, w=2.5)               # 플랩(액센트)
    P.line(48, 52, 78, 28, color=ACC, w=2.5)

def i_calendar(P):
    P.box(20, 22, 56, 54, fill=WHITE, line=STROKE)          # 몸체
    P.box(20, 22, 56, 14, fill=NAVY, line=None)            # 헤더(네이비)
    P.box(32, 14, 5, 12, fill=NAVY, line=None)             # 고리
    P.box(59, 14, 5, 12, fill=NAVY, line=None)
    P.box(42, 50, 12, 12, fill=ACC, line=None)             # 강조된 날짜(액센트)

def i_flag(P):
    P.box(28, 14, 5, 68, fill=NAVY, line=None)             # 깃대
    P.box(33, 16, 40, 26, fill=ACC, line=None)             # 깃발(액센트)
    P.oval(24, 78, 11, 8, fill=NAVY, line=None)            # 받침

def i_bulb(P):
    P.oval(30, 12, 36, 36, fill=WHITE, line=STROKE)         # 전구
    P.box(40, 46, 16, 10, fill=NAVY, line=None)            # 소켓
    P.box(42, 56, 12, 7, fill=NAVY, line=None)
    P.line(42, 30, 48, 40, color=ACC, w=2.5)               # 필라멘트(액센트)
    P.line(48, 40, 54, 30, color=ACC, w=2.5)
    P.line(20, 18, 26, 24, color=ACC, w=2)                 # 광선(액센트)
    P.line(48, 6, 48, 12, color=ACC, w=2)
    P.line(76, 18, 70, 24, color=ACC, w=2)

def i_clock(P):
    P.box(20, 18, 56, 56, fill=WHITE, line=STROKE, lw=2.5,
          shape=MSO_SHAPE.DONUT)                             # 시계 링
    P.line(48, 46, 48, 30, color=NAVY, w=3.5)              # 시침
    P.line(48, 46, 62, 52, color=ACC, w=2.5)               # 분침(액센트)
    P.oval(44, 42, 8, 8, fill=NAVY, line=None)             # 중심

def i_gear(P):
    P.box(18, 16, 60, 60, fill=NAVY, line=None,
          shape=MSO_SHAPE.GEAR_6)                            # 톱니
    P.oval(38, 36, 20, 20, fill=WHITE, line=None)          # 중앙 구멍
    P.oval(44, 42, 8, 8, fill=ACC, line=None)             # 중심점(액센트)

def i_cloud(P):
    P.box(14, 28, 68, 44, fill=WHITE, line=STROKE,
          shape=MSO_SHAPE.CLOUD)                             # 구름
    P.box(43, 40, 10, 24, fill=ACC, line=None,
          shape=MSO_SHAPE.UP_ARROW)                          # 업로드(액센트)

def i_link(P):
    P.box(18, 38, 38, 22, fill=None, line=NAVY, lw=4.5,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)                 # 링1
    P.box(40, 38, 38, 22, fill=None, line=ACC, lw=4.5,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)                 # 링2(액센트)


# =========================================================
# 아이콘 레지스트리
# =========================================================
ICONS = [
    ("ICN-001", "document",  "문서",       i_document,  ["문서", "리포트", "자료"]),
    ("ICN-002", "barchart",  "막대그래프", i_barchart,  ["차트", "통계", "실적"]),
    ("ICN-003", "growth",    "성장",       i_growth,    ["성장", "상향", "증가"]),
    ("ICN-004", "magnifier", "분석",       i_magnifier, ["분석", "돋보기", "탐색"]),
    ("ICN-005", "checkcircle","검증",      i_checkcircle,["검증", "완료", "승인"]),
    ("ICN-006", "pin",       "위치",       i_pin,       ["위치", "지역", "현장"]),
    ("ICN-007", "database",  "데이터베이스", i_database,  ["데이터", "저장", "DB"]),
    ("ICN-008", "network",   "네트워크",   i_network,   ["연계", "네트워크", "협력"]),
    ("ICN-009", "shield",    "보안",       i_shield,    ["보안", "보호", "안전"]),
    ("ICN-010", "target",    "목표",       i_target,    ["목표", "타깃", "정확도"]),
    ("ICN-011", "monitor",   "모니터",     i_monitor,   ["시스템", "화면", "플랫폼"]),
    ("ICN-012", "mobile",    "모바일",     i_mobile,    ["모바일", "앱", "서비스"]),
    ("ICN-013", "mail",      "메일",       i_mail,      ["메일", "소통", "안내"]),
    ("ICN-014", "calendar",  "일정",       i_calendar,  ["일정", "달력", "기간"]),
    ("ICN-015", "flag",      "목표지점",   i_flag,      ["마일스톤", "깃발", "달성"]),
    ("ICN-016", "bulb",      "아이디어",   i_bulb,      ["아이디어", "혁신", "기획"]),
    ("ICN-017", "clock",     "시간",       i_clock,     ["시간", "효율", "신속"]),
    ("ICN-018", "gear",      "설정",       i_gear,      ["설정", "운영", "관리"]),
    ("ICN-019", "cloud",     "클라우드",   i_cloud,     ["클라우드", "업로드", "인프라"]),
    ("ICN-020", "link",      "연결",       i_link,      ["연계", "링크", "통합"]),
]


# ---------- 시트 배치 ----------
def title_block(slide, text, sub):
    c.add_text(slide, 0.55, 0.42, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    c.add_text(slide, 0.55, 0.92, 12.2, 0.3, sub, size=11,
               color=R("muted_text"), align=PP_ALIGN.LEFT)


def place_icon(slide, meta, col, row):
    asset_id, glyph, label, fn, _tags = meta
    S = 1.25
    left, top = 0.75, 1.55
    colw, rowh = (13.333 - 1.5) / 4.0, 2.75
    cellx = left + col * colw
    ix = cellx + (colw - S) / 2.0
    iy = top + row * rowh + 0.30
    # 개별 ID 캡션(아이콘 상단) — 라이브러리 라벨이므로 그룹 밖(자산에 포함 안 함)
    c.add_text(slide, ix, iy - 0.30, S, 0.24, asset_id, size=9, bold=True,
               color=CC["gray_500"], align=PP_ALIGN.CENTER)
    # 아이콘 도형 조합 — Painter 가 그린 모든 도형을 수집
    P = Painter(slide, ix, iy, S)
    fn(P)
    # 하단 한글 라벨 — 자산의 일부이므로 그룹 안에 포함
    lbl = c.add_text(slide, cellx, iy + S + 0.06, colw, 0.32, label, size=12, bold=True,
                     color=R("body_text"), align=PP_ALIGN.CENTER)
    # 자체완결 그룹화 → 그룹 자체가 asset:ICN-0NN 앵커(투명 컨테이너 폐기)
    c.group_asset(slide, P.shapes + [lbl], asset_id)


# =========================================================
# 빌드
# =========================================================
FILE_REL = "decks/11_icons/ICN_flat-set_v1.pptx"
entries = []
p = c.new_deck()

# 슬라이드당 8개(4x2). 20개 → 8/8/4 = 3 슬라이드
chunks = [ICONS[0:8], ICONS[8:16], ICONS[16:20]]
slide_of = {}
for si, chunk in enumerate(chunks):
    s = c.blank_slide(p)
    rng = "%s ~ %s" % (chunk[0][0], chunk[-1][0])
    c.id_caption(s, rng)
    title_block(s, "아이콘 세트 · ICN (%d/%d)" % (si + 1, len(chunks)),
                "제안서용 플랫 아이콘 · 네이티브 도형 조합 · 2색(네이비+시안) · 편집 가능")
    for k, meta in enumerate(chunk):
        place_icon(s, meta, col=k % 4, row=k // 4)
        slide_of[meta[0]] = si + 1  # 1-based

out = c.save_deck(p, FILE_REL)

# ---------- 메타데이터 ----------
for asset_id, glyph, label, fn, tags in ICONS:
    full_tags = ["아이콘"] + tags + ["플랫", "2색"]
    entries.append(c.entry(
        asset_id, "ICN", "아이콘 · %s" % label, FILE_REL, slide_of[asset_id],
        full_tags,
        {"glyph": glyph, "colors": 2},
        {"label": label},
        ["color", "label", "size"],
        recommended_use=["항목 아이콘", "불릿 강조", "섹션 상징"],
    ))

frag = c.write_fragment("ICN", entries)
print("SAVED:", out)
print("FRAGMENT:", frag)
print("COUNT:", len(entries))
print("SLIDES:", len(chunks))
