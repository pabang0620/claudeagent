# -*- coding: utf-8 -*-
"""CMP 대량 확충 배치 생성기 — CMP-019 ~ CMP-200 (182개).
구조 패밀리 × 색 테마 × 항목수 스윕. 1슬라이드 = 1에셋.
표 계열: graphicFrame(asset:<ID>) 앵커. 다중도형/카드: c.group_asset 자체완결 그룹.
페이지이미지·SmartArt 금지. 색은 c.role/c.C만.
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

R = c.role
CC = c.C

# ---------------------------------------------------------------------------
# 색 테마 (accent / accent2 / header 다크) — 색 스윕용
# ---------------------------------------------------------------------------
THEMES = [
    {"key": "blue",   "accent": R("accent_primary"),   "accent2": R("accent_secondary"), "header": CC["navy_800"]},
    {"key": "teal",   "accent": R("accent_secondary"), "accent2": R("accent_primary"),   "header": CC["navy_600"]},
    {"key": "cyan",   "accent": R("accent_point"),     "accent2": R("accent_secondary"), "header": CC["navy_900"]},
    {"key": "purple", "accent": R("sub_header"),       "accent2": R("accent_primary"),   "header": CC["navy_800"]},
    {"key": "navy",   "accent": CC["navy_600"],        "accent2": R("accent_point"),     "header": CC["navy_900"]},
]
NTH = len(THEMES)

# ---------------------------------------------------------------------------
# 콘텐츠 풀 (리터럴 중복 방지 — 오프셋 슬라이스로 변주)
# ---------------------------------------------------------------------------
EVAL_ITEMS = ["글로벌 대응 역량", "협업 프로세스", "AI 활용 성숙도", "품질 관리 체계",
              "가격 경쟁력", "납기 대응력", "기술 전문성", "데이터 보안 수준",
              "사용자 경험", "시스템 확장성", "유지보수 편의", "고객 지원 체계",
              "시장 이해도", "운영 효율성", "혁신 역량", "브랜드 신뢰도",
              "현지화 수준", "리스크 관리", "규제 준수", "연동 용이성"]
CRITERIA = ["기술성", "경제성", "실현성", "확장성", "안정성", "효율성", "혁신성", "지속성"]
COMPETITORS = ["자사", "경쟁사 A", "경쟁사 B", "경쟁사 C", "경쟁사 D"]
OPTION_COLS = ["기본", "표준", "프리미엄", "엔터프라이즈", "라이트", "프로"]
TIER_NAMES = [("Basic", "기본형"), ("Standard", "표준형"), ("Premium", "프리미엄"),
              ("Enterprise", "확장형"), ("Starter", "스타터"), ("Pro", "프로")]
TIER_PRICES = ["50", "100", "180", "250", "320", "맞춤"]
TIER_FEATS = ["핵심 기능 제공", "표준 기술 지원", "기본 리포트", "AI 활용 지원",
              "글로벌 다국어 대응", "우선 지원 채널", "맞춤 개발 포함", "전담 협업팀",
              "무제한 확장", "고급 분석 리포트", "SLA 보장", "전용 인프라"]
QUAD_TITLES = ["집중 육성", "선택 확대", "현상 유지", "재검토", "핵심 과제", "전략 과제",
               "보완 과제", "관망 과제", "즉시 실행", "단계 확대", "방어 영역", "축소 검토"]
QUAD_DESC = ["역량 우위 영역", "강화 필요 영역", "안정화 영역", "우선순위 후순위",
             "즉시 착수 대상", "단계적 확대", "리스크 관리", "관망 대상",
             "최우선 투자", "점진 투자", "현상 유지", "축소 대상"]
AXES = [("시장 매력도", "실행 역량"), ("기대 효과", "투자 규모"), ("성장성", "수익성"),
        ("중요도", "긴급도"), ("전략 적합성", "실현 가능성")]
SWOT_ITEMS = {
    "S": ["글로벌 네트워크 보유", "검증된 기술 역량", "풍부한 수행 실적", "강력한 파트너십"],
    "W": ["AI 인력 부족", "브랜드 인지도 열위", "자금 조달 제약", "레거시 의존"],
    "O": ["해외 시장 확대", "정부 지원 정책", "디지털 전환 수요", "신규 채널 부상"],
    "T": ["경쟁 심화", "환율 변동성", "규제 강화 흐름", "기술 표준 변화"],
}
PROS = ["글로벌 표준 대응 체계 확보", "AI 기반 품질 자동화", "통합 협업 플랫폼 운영",
        "데이터 기반 의사결정", "확장 가능한 아키텍처", "검증된 보안 체계"]
CONS = ["초기 도입 비용 부담", "내부 역량 확보 기간 필요", "레거시 시스템 연동 복잡",
        "변화 관리 리스크", "운영 인력 추가 필요", "학습 곡선 존재"]
RANK_TARGETS = ["자사 컨소시엄", "경쟁사 A", "경쟁사 B", "경쟁사 C", "경쟁사 D",
                "경쟁사 E", "경쟁사 F"]
RANK_NOTE = {1: "선정 우선", 2: "차순위", 3: "예비"}
BA_BEFORE = ["글로벌 대응 체계 부재", "수작업 중심 품질 관리", "협업 프로세스 분산",
             "AI 미적용", "데이터 사일로", "대응 속도 지연"]
BA_AFTER = ["글로벌 표준 대응 체계", "AI 기반 품질 자동화", "통합 협업 플랫폼",
            "데이터 기반 의사결정", "단일 데이터 허브", "실시간 즉시 대응"]
ALTS = ["대안 A (자체구축)", "대안 B (협업제휴)", "대안 C (외주위탁)", "대안 D (혼합형)"]
PROP_HEADS = ["추진 배경", "기대 효과", "핵심 전략", "실행 방안", "차별화 요소",
              "리스크 대응", "운영 계획", "성과 지표"]
PROP_SENTS = ["글로벌 표준에 부합하는 대응 체계를 구축한다", "AI 기반 자동화로 품질 편차를 최소화한다",
              "통합 협업 플랫폼으로 커뮤니케이션 비용을 절감한다", "데이터 기반 의사결정 체계를 정착시킨다",
              "현지 파트너십을 통해 시장 진입 장벽을 낮춘다", "단계별 마일스톤으로 리스크를 관리한다",
              "전담 조직을 구성해 실행력을 확보한다", "정량 지표로 성과를 추적한다",
              "사용자 피드백을 신속히 반영한다", "보안·규제 요건을 선제적으로 충족한다",
              "확장 가능한 아키텍처로 미래 수요에 대응한다", "비용 효율적 운영 모델을 적용한다"]
PROP_TITLES = ["사업 추진 제안", "핵심 전략 요약", "실행 계획 개요", "차별화 전략",
               "기대 성과 제안", "운영 방안 요약"]


def pick(pool, start, n):
    return [pool[(start + i) % len(pool)] for i in range(n)]


# ---------------------------------------------------------------------------
# 공통 표/텍스트 헬퍼
# ---------------------------------------------------------------------------
def add_table(slide, x, y, w, h, nrows, ncols, col_widths=None, row_h=None):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    return gf, tbl


def set_cell(cell, text, size=11, bold=False, color=None, fill=None,
             align=PP_ALIGN.CENTER, font=None):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill is not None else R("row_base")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(5)
    cell.margin_top = cell.margin_bottom = Pt(3)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))


def title_block(slide, text, sub=None):
    c.add_text(slide, 0.55, 0.52, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.02, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)


def _bullets(slide, x, y, w, items, dot_color, step=0.42, size=11):
    """불릿 리스트 그려서 shape 리스트 반환 (group 수집용)."""
    out = []
    for i, it in enumerate(items):
        yy = y + i * step
        d = c.add_box(slide, x, yy + 0.11, 0.12, 0.12, fill=dot_color, line=None,
                      shape=MSO_SHAPE.OVAL)
        t = c.add_text(slide, x + 0.28, yy - 0.02, w - 0.28, step, it, size=size,
                       color=R("body_text"), align=PP_ALIGN.LEFT)
        out += [d, t]
    return out


# ---------------------------------------------------------------------------
# 패밀리 그리기 함수 — 각기 metadata dict 반환
# ---------------------------------------------------------------------------
def fam_matrix(slide, aid, vi):
    th = THEMES[vi % NTH]
    axis_x, axis_y = AXES[vi % len(AXES)]
    titles = pick(QUAD_TITLES, vi, 4)
    descs = pick(QUAD_DESC, vi, 4)
    palette = [th["accent"], th["accent2"], R("muted_text"), th["header"]]
    quads = [{"title": t, "desc": d} for t, d in zip(titles, descs)]
    title_block(slide, "2x2 전략 포지셔닝 매트릭스", "축 라벨과 4개 사분면을 자유 편집 (%s)" % th["key"])
    shapes = []
    x0, x1, y0, y1 = 2.95, 11.35, 1.55, 6.15
    gap = 0.30
    bw = (x1 - x0 - gap) / 2
    bh = (y1 - y0 - gap) / 2
    cx = x0 + bw + gap / 2
    cy = y0 + bh + gap / 2
    coords = [(x0, y0), (cx + gap / 2, y0), (x0, cy + gap / 2), (cx + gap / 2, cy + gap / 2)]
    for (bx, by), q, fill in zip(coords, quads, palette):
        box = c.add_box(slide, bx, by, bw, bh, fill=fill, line=R("border"), line_w=1.0,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = q["title"]
        c.set_kfont(r, c.FONT_H, 15, True, R("header_text"))
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = q["desc"]
        c.set_kfont(r2, c.FONT_B, 10, False, R("header_text"))
        shapes.append(box)
    shapes.append(c.connector(slide, x0, cy, x1, cy, color=R("muted_text"), w=1.5))
    shapes.append(c.connector(slide, cx, y0, cx, y1, color=R("muted_text"), w=1.5))
    xl = c.add_text(slide, x0, y1 + 0.12, x1 - x0, 0.35, "%s  →" % axis_x, size=12, bold=True,
                    color=th["accent"], align=PP_ALIGN.CENTER)
    yl = c.add_text(slide, 0.1, cy - 0.18, y1 - y0, 0.36, "%s  →" % axis_y, size=12, bold=True,
                    color=th["accent"], align=PP_ALIGN.CENTER)
    yl.left = Inches((x0 - 0.45) - (y1 - y0) / 2)
    yl.top = Inches(cy - 0.18)
    yl.rotation = 270
    shapes += [xl, yl]
    c.group_asset(slide, shapes, aid)
    return {"name": "2x2 전략 포지셔닝 매트릭스", "tags": ["매트릭스", "2x2", "사분면", "포지셔닝", th["key"]],
            "params": {"family": "matrix2x2", "quadrants": 4, "axis": 2, "theme": th["key"]},
            "bindings": {"quadrants": ["q1", "q2", "q3", "q4"], "axis": ["x", "y"]},
            "editable": ["quad_titles", "quad_desc", "axis_labels", "quad_fill"],
            "use": ["전략 포지셔닝", "우선순위", "경쟁력 비교"]}


def fam_eval(slide, aid, vi):
    th = THEMES[vi % NTH]
    ncols = 3 + (vi % 2)   # 3 또는 4
    nrows = 4 + (vi % 4)   # 4~7
    items = pick(EVAL_ITEMS, vi, nrows)
    GOOD = ("○ 우수", R("accent_secondary"))
    MID = ("△ 보통", R("muted_text"))
    BAD = ("✕ 미흡", R("warn"))
    heads = ["평가 항목", "현행", "개선"] + (["목표"] if ncols == 4 else [])
    title_block(slide, "%d열 평가표 (현행 vs 개선)" % ncols, "항목별 현행·개선 상태를 기호로 표기 (%s)" % th["key"])
    x, y, w = 1.1, 1.55, 11.1
    first_w = 4.7 if ncols == 3 else 4.1
    rest_w = (w - first_w) / (ncols - 1)
    col_widths = [first_w] + [rest_w] * (ncols - 1)
    tn = nrows + 1
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, ncols, col_widths=col_widths, row_h=[rh] * tn)
    c.name_asset(gf, aid)
    for j, hd in enumerate(heads):
        fill = th["header"] if j != ncols - 1 else th["accent"]
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, it in enumerate(items, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), it, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        cur = BAD if (i + vi) % 2 == 0 else MID
        set_cell(tbl.cell(i, 1), "%s" % cur[0], size=11, bold=True, color=cur[1], fill=base)
        set_cell(tbl.cell(i, 2), GOOD[0], size=11, bold=True, color=GOOD[1], fill=base)
        if ncols == 4:
            tgt = GOOD if (i + vi) % 3 else MID
            set_cell(tbl.cell(i, 3), tgt[0], size=11, bold=True, color=tgt[1], fill=base)
    return {"name": "%d열 평가표 (현행 vs 개선)" % ncols,
            "tags": ["평가표", "비교표", "현행개선", "기호평가", th["key"]],
            "params": {"family": "evaltable", "cols": ncols, "rows": nrows, "symbols": ["○", "△", "✕"], "theme": th["key"]},
            "bindings": {"rows": items, "cols": heads},
            "editable": ["rows", "cell_symbols", "cell_color", "headers"],
            "use": ["개선 효과 제시", "현황 진단", "평가 기준"]}


def fam_versus(slide, aid, vi):
    th = THEMES[vi % NTH]
    ncomp = 2 + (vi % 3)   # 경쟁 대상 수 → 열: 자사 + 경쟁 (2~4)
    nrows = 5 + (vi % 3)   # 5~7
    items = pick(EVAL_ITEMS, vi, nrows)
    cols = ["평가 항목"] + pick(COMPETITORS, 0, ncomp + 1)[:ncomp + 1]
    cols = ["평가 항목"] + COMPETITORS[:ncomp + 1]
    ncol = len(cols)
    GOOD = ("○", R("accent_secondary"))
    MID = ("△", R("muted_text"))
    BAD = ("✕", R("warn"))
    glyphs = [GOOD, MID, BAD]
    own_col = 1
    title_block(slide, "우열 비교표 (자사 vs 경쟁)", "자사 열 강조 · ○△✕ 기호로 우열 표기 (%s)" % th["key"])
    x, y, w = 1.1, 1.55, 11.1
    first_w = 4.4
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    tn = nrows + 1
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, ncol, col_widths=col_widths, row_h=[rh] * tn)
    c.name_asset(gf, aid)
    for j, hd in enumerate(cols):
        fill = th["accent"] if j == own_col else th["header"]
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, it in enumerate(items, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), it, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j in range(1, ncol):
            if j == own_col:
                g = GOOD if (i + vi) % 4 else MID
                cfill = CC["gray_050"]
            else:
                g = glyphs[(i + j + vi) % 3]
                cfill = base
            set_cell(tbl.cell(i, j), g[0], size=14, bold=True, color=g[1], fill=cfill)
    return {"name": "우열 비교표 (자사 vs 경쟁 %d사)" % ncomp,
            "tags": ["비교표", "우열비교", "자사강조", "경쟁분석", th["key"]],
            "params": {"family": "versus", "cols": ncol, "rows": nrows, "highlight_col": "자사",
                       "competitors": ncomp, "symbols": ["○", "△", "✕"], "theme": th["key"]},
            "bindings": {"rows": items, "cols": cols},
            "editable": ["rows", "cell_symbols", "highlight_col", "headers"],
            "use": ["경쟁력 비교", "차별성 강조", "제안 우위"]}


WEIGHTS = {4: [30, 30, 20, 20], 5: [30, 20, 20, 15, 15], 6: [25, 20, 20, 15, 10, 10]}


def fam_scorecard(slide, aid, vi):
    th = THEMES[vi % NTH]
    nrows = 4 + (vi % 3)   # 4~6
    items = pick(EVAL_ITEMS, vi, nrows)
    wts = WEIGHTS[nrows]
    rows = []
    conv_sum = 0.0
    for i, (it, wt) in enumerate(zip(items, wts)):
        sc = 85 + ((vi + i) % 11)
        conv = round(wt * sc / 100.0, 1)
        conv_sum += conv
        rows.append((it, "%d%%" % wt, "%d" % sc, "%.1f" % conv))
    title_block(slide, "가중치 스코어카드 (환산점수)", "항목별 가중치 × 점수 = 환산점수, 종합 합산 (%s)" % th["key"])
    x, y, w = 1.6, 1.6, 10.1
    tn = nrows + 2
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, 4, col_widths=[4.3, 1.9, 1.9, 2.0], row_h=[rh] * tn)
    c.name_asset(gf, aid)
    for j, hd in enumerate(["평가 항목", "가중치", "점수(/100)", "환산점수"]):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=th["header"],
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (it, wt, sc, cv) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), it, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        set_cell(tbl.cell(i, 1), wt, size=11, fill=base, color=R("muted_text"))
        set_cell(tbl.cell(i, 2), sc, size=11, fill=base, color=R("body_text"))
        set_cell(tbl.cell(i, 3), cv, size=11, bold=True, fill=CC["gray_050"], color=th["accent"])
    li = tn - 1
    set_cell(tbl.cell(li, 0), "종합 환산점수", size=12, bold=True, color=R("header_text"), fill=th["header"], align=PP_ALIGN.LEFT)
    set_cell(tbl.cell(li, 1), "100%", size=12, bold=True, color=R("header_text"), fill=th["header"])
    set_cell(tbl.cell(li, 2), "—", size=12, bold=True, color=R("header_text"), fill=th["header"])
    set_cell(tbl.cell(li, 3), "%.1f" % conv_sum, size=13, bold=True, color=R("accent_point"), fill=th["header"])
    return {"name": "가중치 스코어카드 (%d항목)" % nrows,
            "tags": ["평가표", "스코어카드", "가중치", "환산점수", th["key"]],
            "params": {"family": "scorecard", "cols": 4, "rows": nrows, "total_row": True,
                       "computed_col": "환산", "theme": th["key"]},
            "bindings": {"rows": items, "cols": ["평가 항목", "가중치", "점수", "환산점수"]},
            "editable": ["rows", "weights", "scores", "converted", "total"],
            "use": ["정량 평가", "가중 채점", "심사 대응"]}


def fam_feature(slide, aid, vi):
    th = THEMES[vi % NTH]
    nopt = 3 + (vi % 2)    # 3~4 옵션
    nrows = 5 + (vi % 3)   # 5~7
    items = pick(EVAL_ITEMS, vi, nrows)
    opts = OPTION_COLS[:nopt]
    cols = ["기능 항목"] + opts
    ncol = len(cols)
    glyph = {"y": ("✓", R("accent_secondary")), "n": ("✗", R("warn")), "p": ("△", R("muted_text"))}
    states = ["y", "p", "n"]
    title_block(slide, "기능 매트릭스 그리드 (제공 여부)", "기능별 옵션 제공 여부 ✓ / △ / ✗ (%s)" % th["key"])
    x, y, w = 1.6, 1.6, 10.1
    first_w = 4.3
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    tn = nrows + 1
    rh = 0.6
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, ncol, col_widths=col_widths, row_h=[rh] * tn)
    c.name_asset(gf, aid)
    for j, hd in enumerate(cols):
        fill = th["accent"] if j == ncol - 1 else th["header"]
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i in range(1, tn):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), items[i - 1], size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j in range(1, ncol):
            # 상위 옵션일수록 제공(✓) 경향
            st = states[max(0, 2 - ((j - 1) + ((i + vi) % 2)))]
            sym, col = glyph[st]
            cfill = CC["gray_050"] if j == ncol - 1 else base
            set_cell(tbl.cell(i, j), sym, size=15, bold=True, color=col, fill=cfill)
    return {"name": "기능 매트릭스 그리드 (%d옵션)" % nopt,
            "tags": ["비교표", "기능매트릭스", "그리드", "제공여부", th["key"]],
            "params": {"family": "featgrid", "cols": ncol, "rows": nrows, "options": nopt,
                       "states": ["y", "p", "n"], "symbols": ["✓", "△", "✗"], "theme": th["key"]},
            "bindings": {"rows": items, "cols": cols},
            "editable": ["rows", "cell_states", "state_symbols", "headers"],
            "use": ["기능 비교", "옵션 제공표", "제안 범위"]}


def fam_swot(slide, aid, vi):
    th = THEMES[vi % NTH]
    ipq = 2 + (vi % 2)   # 사분면당 2~3
    title_block(slide, "SWOT 2x2 분석", "강점·약점·기회·위협 4분면 정리 (%s)" % th["key"])
    quads = [
        {"letter": "S", "label": "강점 (Strength)", "color": th["accent"], "items": SWOT_ITEMS["S"][:ipq]},
        {"letter": "W", "label": "약점 (Weakness)", "color": R("warn"), "items": SWOT_ITEMS["W"][:ipq]},
        {"letter": "O", "label": "기회 (Opportunity)", "color": th["accent2"], "items": SWOT_ITEMS["O"][:ipq]},
        {"letter": "T", "label": "위협 (Threat)", "color": th["header"], "items": SWOT_ITEMS["T"][:ipq]},
    ]
    shapes = []
    x0, x1, y0, y1 = 1.35, 11.95, 1.75, 6.45
    gap = 0.3
    bw = (x1 - x0 - gap) / 2
    bh = (y1 - y0 - gap) / 2
    coords = [(x0, y0), (x0 + bw + gap, y0), (x0, y0 + bh + gap), (x0 + bw + gap, y0 + bh + gap)]
    for (bx, by), q in zip(coords, quads):
        panel = c.add_box(slide, bx, by, bw, bh, fill=CC["gray_050"], line=q["color"], line_w=1.5,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(panel)
        badge = c.add_box(slide, bx + 0.16, by + 0.16, 0.6, 0.6, fill=q["color"], line=None, shape=MSO_SHAPE.OVAL)
        c.set_shape_text(badge, q["letter"], size=20, bold=True, color=R("header_text"), font=c.FONT_H)
        shapes.append(badge)
        lab = c.add_text(slide, bx + 0.86, by + 0.22, bw - 1.0, 0.5, q["label"], size=14, bold=True,
                         color=q["color"], align=PP_ALIGN.LEFT, font=c.FONT_H)
        shapes.append(lab)
        shapes += _bullets(slide, bx + 0.28, by + 0.92, bw - 0.55, q["items"], q["color"], step=0.42, size=10.5)
    c.group_asset(slide, shapes, aid)
    return {"name": "SWOT 2x2 분석 (%d항목)" % ipq,
            "tags": ["매트릭스", "2x2", "SWOT", "사분면", th["key"]],
            "params": {"family": "swot", "quadrants": 4, "items_per_quad": ipq, "layout": "SWOT", "theme": th["key"]},
            "bindings": {"quadrants": ["S", "W", "O", "T"], "axis": ["internal", "external"]},
            "editable": ["quad_letters", "quad_labels", "quad_items", "quad_color"],
            "use": ["SWOT 분석", "전략 진단", "환경 분석"]}


def fam_harvey(slide, aid, vi):
    th = THEMES[vi % NTH]
    ncomp = 3 + (vi % 2)   # 대상 열 3~4
    nrows = 5 + (vi % 3)   # 5~7
    items = pick(EVAL_ITEMS, vi, nrows)
    cols = ["평가 항목"] + COMPETITORS[:ncomp]
    ncol = len(cols)
    HB = [("●", th["accent"]), ("◕", th["accent"]), ("◐", th["accent2"]),
          ("◔", R("muted_text")), ("○", R("muted_text"))]
    title_block(slide, "하비볼 비교표 (충족도 우열)", "●(우수) ◕ ◐ ◔ ○(미흡) 채움 정도로 우열 (%s)" % th["key"])
    x, y, w = 1.1, 1.7, 11.1
    first_w = 4.1
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    tn = nrows + 1
    rh = 0.62
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, ncol, col_widths=col_widths, row_h=[rh] * tn)
    c.name_asset(gf, aid)
    for j, hd in enumerate(cols):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=th["header"],
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, it in enumerate(items, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), it, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j in range(1, ncol):
            gi = ((j - 1) + i + vi) % 5
            if j == 1:  # 자사 열은 상위 편중
                gi = min(gi, 1)
            sym, col = HB[gi]
            set_cell(tbl.cell(i, j), sym, size=17, bold=True, color=col, fill=base)
    return {"name": "하비볼 비교표 (%d사)" % ncomp,
            "tags": ["비교표", "하비볼", "우열비교", "충족도", th["key"]],
            "params": {"family": "harvey", "cols": ncol, "rows": nrows, "competitors": ncomp,
                       "glyphs": ["●", "◕", "◐", "◔", "○"], "theme": th["key"]},
            "bindings": {"rows": items, "cols": cols},
            "editable": ["rows", "cell_glyphs", "glyph_color", "headers"],
            "use": ["충족도 평가", "경쟁력 비교", "제안 우위"]}


def fam_tier(slide, aid, vi):
    th = THEMES[vi % NTH]
    n = 3 + (vi % 2)   # 3~4 등급
    rec = 1            # 추천 인덱스
    tiers = []
    for i in range(n):
        nm, badge = TIER_NAMES[i % len(TIER_NAMES)]
        col = th["accent"] if i == rec else (R("muted_text") if i == 0 else th["accent2"])
        tiers.append({"name": nm, "badge": badge, "price": TIER_PRICES[i % len(TIER_PRICES)],
                      "unit": "만원 / 월" if TIER_PRICES[i % len(TIER_PRICES)] != "맞춤" else "별도 견적",
                      "color": col, "feats": pick(TIER_FEATS, i * 2 + vi, 3 + (i % 2)),
                      "recommended": i == rec})
    title_block(slide, "%d단계 요금/등급 카드" % n, "등급 뱃지 · 대형 가격 · ✓ 기능 · 추천 강조 (%s)" % th["key"])
    shapes = []
    x_start, total_w, gap = 1.1, 11.1, 0.4
    cw = (total_w - gap * (n - 1)) / n
    base_y, base_h = 1.95, 4.55
    for i, t in enumerate(tiers):
        is_rec = t["recommended"]
        cx = x_start + i * (cw + gap)
        cy = base_y - (0.22 if is_rec else 0.0)
        ch = base_h + (0.44 if is_rec else 0.0)
        card = c.add_box(slide, cx, cy, cw, ch, fill=R("row_base"),
                         line=R("accent_point") if is_rec else R("border"),
                         line_w=2.5 if is_rec else 1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(card)
        bw = min(1.9, cw - 0.4)
        badge = c.add_box(slide, cx + (cw - bw) / 2, cy + 0.22, bw, 0.38, fill=t["color"], line=None,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.set_shape_text(badge, t["badge"], size=10, bold=True, color=R("header_text"))
        shapes.append(badge)
        shapes.append(c.add_text(slide, cx + 0.12, cy + 0.72, cw - 0.24, 0.4, t["name"], size=15, bold=True,
                                 color=R("header_fill"), align=PP_ALIGN.CENTER, font=c.FONT_H))
        shapes.append(c.add_text(slide, cx + 0.12, cy + 1.14, cw - 0.24, 0.7, t["price"], size=32, bold=True,
                                 color=t["color"], align=PP_ALIGN.CENTER, font=c.FONT_H))
        shapes.append(c.add_text(slide, cx + 0.12, cy + 1.86, cw - 0.24, 0.28, t["unit"], size=10,
                                 color=R("muted_text"), align=PP_ALIGN.CENTER))
        shapes.append(c.add_box(slide, cx + 0.35, cy + 2.25, cw - 0.7, 0.02, fill=R("border"), line=None))
        fy = cy + 2.42
        for feat in t["feats"]:
            shapes.append(c.add_text(slide, cx + 0.28, fy, 0.3, 0.34, "✓", size=12, bold=True,
                                     color=t["color"], align=PP_ALIGN.CENTER))
            shapes.append(c.add_text(slide, cx + 0.56, fy, cw - 0.7, 0.34, feat, size=10, color=R("body_text"),
                                     align=PP_ALIGN.LEFT))
            fy += 0.4
        if is_rec:
            rb = min(1.5, cw - 0.3)
            ribbon = c.add_box(slide, cx + cw - rb - 0.1, cy - 0.02, rb, 0.4, fill=R("warn"), line=None,
                               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            c.set_shape_text(ribbon, "★ 추천", size=10, bold=True, color=R("header_text"))
            shapes.append(ribbon)
    c.group_asset(slide, shapes, aid)
    return {"name": "%d단계 요금/등급 카드" % n,
            "tags": ["비교표", "요금카드", "등급카드", "플랜비교", th["key"]],
            "params": {"family": "tiercard", "columns": n, "recommended_idx": rec, "price_emphasis": True, "theme": th["key"]},
            "bindings": {"rows": [t["name"] for t in tiers], "cols": [t["badge"] for t in tiers]},
            "editable": ["tier_names", "badges", "prices", "units", "feats", "tier_color", "recommended_badge"],
            "use": ["요금제 비교", "등급 제안", "가격 정책"]}


def fam_proscons(slide, aid, vi):
    th = THEMES[vi % NTH]
    nit = 3 + (vi % 2)   # 3~4
    pros = pick(PROS, vi, nit)
    cons = pick(CONS, vi, nit)
    title_block(slide, "장단점 2열 (Pros / Cons)", "장점·단점을 색으로 구분한 2열 대비 (%s)" % th["key"])
    shapes = []
    y0, ch, cw, gap = 1.85, 4.5, 5.35, 0.4
    x_p = 1.1
    x_c = x_p + cw + gap
    cfg = [(x_p, "장점 (Pros)", "＋", th["accent"], CC["gray_050"], pros),
           (x_c, "단점 (Cons)", "－", R("warn"), CC["beige_100"], cons)]
    for cx, head, mark, accent, bg, its in cfg:
        card = c.add_box(slide, cx, y0, cw, ch, fill=bg, line=accent, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(card)
        hd = c.add_box(slide, cx, y0, cw, 0.8, fill=accent, line=None, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        c.set_shape_text(hd, head, size=15, bold=True, color=R("header_text"))
        shapes.append(hd)
        fy = y0 + 1.05
        for it in its:
            badge = c.add_box(slide, cx + 0.3, fy + 0.04, 0.34, 0.34, fill=accent, line=None, shape=MSO_SHAPE.OVAL)
            c.set_shape_text(badge, mark, size=13, bold=True, color=R("header_text"))
            tx = c.add_text(slide, cx + 0.78, fy - 0.02, cw - 1.05, 0.5, it, size=11.5, color=R("body_text"),
                            align=PP_ALIGN.LEFT)
            shapes += [badge, tx]
            fy += 0.78
    c.group_asset(slide, shapes, aid)
    return {"name": "장단점 2열 (%d항목)" % nit,
            "tags": ["비교표", "장단점", "proscons", "2열대비", th["key"]],
            "params": {"family": "proscons", "columns": 2, "items": nit, "theme": th["key"]},
            "bindings": {"quadrants": ["pros", "cons"], "axis": ["state"]},
            "editable": ["pros_items", "cons_items", "pros_color", "cons_color"],
            "use": ["장단점 정리", "의사결정 근거", "리스크 제시"]}


def fam_ranking(slide, aid, vi):
    th = THEMES[vi % NTH]
    nrows = 4 + (vi % 4)   # 4~7
    rank_fill = {1: th["accent"], 2: R("accent_point"), 3: th["accent2"]}
    tgts = pick(RANK_TARGETS, vi, nrows)
    title_block(slide, "순위 비교표 (종합 순위)", "종합 점수 기준 순위 뱃지 · 상위 3위 강조 (%s)" % th["key"])
    x, y, w = 1.6, 1.7, 10.1
    tn = nrows + 1
    rh = 0.62
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, 4, col_widths=[1.4, 4.5, 2.1, 2.1], row_h=[rh] * tn)
    c.name_asset(gf, aid)
    for j, hd in enumerate(["순위", "평가 대상", "종합 점수", "비고"]):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=th["header"],
                 align=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER)
    for i, tgt in enumerate(tgts, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        score = 90 - (i - 1) * 4 - (vi % 3)
        if i in rank_fill:
            set_cell(tbl.cell(i, 0), "%d위" % i, size=13, bold=True, color=R("header_text"), fill=rank_fill[i])
        else:
            set_cell(tbl.cell(i, 0), "%d위" % i, size=11, bold=True, color=R("muted_text"), fill=base)
        set_cell(tbl.cell(i, 1), tgt, size=11, bold=(i <= 3), fill=base, align=PP_ALIGN.LEFT)
        set_cell(tbl.cell(i, 2), "%.1f" % score, size=11, bold=True, fill=base, color=th["accent"])
        set_cell(tbl.cell(i, 3), RANK_NOTE.get(i, "탈락"), size=10, fill=base, color=R("muted_text"))
    return {"name": "순위 비교표 (%d위)" % nrows,
            "tags": ["비교표", "순위표", "랭킹", "순위뱃지", th["key"]],
            "params": {"family": "ranking", "cols": 4, "rows": nrows, "rank_badge": True,
                       "top_n_highlight": 3, "theme": th["key"]},
            "bindings": {"rows": ["%d위" % i for i in range(1, nrows + 1)],
                         "cols": ["순위", "평가 대상", "종합 점수", "비고"]},
            "editable": ["rows", "rank", "target", "score", "note", "rank_fill"],
            "use": ["순위 제시", "종합 평가", "심사 결과"]}


def fam_decision(slide, aid, vi):
    th = THEMES[vi % NTH]
    ncrit = 3 + (vi % 3)   # 성능 기준 3~5
    nopt = 3 + (vi % 2)    # 대안 3~4
    crits = pick(CRITERIA, vi, ncrit)
    criteria = crits + ["종합점수"]
    cols = ["대안"] + criteria
    ncol = len(cols)
    opts = ALTS[:nopt]
    rows = []
    for oi, nm in enumerate(opts):
        scores = []
        tot = 0
        for ci in range(ncrit):
            sv = 7 + ((vi + oi + ci) % 3)
            tot += sv
            scores.append("%d" % sv)
        scores.append("%d" % tot)
        rows.append((nm, scores, tot))
    best_idx = max(range(len(rows)), key=lambda k: rows[k][2])
    title_block(slide, "의사결정 매트릭스 (대안 평가)", "기준별 점수·종합점수로 최적 대안 선정 (%s)" % th["key"])
    x, y, w = 1.1, 1.7, 11.1
    first_w = 3.0
    rest_w = (w - first_w) / (ncol - 1)
    col_widths = [first_w] + [rest_w] * (ncol - 1)
    tn = len(rows) + 1
    rh = 0.66
    gf, tbl = add_table(slide, x, y, w, rh * tn, tn, ncol, col_widths=col_widths, row_h=[rh] * tn)
    c.name_asset(gf, aid)
    last = ncol - 1
    for j, hd in enumerate(cols):
        fill = th["accent"] if j == last else th["header"]
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"), fill=fill,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (nm, scores, _t) in enumerate(rows, start=1):
        is_best = (i - 1) == best_idx
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), nm, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
        for j, sv in enumerate(scores, start=1):
            if j == last:
                cfill = th["accent"] if is_best else CC["gray_050"]
                col = R("header_text") if is_best else th["accent"]
                set_cell(tbl.cell(i, j), sv, size=13, bold=True, color=col, fill=cfill)
            else:
                set_cell(tbl.cell(i, j), sv, size=11, fill=base, color=R("body_text"))
    return {"name": "의사결정 매트릭스 (%d기준×%d대안)" % (ncrit, nopt),
            "tags": ["매트릭스", "의사결정", "대안평가", "점수표", th["key"]],
            "params": {"family": "decision", "criteria": ncrit, "options": nopt,
                       "highlight_col": "종합점수", "best_row": True, "theme": th["key"]},
            "bindings": {"rows": opts, "cols": cols},
            "editable": ["options", "criteria", "scores", "best_idx", "highlight_col"],
            "use": ["대안 비교", "의사결정 근거", "선정 논리"]}


def fam_beforeafter(slide, aid, vi):
    th = THEMES[vi % NTH]
    nit = 3 + (vi % 3)   # 3~5
    before = pick(BA_BEFORE, vi, nit)
    after = pick(BA_AFTER, vi, nit)
    title_block(slide, "Before / After 대비 카드", "개선 전·후를 2열 카드로 대비 (%s)" % th["key"])
    shapes = []
    y0, ch, cw, gap_mid = 1.65, 4.6, 5.35, 1.1
    x_b = 1.1
    x_a = x_b + cw + gap_mid
    cb = c.add_box(slide, x_b, y0, cw, ch, fill=CC["gray_100"], line=R("border"), line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    hb = c.add_box(slide, x_b, y0, cw, 0.75, fill=R("muted_text"), line=None, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    c.set_shape_text(hb, "Before (현행)", size=15, bold=True, color=R("header_text"))
    shapes += [cb, hb]
    shapes += _bullets(slide, x_b + 0.35, y0 + 1.0, cw - 0.7, before, R("muted_text"), step=0.62, size=11)
    ca = c.add_box(slide, x_a, y0, cw, ch, fill=CC["gray_050"], line=th["accent"], line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ha = c.add_box(slide, x_a, y0, cw, 0.75, fill=th["accent"], line=None, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    c.set_shape_text(ha, "After (개선)", size=15, bold=True, color=R("header_text"))
    shapes += [ca, ha]
    shapes += _bullets(slide, x_a + 0.35, y0 + 1.0, cw - 0.7, after, th["accent2"], step=0.62, size=11)
    ar = c.add_box(slide, x_b + cw + 0.15, y0 + ch / 2 - 0.4, gap_mid - 0.3, 0.8, fill=R("accent_point"),
                   line=None, shape=MSO_SHAPE.RIGHT_ARROW)
    c.set_shape_text(ar, "개선", size=11, bold=True, color=R("header_text"))
    shapes.append(ar)
    c.group_asset(slide, shapes, aid)
    return {"name": "Before / After 대비 카드 (%d항목)" % nit,
            "tags": ["비교표", "전후비교", "대비카드", "개선효과", th["key"]],
            "params": {"family": "beforeafter", "columns": 2, "items": nit, "theme": th["key"]},
            "bindings": {"quadrants": ["before", "after"], "axis": ["state"]},
            "editable": ["before_head", "before_items", "after_head", "after_items", "card_color"],
            "use": ["개선 효과 제시", "전후 대비", "성과 강조"]}


def fam_textprop(slide, aid, vi):
    th = THEMES[vi % NTH]
    nsec = 2 + (vi % 2)   # 섹션 2~3
    nb = 3 if nsec == 2 else 2
    heading = PROP_TITLES[vi % len(PROP_TITLES)]
    shapes = []
    px, pw = 1.1, 11.1
    py, ph = 1.5, 5.5
    panel = c.add_box(slide, px, py, pw, ph, fill=CC["gray_050"], line=R("border"), line_w=1.0,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shapes.append(panel)
    hbar = c.add_box(slide, px, py, pw, 0.7, fill=th["header"], line=None, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    c.set_shape_text(hbar, heading, size=16, bold=True, color=R("header_text"), font=c.FONT_H)
    shapes.append(hbar)
    cy = py + 0.9
    sh = (ph - 0.9) / nsec
    for si in range(nsec):
        sub = PROP_HEADS[(vi + si) % len(PROP_HEADS)]
        tab = c.add_box(slide, px + 0.35, cy + 0.05, 0.12, 0.34, fill=th["accent"], line=None)
        subt = c.add_text(slide, px + 0.6, cy, pw - 1.0, 0.4, sub, size=13, bold=True,
                          color=th["accent"], align=PP_ALIGN.LEFT, font=c.FONT_H)
        shapes += [tab, subt]
        sents = pick(PROP_SENTS, vi + si * 3, nb)
        shapes += _bullets(slide, px + 0.6, cy + 0.5, pw - 1.1, sents, th["accent2"], step=0.42, size=11)
        cy += sh
    c.group_asset(slide, shapes, aid)
    return {"name": "텍스트 제안서 블록 (%d섹션)" % nsec,
            "tags": ["제안서", "텍스트블록", "불릿", "요약패널", th["key"]],
            "params": {"family": "textprop", "sections": nsec, "bullets": nb, "theme": th["key"]},
            "bindings": {"rows": PROP_HEADS[:nsec], "cols": ["heading"]},
            "editable": ["heading", "section_subheads", "section_bullets", "accent_color"],
            "use": ["제안 요약", "본문 채움", "설명 슬라이드"]}


# ---------------------------------------------------------------------------
# 드라이버 — 패밀리별 파일 분산(≤25슬라이드), ID 순차 부여
# ---------------------------------------------------------------------------
FAMILIES = [
    ("matrix2x2",   14, fam_matrix),
    ("evaltable",   14, fam_eval),
    ("versus",      16, fam_versus),
    ("scorecard",   14, fam_scorecard),
    ("featgrid",    16, fam_feature),
    ("swot",        12, fam_swot),
    ("harvey",      16, fam_harvey),
    ("tiercard",    14, fam_tier),
    ("proscons",    12, fam_proscons),
    ("ranking",     14, fam_ranking),
    ("decision",    14, fam_decision),
    ("beforeafter", 12, fam_beforeafter),
    ("textprop",    14, fam_textprop),
]

entries = []
files = []
gid = 19
CHUNK = 25

for fam_name, count, fn in FAMILIES:
    made = 0
    file_idx = 1
    while made < count:
        chunk = min(CHUNK, count - made)
        prs = c.new_deck()
        rel = "decks/04_compare/CMP_bulk_%s_v%d.pptx" % (fam_name, file_idx)
        for k in range(chunk):
            vi = made + k
            aid = "CMP-%03d" % gid
            s = c.blank_slide(prs)
            c.id_caption(s, aid)
            meta = fn(s, aid, vi)
            entries.append(c.entry(aid, "CMP", meta["name"], rel, k + 1,
                                   meta["tags"], meta["params"], meta["bindings"], meta["editable"],
                                   recommended_use=meta["use"]))
            gid += 1
        c.save_deck(prs, rel)
        files.append(rel)
        made += chunk
        file_idx += 1

frag = c.write_fragment("CMP_bulk", entries)

print("ENTRIES", len(entries))
print("ID_RANGE", entries[0]["id"], "-", entries[-1]["id"])
print("FILES", len(files))
for f in files:
    print("FILE", f)
print("FRAG", frag)
