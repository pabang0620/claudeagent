# -*- coding: utf-8 -*-
"""
06_org (ORG) v2 변형팩 — 조직/체계 6종 (ORG-009 ~ ORG-014)
v1(위계/네트워크)과 다른 축: RACI · 기능별 매트릭스조직 · 파트너 생태계 ·
위원회 구조 · 컨소시엄 역할분담표 · 에스컬레이션 체계.

규칙:
- 색은 c.role / c.C 만. 한글은 c 헬퍼(set_kfont 경유).
- 다중도형 에셋 = c.group_asset(slide, shapes, 'ORG-0NN')  (그룹명이 앵커).
- 단일 표 에셋 = graphicFrame 에 c.name_asset 만 (그룹 미사용).
- 페이지이미지·SmartArt 금지. 1슬라이드 1에셋. id_caption 은 그룹 제외.

파일 분산:
  ORG_matrix_v2.pptx     : RACI(009) / 기능별매트릭스(010) / 컨소시엄분담표(013)
  ORG_ecosystem_v2.pptx  : 파트너생태계(011) / 위원회구조(012) / 에스컬레이션(014)
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

R = c.role
CC = c.C
RR = MSO_SHAPE.ROUNDED_RECTANGLE
RECT = MSO_SHAPE.RECTANGLE
OVAL = MSO_SHAPE.OVAL
UPARROW = MSO_SHAPE.UP_ARROW

F1 = "decks/06_org/ORG_matrix_v2.pptx"
F2 = "decks/06_org/ORG_ecosystem_v2.pptx"

entries = []
def E(asset_id, name, file_rel, slide_idx, tags, params, bindings, editable, rec):
    entries.append(c.entry(asset_id, "ORG", name, file_rel, slide_idx,
                           tags, params, bindings, editable, recommended_use=rec))

# ─────────────────────────────────────────────────────────────
# 저수준 헬퍼 (좌표 inch)
# ─────────────────────────────────────────────────────────────
def cx(sp):  return sp.left / 914400.0 + sp.width / 914400.0 / 2
def cyf(sp): return sp.top / 914400.0 + sp.height / 914400.0 / 2
def topf(sp):return sp.top / 914400.0
def botf(sp):return sp.top / 914400.0 + sp.height / 914400.0
def rightf(sp): return sp.left / 914400.0 + sp.width / 914400.0
def leftf(sp):  return sp.left / 914400.0

def node(slide, x, y, w, h, text, fill, txt=None, size=12, bold=True,
         line=None, shape=RR, radius=0.10):
    sp = c.add_box(slide, x, y, w, h, fill=fill, line=line, line_w=1.0, shape=shape)
    if shape == RR:
        sp.adjustments[0] = radius
    c.set_shape_text(sp, text, size=size, bold=bold,
                     color=txt or R("header_text"), align=PP_ALIGN.CENTER, font=c.FONT_H)
    return sp

def conn(slide, x1, y1, x2, y2, color=None, w=1.5, dashed=False, head=False, tail=False):
    cn = c.connector(slide, x1, y1, x2, y2, color=color or R("muted_text"), w=w)
    ln = cn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn

def tree(slide, parent, children, color=None, bus_gap=0.42, w=1.6):
    """부모 하단 → 수평 버스 → 각 자식 상단. 커넥터 리스트 반환."""
    col = color or R("muted_text")
    out = []
    py = botf(parent); busy = py + bus_gap; pcx = cx(parent)
    xs = [cx(ch) for ch in children]
    out.append(conn(slide, pcx, py, pcx, busy, color=col, w=w))
    out.append(conn(slide, min(xs + [pcx]), busy, max(xs + [pcx]), busy, color=col, w=w))
    for ch in children:
        out.append(conn(slide, cx(ch), busy, cx(ch), topf(ch), color=col, w=w))
    return out

def slide_title(slide, text, sub=None):
    """슬라이드 상단 제목(에셋 그룹 미포함 — 라이브러리 식별용 데코)."""
    c.add_text(slide, 0.55, 0.5, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.0, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)

# ── 표 헬퍼 ──
def add_table(slide, x, y, w, h, nrows, ncols, col_widths=None, row_h=None):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    return gf, tbl

def cell(cl, text, size=11, bold=False, color=None, fill=None,
         align=PP_ALIGN.CENTER, font=None):
    cl.fill.solid()
    cl.fill.fore_color.rgb = fill if fill is not None else R("row_base")
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_left = cl.margin_right = Pt(5)
    cl.margin_top = cl.margin_bottom = Pt(3)
    tf = cl.text_frame; tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))

# RACI 마커 색
MARK = {
    "A": (R("warn"), True),
    "R": (R("accent_primary"), True),
    "C": (R("accent_secondary"), False),
    "I": (R("muted_text"), False),
    "": (R("muted_text"), False),
}


# ═════════════════════════════════════════════════════════════
# FILE 1 — ORG_matrix_v2
# ═════════════════════════════════════════════════════════════
prs = c.new_deck()

# ── ORG-009 : RACI 책임 배분표 (업무 × 역할) ────────────────
s = c.blank_slide(prs); c.id_caption(s, "ORG-009 · RACI 책임 배분표")
slide_title(s, "RACI 책임 배분 매트릭스", "핵심 업무별 역할 책임(R/A/C/I) 정의 — 제안 수행조직")
roles = ["총괄PM", "기획팀", "개발팀", "현지화팀", "품질(QA)"]
raci_rows = [
    ("사업 착수·킥오프",   ["A", "R", "C", "I", "I"]),
    ("요구사항 정의",       ["A", "R", "C", "C", "I"]),
    ("플랫폼 설계·개발",    ["A", "C", "R", "I", "C"]),
    ("콘텐츠 현지화",       ["A", "I", "C", "R", "C"]),
    ("통합 QA 테스트",      ["A", "I", "C", "C", "R"]),
    ("발주처 보고·검수",    ["A", "R", "I", "I", "I"]),
    ("안정화·운영 이관",    ["A", "I", "R", "C", "R"]),
]
ncol = 1 + len(roles)
nrow = 1 + len(raci_rows) + 1  # header + 업무 + legend
cw = [2.6] + [1.85] * len(roles)
rh = [0.5] * nrow
gf, tbl = add_table(s, 0.7, 1.55, sum(cw), sum(rh), nrow, ncol, cw, rh)
c.name_asset(gf, "ORG-009")
cell(tbl.cell(0, 0), "핵심 업무", size=12, bold=True, color=R("header_text"),
     fill=R("header_fill"), align=PP_ALIGN.LEFT)
for j, rn in enumerate(roles):
    cell(tbl.cell(0, j + 1), rn, size=12, bold=True, color=R("header_text"), fill=R("header_fill"))
for i, (act, marks) in enumerate(raci_rows):
    base = R("row_stripe") if i % 2 else R("row_base")
    cell(tbl.cell(i + 1, 0), act, size=11, bold=True, fill=base, align=PP_ALIGN.LEFT)
    for j, m in enumerate(marks):
        col, bd = MARK[m]
        cell(tbl.cell(i + 1, j + 1), m, size=13, bold=bd, color=col, fill=base)
li = nrow - 1
tbl.cell(li, 0).merge(tbl.cell(li, ncol - 1))
cell(tbl.cell(li, 0), "R 실무 담당(Responsible)   ·   A 최종 책임(Accountable)   ·   C 자문(Consulted)   ·   I 결과 공유(Informed)",
     size=10, bold=True, color=R("muted_text"), fill=R("panel_bg"))
E("ORG-009", "RACI 책임 배분표 (업무×역할, R/A/C/I)", F1, 1,
  ["RACI", "책임배분", "역할분담", "매트릭스", "표"],
  {"rows": len(raci_rows), "roles": len(roles), "codes": ["R", "A", "C", "I"]},
  {"roles": roles,
   "activities": [{"task": a, "assign": dict(zip(roles, m))} for a, m in raci_rows],
   "legend": {"R": "실무 담당", "A": "최종 책임", "C": "자문", "I": "결과 공유"}},
  ["cell-text", "add-row", "add-col", "marker", "color"],
  ["RACI", "책임배분(R&R)", "역할 책임 정의"])

# ── ORG-010 : 기능별 매트릭스 조직 (기능축 × 사업축) ─────────
s = c.blank_slide(prs); c.id_caption(s, "ORG-010 · 기능별 매트릭스 조직")
slide_title(s, "기능별 매트릭스 조직", "기능(수직) × 사업(수평) 이중 보고 체계 — 매트릭스 조직")
shp_conn = []; shp_box = []; shp_dot = []
proj = [("사업 A\n해외 A권역", 6.2), ("사업 B\n해외 B권역", 8.5), ("사업 C\n신흥시장", 10.8)]
func = [("기획본부", 3.5), ("개발본부", 4.25), ("현지화본부", 5.0), ("품질본부", 5.75)]
pfill = [CC["navy_600"], R("accent_primary"), R("accent_point")]
ffill = [R("accent_secondary"), CC["navy_800"], R("sub_header"), CC["navy_600"]]
# 총괄 PMO
pm = node(s, 5.4, 1.2, 2.6, 0.7, "총괄 PMO\n(사업 총괄)", R("header_fill"), size=13)
shp_box.append(pm)
# 사업축(수평) 헤더
pboxes = []
for (lab, px), fc in zip(proj, pfill):
    b = node(s, px - 0.95, 2.35, 1.9, 0.62, lab, fc, size=11)
    pboxes.append(b); shp_box.append(b)
# 기능축(수직) 헤더 (좌측)
fboxes = []
for (lab, fy), fc in zip(func, ffill):
    b = node(s, 0.6, fy - 0.28, 2.0, 0.56, lab, fc, size=11.5)
    fboxes.append(b); shp_box.append(b)
# 솔리드: PMO → 사업 헤더 (트리)
shp_conn += tree(s, pm, pboxes, color=R("muted_text"), bus_gap=0.28, w=1.6)
# 솔리드: PMO → 좌측 기능 스파인
spine_x = 2.6
shp_conn.append(conn(s, leftf(pm), cyf(pm), spine_x, cyf(pm), color=R("muted_text"), w=1.6))
shp_conn.append(conn(s, spine_x, cyf(pm), spine_x, func[-1][1], color=R("muted_text"), w=1.6))
# 점선 그리드: 사업 수직선(아래로) + 기능 수평선(오른쪽으로)
for (lab, px), fc in zip(proj, pfill):
    shp_conn.append(conn(s, px, 2.97, px, 6.05, color=R("border"), w=1.25, dashed=True))
for (lab, fy), fc in zip(func, ffill):
    shp_conn.append(conn(s, spine_x, fy, 11.25, fy, color=R("border"), w=1.25, dashed=True))
# 교차 담당자 노드(작은 원)
for (plab, px) in proj:
    for (flab, fy) in func:
        d = c.add_box(s, px - 0.16, fy - 0.16, 0.32, 0.32, fill=R("accent_primary"),
                      line=CC["white"], line_w=1.0, shape=OVAL)
        shp_dot.append(d)
c.group_asset(s, shp_conn + shp_box + shp_dot, "ORG-010")
E("ORG-010", "기능별 매트릭스 조직 (기능축×사업축)", F1, 2,
  ["조직도", "매트릭스", "이중보고", "기능별", "사업별"],
  {"functions": len(func), "projects": len(proj), "orient": "matrix-grid"},
  {"pmo": "총괄 PMO",
   "functions": [f[0] for f in func],
   "projects": [p[0].split("\n")[0] for p in proj],
   "cells": len(func) * len(proj)},
  ["node-text", "add-function", "add-project", "color", "connector"],
  ["매트릭스 조직", "이중 보고 체계", "수행조직"])

# ── ORG-013 : 컨소시엄 역할분담표 (주관/참여 × 업무) ────────
s = c.blank_slide(prs); c.id_caption(s, "ORG-013 · 컨소시엄 역할분담표")
slide_title(s, "컨소시엄 역할분담표", "주관사·참여사 × 업무분장(WBS)별 담당·분담율")
wps = ["WP1\n플랫폼 개발", "WP2\n현지화", "WP3\n마케팅·배급", "WP4\n운영·유지보수"]
comp_rows = [
    ("주관사 ㈜가",   ["◎", "●", "–", "◎"], "40%"),
    ("참여사 ㈜나",   ["●", "◎", "–", "●"], "25%"),
    ("참여사 ㈜다",   ["–", "●", "◎", "–"], "20%"),
    ("참여사 ㈜라",   ["–", "–", "●", "●"], "15%"),
]
MK2 = {"◎": (R("accent_primary"), True), "●": (R("accent_secondary"), True), "–": (R("muted_text"), False)}
nc = 1 + len(wps) + 1
nr = 1 + len(comp_rows) + 1
cw2 = [2.0] + [2.0] * len(wps) + [1.4]
rh2 = [0.7] + [0.62] * len(comp_rows) + [0.55]
gf2, t2 = add_table(s, 0.9, 1.6, sum(cw2), sum(rh2), nr, nc, cw2, rh2)
c.name_asset(gf2, "ORG-013")
cell(t2.cell(0, 0), "구분", size=12, bold=True, color=R("header_text"),
     fill=R("header_fill"), align=PP_ALIGN.LEFT)
for j, wp in enumerate(wps):
    cell(t2.cell(0, j + 1), wp, size=11, bold=True, color=R("header_text"), fill=R("header_fill"))
cell(t2.cell(0, nc - 1), "분담율", size=12, bold=True, color=R("header_text"), fill=R("header_fill"))
for i, (nm, marks, share) in enumerate(comp_rows):
    base = R("row_stripe") if i % 2 else R("row_base")
    lead = i == 0
    cell(t2.cell(i + 1, 0), nm, size=11, bold=True,
         color=R("header_fill") if lead else R("body_text"), fill=base, align=PP_ALIGN.LEFT)
    for j, m in enumerate(marks):
        col, bd = MK2[m]
        cell(t2.cell(i + 1, j + 1), m, size=14, bold=bd, color=col, fill=base)
    cell(t2.cell(i + 1, nc - 1), share, size=12, bold=True, color=R("accent_primary"), fill=base)
si = nr - 1
t2.cell(si, 1).merge(t2.cell(si, nc - 2))
cell(t2.cell(si, 0), "합계", size=12, bold=True, color=R("header_text"), fill=CC["navy_600"])
cell(t2.cell(si, 1), "◎ 주관 담당    ●  공동 참여    –  미참여", size=10.5, bold=True,
     color=R("muted_text"), fill=R("panel_bg"))
cell(t2.cell(si, nc - 1), "100%", size=12, bold=True, color=R("header_text"), fill=CC["navy_600"])
E("ORG-013", "컨소시엄 역할분담표 (주관/참여사×업무)", F1, 3,
  ["컨소시엄", "역할분담", "업무분장", "WBS", "표"],
  {"companies": len(comp_rows), "work_packages": len(wps), "show_share": True},
  {"work_packages": [w.replace("\n", " ") for w in wps],
   "members": [{"name": n, "assign": dict(zip([w.split("\n")[0] for w in wps], m)), "share": sh}
               for n, m, sh in comp_rows],
   "legend": {"◎": "주관 담당", "●": "공동 참여", "–": "미참여"}},
  ["cell-text", "add-row", "add-col", "marker", "share", "color"],
  ["컨소시엄", "역할분담", "업무분장"])

c.save_deck(prs, F1)

# ═════════════════════════════════════════════════════════════
# FILE 2 — ORG_ecosystem_v2
# ═════════════════════════════════════════════════════════════
prs = c.new_deck()

# ── ORG-011 : 파트너 생태계 (중앙 + 계층별 파트너군) ────────
s = c.blank_slide(prs); c.id_caption(s, "ORG-011 · 파트너 생태계")
slide_title(s, "파트너 생태계 지도", "중앙 플랫폼 + 계층별(1~3계층) 파트너군 협력 구조")
e_conn = []; e_panel = []; e_band = []; e_chip = []
tiers = [
    # (헤더, 밴드색, 패널좌표 x,y,w,h, chips[])
    ("1계층 · 핵심 파트너", R("accent_primary"), (4.85, 0.95, 3.6, 1.55),
     ["글로벌 OTT", "현지 메이저 배급사"]),
    ("2계층 · 전략 파트너", R("accent_secondary"), (0.6, 5.05, 4.15, 1.9),
     ["제작 스튜디오", "마케팅 대행사", "전문 번역사"]),
    ("3계층 · 협력 파트너", R("sub_header"), (8.6, 5.05, 4.15, 1.9),
     ["법무·자문", "결제 PG사", "데이터 분석"]),
]
# 허브
hub = node(s, 5.47, 3.15, 2.4, 1.4, "K-콘텐츠\n통합 플랫폼", R("header_fill"), size=14, shape=OVAL)
hub_cx, hub_cy = cx(hub), cyf(hub)
panel_centers = []
for header, band_c, (px, py, pw, ph), chips in tiers:
    panel = c.add_box(s, px, py, pw, ph, fill=R("panel_bg"), line=R("border"), line_w=1.0, shape=RR)
    panel.adjustments[0] = 0.06
    e_panel.append(panel)
    band = c.add_box(s, px, py, pw, 0.42, fill=band_c, line=None, shape=RECT)
    c.set_shape_text(band, header, size=11.5, bold=True, color=R("header_text"),
                     align=PP_ALIGN.CENTER, font=c.FONT_H)
    e_band.append(band)
    n = len(chips); chw = (pw - 0.3 - 0.15 * (n - 1)) / n
    for k, ch in enumerate(chips):
        chx = px + 0.15 + k * (chw + 0.15)
        chip = node(s, chx, py + 0.58, chw, ph - 0.75, ch, CC["white"],
                    txt=R("body_text"), size=10.5, bold=True, line=band_c, shape=RR, radius=0.12)
        e_chip.append(chip)
    panel_centers.append((px + pw / 2, py + ph / 2))
for pcx, pcy in panel_centers:
    e_conn.append(conn(s, hub_cx, hub_cy, pcx, pcy, color=R("accent_primary"), w=1.75))
# z: 커넥터 → 패널 → 밴드 → 칩 → 허브(최상단)
c.group_asset(s, e_conn + e_panel + e_band + e_chip + [hub], "ORG-011")
E("ORG-011", "파트너 생태계 (중앙+계층별 파트너군)", F2, 1,
  ["생태계", "파트너", "네트워크", "계층", "협업"],
  {"hub": 1, "tiers": 3, "orient": "ecosystem"},
  {"hub": "K-콘텐츠 통합 플랫폼",
   "tiers": [{"tier": h, "partners": ch} for h, _, _, ch in tiers]},
  ["node-text", "add-partner", "add-tier", "color", "connector"],
  ["파트너 생태계", "협업 네트워크", "협력 구조"])

# ── ORG-012 : 위원회 구조 (위원장 - 분과위 3) ───────────────
s = c.blank_slide(prs); c.id_caption(s, "ORG-012 · 위원회 구조")
slide_title(s, "추진위원회 구조", "위원장 · 사무국(간사) · 분과위원회 3 운영 체계")
w_conn = []; w_box = []; w_panel = []
chair = node(s, 5.4, 1.15, 2.55, 0.72, "추진위원장", R("header_fill"), size=14)
w_box.append(chair)
# 사무국(간사) — 점선 우측
sec = node(s, 9.7, 1.2, 2.1, 0.62, "사무국(간사)", R("muted_text"), size=11.5)
w_box.append(sec)
w_conn.append(conn(s, rightf(chair), cyf(chair), leftf(sec), cyf(sec),
                   color=R("muted_text"), w=1.4, dashed=True))
subs = [
    ("기술 분과위", R("accent_primary"), 2.75, ["시스템·인프라 심의", "보안·품질 기준", "기술 자문"]),
    ("콘텐츠 분과위", R("accent_secondary"), 6.667, ["콘텐츠 기획 심의", "현지화 방향", "저작권 검토"]),
    ("운영 분과위", R("sub_header"), 10.58, ["사업 운영 점검", "예산·일정 조정", "성과 평가"]),
]
sboxes = []
for lab, fc, scx, bullets in subs:
    b = node(s, scx - 1.6, 2.95, 3.2, 0.68, lab, fc, size=12.5)
    sboxes.append(b); w_box.append(b)
    panel = c.add_box(s, scx - 1.6, 3.8, 3.2, 1.85, fill=R("panel_bg"),
                      line=R("border"), line_w=0.9, shape=RR)
    panel.adjustments[0] = 0.05
    txt = "\n".join("· " + t for t in bullets)
    c.set_shape_text(panel, txt, size=11, bold=False, color=R("body_text"),
                     align=PP_ALIGN.LEFT, font=c.FONT_B)
    panel.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for p in panel.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
    w_panel.append(panel)
w_conn += tree(s, chair, sboxes, color=R("muted_text"), bus_gap=0.5, w=1.6)
# 분과 → 패널 짧은 연결
for b, panel in zip(sboxes, w_panel):
    w_conn.append(conn(s, cx(b), botf(b), cx(b), topf(panel), color=R("border"), w=1.2))
c.group_asset(s, w_conn + w_box + w_panel, "ORG-012")
E("ORG-012", "위원회 구조 (위원장-분과위 3)", F2, 2,
  ["위원회", "거버넌스", "분과위", "심의", "조직도"],
  {"chair": 1, "secretariat": 1, "subcommittees": 3},
  {"chair": "추진위원장", "secretariat": "사무국(간사)",
   "subcommittees": [{"name": n, "duties": b} for n, _, _, b in subs]},
  ["node-text", "add-subcommittee", "color", "connector"],
  ["위원회 구조", "거버넌스", "심의 체계"])

# ── ORG-014 : 에스컬레이션 체계 (단계별 대응 상향) ──────────
s = c.blank_slide(prs); c.id_caption(s, "ORG-014 · 에스컬레이션 체계")
slide_title(s, "이슈 에스컬레이션 체계", "미해결 이슈의 단계별 대응 주체 상향(escalation) 흐름")
x_conn = []; x_box = []; x_ann = []
levels = [
    # (라벨, 색, y, SLA 주석)
    ("1단계 · 담당자 대응", R("accent_secondary"), 5.9, "대응시간 4시간 이내 · 현장 1차 조치"),
    ("2단계 · 팀장 대응",   R("accent_primary"), 4.7, "4시간 초과 시 · 자원 재배치 결정"),
    ("3단계 · 총괄PM 대응", R("sub_header"), 3.5, "24시간 초과 시 · 범위·일정 조정"),
    ("4단계 · 운영위원회 대응", R("warn"), 2.3, "중대 이슈 · 계약·정책 의사결정"),
]
lb_x, lb_w, lb_h = 3.3, 5.3, 0.82
lboxes = []
for lab, fc, ly, sla in levels:
    b = node(s, lb_x, ly, lb_w, lb_h, lab, fc, size=13)
    lboxes.append(b); x_box.append(b)
    ann = c.add_box(s, lb_x + lb_w + 0.35, ly + 0.06, 3.55, lb_h - 0.12,
                    fill=R("panel_bg"), line=R("border"), line_w=0.8, shape=RR)
    ann.adjustments[0] = 0.08
    c.set_shape_text(ann, sla, size=10.5, bold=False, color=R("muted_text"),
                     align=PP_ALIGN.LEFT, font=c.FONT_B)
    x_ann.append(ann)
# 상향 화살표 커넥터 (하위 top → 상위 bottom, 화살표 위로)
for lower, upper in zip(lboxes[:-1], lboxes[1:]):
    x_conn.append(conn(s, cx(lower), topf(lower), cx(upper), botf(upper),
                       color=R("warn"), w=2.25, head=True))
# 좌측 상향 대형 화살표(데코)
big = c.add_box(s, 1.15, 2.25, 1.35, 4.55, fill=R("accent_primary"), line=None, shape=UPARROW)
c.set_shape_text(big, "대응\n상향", size=13, bold=True, color=R("header_text"),
                 align=PP_ALIGN.CENTER, font=c.FONT_H)
x_box.append(big)
c.group_asset(s, x_conn + x_box + x_ann, "ORG-014")
E("ORG-014", "에스컬레이션 체계 (단계별 대응 상향)", F2, 3,
  ["에스컬레이션", "이슈대응", "단계", "상향", "체계도"],
  {"levels": 4, "direction": "upward"},
  {"levels": [{"step": i + 1, "owner": lab.split("· ")[1], "sla": sla}
              for i, (lab, _, _, sla) in enumerate(levels)]},
  ["node-text", "add-level", "sla-text", "color", "connector"],
  ["에스컬레이션", "이슈 대응 체계", "리스크 관리"])

c.save_deck(prs, F2)

# ─────────────────────────────────────────────────────────────
frag = c.write_fragment("ORG_v2", entries)
print("SAVED:", F1, "|", F2)
print("FRAGMENT:", frag)
print("ENTRIES:", len(entries), [e["id"] for e in entries])
