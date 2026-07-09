# -*- coding: utf-8 -*-
"""
KPI 변형팩 v2 생성기 (KPI-011 ~ KPI-018)
v1(KPI-001~010)과 형태를 달리한 8개 KPI 에셋. 네이티브 도형+텍스트만.
다중도형 에셋은 c.group_asset 으로 asset:<ID> 그룹 1개로 묶는다(id_caption/제목 제외).
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'lib'))
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RR = MSO_SHAPE.ROUNDED_RECTANGLE
OVAL = MSO_SHAPE.OVAL
PIE = MSO_SHAPE.PIE
RADIUS = c.TOKENS["geom"]["round_radius"]
SW = 13.333


def sn(slide, x, y, w, num, unit="", num_color=None, unit_color=None,
       num_size=40, unit_size=16, align=PP_ALIGN.CENTER, h=0.9):
    """큰 숫자 + 단위 텍스트박스. 반환 tb(그룹에 append)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(num)
    c.set_kfont(r, c.FONT_H, num_size, True, num_color or c.role("accent_primary"))
    if unit:
        ru = p.add_run(); ru.text = " " + unit
        c.set_kfont(ru, c.FONT_H, unit_size, True, unit_color or c.role("muted_text"))
    return tb


def title(slide, text):
    c.add_text(slide, 0.15, 0.5, 11, 0.4, text, size=12, bold=True,
               color=c.role("muted_text"))


ACC = [c.role("accent_primary"), c.role("accent_secondary"),
       c.role("accent_point"), c.role("sub_header")]

# ============================================================
# FILE 1: KPI_gauge_v2.pptx  (링/게이지/배지/비교 = 진행·지표 시각화)
# ============================================================
prs1 = c.new_deck()

# ---- KPI-011: 링 프로그레스 세트 (얇은 도넛 링 3, % 중앙) ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-011")
title(s, "링 프로그레스 세트 (얇은 도넛 링 + % 강조)")
RINGS = [
    {"pct": 95, "unit": "%", "label": "목표 달성률", "col": 0},
    {"pct": 82, "unit": "%", "label": "고객 만족도", "col": 1},
    {"pct": 68, "unit": "%", "label": "수출 전환율", "col": 3},
]
A = []
gd = 2.4; gy = 2.0; ggap = 0.8
tot = gd*3 + ggap*2; gx0 = (SW - tot)/2
for i, cd in enumerate(RINGS):
    x = gx0 + i*(gd+ggap); frac = cd["pct"]/100.0
    accent = ACC[cd["col"]]
    A.append(c.add_box(s, x, gy, gd, gd, fill=c.role("border"), shape=OVAL))
    pie = c.add_box(s, x, gy, gd, gd, fill=accent, shape=PIE)
    pie.adjustments[0] = -90.0; pie.adjustments[1] = -90.0 + 360.0*frac
    A.append(pie)
    inr = gd*0.74  # 얇은 링
    A.append(c.add_box(s, x+(gd-inr)/2, gy+(gd-inr)/2, inr, inr,
                       fill=c.role("row_base"), shape=OVAL))
    A.append(sn(s, x, gy+gd/2-0.45, gd, str(cd["pct"]), cd["unit"],
                num_color=c.role("header_fill"), unit_color=accent, num_size=30, unit_size=14))
    A.append(c.add_text(s, x, gy+gd+0.12, gd, 0.5, cd["label"], size=13, bold=True,
                        color=c.role("body_text"), align=PP_ALIGN.CENTER))
c.group_asset(s, A, "KPI-011")

# ---- KPI-015: 배지형 지표 (둥근 배지 안 숫자 4) ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-015")
title(s, "배지형 지표 (둥근 배지 안 숫자)")
BADGE = [
    {"value": "98%", "label": "정시 납품률", "col": 0},
    {"value": "4.8", "label": "이용자 만족도", "col": 1},
    {"value": "120", "label": "협력 기업수", "col": 2},
    {"value": "32",  "label": "수출 대상국", "col": 3},
]
A = []
bd = 1.8; by = 2.2; bgap = 0.55
tot = bd*4 + bgap*3; bx0 = (SW - tot)/2
for i, cd in enumerate(BADGE):
    x = bx0 + i*(bd+bgap); accent = ACC[cd["col"]]
    ring = c.add_box(s, x-0.09, by-0.09, bd+0.18, bd+0.18, fill=None,
                     line=accent, line_w=1.5, shape=OVAL)
    A.append(ring)
    badge = c.add_box(s, x, by, bd, bd, fill=accent, shape=OVAL)
    c.set_shape_text(badge, cd["value"], size=30, bold=True, color=c.C["white"])
    A.append(badge)
    A.append(c.add_text(s, x-0.2, by+bd+0.15, bd+0.4, 0.5, cd["label"], size=13, bold=True,
                        color=c.role("body_text"), align=PP_ALIGN.CENTER))
c.group_asset(s, A, "KPI-015")

# ---- KPI-016: 반원 게이지 row (semicircle gauge 3) ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-016")
title(s, "반원 게이지 (semicircle gauge)")
SEMI = [
    {"pct": 88, "label": "예산 집행률", "col": 0},
    {"pct": 73, "label": "사업 목표 달성", "col": 1},
    {"pct": 95, "label": "시스템 가동률", "col": 3},
]
A = []
gd = 2.8; gy = 2.1; ggap = 0.9
tot = gd*3 + ggap*2; gx0 = (SW - tot)/2
for i, cd in enumerate(SEMI):
    x = gx0 + i*(gd+ggap); frac = cd["pct"]/100.0
    accent = ACC[cd["col"]]; cyl = gy + gd/2  # 지름선 y
    # 1) 회색 풀 원 (트랙)
    A.append(c.add_box(s, x, gy, gd, gd, fill=c.role("border"), shape=OVAL))
    # 2) 컬러 채움 PIE (서→우, 상단 반원 좌측부터)
    pie = c.add_box(s, x, gy, gd, gd, fill=accent, shape=PIE)
    pie.adjustments[0] = 180.0; pie.adjustments[1] = 180.0 + 180.0*frac
    A.append(pie)
    # 3) 안쪽 흰 원 → 링 홀
    inr = gd*0.58
    A.append(c.add_box(s, x+(gd-inr)/2, gy+(gd-inr)/2, inr, inr,
                       fill=c.role("row_base"), shape=OVAL))
    # 4) 하단 반원 흰 사각으로 클립 (지름선 아래 제거)
    A.append(c.add_box(s, x-0.05, cyl, gd+0.1, gd/2+0.1, fill=c.role("row_base")))
    # 5) % 숫자 (지름선 바로 위)
    A.append(sn(s, x, cyl-0.72, gd, str(cd["pct"]), "%",
                num_color=c.role("header_fill"), unit_color=accent, num_size=28, unit_size=13))
    # 6) 라벨 (지름선 아래)
    A.append(c.add_text(s, x, cyl+0.12, gd, 0.5, cd["label"], size=13, bold=True,
                        color=c.role("body_text"), align=PP_ALIGN.CENTER))
c.group_asset(s, A, "KPI-016")

# ---- KPI-017: 비교 KPI (전년 vs 올해 2숫자 대비 3) ----
s = c.blank_slide(prs1); c.id_caption(s, "KPI-017")
title(s, "비교 KPI (전년 vs 올해 대비)")
CMP = [
    {"label": "수출상담액", "prev": "33", "cur": "42", "unit": "만$", "col": 0},
    {"label": "참여기업",   "prev": "28", "cur": "40", "unit": "개사", "col": 1},
    {"label": "이용자 만족도","prev": "4.2","cur": "4.6","unit": "점",  "col": 3},
]
A = []
cw = (SW - 1.0 - 0.4*2)/3; x0 = 0.5; y = 1.7; ch = 3.6
for i, cd in enumerate(CMP):
    x = x0 + i*(cw+0.4); accent = ACC[cd["col"]]
    bg = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = RADIUS; A.append(bg)
    A.append(c.add_box(s, x, y, cw, 0.55, fill=accent, shape=RR))
    A.append(c.add_text(s, x, y, cw, 0.55, cd["label"], size=14, bold=True,
                        color=c.C["white"], align=PP_ALIGN.CENTER))
    # 전년 (좌, 흐림)
    A.append(c.add_text(s, x+0.1, y+0.75, cw-0.2, 0.35, "전년", size=11, bold=True,
                        color=c.role("muted_text"), align=PP_ALIGN.CENTER))
    A.append(sn(s, x+0.1, y+1.05, cw-0.2, cd["prev"], cd["unit"],
                num_color=c.role("muted_text"), unit_color=c.role("muted_text"),
                num_size=26, unit_size=12))
    # 화살표
    A.append(c.add_text(s, x, y+1.95, cw, 0.4, "▼", size=16, bold=True,
                        color=accent, align=PP_ALIGN.CENTER))
    # 올해 (강조)
    A.append(c.add_text(s, x+0.1, y+2.25, cw-0.2, 0.35, "올해", size=11, bold=True,
                        color=accent, align=PP_ALIGN.CENTER))
    A.append(sn(s, x+0.1, y+2.55, cw-0.2, cd["cur"], cd["unit"],
                num_color=c.role("header_fill"), unit_color=accent,
                num_size=34, unit_size=15))
c.group_asset(s, A, "KPI-017")

out1 = c.save_deck(prs1, "decks/02_kpi/KPI_gauge_v2.pptx")


# ============================================================
# FILE 2: KPI_cards_v2.pptx  (스트립/빅넘버/트렌드/스탯밴드)
# ============================================================
prs2 = c.new_deck()

# ---- KPI-012: 아이콘+지표 가로 스트립 (아이콘원+숫자+라벨 4행) ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-012")
title(s, "아이콘 + 지표 가로 스트립 (4행)")
STRIP = [
    {"value": "1,000", "unit": "건",  "label": "검증완료 기업정보", "col": 0},
    {"value": "40",    "unit": "개사", "label": "참여기업",         "col": 1},
    {"value": "42",    "unit": "만$",  "label": "수출상담액",       "col": 2},
    {"value": "4.6",   "unit": "점",   "label": "이용자 만족도",     "col": 3},
]
A = []
rw = 11.0; rh = 1.0; rx = (SW-rw)/2; ry0 = 1.3; rgap = 0.18
for i, cd in enumerate(STRIP):
    yy = ry0 + i*(rh+rgap); accent = ACC[cd["col"]]
    bg = c.add_box(s, rx, yy, rw, rh, fill=c.role("panel_bg"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = 0.14; A.append(bg)
    icd = 0.66
    ic = c.add_box(s, rx+0.3, yy+(rh-icd)/2, icd, icd, fill=accent, shape=OVAL)
    c.set_shape_text(ic, "ICON", size=8, bold=True, color=c.C["white"])
    A.append(ic)
    A.append(sn(s, rx+1.25, yy, 2.6, cd["value"], cd["unit"],
                num_color=accent, unit_color=c.role("muted_text"),
                num_size=28, unit_size=13, align=PP_ALIGN.LEFT, h=rh))
    A.append(c.add_text(s, rx+4.2, yy, rw-4.5, rh, cd["label"], size=15, bold=True,
                        color=c.role("body_text"), align=PP_ALIGN.LEFT))
c.group_asset(s, A, "KPI-012")

# ---- KPI-013: 큰숫자-좌 / 라벨-우 카드 3행 ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-013")
title(s, "큰숫자 좌 · 라벨 우 카드")
BIG = [
    {"value": "1,000", "unit": "건",  "label": "검증완료 기업정보 DB 구축", "col": 0},
    {"value": "42",    "unit": "만$", "label": "누적 수출상담 성사액",     "col": 1},
    {"value": "95",    "unit": "%",   "label": "연간 목표 대비 달성률",     "col": 3},
]
A = []
rw = 11.0; rh = 1.5; rx = (SW-rw)/2; ry0 = 1.4; rgap = 0.35
for i, cd in enumerate(BIG):
    yy = ry0 + i*(rh+rgap); accent = ACC[cd["col"]]
    bg = c.add_box(s, rx, yy, rw, rh, fill=c.role("row_base"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = 0.1; A.append(bg)
    A.append(c.add_box(s, rx, yy, 0.16, rh, fill=accent))
    A.append(sn(s, rx+0.5, yy, 4.0, cd["value"], cd["unit"],
                num_color=accent, unit_color=c.role("muted_text"),
                num_size=40, unit_size=17, align=PP_ALIGN.LEFT, h=rh))
    A.append(c.add_box(s, rx+4.8, yy+0.3, 0.02, rh-0.6, fill=c.role("border")))
    A.append(c.add_text(s, rx+5.15, yy, rw-5.4, rh, cd["label"], size=17, bold=True,
                        color=c.role("body_text"), align=PP_ALIGN.LEFT))
c.group_asset(s, A, "KPI-013")

# ---- KPI-014: 트렌드 카드 (미니 막대 스파크라인 + 숫자 3) ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-014")
title(s, "트렌드 카드 (미니 막대 + 숫자)")
TREND = [
    {"value": "1,000", "unit": "건",  "label": "기업정보",   "delta": "+18%", "col": 0},
    {"value": "42",    "unit": "만$", "label": "수출상담액", "delta": "+27%", "col": 1},
    {"value": "4.6",   "unit": "점",  "label": "이용자 만족도","delta": "+5%",  "col": 3},
]
A = []
cw = (SW - 1.0 - 0.4*2)/3; x0 = 0.5; y = 1.7; ch = 3.3
heights = [0.5, 0.72, 0.62, 0.95, 1.25]
for i, cd in enumerate(TREND):
    x = x0 + i*(cw+0.4); accent = ACC[cd["col"]]
    bg = c.add_box(s, x, y, cw, ch, fill=c.role("row_base"),
                   line=c.role("border"), line_w=1.0, shape=RR)
    bg.adjustments[0] = RADIUS; A.append(bg)
    A.append(c.add_text(s, x+0.25, y+0.3, cw-0.5, 0.4, cd["label"], size=14, bold=True,
                        color=c.role("muted_text"), align=PP_ALIGN.LEFT))
    A.append(sn(s, x+0.25, y+0.7, cw-0.5, cd["value"], cd["unit"],
                num_color=c.role("header_fill"), unit_color=c.role("muted_text"),
                num_size=36, unit_size=15, align=PP_ALIGN.LEFT))
    A.append(c.add_text(s, x+0.25, y+1.55, cw-0.5, 0.35, "▲ %s 전년比" % cd["delta"],
                        size=12, bold=True, color=c.role("accent_secondary"),
                        align=PP_ALIGN.LEFT))
    # 미니 막대 스파크라인 (하단)
    nb = len(heights); bw = 0.34; bgp = 0.16
    tot = nb*bw + (nb-1)*bgp; bx = x + (cw-tot)/2
    base = y + ch - 0.4
    for j, hh in enumerate(heights):
        col = accent if j == nb-1 else c.C["gray_300"]
        A.append(c.add_box(s, bx + j*(bw+bgp), base-hh, bw, hh, fill=col))
c.group_asset(s, A, "KPI-014")

# ---- KPI-018: 통계 밴드 (가로 긴 띠 4지표 구분) ----
s = c.blank_slide(prs2); c.id_caption(s, "KPI-018")
title(s, "통계 밴드 (가로 띠 4지표)")
BAND = [
    {"value": "1,000", "unit": "건",  "label": "검증 기업정보"},
    {"value": "40",    "unit": "개사", "label": "참여기업"},
    {"value": "42",    "unit": "만$",  "label": "수출상담액"},
    {"value": "4.6",   "unit": "점",   "label": "이용자 만족도"},
]
A = []
bw = 12.13; bh = 2.0; bx = (SW-bw)/2; byy = 2.6
band = c.add_box(s, bx, byy, bw, bh, fill=c.role("header_fill"), shape=RR)
band.adjustments[0] = 0.06; A.append(band)
seg = bw/4
for i, cd in enumerate(BAND):
    sx = bx + i*seg
    if i > 0:
        A.append(c.add_box(s, sx, byy+0.35, 0.02, bh-0.7, fill=c.role("accent_point")))
    A.append(sn(s, sx, byy+0.4, seg, cd["value"], cd["unit"],
                num_color=c.C["white"], unit_color=c.role("accent_point"),
                num_size=34, unit_size=15))
    A.append(c.add_text(s, sx, byy+bh-0.65, seg, 0.45, cd["label"], size=14, bold=True,
                        color=c.C["gray_100"], align=PP_ALIGN.CENTER))
c.group_asset(s, A, "KPI-018")

out2 = c.save_deck(prs2, "decks/02_kpi/KPI_cards_v2.pptx")

f1 = "decks/02_kpi/KPI_gauge_v2.pptx"
f2 = "decks/02_kpi/KPI_cards_v2.pptx"


# ============================================================
# 매니페스트 조각
# ============================================================
def E(aid, name, file_rel, slide_idx, cards, count, tags, editable, params, ruse):
    return c.entry(asset_id=aid, category="KPI", name=name, file_rel=file_rel,
                   slide_idx=slide_idx, tags=tags, params=params,
                   bindings={"cards": cards, "count": count},
                   editable=editable, recommended_use=ruse)


def cl(items, keys):
    return [{k: d[k] for k in keys} for d in items]


entries = [
    E("KPI-011", "링 프로그레스 세트 3종", f1, 1,
      [{"value": str(d["pct"]), "unit": d["unit"], "label": d["label"]} for d in RINGS], 3,
      ["KPI","성과지표","링","프로그레스","도넛","비율강조"],
      ["value","label","ring-color","pct"],
      {"rings": 3, "style": "ring-progress", "shape": "native-pie", "num_size": 30},
      ["달성률","목표대비","진행률"]),
    E("KPI-015", "배지형 지표 4종", f1, 2,
      cl(BADGE, ["value","label"]), 4,
      ["KPI","성과지표","배지","원형","핵심수치"],
      ["value","label","badge-color"],
      {"badges": 4, "style": "circle-badge", "num_size": 30},
      ["핵심지표","실적강조","한눈요약"]),
    E("KPI-016", "반원 게이지 3종", f1, 3,
      [{"value": str(d["pct"]), "unit": "%", "label": d["label"]} for d in SEMI], 3,
      ["KPI","성과지표","게이지","반원","semicircle","달성률"],
      ["value","label","gauge-color","pct"],
      {"gauges": 3, "style": "semicircle-gauge", "shape": "native-pie", "num_size": 28},
      ["달성률","목표대비","가동률"]),
    E("KPI-017", "비교 KPI 전년vs올해 3종", f1, 4,
      [{"value": d["cur"], "unit": d["unit"], "label": d["label"], "prev": d["prev"]} for d in CMP], 3,
      ["KPI","성과지표","비교","전년대비","2숫자대비"],
      ["value","prev","unit","label","accent-color"],
      {"cards": 3, "style": "compare-prev-cur", "num_size": 34},
      ["증감비교","전년대비","성과대비"]),
    E("KPI-012", "아이콘+지표 가로 스트립 4행", f2, 1,
      cl(STRIP, ["value","unit","label"]), 4,
      ["KPI","성과지표","가로스트립","아이콘","4행","리스트"],
      ["value","unit","label","icon","icon-color"],
      {"rows": 4, "style": "icon-strip", "icon_placeholder": True, "num_size": 28},
      ["성과목표","실적강조","서비스지표"]),
    E("KPI-013", "큰숫자 좌·라벨 우 카드 3행", f2, 2,
      cl(BIG, ["value","unit","label"]), 3,
      ["KPI","성과지표","빅넘버","라벨우측","가로카드"],
      ["value","unit","label","accent-color"],
      {"rows": 3, "style": "bignum-left-label-right", "num_size": 40},
      ["대표성과","핵심지표","실적강조"]),
    E("KPI-014", "트렌드 카드 3종 (미니막대)", f2, 3,
      [{"value": d["value"], "unit": d["unit"], "label": d["label"], "delta": d["delta"]} for d in TREND], 3,
      ["KPI","성과지표","트렌드","스파크라인","미니막대","증감"],
      ["value","unit","label","delta","trend-color"],
      {"cards": 3, "style": "trend-sparkbar", "bars": 5, "num_size": 36},
      ["성과추이","증감비교","성장강조"]),
    E("KPI-018", "통계 밴드 가로띠 4지표", f2, 4,
      cl(BAND, ["value","unit","label"]), 4,
      ["KPI","성과지표","통계밴드","가로띠","4지표","구분선"],
      ["value","unit","label","band-color"],
      {"segments": 4, "style": "stat-band", "num_size": 34},
      ["종합성과","핵심지표","한눈요약"]),
]

frag = c.write_fragment("KPI_v2", entries)
print("OUT1:", out1)
print("OUT2:", out2)
print("FRAG:", frag)
print("ENTRIES:", len(entries))
