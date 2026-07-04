# -*- coding: utf-8 -*-
"""TBL 변형팩 v2 — TBL-013 ~ TBL-022 (표 에셋 10종).
2개 파일:
  decks/01_tables/TBL_finance_v2.pptx  : 013,015,016,019,022 (재무/수치 계열)
  decks/01_tables/TBL_layout_v2.pptx   : 014,017,018,020,021 (레이아웃/상태 계열)
색은 c.role / c.C 만 사용. 한글은 c.set_kfont 경유. 페이지이미지·SmartArt 금지.
"""
import sys
sys.path.insert(0, '/home/pabang/myapp/.claude/pptx-asset-library/generators/lib')
import common as c
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

R = c.role
CC = c.C
NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid (테마 테두리/밴딩 제거)

# ---------- 로컬 테이블 헬퍼 ----------
def add_table(slide, x, y, w, h, nrows, ncols, col_widths=None, row_h=None):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    # 테마 스타일 제거 → 우리가 직접 채움/테두리 제어
    tblPr = tbl._tbl.tblPr
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        sid = etree.SubElement(tblPr, qn('a:tableStyleId'))
    sid.text = NO_STYLE
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    return gf, tbl

def set_cell(cell, text, size=11, bold=False, color=None, fill=None,
             align=PP_ALIGN.CENTER, font=None, ml=5):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = R("row_base")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(ml)
    cell.margin_top = cell.margin_bottom = Pt(2)
    tf = cell.text_frame; tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        c.set_kfont(r, font or c.FONT_B, size, bold, color or R("body_text"))

def cell_borders(cell, spec):
    """spec: {'L'|'R'|'T'|'B': (color_or_None, w_pt)}. color None → 투명(noFill)."""
    tcPr = cell._tc.get_or_add_tcPr()
    # B,T,R,L 순으로 앞쪽 insert → 최종 스키마 순서 L,R,T,B
    for edge, tag in (('B', 'a:lnB'), ('T', 'a:lnT'), ('R', 'a:lnR'), ('L', 'a:lnL')):
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)
        if edge not in spec:
            continue
        color, wpt = spec[edge]
        ln = etree.Element(qn(tag))
        ln.set('w', str(int(wpt * 12700))); ln.set('cap', 'flat')
        if color is None:
            etree.SubElement(ln, qn('a:noFill'))
        else:
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr')).set('val', str(color))
        tcPr.insert(0, ln)

def grid(tbl, color, wpt=0.75):
    for row in tbl.rows:
        for cell in row.cells:
            cell_borders(cell, {'L': (color, wpt), 'R': (color, wpt),
                                'T': (color, wpt), 'B': (color, wpt)})

def title_block(slide, text, sub=None):
    c.add_text(slide, 0.55, 0.5, 12.2, 0.5, text, size=18, bold=True,
               color=R("header_fill"), align=PP_ALIGN.LEFT, font=c.FONT_H)
    if sub:
        c.add_text(slide, 0.55, 1.0, 12.2, 0.32, sub, size=11,
                   color=R("muted_text"), align=PP_ALIGN.LEFT)

L = PP_ALIGN.LEFT; Rg = PP_ALIGN.RIGHT; Ct = PP_ALIGN.CENTER

# =========================================================
# TBL-013  재무/합계표 (소계·합계행 강조)
# =========================================================
def t013(slide):
    title_block(slide, "연도별 손익 요약표", "단위: 백만원 · 소계/영업이익 행 강조")
    x, y, w = 1.0, 1.55, 11.33
    heads = ["구분", "2024년", "2025년", "2026년(E)", "증감률"]
    # (label, v1, v2, v3, rate, kind)  kind: base/sub/total
    rows = [
        ("국내 매출", "1,200", "1,480", "1,760", "+18.9%", "base"),
        ("해외 매출", "320", "610", "980", "+60.7%", "base"),
        ("소계 (매출)", "1,520", "2,090", "2,740", "+31.1%", "sub"),
        ("인건비", "680", "790", "910", "+15.2%", "base"),
        ("운영비", "240", "310", "380", "+22.6%", "base"),
        ("소계 (비용)", "920", "1,100", "1,290", "+17.3%", "sub"),
        ("영업이익", "600", "990", "1,450", "+46.5%", "total"),
    ]
    nrows = len(rows) + 1
    rh = 0.56
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 5,
                        col_widths=[3.13, 2.05, 2.05, 2.05, 2.05], row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=L if j == 0 else Ct)
    for i, (lab, a, b, d, rate, kind) in enumerate(rows, start=1):
        if kind == "total":
            fill = CC["navy_800"]; tcol = R("header_text"); bold = True
        elif kind == "sub":
            fill = CC["gray_300"]; tcol = R("header_fill"); bold = True
        else:
            fill = R("row_stripe") if i % 2 == 0 else R("row_base"); tcol = R("body_text"); bold = False
        set_cell(tbl.cell(i, 0), lab, size=11, bold=bold or kind != "base", color=tcol, fill=fill, align=L)
        vals = [a, b, d]
        for k, v in enumerate(vals, start=1):
            set_cell(tbl.cell(i, k), v, size=11, bold=bold, color=tcol, fill=fill, align=Rg)
        rcol = tcol if kind != "base" else R("accent_secondary")
        set_cell(tbl.cell(i, 4), rate, size=11, bold=True, color=rcol, fill=fill, align=Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, "TBL-013")
    c.id_caption(slide, "TBL-013")

# =========================================================
# TBL-015  고밀도 컴팩트표 (작은 폰트 다열)
# =========================================================
def t015(slide):
    title_block(slide, "분기 KPI 상세 실적표", "단위: 건/백만원 · 8열 고밀도 컴팩트")
    x, y, w = 0.7, 1.5, 11.93
    heads = ["지표", "Q1", "Q2", "Q3", "Q4", "연간", "목표", "달성률"]
    rows = [
        ("신규 계약", "42", "55", "61", "73", "231", "220", "105%"),
        ("리텐션(%)", "88", "89", "91", "92", "90", "88", "102%"),
        ("MAU(천)", "120", "138", "156", "181", "181", "170", "106%"),
        ("ARPU", "3.2", "3.4", "3.6", "3.9", "3.5", "3.4", "103%"),
        ("매출", "1,180", "1,340", "1,510", "1,760", "5,790", "5,500", "105%"),
        ("CAC", "48", "45", "43", "41", "44", "46", "96%"),
        ("이탈률(%)", "4.1", "3.8", "3.5", "3.1", "3.6", "4.0", "90%"),
        ("NPS", "38", "42", "46", "51", "51", "45", "113%"),
    ]
    nrows = len(rows) + 1
    rh = 0.34
    cw = [2.13] + [1.4] * 7
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 8, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=9, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=L if j == 0 else Ct, ml=3)
    for i, row in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), row[0], size=9, bold=True, fill=base, align=L, ml=3)
        for j in range(1, 8):
            col = R("accent_primary") if j == 7 else R("body_text")
            set_cell(tbl.cell(i, j), row[j], size=9, bold=(j == 7), color=col, fill=base, align=Rg, ml=3)
    grid(tbl, R("border"), 0.5)
    c.name_asset(gf, "TBL-015")
    c.id_caption(slide, "TBL-015")

# =========================================================
# TBL-016  셀내 KPI 표 (숫자 + 증감칩)  ── 결합형(group_asset)
# =========================================================
def t016(slide):
    title_block(slide, "핵심 성과지표(KPI) 대시보드 표", "값 + 전기대비 증감칩 오버레이")
    x, y, w = 1.0, 1.6, 11.33
    heads = ["KPI 지표", "현재 값", "목표", "전기 대비"]
    # (지표, 값, 목표, delta_text, up_bool)
    rows = [
        ("월 거래액", "₩ 4.82억", "4.5억", "12.4%", True),
        ("활성 사용자", "18.1만", "17.0만", "6.5%", True),
        ("전환율", "3.9%", "3.5%", "0.4%p", True),
        ("고객이탈률", "3.1%", "4.0%", "0.9%p", False),
        ("객단가", "₩ 26,700", "25,000", "6.8%", True),
    ]
    nrows = len(rows) + 1
    rh = 0.66
    cw = [4.03, 2.9, 2.1, 2.3]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 4, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=L if j == 0 else Ct)
    for i, (lab, val, tgt, delta, up) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), lab, size=12, bold=True, fill=base, align=L)
        set_cell(tbl.cell(i, 1), val, size=15, bold=True, color=R("header_fill"), fill=base, align=Ct)
        set_cell(tbl.cell(i, 2), tgt, size=11, color=R("muted_text"), fill=base, align=Ct)
        set_cell(tbl.cell(i, 3), "", fill=base)
    grid(tbl, R("border"), 0.75)
    # 증감칩 오버레이 (마지막 열)
    chips = []
    chip_w, chip_h = 1.7, 0.42
    col3_x = x + sum(cw[:3])
    for i, (lab, val, tgt, delta, up) in enumerate(rows, start=1):
        cx = col3_x + (cw[3] - chip_w) / 2
        cy = y + rh * i + (rh - chip_h) / 2
        col = R("accent_secondary") if up else R("warn")
        chip = c.add_box(slide, cx, cy, chip_w, chip_h, fill=col, line=None,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        arrow = "▲" if up else "▼"
        c.set_shape_text(chip, "%s %s" % (arrow, delta), size=11, bold=True, color=R("header_text"))
        chips.append(chip)
    c.group_asset(slide, [gf] + chips, "TBL-016")
    c.id_caption(slide, "TBL-016")

# =========================================================
# TBL-019  예산집행표 (항목·예산·집행·잔액)
# =========================================================
def t019(slide):
    title_block(slide, "사업비 집행 현황표", "단위: 천원 · 집행률 및 합계 행 강조")
    x, y, w = 1.0, 1.55, 11.33
    heads = ["사업 항목", "예산액", "집행액", "잔액", "집행률"]
    rows = [
        ("인건비", "180,000", "142,500", "37,500", "79.2%"),
        ("콘텐츠 제작", "240,000", "205,800", "34,200", "85.8%"),
        ("마케팅·홍보", "120,000", "96,400", "23,600", "80.3%"),
        ("시스템 구축", "160,000", "88,200", "71,800", "55.1%"),
        ("일반 운영비", "60,000", "47,900", "12,100", "79.8%"),
    ]
    total = ("합계", "760,000", "580,800", "179,200", "76.4%")
    nrows = len(rows) + 2
    rh = 0.58
    cw = [3.33, 2.0, 2.0, 2.0, 2.0]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 5, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=L if j == 0 else Ct)
    for i, (lab, bud, exe, rem, rate) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), lab, size=11, bold=True, fill=base, align=L)
        set_cell(tbl.cell(i, 1), bud, size=11, fill=base, align=Rg, color=R("muted_text"))
        set_cell(tbl.cell(i, 2), exe, size=11, bold=True, fill=base, align=Rg, color=R("accent_primary"))
        set_cell(tbl.cell(i, 3), rem, size=11, fill=base, align=Rg)
        set_cell(tbl.cell(i, 4), rate, size=11, bold=True, fill=base, align=Rg, color=R("header_fill"))
    li = nrows - 1
    for j, v in enumerate(total):
        set_cell(tbl.cell(li, j), v, size=12, bold=True, color=R("header_text"),
                 fill=CC["navy_800"], align=L if j == 0 else Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, "TBL-019")
    c.id_caption(slide, "TBL-019")

# =========================================================
# TBL-022  우측정렬 수치표 (재무 정렬 · 음수 적색)
# =========================================================
def t022(slide):
    title_block(slide, "요약 재무상태표", "단위: 백만원 · 수치 우측정렬, 감소 항목 적색")
    x, y, w = 1.0, 1.55, 11.33
    heads = ["계정 과목", "당기", "전기", "증감액", "증감률"]
    # (계정, 당기, 전기, 증감, 율, neg_bool)
    rows = [
        ("유동자산", "3,420", "2,980", "+440", "+14.8%", False),
        ("비유동자산", "5,180", "5,360", "-180", "-3.4%", True),
        ("자산 총계", "8,600", "8,340", "+260", "+3.1%", False),
        ("유동부채", "2,110", "2,450", "-340", "-13.9%", True),
        ("비유동부채", "1,780", "1,690", "+90", "+5.3%", False),
        ("자본 총계", "4,710", "4,200", "+510", "+12.1%", False),
    ]
    nrows = len(rows) + 1
    rh = 0.56
    cw = [3.33, 2.0, 2.0, 2.0, 2.0]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 5, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=L if j == 0 else Rg)
    for i, (lab, cur, pre, dv, rate, neg) in enumerate(rows, start=1):
        emph = lab.endswith("총계")
        base = CC["gray_100"] if emph else (R("row_stripe") if i % 2 == 0 else R("row_base"))
        set_cell(tbl.cell(i, 0), lab, size=11, bold=emph, fill=base, align=L, color=R("header_fill") if emph else R("body_text"))
        set_cell(tbl.cell(i, 1), cur, size=11, bold=emph, fill=base, align=Rg)
        set_cell(tbl.cell(i, 2), pre, size=11, fill=base, align=Rg, color=R("muted_text"))
        dcol = R("warn") if neg else R("accent_secondary")
        set_cell(tbl.cell(i, 3), dv, size=11, bold=True, color=dcol, fill=base, align=Rg)
        set_cell(tbl.cell(i, 4), rate, size=11, bold=True, color=dcol, fill=base, align=Rg)
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, "TBL-022")
    c.id_caption(slide, "TBL-022")

# =========================================================
# TBL-014  보더리스 미니멀표 (헤더 밑줄만)
# =========================================================
def t014(slide):
    title_block(slide, "사업 개요 요약표", "보더리스 · 헤더 하단 밑줄만 표기")
    x, y, w = 1.0, 1.7, 11.33
    heads = ["구분", "주요 내용", "담당 부서"]
    rows = [
        ("사업명", "콘텐츠 해외진출 통합 플랫폼 구축", "사업기획팀"),
        ("사업 기간", "착수 후 12개월 (단계별 마일스톤 운영)", "PMO"),
        ("추진 목표", "글로벌 유통 채널 15개국 확장 및 현지화", "해외사업팀"),
        ("핵심 산출물", "플랫폼 v1.0 · 현지화 가이드 · 운영 매뉴얼", "개발팀"),
        ("기대 효과", "해외 매출 비중 35% 달성, 신규 파트너 40개사", "전략기획실"),
    ]
    nrows = len(rows) + 1
    rh = 0.7
    cw = [2.4, 6.53, 2.4]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 3, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=13, bold=True, color=R("header_fill"),
                 fill=R("row_base"), align=L)
        # 헤더 하단 밑줄(accent) 만, 나머지 투명
        cell_borders(tbl.cell(0, j), {'B': (R("accent_primary"), 2.0),
                                      'L': (None, 0.5), 'R': (None, 0.5), 'T': (None, 0.5)})
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            col = R("header_fill") if j == 0 else R("body_text")
            set_cell(tbl.cell(i, j), v, size=11, bold=(j == 0), color=col, fill=R("row_base"),
                     align=L)
            # 본문은 아주 옅은 하단 헤어라인만 (미니멀 구분)
            hair = R("border") if i < len(rows) else None
            cell_borders(tbl.cell(i, j), {'B': (hair, 0.5) if hair else (None, 0.5),
                                          'L': (None, 0.5), 'R': (None, 0.5), 'T': (None, 0.5)})
    c.name_asset(gf, "TBL-014")
    c.id_caption(slide, "TBL-014")

# =========================================================
# TBL-017  상태표 (신호등 색 셀)
# =========================================================
def t017(slide):
    title_block(slide, "과업 진행 상태 점검표", "상태 셀 신호등 색상 (완료·진행·지연)")
    x, y, w = 1.0, 1.55, 11.33
    heads = ["과업 항목", "상태", "담당", "비고"]
    # (항목, 상태라벨, 상태종류, 담당, 비고)  종류: ok/go/warn
    rows = [
        ("요구사항 정의", "완료", "ok", "기획팀", "확정본 배포"),
        ("UI/UX 설계", "완료", "ok", "디자인팀", "시안 승인"),
        ("백엔드 개발", "진행중", "go", "개발팀", "68% 진척"),
        ("현지화 번역", "진행중", "go", "해외팀", "3개국 완료"),
        ("통합 테스트", "지연", "warn", "QA팀", "일정 재조정"),
        ("오픈 준비", "대기", "wait", "PMO", "선행과업 대기"),
    ]
    signal = {"ok": R("accent_secondary"), "go": R("accent_primary"),
              "warn": R("warn"), "wait": R("muted_text")}
    nrows = len(rows) + 1
    rh = 0.56
    cw = [3.93, 1.8, 2.2, 3.4]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 4, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=L if j in (0, 3) else Ct)
    for i, (item, slab, kind, owner, note) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), item, size=11, bold=True, fill=base, align=L)
        set_cell(tbl.cell(i, 1), slab, size=11, bold=True, color=R("header_text"),
                 fill=signal[kind], align=Ct)
        set_cell(tbl.cell(i, 2), owner, size=11, fill=base, align=Ct, color=R("muted_text"))
        set_cell(tbl.cell(i, 3), note, size=11, fill=base, align=L, color=R("body_text"))
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, "TBL-017")
    c.id_caption(slide, "TBL-017")

# =========================================================
# TBL-018  추진일정표 (단계·활동·산출물)
# =========================================================
def t018(slide):
    title_block(slide, "단계별 추진 일정표", "단계 · 주요 활동 · 산출물 · 일정")
    x, y, w = 1.0, 1.55, 11.33
    heads = ["단계", "주요 활동", "산출물", "일정"]
    # (단계, 활동, 산출물, 일정, phase_color)
    rows = [
        ("1단계\n기획", "요구사항 분석 · 벤치마킹 · 범위 확정", "요구정의서 / WBS", "M1~M2", CC["navy_800"]),
        ("2단계\n설계", "정보구조 · 화면 · DB 스키마 설계", "설계서 / 프로토타입", "M3~M4", CC["blue_500"]),
        ("3단계\n구축", "핵심 기능 개발 · 현지화 · 연동", "플랫폼 v0.9", "M5~M9", CC["teal_500"]),
        ("4단계\n검증", "통합·부하 테스트 · 결함 조치", "테스트 결과서", "M10~M11", CC["cyan_500"]),
        ("5단계\n안정화", "오픈 · 운영 이관 · 교육", "운영 매뉴얼", "M12", CC["purple_600"]),
    ]
    nrows = len(rows) + 1
    rh = 0.72
    cw = [1.7, 5.13, 2.7, 1.8]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 4, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=Ct if j != 1 else L)
    for i, (stage, act, out, sched, pcol) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), stage, size=11, bold=True, color=R("header_text"),
                 fill=pcol, align=Ct)
        set_cell(tbl.cell(i, 1), act, size=11, fill=base, align=L)
        set_cell(tbl.cell(i, 2), out, size=11, bold=True, fill=base, align=Ct, color=R("accent_primary"))
        set_cell(tbl.cell(i, 3), sched, size=11, bold=True, fill=base, align=Ct, color=R("header_fill"))
    grid(tbl, R("border"), 0.75)
    c.name_asset(gf, "TBL-018")
    c.id_caption(slide, "TBL-018")

# =========================================================
# TBL-020  체크리스트표 (✓ 도형)  ── 결합형(group_asset)
# =========================================================
def t020(slide):
    title_block(slide, "오픈 전 점검 체크리스트", "✓ 도형 오버레이 · 완료 여부 표기")
    x, y, w = 1.0, 1.6, 11.33
    heads = ["확인", "점검 항목", "상태"]
    # (항목, checked, 상태text)
    rows = [
        ("보안 취약점 점검 및 조치 완료", True, "완료"),
        ("부하 테스트 (동시접속 5,000) 통과", True, "완료"),
        ("현지화 콘텐츠 3개국 검수 완료", True, "완료"),
        ("결제 모듈 실거래 검증", False, "진행중"),
        ("운영 매뉴얼 · 장애 대응 체계 수립", False, "예정"),
    ]
    nrows = len(rows) + 1
    rh = 0.66
    cw = [1.3, 7.53, 2.5]
    gf, tbl = add_table(slide, x, y, w, rh * nrows, nrows, 3, col_widths=cw, row_h=[rh] * nrows)
    for j, hd in enumerate(heads):
        set_cell(tbl.cell(0, j), hd, size=12, bold=True, color=R("header_text"),
                 fill=R("header_fill"), align=Ct if j != 1 else L)
    for i, (item, checked, note) in enumerate(rows, start=1):
        base = R("row_stripe") if i % 2 == 0 else R("row_base")
        set_cell(tbl.cell(i, 0), "", fill=base)
        set_cell(tbl.cell(i, 1), item, size=11, fill=base, align=L)
        ncol = R("accent_secondary") if checked else R("muted_text")
        set_cell(tbl.cell(i, 2), note, size=11, bold=True, color=ncol, fill=base, align=Ct)
    grid(tbl, R("border"), 0.75)
    boxes = []
    for i, (item, checked, note) in enumerate(rows, start=1):
        cbx = x + (cw[0] - 0.34) / 2
        cby = y + rh * i + (rh - 0.34) / 2
        fill = R("accent_secondary") if checked else R("row_base")
        line = R("accent_secondary") if checked else R("border")
        cb = c.add_box(slide, cbx, cby, 0.34, 0.34, fill=fill, line=line, line_w=1.5,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        if checked:
            c.set_shape_text(cb, "✓", size=14, bold=True, color=R("header_text"))
        boxes.append(cb)
    c.group_asset(slide, [gf] + boxes, "TBL-020")
    c.id_caption(slide, "TBL-020")

# =========================================================
# TBL-021  2톤 카드형 표 (행이 카드처럼)  ── 도형 결합형(group_asset)
# =========================================================
def t021(slide):
    title_block(slide, "핵심 성과 하이라이트 (카드형 표)", "행이 카드 · 좌측 라벨/우측 내용 2톤")
    x, w = 1.0, 11.33
    y0 = 1.65
    card_h, gap = 0.86, 0.16
    label_w = 3.0
    # (라벨, 지표값, 상세, 라벨색)
    rows = [
        ("글로벌 진출", "15개국", "누적 해외 매출 240억 · 현지 파트너 38개사", CC["navy_800"]),
        ("플랫폼 성장", "MAU 18.1만", "전년 대비 +64% · 리텐션 91%", CC["blue_500"]),
        ("콘텐츠 유통", "2,400편", "12개 언어 현지화 · 월 신규 180편", CC["teal_500"]),
        ("운영 효율", "-32% 비용", "자동화로 처리시간 절감 · SLA 99.9%", CC["purple_600"]),
    ]
    shapes = []
    for i, (lab, metric, desc, lc) in enumerate(rows):
        cy = y0 + i * (card_h + gap)
        # 카드 바탕(연한 톤)
        card = c.add_box(slide, x, cy, w, card_h, fill=CC["gray_050"], line=R("border"),
                         line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shapes.append(card)
        # 좌측 라벨 패널(진한 톤) — 2톤
        panel = c.add_box(slide, x, cy, label_w, card_h, fill=lc, line=None,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.set_shape_text(panel, lab, size=13, bold=True, color=R("header_text"))
        shapes.append(panel)
        # 우측 값(강조)
        vb = c.add_text(slide, x + label_w + 0.2, cy + 0.06, 2.9, card_h - 0.12, metric,
                        size=17, bold=True, color=lc, align=L, anchor=MSO_ANCHOR.MIDDLE, font=c.FONT_H)
        shapes.append(vb)
        # 우측 상세
        db = c.add_text(slide, x + label_w + 3.15, cy + 0.06, w - label_w - 3.35, card_h - 0.12,
                        desc, size=11, color=R("body_text"), align=L, anchor=MSO_ANCHOR.MIDDLE)
        shapes.append(db)
    c.group_asset(slide, shapes, "TBL-021")
    c.id_caption(slide, "TBL-021")

# =========================================================
# 빌드
# =========================================================
def build():
    # ---- 파일 A: 재무/수치 계열 ----
    prsA = c.new_deck()
    for fn in (t013, t015, t016, t019, t022):
        fn(c.blank_slide(prsA))
    pA = c.save_deck(prsA, "decks/01_tables/TBL_finance_v2.pptx")

    # ---- 파일 B: 레이아웃/상태 계열 ----
    prsB = c.new_deck()
    for fn in (t014, t017, t018, t020, t021):
        fn(c.blank_slide(prsB))
    pB = c.save_deck(prsB, "decks/01_tables/TBL_layout_v2.pptx")

    # ---- 매니페스트 조각 ----
    fA = "decks/01_tables/TBL_finance_v2.pptx"
    fB = "decks/01_tables/TBL_layout_v2.pptx"
    E = c.entry
    entries = [
        E("TBL-013", "TBL", "재무 손익 요약표 (소계·합계행 강조)", fA, 1,
          ["table", "finance", "subtotal", "total", "재무", "합계표"],
          {"cols": 5, "rows": 8}, ["구분", "연도값", "증감률"],
          ["header", "row_label", "cell_value"],
          recommended_use=["손익요약", "재무보고", "연도별 실적"]),
        E("TBL-015", "TBL", "고밀도 컴팩트표 (8열 소폰트)", fA, 2,
          ["table", "compact", "dense", "고밀도", "다열"],
          {"cols": 8, "rows": 8, "font_pt": 9}, ["지표", "분기값", "달성률"],
          ["header", "cell_value"],
          recommended_use=["KPI 상세", "분기실적", "다지표 요약"]),
        E("TBL-016", "TBL", "셀내 KPI 표 (값+증감칩)", fA, 3,
          ["table", "kpi", "chip", "delta", "증감칩", "결합형"],
          {"cols": 4, "rows": 5, "combined": True}, ["지표", "값", "증감칩"],
          ["header", "kpi_value", "delta_chip"],
          recommended_use=["KPI 대시보드", "성과요약"]),
        E("TBL-019", "TBL", "예산집행표 (예산·집행·잔액·합계강조)", fA, 4,
          ["table", "budget", "execution", "예산", "집행", "합계"],
          {"cols": 5, "rows": 6}, ["항목", "금액", "집행률"],
          ["header", "row_label", "amount", "total_row"],
          recommended_use=["예산집행", "사업비 현황"]),
        E("TBL-022", "TBL", "우측정렬 수치표 (재무·음수적색)", fA, 5,
          ["table", "numeric", "right-align", "finance", "재무정렬"],
          {"cols": 5, "rows": 6, "align": "right"}, ["계정", "수치", "증감"],
          ["header", "account", "figure"],
          recommended_use=["재무상태표", "회계 요약", "수치정렬"]),
        E("TBL-014", "TBL", "보더리스 미니멀표 (헤더 밑줄만)", fB, 1,
          ["table", "borderless", "minimal", "미니멀", "밑줄"],
          {"cols": 3, "rows": 5, "borderless": True}, ["구분", "내용", "담당"],
          ["header", "row_label", "content"],
          recommended_use=["사업개요", "깔끔한 요약", "본문 삽입표"]),
        E("TBL-017", "TBL", "상태표 (신호등 색 셀)", fB, 2,
          ["table", "status", "traffic-light", "상태", "신호등"],
          {"cols": 4, "rows": 6}, ["항목", "상태", "담당"],
          ["header", "item", "status_cell"],
          recommended_use=["진행상태", "과업점검", "RAG 상태"]),
        E("TBL-018", "TBL", "추진일정표 (단계·활동·산출물)", fB, 3,
          ["table", "schedule", "phase", "추진일정", "산출물"],
          {"cols": 4, "rows": 5}, ["단계", "활동", "산출물", "일정"],
          ["header", "phase", "activity", "deliverable"],
          recommended_use=["추진일정", "단계별 계획", "WBS 요약"]),
        E("TBL-020", "TBL", "체크리스트표 (✓ 도형)", fB, 4,
          ["table", "checklist", "check", "체크리스트", "결합형"],
          {"cols": 3, "rows": 5, "combined": True}, ["점검항목", "체크", "상태"],
          ["header", "check_shape", "item", "status"],
          recommended_use=["점검표", "오픈 체크리스트", "완료여부"]),
        E("TBL-021", "TBL", "2톤 카드형 표 (행=카드)", fB, 5,
          ["table", "card", "two-tone", "카드형", "도형결합"],
          {"rows": 4, "shape_based": True}, ["라벨", "값", "상세"],
          ["label_panel", "metric", "detail"],
          recommended_use=["성과 하이라이트", "카드형 요약", "강조 표"]),
    ]
    frag = c.write_fragment("TBL_v2", entries)
    print("SAVED:", pA)
    print("SAVED:", pB)
    print("FRAGMENT:", frag)

if __name__ == "__main__":
    build()
